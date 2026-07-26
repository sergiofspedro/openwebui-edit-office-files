    async def generate_document(self, content: str, title: str = "Document", theme: str = "professional", __user__=None, __request__=None) -> str:
        """Generate a professional Word document with modern styling and themes.

        Args:
            content: Markdown-formatted content
            title: Document title (used for cover page and filename)
            theme: Visual theme - professional, modern, creative, corporate, minimal, elegant
        Returns:
            Markdown link to the generated file
        """
        try:
            from docx import Document
            from docx.shared import Inches, Pt, Cm, RGBColor, Emu
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.table import WD_TABLE_ALIGNMENT
            from docx.oxml.ns import qn, nsdecls
            from docx.oxml import parse_xml
            import datetime

            doc = Document()

            # --- Page Setup ---
            section = doc.sections[0]
            section.top_margin = Cm(2.0)
            section.bottom_margin = Cm(2.0)
            section.left_margin = Cm(2.5)
            section.right_margin = Cm(2.5)

            # --- Color Themes ---
            themes = {
                "professional": {"primary": "1F4E79", "accent": "2E75B6", "light": "D6E4F0", "text": "333333"},
                "modern": {"primary": "2D3436", "accent": "6C5CE7", "light": "DFE6E9", "text": "2D3436"},
                "creative": {"primary": "E17055", "accent": "FDCB6E", "light": "FFF3E0", "text": "2D3436"},
                "corporate": {"primary": "003366", "accent": "CC0000", "light": "E8EEF4", "text": "1A1A1A"},
                "minimal": {"primary": "000000", "accent": "666666", "light": "F5F5F5", "text": "333333"},
                "elegant": {"primary": "4A235A", "accent": "8E44AD", "light": "F3E5F5", "text": "1A1A1A"},
            }
            colors = themes.get(theme, themes["professional"])

            def hex_to_rgb(hex_color):
                return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

            # --- Default Font ---
            style = doc.styles['Normal']
            font = style.font
            font.name = 'Calibri'
            font.size = Pt(11)
            font.color.rgb = hex_to_rgb(colors["text"])
            style.paragraph_format.space_after = Pt(8)
            style.paragraph_format.line_spacing = 1.15

            # --- Heading Styles ---
            for i in range(1, 4):
                heading_style = doc.styles['Heading %d' % i]
                heading_style.font.name = 'Calibri'
                heading_style.font.color.rgb = hex_to_rgb(colors["primary"])
                if i == 1:
                    heading_style.font.size = Pt(24)
                    heading_style.paragraph_format.space_before = Pt(24)
                    heading_style.paragraph_format.space_after = Pt(12)
                elif i == 2:
                    heading_style.font.size = Pt(18)
                    heading_style.paragraph_format.space_before = Pt(18)
                    heading_style.paragraph_format.space_after = Pt(8)
                else:
                    heading_style.font.size = Pt(14)
                    heading_style.paragraph_format.space_before = Pt(12)
                    heading_style.paragraph_format.space_after = Pt(6)

            # --- Cover Page ---
            cover_table = doc.add_table(rows=1, cols=1)
            cover_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            cell = cover_table.rows[0].cells[0]
            cell.width = Inches(6.5)
            shading_elm = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["primary"]}"/>')
            cell._tc.get_or_add_tcPr().append(shading_elm)

            p = cell.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_before = Pt(40)
            p.paragraph_format.space_after = Pt(8)
            run = p.add_run(_format_text(title))
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.name = 'Calibri'

            p = cell.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_after = Pt(30)
            run = p.add_run(datetime.datetime.now().strftime("%B %d, %Y"))
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(200, 200, 200)
            run.font.name = 'Calibri'

            doc.add_paragraph()

            # --- Process Content ---
            lines = content.split('\n')
            table_rows = []
            callout_lines = []

            for line in lines:
                line = line.strip()
                if not line:
                    if callout_lines:
                        _add_callout_box(doc, callout_lines, colors, hex_to_rgb)
                        callout_lines = []
                    continue

                if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
                    level = line.count('#')
                    text = _format_text(line.lstrip('#').strip())
                    doc.add_heading(text, level=min(level, 3))
                    continue

                if line.startswith('> '):
                    callout_lines.append(line[2:])
                    continue
                elif callout_lines:
                    _add_callout_box(doc, callout_lines, colors, hex_to_rgb)
                    callout_lines = []

                if line.startswith('|') and line.endswith('|'):
                    cells = [c.strip() for c in line.split('|')[1:-1]]
                    if all(c.startswith('---') for c in cells):
                        continue
                    table_rows.append(cells)
                    continue
                elif table_rows:
                    _add_professional_table(doc, table_rows, colors, hex_to_rgb)
                    table_rows = []

                if line.startswith('- ') or line.startswith('* '):
                    text = _format_text(line[2:].strip())
                    doc.add_paragraph(text, style='List Bullet')
                    continue

                if re.match(r'^\d+\.\s', line):
                    text = _format_text(re.sub(r'^\d+\.\s', '', line))
                    doc.add_paragraph(text, style='List Number')
                    continue

                text = _format_text(line)
                doc.add_paragraph(text)

            if table_rows:
                _add_professional_table(doc, table_rows, colors, hex_to_rgb)

            # --- Footer with page numbers ---
            footer = section.footer
            footer.is_linked_to_previous = False
            p = footer.paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run("Page ")
            run.font.size = Pt(8)
            run.font.color.rgb = hex_to_rgb(colors["accent"])
            fldChar1 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
            run._r.append(fldChar1)
            instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
            run._r.append(instrText)
            fldChar2 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
            run._r.append(fldChar2)

            out = io.BytesIO()
            doc.save(out)
            out.seek(0)
            url, fname = await self._save_and_link(out.getvalue(), "%s.docx" % title, __request__)
            if url:
                return "Document created: [%s](%s)" % (fname, url)
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def generate_slides(self, content: str, title: str = "Presentation", theme: str = "modern", __user__=None, __request__=None) -> str:
        """Generate professional PowerPoint slides with modern design.

        Args:
            content: Markdown content (headings become slides, bullets become content)
            title: Presentation title for the first slide
            theme: Visual theme - modern, light, dark, corporate, creative, minimal
        Returns:
            Markdown link to the generated file
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu, Cm
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
            import datetime

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # --- Color Themes ---
            themes = {
                "modern": {"bg": "1A1A2E", "accent": "E94560", "text": "FFFFFF", "subtitle": "B0B0B0", "card": "16213E"},
                "light": {"bg": "FFFFFF", "accent": "2563EB", "text": "1E293B", "subtitle": "64748B", "card": "F1F5F9"},
                "dark": {"bg": "0F172A", "accent": "38BDF8", "text": "F8FAFC", "subtitle": "94A3B8", "card": "1E293B"},
                "corporate": {"bg": "FFFFFF", "accent": "003366", "text": "1A1A1A", "subtitle": "666666", "card": "F0F4F8"},
                "creative": {"bg": "FFF8F0", "accent": "FF6B6B", "text": "2D3436", "subtitle": "636E72", "card": "FFEAA7"},
                "minimal": {"bg": "FAFAFA", "accent": "000000", "text": "333333", "subtitle": "999999", "card": "F0F0F0"},
            }
            colors = themes.get(theme, themes["modern"])

            def hex_to_rgb(h):
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

            def set_slide_bg(slide, color_hex):
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(color_hex)

            def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
                txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = _format_text(text)
                p.font.size = Pt(font_size)
                p.font.bold = bold
                p.font.color.rgb = hex_to_rgb(color or colors["text"])
                p.font.name = 'Calibri'
                p.alignment = alignment
                return txBox

            def add_accent_bar(slide, left, top, width, height, color_hex):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
                shape.line.fill.background()
                return shape

            # --- Slide 1: Title Slide ---
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide, colors["bg"])

            add_accent_bar(slide, 0, 3.0, 13.333, 0.06, colors["accent"])
            add_text_box(slide, 1.5, 1.5, 10, 1.5, _format_text(title), font_size=44, bold=True, color=colors["text"])
            add_text_box(slide, 1.5, 3.3, 10, 0.8, "Generated %s" % datetime.datetime.now().strftime("%B %d, %Y"), font_size=18, color=colors["subtitle"])

            # --- Process content into slides ---
            lines = content.split('\n')
            current_slide_lines = []
            slide_count = 1

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if line.startswith('# ') or line.startswith('## '):
                    if current_slide_lines:
                        _render_content_slide(prs, current_slide_lines, colors, hex_to_rgb, add_text_box, add_accent_bar, set_slide_bg, slide_count)
                        slide_count += 1
                    current_slide_lines = [line]
                else:
                    current_slide_lines.append(line)

            if current_slide_lines:
                _render_content_slide(prs, current_slide_lines, colors, hex_to_rgb, add_text_box, add_accent_bar, set_slide_bg, slide_count)

            out = io.BytesIO()
            prs.save(out)
            out.seek(0)
            url, fname = await self._save_and_link(out.getvalue(), "%s.pptx" % title, __request__)
            if url:
                return "Presentation created: [%s](%s)" % (fname, url)
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def generate_spreadsheet(self, content: str, title: str = "Spreadsheet", theme: str = "professional", __user__=None, __request__=None) -> str:
        """Generate a professional Excel spreadsheet with modern styling.

        Args:
            content: CSV or tab-delimited data (first row = headers, rest = data)
            title: Spreadsheet title / filename
            theme: Visual theme - professional, modern, corporate, minimal, colorful, pastel
        Returns:
            Markdown link to the generated file
        """
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
            from openpyxl.worksheet.table import Table, TableStyleInfo
            import io

            # --- Color Themes ---
            themes = {
                "professional": {"header": "1F4E79", "header_font": "FFFFFF", "alt_row": "F0F5FA", "border": "CCCCCC", "accent": "2E75B6"},
                "modern": {"header": "2D3436", "header_font": "FFFFFF", "alt_row": "F5F5F5", "border": "DFE6E9", "accent": "6C5CE7"},
                "corporate": {"header": "003366", "header_font": "FFFFFF", "alt_row": "E8EEF4", "border": "B0C4DE", "accent": "CC0000"},
                "minimal": {"header": "333333", "header_font": "FFFFFF", "alt_row": "F8F8F8", "border": "DDDDDD", "accent": "666666"},
                "colorful": {"header": "E17055", "header_font": "FFFFFF", "alt_row": "FFF3E0", "border": "FABEAB", "accent": "FDCB6E"},
                "pastel": {"header": "6C5CE7", "header_font": "FFFFFF", "alt_row": "F3E5F5", "border": "CE93D8", "accent": "A29BFE"},
            }
            colors = themes.get(theme, themes["professional"])

            header_fill = PatternFill(start_color=colors["header"], end_color=colors["header"], fill_type="solid")
            header_font = Font(name='Calibri', size=11, bold=True, color=colors["header_font"])
            alt_fill = PatternFill(start_color=colors["alt_row"], end_color=colors["alt_row"], fill_type="solid")
            body_font = Font(name='Calibri', size=11)
            thin_border = Border(
                left=Side(style='thin', color=colors["border"]),
                right=Side(style='thin', color=colors["border"]),
                top=Side(style='thin', color=colors["border"]),
                bottom=Side(style='thin', color=colors["border"])
            )

            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = title[:31]

            # Parse CSV or tab-delimited content
            lines = [l.strip() for l in content.split('\n') if l.strip()]
            data = []
            for line in lines:
                sep = '\t' if '\t' in line else ','
                cells = [c.strip().strip('"') for c in line.split(sep)]
                data.append(cells)

            if not data:
                return json.dumps({"error": "No data provided"})

            # Write data
            for row in data:
                ws.append([_format_text(c) if isinstance(c, str) else c for c in row])

            # Style header row
            for cell in ws[1]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border

            # Zebra striping for data rows
            for i, row in enumerate(ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=ws.max_column), start=2):
                if i % 2 == 0:
                    for cell in row:
                        cell.fill = alt_fill
                        cell.border = thin_border
                        cell.font = body_font
                else:
                    for cell in row:
                        cell.border = thin_border
                        cell.font = body_font

            # Auto-fit columns
            for col in ws.columns:
                max_len = 0
                col_letter = get_column_letter(col[0].column)
                for cell in col:
                    try:
                        max_len = max(max_len, len(str(cell.value or '')))
                    except:
                        pass
                ws.column_dimensions[col_letter].width = min(max_len + 3, 40)

            # Freeze top row
            ws.freeze_panes = 'A2'

            # Add auto-filter
            if ws.max_row > 1:
                ws.auto_filter.ref = "A1:%s%d" % (get_column_letter(ws.max_column), ws.max_row)

            # Add Excel Table if data exists
            if ws.max_row > 1 and ws.max_column > 0:
                try:
                    tab = Table(
                        displayName="Table_" + ws.title.replace(' ', '')[:20],
                        ref="A1:%s%d" % (get_column_letter(ws.max_column), ws.max_row)
                    )
                    style = TableStyleInfo(
                        name="TableStyleMedium2",
                        showFirstColumn=False,
                        showLastColumn=False,
                        showRowStripes=True,
                        showColumnStripes=False
                    )
                    tab.tableStyleInfo = style
                    ws.add_table(tab)
                except Exception:
                    pass

            out = io.BytesIO()
            wb.save(out)
            out.seek(0)
            url, fname = await self._save_and_link(out.getvalue(), "%s.xlsx" % title, __request__)
            if url:
                return "Spreadsheet created: [%s](%s)" % (fname, url)
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

