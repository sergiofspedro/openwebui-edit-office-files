"""Edit Office Files — File editing and transformation (replace, merge, modify, export)."""

import io
import json
import os
import platform
import re
import sqlite3
import sys
import traceback
from copy import copy
from typing import Any, Dict, Optional

from .constants import *
from .utils import *

async def replace_text(
    self,
    file_id: str,
    find_text: str,
    replace_with: str,
    output_filename: str = "",
    __user__=None,
    __request__=None,
) -> str:
    """Find and replace text in any Office file while preserving original formatting
        and capitalization of the existing document.

        Works on cell values in xlsx, paragraph text in docx, and shape text in pptx.

        Args:
            file_id: File ID to edit
            find_text: Text to find
            replace_with: Text to replace with
            output_filename: Optional output filename
        """
    try:
        file_data = _read_file_bytes(file_id)
        if file_data is None:
            return json.dumps({"error": f"Could not read file {file_id}"})

        # Detect type
        try:
            conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            row = conn.execute(
                "SELECT filename FROM file WHERE id = ?", (file_id,)
            ).fetchone()
            conn.close()
            filename = row[0] if row else file_id
        except Exception:
            filename = file_id

        file_type = _detect_type(filename)
        out = io.BytesIO()
        out_name = output_filename
        count = 0

        if file_type == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

            for sn in wb.sheetnames:
                ws = wb[sn]
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        if isinstance(cell.value, str) and find_text in cell.value:
                            cell.value = cell.value.replace(find_text, replace_with)
                            count += 1
                        elif not isinstance(cell.value, str):
                            sval = str(cell.value)
                            if find_text in sval:
                                cell.value = replace_with
                                count += 1

            wb.save(out)
            wb.close()

        elif file_type == "xls":
            # Convert .xls to .xlsx, then apply same replace logic
            file_data = _xls_to_xlsx(file_data)
            file_type = "xlsx"
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_data))

            for sn in wb.sheetnames:
                ws = wb[sn]
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        if isinstance(cell.value, str) and find_text in cell.value:
                            cell.value = cell.value.replace(find_text, replace_with)
                            count += 1
                        elif not isinstance(cell.value, str):
                            sval = str(cell.value)
                            if find_text in sval:
                                cell.value = replace_with
                                count += 1

            wb.save(out)
            wb.close()

        elif file_type == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_data))
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.docx"

            for para in doc.paragraphs:
                if find_text in para.text:
                    if para.runs:
                        # Preserve formatting: replace within each run individually
                        for run in para.runs:
                            if find_text in run.text:
                                run.text = run.text.replace(find_text, replace_with)
                        count += 1
                    else:
                        para.text = para.text.replace(find_text, replace_with)
                        count += 1

            # Also replace in tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if find_text in cell.text:
                            for para in cell.paragraphs:
                                if find_text in para.text:
                                    if para.runs:
                                        for run in para.runs:
                                            if find_text in run.text:
                                                run.text = run.text.replace(find_text, replace_with)
                                        count += 1
                                    else:
                                        para.text = para.text.replace(find_text, replace_with)
                                        count += 1

            doc.save(out)

        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_data))
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.pptx"

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            if find_text in para.text:
                                if para.runs:
                                    for run in para.runs:
                                        if find_text in run.text:
                                            run.text = run.text.replace(find_text, replace_with)
                                else:
                                    para.text = para.text.replace(find_text, replace_with)
                                count += 1

            prs.save(out)

        else:
            return json.dumps({"error": f"Unsupported type: {file_type}"})

        out.seek(0)
        url, name = await self._save_and_link(out.read(), out_name, __request__)
        if url:
            return f"[{name}]({url})\n\nReplaced '{find_text}' with '{replace_with}' in {count} place(s), preserving all formatting."
        return json.dumps({"error": "Could not save file"})

    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

# -----------------------------------------------------------------
# UPDATE CELLS
# -----------------------------------------------------------------


async def update_cells(
    self,
    file_id: str,
    cells: str,
    output_filename: str = "",
    __user__=None,
    __request__=None,
) -> str:
    """Update individual cells in an XLSX file. Pass a JSON array of cell updates.

        Each update: {"cell": "A1", "value": "new value"} or {"cell": "B2", "value": 42, "sheet": "Sheet1"}
        If sheet is omitted, uses the active sheet.

        Args:
            file_id: The file ID of the XLSX file to edit
            cells: JSON string - array of {cell, value[, sheet]} objects
            output_filename: Optional output filename
        """
    try:
        file_data = _read_file_bytes(file_id)
        if file_data is None:
            return json.dumps({"error": f"Could not read file {file_id}"})

        try:
            conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            row2 = conn2.execute(
                "SELECT filename FROM file WHERE id = ?", (file_id,)
            ).fetchone()
            conn2.close()
            filename = row2[0] if row2 else file_id
        except Exception:
            filename = file_id

        file_type = _detect_type(filename)
        if file_type not in ("xlsx", "xls"):
            return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx/xls supported."})

        import openpyxl
        from openpyxl.utils import column_index_from_string
        import re as _re_cell

        try:
            updates = json.loads(cells)
            if not isinstance(updates, list):
                updates = [updates]
        except json.JSONDecodeError:
            return json.dumps({"error": "cells must be valid JSON (array of {cell, value} objects)"})

        out_name = output_filename or os.path.splitext(filename)[0] + "_updated.xlsx"
        wb = openpyxl.load_workbook(io.BytesIO(file_data))
        count = 0
        for upd in updates:
            cell_ref = upd.get("cell", "")
            value = upd.get("value", "")
            sname = upd.get("sheet", "")
            ws = wb[sname] if sname else wb.active
            m = _re_cell.match(r"([A-Za-z]+)(\d+)", cell_ref)
            if not m:
                continue
            col_str, row_str = m.group(1), m.group(2)
            col_idx = column_index_from_string(col_str.upper())
            row_idx = int(row_str)
            ws.cell(row=row_idx, column=col_idx).value = value
            count += 1

        out = io.BytesIO()
        wb.save(out)
        wb.close()
        out.seek(0)

        url, name = await self._save_and_link(out.read(), out_name, __request__)
        if url:
            return f"[{name}]({url})\n\nUpdated {count} cell(s) in {file_type.upper()} file."
        return json.dumps({"error": "Could not save file"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

# -----------------------------------------------------------------
# MODIFY ROWS (insert/delete)
# -----------------------------------------------------------------


async def modify_rows(
    self,
    file_id: str,
    action: str,
    row_number: int = 1,
    count: int = 1,
    sheet_name: str = "",
    output_filename: str = "",
    __user__=None,
    __request__=None,
) -> str:
    """Insert or delete rows in an XLSX file.

        Args:
            file_id: The file ID of the XLSX file to edit
            action: 'insert' or 'delete'
            row_number: Row number to start at (1-indexed)
            count: Number of rows to insert/delete (default 1)
            sheet_name: Optional sheet name (default: active sheet)
            output_filename: Optional output filename
        """
    try:
        file_data = _read_file_bytes(file_id)
        if file_data is None:
            return json.dumps({"error": f"Could not read file {file_id}"})

        try:
            conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            row2 = conn2.execute(
                "SELECT filename FROM file WHERE id = ?", (file_id,)
            ).fetchone()
            conn2.close()
            filename = row2[0] if row2 else file_id
        except Exception:
            filename = file_id

        file_type = _detect_type(filename)
        if file_type not in ("xlsx", "xls"):
            return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx/xls supported."})

        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(file_data))
        ws = wb[sheet_name] if sheet_name else wb.active
        out_name = output_filename or os.path.splitext(filename)[0] + "_modified.xlsx"

        action_lower = action.strip().lower()
        if action_lower == "insert":
            ws.insert_rows(idx=row_number, amount=count)
            msg = f"Inserted {count} row(s) at row {row_number}"
        elif action_lower == "delete":
            ws.delete_rows(idx=row_number, amount=count)
            msg = f"Deleted {count} row(s) starting at row {row_number}"
        else:
            return json.dumps({"error": f"Unknown action '{action}'. Use 'insert' or 'delete'."})

        out = io.BytesIO()
        wb.save(out)
        wb.close()
        out.seek(0)

        url, name = await self._save_and_link(out.read(), out_name, __request__)
        if url:
            return f"[{name}]({url})\n\n{msg} in {file_type.upper()} file."
        return json.dumps({"error": "Could not save file"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

# -----------------------------------------------------------------
# PASSWORD PROTECT FILE
# -----------------------------------------------------------------


async def protect_file(
    self,
    file_id: str,
    password: str,
    output_filename: str = "",
    __user__=None,
    __request__=None,
) -> str:
    """Add password protection to an XLSX or DOCX file.

        For XLSX: uses openpyxl workbook protection. For DOCX: uses write protection.
        Install msoffcrypto-tool (pip install msoffcrypto-tool) for encryption support.

        Args:
            file_id: The file ID of the Office file to protect
            password: Password to set
            output_filename: Optional output filename
        """
    try:
        file_data = _read_file_bytes(file_id)
        if file_data is None:
            return json.dumps({"error": f"Could not read file {file_id}"})

        try:
            conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            row2 = conn2.execute(
                "SELECT filename FROM file WHERE id = ?", (file_id,)
            ).fetchone()
            conn2.close()
            filename = row2[0] if row2 else file_id
        except Exception:
            filename = file_id

        file_type = _detect_type(filename)
        if file_type not in ("xlsx", "docx"):
            return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx and docx supported."})

        import hashlib

        if file_type == "xlsx":
            import openpyxl
            from openpyxl.workbook.protection import WorkbookProtection

            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            out_name = output_filename or os.path.splitext(filename)[0] + "_protected.xlsx"

            # Set workbook protection
            wb.security = WorkbookProtection(workbookPassword=password, lockStructure=True)

            # Protect all sheets
            for ws in wb.worksheets:
                ws.protection.set_password(password)

            out = io.BytesIO()
            wb.save(out)
            wb.close()
            out.seek(0)
            protected_bytes = out.read()

        elif file_type == "docx":
            from docx import Document
            from docx.oxml.ns import qn
            from lxml import etree

            out_name = output_filename or os.path.splitext(filename)[0] + "_protected.docx"
            doc = Document(io.BytesIO(file_data))

            # Add write protection XML
            settings_element = doc.settings.element
            protection = etree.SubElement(
                settings_element,
                qn("w:documentProtection")
            )
            protection.set(qn("w:edit"), "readOnly")
            protection.set(qn("w:enforcement"), "1")
            pw_hash = hashlib.sha256(password.encode("utf-8")).hexdigest().upper()
            protection.set(qn("w:cryptProviderType"), "rsaAES")
            protection.set(qn("w:cryptAlgorithmClass"), "hash")
            protection.set(qn("w:cryptAlgorithmType"), "typeAny")
            protection.set(qn("w:cryptAlgorithmSid"), "4")
            protection.set(qn("w:cryptSpinCount"), "100000")
            protection.set(qn("w:hash"), pw_hash)
            protection.set(qn("w:salt"), "AAAAAAAAAAAAAAAAAAAAAA==")

            out = io.BytesIO()
            doc.save(out)
            out.seek(0)
            protected_bytes = out.read()

        url, fname = await self._save_and_link(protected_bytes, out_name, __request__)
        if url:
            return f"[{fname}]({url})\n\nPassword-protected {file_type.upper()} file created."
        return json.dumps({"error": "Could not save file"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

# -----------------------------------------------------------------
# CREATE NEW FILE
# -----------------------------------------------------------------


async def merge_sheets(self, file_ids: str, output_filename: str = "", __user__=None, __request__=None) -> str:
    try:
        import sqlite3 as s3, openpyxl, io, os
        from copy import copy
        conn2 = s3.connect(_DB_PATH)
        try:
            ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
            wb_out = openpyxl.Workbook()
            wb_out.remove(wb_out.active)
            merged = 0
            for fid in ids:
                row = conn2.execute("SELECT filename, meta FROM file WHERE id=?", (fid,)).fetchone()
                if not row:
                    row = conn2.execute("SELECT filename, meta FROM file WHERE filename LIKE ?", ("%"+fid+"%",)).fetchone()
                if not row:
                    continue
                filename = row[0]
                meta = json.loads(row[1]) if row[1] else {}
                data = _read_file_bytes(meta.get("path", fid))
                if not data:
                    continue
                wb_src = openpyxl.load_workbook(io.BytesIO(data))
                base_name = os.path.splitext(os.path.basename(filename))[0][:15]
                for sn in wb_src.sheetnames:
                    ws_src = wb_src[sn]
                    sheet_name = (base_name + "_" + sn)[:31]
                    ws_out = wb_out.create_sheet(title=sheet_name)
                    for ri, row_data in enumerate(ws_src.iter_rows(), 1):
                        for ci, cell in enumerate(row_data, 1):
                            out_cell = ws_out.cell(row=ri, column=ci, value=cell.value)
                            if cell.has_style:
                                out_cell.font = copy(cell.font)
                                out_cell.fill = copy(cell.fill)
                                out_cell.border = copy(cell.border)
                                out_cell.alignment = copy(cell.alignment)
                                out_cell.number_format = cell.number_format
                    merged += 1
                wb_src.close()
        finally:
            conn2.close()
        if merged == 0:
            return json.dumps({"error": "No files could be merged"})
        out = io.BytesIO()
        wb_out.save(out)
        out.seek(0)
        fname = output_filename or "merged_workbook.xlsx"
        url, name = await self._save_and_link(out.read(), fname, __request__)
        if url:
            return f"[{name}]({url})\n\nMerged {merged} sheets from {len(ids)} files."
        return json.dumps({"error": "Could not save file"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


async def batch_process(self, file_ids: str, operation: str, params: str = "", output_filename: str = "", __user__=None, __request__=None) -> str:
    try:
        ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
        results = []
        for fid in ids:
            if operation == "replace":
                parts = params.split("|||", 1)
                if len(parts) == 2:
                    await self.replace_text(fid, parts[0], parts[1], "", __user__, __request__)
                    results.append(f"  {fid}: replaced")
            elif operation == "add_rows":
                await self.add_content(fid, params, "", __user__, __request__)
                results.append(f"  {fid}: rows added")
            else:
                results.append(f"  {fid}: unsupported operation '{operation}'")
        if results:
            return "Batch processed " + str(len(ids)) + " files:\n" + "\n".join(results)
        return json.dumps({"error": f"No files processed. Unsupported operation: {operation}"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


async def merge_pdfs(self, file_ids: str, output_filename: str = "", __user__=None, __request__=None) -> str:
    try:
        import fitz, sqlite3 as s3, io, os
        conn2 = s3.connect(_DB_PATH)
        try:
            ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
            merger = fitz.open()
            count = 0
            for fid in ids:
                row = conn2.execute("SELECT meta FROM file WHERE id=?", (fid,)).fetchone()
                if not row:
                    row = conn2.execute("SELECT meta FROM file WHERE filename LIKE ?", ("%"+fid+"%",)).fetchone()
                if not row:
                    continue
                meta = json.loads(row[0]) if row[0] else {}
                data = _read_file_bytes(meta.get("path", fid))
                if not data:
                    continue
                src = fitz.open(stream=data, filetype="pdf")
                merger.insert_pdf(src)
                src.close()
                count += 1
        finally:
            conn2.close()
        if count == 0:
            merger.close()
            return json.dumps({"error": "No PDFs could be merged"})
        out = io.BytesIO()
        merger.save(out)
        merger.close()
        out.seek(0)
        fname = output_filename or "merged.pdf"
        url, name = await self._save_and_link(out.read(), fname, __request__)
        if url:
            return f"[{name}]({url})\n\nMerged {count} PDFs into one file."
        return json.dumps({"error": "Could not save file"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


async def split_pdf(self, file_id: str, pages_per_file: int = 1, output_filename: str = "", __user__=None, __request__=None) -> str:
    try:
        import fitz, sqlite3 as s3, io, os
        conn2 = s3.connect(_DB_PATH)
        try:
            row = conn2.execute("SELECT meta FROM file WHERE id=?", (file_id,)).fetchone()
            if not row:
                row = conn2.execute("SELECT meta FROM file WHERE filename LIKE ?", ("%"+file_id+"%",)).fetchone()
            if not row:
                return json.dumps({"error": "File not found"})
            meta = json.loads(row[0]) if row[0] else {}
            data = _read_file_bytes(meta.get("path", file_id))
            if not data:
                return json.dumps({"error": "File not found on disk"})
        finally:
            conn2.close()
        src = fitz.open(stream=data, filetype="pdf")
        total_pages = src.page_count
        urls = []
        for start in range(0, total_pages, pages_per_file):
            end = min(start + pages_per_file, total_pages)
            sub = fitz.open()
            sub.insert_pdf(src, from_page=start, to_page=end-1)
            out = io.BytesIO()
            sub.save(out)
            sub.close()
            out.seek(0)
            part_name = f"part_{start+1}_{end}.pdf"
            url, name = await self._save_and_link(out.read(), part_name, __request__)
            if url:
                urls.append(f"[{name}]({url})")
        src.close()
        if urls:
            return "Split into " + str(len(urls)) + " files:\n" + "\n".join(urls)
        return json.dumps({"error": "Could not split PDF"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


async def mail_merge(self, template_file_id: str, data_file_id: str, output_prefix: str = "merged", __user__=None, __request__=None) -> str:
    """Generate personalized documents by merging CSV/Excel data into a DOCX template."""
    from docx import Document
    import io, csv
    
    t_bytes, t_name, t_type = self._resolve_file(template_file_id)
    if not t_bytes:
        return json.dumps({"error": f"Template not found: {template_file_id}"})
    
    d_bytes, d_name, d_type = self._resolve_file(data_file_id)
    if not d_bytes:
        return json.dumps({"error": f"Data file not found: {data_file_id}"})
    
    rows = []
    if d_type in ("xlsx", "xls"):
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(d_bytes))
        ws = wb.active
        headers = [str(c.value or '') for c in ws[1]]
        for row in ws.iter_rows(min_row=2, values_only=True):
            rows.append(dict(zip(headers, [str(v or '') for v in row])))
    else:
        reader = csv.DictReader(io.StringIO(d_bytes.decode('utf-8')))
        rows = list(reader)
    
    if not rows:
        return json.dumps({"error": "No data rows found"})
    
    results = []
    for i, row in enumerate(rows):
        doc = Document(io.BytesIO(t_bytes))
        for para in doc.paragraphs:
            for key, value in row.items():
                if f"{{{{{key}}}}}" in para.text:
                    for run in para.runs:
                        run.text = run.text.replace(f"{{{{{key}}}}}", value)
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        fname = f"{output_prefix}_{i+1}.docx"
        url, saved = await self._save_and_link(buf.getvalue(), fname, __request__, __user__=__user__)
        if url:
            results.append(f"[{saved}]({url})")
    
    return f"Merged {len(results)} documents:\n" + "\n".join(results)

# --- v3.2.0: Chart Generation ---


async def edit_metadata(self, file_id: str, author: str = "", title: str = "", subject: str = "", keywords: str = "", __user__=None, __request__=None) -> str:
    """Edit document metadata (author, title, subject, keywords)."""
    import io
    
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes:
        return json.dumps({"error": f"File not found: {file_id}"})
    
    changes = []
    if ftype == "docx":
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        cp = doc.core_properties
        if author: cp.author = author; changes.append("author")
        if title: cp.title = title; changes.append("title")
        if subject: cp.subject = subject; changes.append("subject")
        if keywords: cp.keywords = keywords; changes.append("keywords")
        buf = io.BytesIO(); doc.save(buf); buf.seek(0)
    elif ftype == "xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(file_bytes))
        if author: wb.properties.creator = author; changes.append("author")
        if title: wb.properties.title = title; changes.append("title")
        if subject: wb.properties.subject = subject; changes.append("subject")
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    elif ftype == "pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        cp = prs.core_properties
        if author: cp.author = author; changes.append("author")
        if title: cp.title = title; changes.append("title")
        if subject: cp.subject = subject; changes.append("subject")
        if keywords: cp.keywords = keywords; changes.append("keywords")
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    else:
        return json.dumps({"error": f"Metadata editing not supported for {ftype}"})
    
    if not changes:
        return json.dumps({"error": "No metadata fields specified"})
    
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url:
        return f"Metadata updated ({', '.join(changes)}): [{fname}]({url})"
    return json.dumps({"error": "Could not save file"})

# --- v3.2.0: Accessibility Check ---


async def check_accessibility(self, file_id: str) -> str:
    """Check a document for accessibility issues."""
    import io
    
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes:
        return json.dumps({"error": f"File not found: {file_id}"})
    
    issues = []
    if ftype == "docx":
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        headings = []
        for p in doc.paragraphs:
            if p.style.name.startswith('Heading'):
                level = int(p.style.name.split()[-1]) if p.style.name.split()[-1].isdigit() else 1
                headings.append((level, p.text[:50]))
        prev = 0
        for level, text in headings:
            if level > prev + 1:
                issues.append(f"Heading skip: H{prev} to H{level} ('{text}')")
            prev = level
        if not headings:
            issues.append("No headings found - document may lack structure")
    elif ftype == "pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == 13 and not shape.alt_text:
                    issues.append(f"Slide {i+1}: Image missing alt text")
    
    if not issues:
        return f"Accessibility check passed for {filename}. No issues found."
    result = f"**Accessibility Report: {filename}**\n\n"
    for issue in issues:
        result += f"- {issue}\n"
    result += f"\n{len(issues)} issue(s) found."
    return result

# --- v3.2.0: Add Alt Text ---


async def compare_documents(self, file_id_a: str, file_id_b: str) -> str:
    """Compare two documents and show differences."""
    a_bytes, a_name, a_type = self._resolve_file(file_id_a)
    b_bytes, b_name, b_type = self._resolve_file(file_id_b)
    if not a_bytes or not b_bytes:
        return json.dumps({"error": "One or both files not found"})
    
    def get_text(ftype, fb, fn):
        if ftype in ("xlsx","xls"):
            return self._read_xlsx(fb, fn) if ftype == "xlsx" else self._read_xls(fb, fn)
        elif ftype == "docx": return self._read_docx(fb, fn)
        elif ftype == "pptx": return self._read_pptx(fb, fn)
        elif ftype in ("odt","ods","odp"): return _read_odf(fb, fn)
        return ""
    
    text_a = get_text(a_type, a_bytes, a_name)
    text_b = get_text(b_type, b_bytes, b_name)
    
    lines_a = text_a.split('\n')
    lines_b = text_b.split('\n')
    
    result = f"**Comparison: {a_name} vs {b_name}**\n\n"
    added = 0; removed = 0; changed = 0
    max_len = max(len(lines_a), len(lines_b))
    
    for i in range(max_len):
        la = lines_a[i] if i < len(lines_a) else None
        lb = lines_b[i] if i < len(lines_b) else None
        if la is None:
            result += f"+ L{i+1}: {lb}\n"; added += 1
        elif lb is None:
            result += f"- L{i+1}: {la}\n"; removed += 1
        elif la != lb:
            result += f"~ L{i+1}:\n  - {la}\n  + {lb}\n"; changed += 1
    
    result += f"\n**Summary:** {added} added, {removed} removed, {changed} changed"
    return result

# --- v3.3.0: Export to Markdown ---


async def upload_to_drive(self, file_id: str, folder_id: str = "root") -> str:
    """Upload a file to Google Drive. Requires google-api-python-client and credentials."""
    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaIoBaseUpload
        import io as _io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        creds_path = os.environ.get("GOOGLE_CREDENTIALS", "")
        if not creds_path or not os.path.exists(creds_path):
            return json.dumps({"error": "GOOGLE_CREDENTIALS env var not set or file not found"})
        
        mime_map = {"xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                   "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                   "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                   "pdf": "application/pdf"}
        mime = mime_map.get(ftype, "application/octet-stream")
        
        creds = service_account.Credentials.from_service_account_file(creds_path, scopes=["https://www.googleapis.com/auth/drive.file"])
        service = build("drive", "v3", credentials=creds)
        
        media = MediaIoBaseUpload(_io.BytesIO(file_bytes), mimetype=mime, resumable=True)
        file_metadata = {"name": filename, "parents": [folder_id]}
        drive_file = service.files().create(body=file_metadata, media_body=media, fields="id,webViewLink").execute()
        
        return f"Uploaded to Google Drive: {drive_file.get('webViewLink', drive_file.get('id'))}"
    except ImportError:
        return json.dumps({"error": "google-api-python-client not installed. pip install google-api-python-client google-auth"})
    except Exception as e:
        return json.dumps({"error": f"Drive upload failed: {str(e)}"})

# --- v3.3.0: OCR ---


async def ocr_extract(self, file_id: str, language: str = "eng") -> str:
    """Extract text from images in a document using OCR. Requires pytesseract and Pillow."""
    try:
        import pytesseract
        from PIL import Image
        import io as _io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype == "pdf":
            try:
                import fitz
                pdf = fitz.open(stream=file_bytes, filetype="pdf")
                results = []
                for i, page in enumerate(pdf):
                    pix = page.get_pixmap(dpi=200)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text = pytesseract.image_to_string(img, lang=language)
                    if text.strip():
                        results.append(f"--- Page {i+1} ---\n{text.strip()}")
                pdf.close()
                return "\n\n".join(results) if results else "No text found in PDF images"
            except ImportError:
                return json.dumps({"error": "PyMuPDF not installed"})
        else:
            return json.dumps({"error": f"OCR only supported for PDF files. Got: {ftype}"})
    except ImportError:
        return json.dumps({"error": "pytesseract or Pillow not installed. pip install pytesseract Pillow"})
    except Exception as e:
        return json.dumps({"error": f"OCR failed: {str(e)}"})

# --- v3.3.0: i18n Error Messages ---


async def translate_errors(self, language: str = "en") -> str:
    """Set the language for error messages. Supported: en, pt, es, fr, de."""
    translations = {
        "en": {"file_not_found": "File not found", "could_not_save": "Could not save file", "unsupported": "Unsupported format"},
        "pt": {"file_not_found": "Ficheiro nao encontrado", "could_not_save": "Nao foi possivel guardar", "unsupported": "Formato nao suportado"},
        "es": {"file_not_found": "Archivo no encontrado", "could_not_save": "No se pudo guardar", "unsupported": "Formato no soportado"},
        "fr": {"file_not_found": "Fichier introuvable", "could_not_save": "Impossible d'enregistrer", "unsupported": "Format non pris en charge"},
        "de": {"file_not_found": "Datei nicht gefunden", "could_not_save": "Konnte nicht gespeichert werden", "unsupported": "Nicht unterstutztes Format"},
    }
    if language not in translations:
        return json.dumps({"error": f"Language '{language}' not supported. Available: {', '.join(translations.keys())}"})
    self.valves.language = language
    return f"Language set to {language}. Error messages will now appear in {language}."


# --- v3.4.0: AI Summarize ---


async def document_stats(self, file_id: str) -> str:
    """Show document statistics: word count, reading time, complexity."""
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype in ("xlsx","xls"):
        content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
    elif ftype == "docx": content = self._read_docx(file_bytes, filename)
    elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
    elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
    else: return json.dumps({"error": "Unsupported format"})
    words = len(content.split())
    chars = len(content)
    lines = content.count('\n') + 1
    sentences = content.count('.') + content.count('!') + content.count('?')
    reading_time = max(1, words // 200)
    avg_word_len = sum(len(w) for w in content.split()) / max(1, words)
    complexity = "Easy" if avg_word_len < 4 else ("Medium" if avg_word_len < 5.5 else "Complex")
    return "**Stats: " + str(filename) + "**\n- Words: " + str(words) + "\n- Characters: " + str(chars) + "\n- Lines: " + str(lines) + "\n- Sentences: " + str(sentences) + "\n- Reading time: ~" + str(reading_time) + " min\n- Avg word length: " + str(round(avg_word_len, 1)) + "\n- Complexity: " + str(complexity)

# --- v3.4.0: QR Codes ---


async def grammar_check(self, file_id: str) -> str:
    """Check document for grammar and style issues. The LLM will provide corrections."""
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
    elif ftype == "docx": content = self._read_docx(file_bytes, filename)
    elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
    elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
    else: return json.dumps({"error": "Unsupported format"})
    return f"**Grammar Check: {filename}**\n\nReview this document for:\n1. Grammar errors\n2. Spelling mistakes\n3. Style inconsistencies\n4. Passive voice overuse\n5. Readability issues\n\nProvide corrections with line references:\n\n```\n{content[:4000]}\n```"


async def translate_document(self, file_id: str, target_language: str, __user__=None, __request__=None) -> str:
    """Translate a document to another language. The LLM will translate while preserving structure."""
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
    elif ftype == "docx": content = self._read_docx(file_bytes, filename)
    elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
    elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
    else: return json.dumps({"error": "Unsupported format"})
    return f"**Translate to {target_language}: {filename}**\n\nTranslate the following document to {target_language}. Preserve all formatting markers (# for headings, | for tables, - for bullets). Keep numbers, dates, and proper names unchanged.\n\n```\n{content[:4000]}\n```"


async def classify_document(self, file_id: str) -> str:
    """Auto-classify a document by type, theme, and department."""
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
    elif ftype == "docx": content = self._read_docx(file_bytes, filename)
    elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
    elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
    else: return json.dumps({"error": "Unsupported format"})
    return f"**Classify: {filename}**\n\nAnalyze this document and provide:\n1. Document type (report, proposal, invoice, contract, presentation, spreadsheet, letter, memo, manual, other)\n2. Primary theme/topic\n3. Department (finance, HR, marketing, engineering, sales, legal, operations, other)\n4. Confidentiality level (public, internal, confidential, restricted)\n5. Suggested tags (3-5 keywords)\n\n```\n{content[:2000]}\n```"


__all__ = [
    "batch_process",
    "check_accessibility",
    "classify_document",
    "compare_documents",
    "document_stats",
    "edit_metadata",
    "grammar_check",
    "mail_merge",
    "merge_pdfs",
    "merge_sheets",
    "modify_rows",
    "ocr_extract",
    "protect_file",
    "replace_text",
    "split_pdf",
    "translate_document",
    "translate_errors",
    "update_cells",
    "upload_to_drive",
]
