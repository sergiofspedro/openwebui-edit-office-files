"""
title: Edit Office Files
author: giofsp
author_url: https://github.com/sergiofspedro
description: Unified tool to read, edit, and create Office files (.xlsx, .xls, .docx, .pptx) preserving original formatting and styles. Supports markdown rendering in DOCX (headings, bold, italic, code, links). Detects highlights, bold, italic formatting. Detects legacy .doc and .ppt. Note: Track changes are not supported.
version: 3.11.2
requirements: openpyxl, python-docx, python-pptx, xlrd, odfpy
"""

import io
import json
import os
import platform
import re
import sqlite3
import sys
import traceback
from copy import copy
from typing import Any, Dict, Optional

from pydantic import BaseModel, Field


from .constants import *
from .utils import *


# =========================================================================
class Tools:
    class Valves(BaseModel):
        base_url: Optional[str] = Field(
            default=None,
            description="Override the base URL for download links. Auto-detected from X-Original-Host header or WEBUI_URL env var if unset.",
        )
        file_url_pattern: Optional[str] = Field(
            default=None,
            description="Custom file download URL pattern. Use {file_id} placeholder. Examples: /api/v1/files/{file_id}/content (default), /api/files/{file_id}. If unset, uses standard Open WebUI URL.",
        )
        templates: Optional[str] = Field(default="{}", description="JSON map of template names to content strings.")
        cleanup_schedule: Optional[str] = Field(default="{}", description="JSON schedule for auto-cleanup.")
        language: Optional[str] = Field(default="en", description="Language for error messages: en, pt, es, fr, de.")
        pass

    def __init__(self):
        self.valves = self.Valves()

    def _resolve_file(self, file_id: str):
        """Resolve a file ID to (bytes, filename, file_type)."""
        path = _resolve_file_path(file_id)
        if not path:
            return None, None, None
        filename = os.path.basename(path)
        ftype = _detect_type(filename)
        file_bytes = _read_file_bytes(file_id)
        return file_bytes, filename, ftype





    async def _save_and_link(self, file_bytes: bytes, filename: str, __request__=None, __user__=None) -> tuple:
        """Save file to Open WebUI uploads dir, register in DB, return download URL."""
        import base64 as _b64
        import hashlib
        import time as _time
        import uuid as _uuid

        ext = os.path.splitext(filename)[1].lower()
        mt = {
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".xls": "application/vnd.ms-excel",
            ".doc": "application/msword",
            ".ppt": "application/vnd.ms-powerpoint",
            ".pdf": "application/pdf",
            ".odt": "application/vnd.oasis.opendocument.text",
            ".ods": "application/vnd.oasis.opendocument.spreadsheet",
            ".odp": "application/vnd.oasis.opendocument.presentation",
        }
        content_type = mt.get(ext, "application/octet-stream")

        try:
            file_id = str(_uuid.uuid4())
            with open(os.path.join(_UPLOAD_DIR, file_id), "wb") as f:
                f.write(file_bytes)

            file_hash = hashlib.sha256(file_bytes).hexdigest()[:16]
            now = int(_time.time())

            # Generate a preview for the data column
            try:
                if content_type.startswith("text/") or ext in (".csv", ".md", ".txt"):
                    preview = file_bytes.decode('utf-8', errors='replace')[:500]
                elif ext in (".xlsx", ".xls"):
                    preview = f"[Excel spreadsheet: {len(file_bytes)} bytes]"
                elif ext in (".docx",):
                    preview = f"[Word document: {len(file_bytes)} bytes]"
                elif ext in (".pptx",):
                    preview = f"[PowerPoint presentation: {len(file_bytes)} bytes]"
                elif ext in (".pdf",):
                    preview = f"[PDF document: {len(file_bytes)} bytes]"
                else:
                    preview = f"[File: {len(file_bytes)} bytes]"
            except Exception:
                preview = "{}"

            conn = sqlite3.connect(_DB_PATH)
            try:
                conn.execute(
                    """INSERT OR REPLACE INTO file
                       (id, user_id, hash, filename, path, data, meta, created_at, updated_at)
                       VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                    (
                        file_id,
                        __user__.get("id", "") if __user__ and isinstance(__user__, dict) else "",
                        file_hash,
                        _encode_filename(filename),
                        os.path.join(_UPLOAD_DIR, file_id),
                        preview,
                        json.dumps({
                            "name": filename,
                            "content_type": content_type,
                            "size": len(file_bytes),
                            "source": "office-plugin",
                            "generated": True,
                        }),
                        now,
                        now,
                    ),
                )
                conn.commit()
            finally:
                conn.close()

            base_url = self.valves.base_url
            if not base_url and __request__:
                try:
                    host = __request__.headers.get("x-original-host")
                    if host:
                        base_url = f"https://{host}"
                except Exception:
                    pass
            if not base_url:
                base_url = os.environ.get("WEBUI_URL", "http://localhost:3000")
            base_url = base_url.rstrip("/")

            # Standard Open WebUI file download URL (works for most versions)
            # Customize via file_url_pattern valve if needed
            if self.valves.file_url_pattern:
                url = f"{base_url}{self.valves.file_url_pattern.replace('{file_id}', file_id)}"
            else:
                url = f"{base_url}/api/v1/files/{file_id}/content"
            return (url, filename)

        except Exception as e:
            print(f"[office] Save failed: {e}", file=sys.stderr)
            try:
                if len(file_bytes) > 1_000_000:  # 1MB limit
                    return (None, None)
                data = _b64.b64encode(file_bytes).decode("ascii")
                return (f"data:{content_type};base64,{data}", filename)
            except Exception:
                return (None, None)

    # -----------------------------------------------------------------
    # READ
    # -----------------------------------------------------------------











    async def auto_backup(self, __user__=None, __request__=None) -> str:
        try:
            import shutil, datetime
            db_path = _DB_PATH
            backup_dir = os.path.join(_get_owui_data_dir(), "backups")
            os.makedirs(backup_dir, exist_ok=True)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_name = f"webui_backup_{timestamp}.db"
            backup_path = os.path.join(backup_dir, backup_name)
            shutil.copy2(db_path, backup_path)
            size_kb = os.path.getsize(backup_path) / 1024
            return json.dumps({"success": True, "backup_path": backup_path, "size_kb": round(size_kb,1), "message": f"Backup: {backup_name} ({size_kb:.1f} KB)"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})





    async def tool_stats(self, __user__=None, __request__=None) -> str:
        try:
            import sqlite3 as s3
            conn2 = s3.connect(_DB_PATH)
            try:
                try:
                    tool_count = conn2.execute("SELECT COUNT(*) FROM tool WHERE is_active=1").fetchone()[0]
                except Exception:
                    tool_count = 0
                try:
                    func_count = conn2.execute("SELECT COUNT(*) FROM function WHERE is_active=1").fetchone()[0]
                except Exception:
                    func_count = 0
                try:
                    model_count = conn2.execute("SELECT COUNT(*) FROM model WHERE is_active=1").fetchone()[0]
                except Exception:
                    model_count = 0
                exports_dir = os.path.join(_get_owui_data_dir(), "exports")
                export_count = len([f for f in os.listdir(exports_dir) if os.path.isfile(os.path.join(exports_dir, f))]) if os.path.exists(exports_dir) else 0
                db_size_kb = os.path.getsize(_DB_PATH) / 1024
            finally:
                conn2.close()
            return json.dumps({
                "tools": tool_count,
                "functions": func_count,
                "models": model_count,
                "exported_files": export_count,
                "db_size_kb": round(db_size_kb, 1)
            }, indent=2)
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    

    async def cleanup_files(self, days_old: int = 30) -> str:
        """Remove generated Office files older than N days. Default 30 days."""
        import time as _time
        
        cutoff = int(_time.time()) - (days_old * 86400)
        
        conn = sqlite3.connect(_DB_PATH)
        try:
            conn.row_factory = sqlite3.Row
            rows = conn.execute(
                "SELECT id, filename, created_at FROM file WHERE meta LIKE '%office-plugin%' AND created_at < ?",
                (cutoff,)
            ).fetchall()
            
            if not rows:
                return f"No office-generated files older than {days_old} days found."
            
            deleted = []
            errors = []
            for row in rows:
                try:
                    fpath = os.path.join(_UPLOAD_DIR, row["id"])
                    if os.path.exists(fpath):
                        os.remove(fpath)
                    conn.execute("DELETE FROM file WHERE id = ?", (row["id"],))
                    deleted.append(row["filename"])
                except Exception as e:
                    errors.append(f"{row['filename']}: {e}")
            
            conn.commit()
        finally:
            conn.close()
        
        result = f"Cleaned up {len(deleted)} file(s) older than {days_old} days:\n"
        for f in deleted:
            result += f"- {f}\n"
        if errors:
            result += f"\nErrors ({len(errors)}):\n"
            for e in errors:
                result += f"- {e}\n"
        return result

    async def convert_format(self, file_id: str, target_format: str, __user__=None, __request__=None) -> str:
        """Convert between Office formats (docx<->odt, xlsx<->ods, pptx<->odp)."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        base_name = os.path.splitext(filename)[0]
        
        if ftype in ("xlsx", "xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx":
            content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx":
            content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt", "ods", "odp"):
            content = _read_odf(file_bytes, filename)
        else:
            return json.dumps({"error": f"Unsupported source format: {ftype}"})
        
        if target_format in ("odt", "ods", "odp"):
            return await self.create_odf(content, base_name, target_format, __user__=__user__, __request__=__request__)
        elif target_format == "docx":
            return await self.generate_document(content, base_name, __user__=__user__, __request__=__request__)
        elif target_format == "pptx":
            return await self.generate_slides(content, base_name, __user__=__user__, __request__=__request__)
        elif target_format == "xlsx":
            return await self.generate_spreadsheet(content, base_name, __user__=__user__, __request__=__request__)
        return json.dumps({"error": f"Unsupported target format: {target_format}"})

    # --- v3.2.0: Template System ---
    async def save_template(self, name: str, content: str) -> str:
        """Save a document template for reuse."""
        templates = json.loads(self.valves.templates or "{}")
        templates[name] = content
        self.valves.templates = json.dumps(templates)
        return f"Template '{name}' saved."

    async def use_template(self, name: str, __user__=None, __request__=None, **kwargs) -> str:
        """Generate a document from a saved template, replacing {placeholders}."""
        templates = json.loads(self.valves.templates or "{}")
        if name not in templates:
            return f"Template '{name}' not found. Available: {', '.join(templates.keys())}"
        content = templates[name]
        for key, value in kwargs.items():
            content = content.replace(f"{{{key}}}", str(value))
        return await self.generate_document(content, name, __user__=__user__, __request__=__request__)

    async def list_templates(self) -> str:
        """List all saved templates."""
        templates = json.loads(self.valves.templates or "{}")
        if not templates:
            return "No templates saved."
        result = "Available templates:\n"
        for name in templates:
            preview = templates[name][:50].replace('\n', ' ')
            result += f"- {name}: {preview}...\n"
        return result

    # --- v3.2.0: Scheduled Cleanup ---
    async def schedule_cleanup(self, days_old: int = 30, interval_hours: int = 24) -> str:
        """Schedule automatic cleanup every N hours. Set interval_hours=0 to disable."""
        schedule = {"days_old": days_old, "interval_hours": interval_hours, "enabled": interval_hours > 0}
        self.valves.cleanup_schedule = json.dumps(schedule)
        if interval_hours > 0:
            return f"Cleanup scheduled: remove files older than {days_old} days, every {interval_hours} hours."
        return "Scheduled cleanup disabled."

    # --- v3.2.0: Mail Merge ---
    async def preview_file(self, file_id: str, max_lines: int = 20) -> str:
        """Show a text preview of any Office file before downloading."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype in ("xlsx", "xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx":
            content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx":
            content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt", "ods", "odp"):
            content = _read_odf(file_bytes, filename)
        else:
            return json.dumps({"error": f"Preview not supported for {ftype}"})
        
        lines = content.split('\n')
        preview = '\n'.join(lines[:max_lines])
        total = len(lines)
        result = f"**{filename}** ({ftype.upper()}, {total} lines)\n\n```\n{preview}\n```"
        if total > max_lines:
            result += f"\n... ({total - max_lines} more lines)"
        return result

    # --- v3.2.0: Metadata Editing ---
    async def _progress(self, current: int, total: int, operation: str = "Processing", __event_emitter__=None) -> None:
        """Emit progress via __event_emitter__ if available."""
        if __event_emitter__ is None:
            return
        try:
            await __event_emitter__({"type": "status", "data": {"description": f"{operation}: {current}/{total}", "done": current >= total}})
        except Exception:
            pass


    # --- v3.3.0: Document Comparison ---
    async def export_to_markdown(self, file_id: str, __user__=None, __request__=None) -> str:
        """Export any Office file to Markdown format."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype in ("xlsx","xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx":
            content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx":
            content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"):
            content = _read_odf(file_bytes, filename)
        else:
            return json.dumps({"error": f"Export not supported for {ftype}"})
        
        base = os.path.splitext(filename)[0]
        md_content = f"# {base}\n\n{content}"
        md_bytes = md_content.encode('utf-8')
        
        url, fname = await self._save_and_link(md_bytes, f"{base}.md", __request__, __user__=__user__)
        if url:
            return f"Exported to Markdown: [{fname}]({url})"
        return json.dumps({"error": "Could not save file"})

    # --- v3.3.0: Import from URL ---
    async def import_from_url(self, url: str, title: str = "Web Document", __user__=None, __request__=None) -> str:
        """Fetch a web page and convert it to a Word document."""
        try:
            import urllib.request as _urllib
            req = _urllib.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            resp = _urllib.urlopen(req, timeout=15)
            html = resp.read().decode('utf-8', errors='replace')
        except Exception as e:
            return json.dumps({"error": f"Failed to fetch URL: {str(e)}"})
        
        # Simple HTML to text
        import re as _re
        text = _re.sub(r'<script[^>]*>.*?</script>', '', html, flags=_re.DOTALL|_re.IGNORECASE)
        text = _re.sub(r'<style[^>]*>.*?</style>', '', text, flags=_re.DOTALL|_re.IGNORECASE)
        text = _re.sub(r'<[^>]+>', '\n', text)
        text = _re.sub(r'\n\s*\n', '\n\n', text)
        text = _re.sub(r'[ \t]+', ' ', text)
        text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
        
        if not text:
            return json.dumps({"error": "No text content extracted from URL"})
        
        return await self.generate_document(text[:50000], title, __user__=__user__, __request__=__request__)

    # --- v3.3.0: File Versioning ---
    async def version_file(self, file_id: str, label: str = "", __user__=None, __request__=None) -> str:
        """Save a versioned copy of a file before editing."""
        import time as _time
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        base, ext = os.path.splitext(filename)
        ts = _time.strftime("%Y%m%d_%H%M%S")
        label_str = f"_{label}" if label else ""
        version_name = f"{base}_v{ts}{label_str}{ext}"
        
        url, fname = await self._save_and_link(file_bytes, version_name, __request__, __user__=__user__)
        if url:
            return f"Version saved: [{fname}]({url})"
        return json.dumps({"error": "Could not save version"})

    # --- v3.3.0: Cloud Storage (Google Drive) ---
    async def ai_summarize(self, file_id: str) -> str:
        """Extract document text for LLM summarization. The LLM will then summarize it."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": "File not found: " + str(file_id)})
        if ftype in ("xlsx","xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": "Unsupported format: " + str(ftype)})
        words = len(content.split())
        preview = content[:3000]
        return "**" + str(filename) + "** (" + str(words) + " words)\n\n" + preview + ("\n\n... (truncated, " + str(words) + " total words)" if len(content) > 3000 else "")

    # --- v3.4.0: Speaker Notes ---
    async def bulk_folder_ops(self, operation: str = "list", pattern: str = "*", __user__=None, __request__=None) -> str:
        """Apply an operation to all files in the uploads folder. Operations: list, delete_old, stats."""
        import glob as _glob, time as _time
        uploads = _UPLOAD_DIR
        if operation == "list":
            files = sorted(_glob.glob(os.path.join(uploads, pattern)))
            if not files: return "No files found matching: " + str(pattern)
            result = "**Files in uploads (" + str(len(files)) + "):**\n"
            for f in files[:50]:
                size = os.path.getsize(f)
                mtime = _time.strftime("%Y-%m-%d %H:%M", _time.localtime(os.path.getmtime(f)))
                result += "- " + os.path.basename(f) + " (" + str(round(size/1024,1)) + " KB, " + str(mtime) + ")\n"
            if len(files) > 50: result += "... and " + str(len(files)-50) + " more"
            return result
        elif operation == "delete_old":
            cutoff = _time.time() - (30 * 86400)
            deleted = 0
            for f in _glob.glob(os.path.join(uploads, pattern)):
                if os.path.getmtime(f) < cutoff:
                    os.remove(f); deleted += 1
            return "Deleted " + str(deleted) + " files older than 30 days."
        elif operation == "stats":
            files = _glob.glob(os.path.join(uploads, pattern))
            total_size = sum(os.path.getsize(f) for f in files)
            return "**Uploads Stats:**\n- Files: " + str(len(files)) + "\n- Total size: " + str(round(total_size/1024/1024,1)) + " MB"
        return json.dumps({"error": "Unknown operation: " + str(operation) + ". Use: list, delete_old, stats"})

    # --- v3.4.0: File Search ---
    async def file_search(self, query: str, file_type: str = "all") -> str:
        """Search for text across all generated files. file_type: all, xlsx, docx, pptx, pdf."""
        import glob as _glob
        results = []
        patterns = {"all": "*.*", "xlsx": "*.xlsx", "docx": "*.docx", "pptx": "*.pptx", "pdf": "*.pdf"}
        pattern = patterns.get(file_type, "*.*")
        for fpath in _glob.glob(os.path.join(_UPLOAD_DIR, pattern)):
            try:
                fname = os.path.basename(fpath)
                ext = os.path.splitext(fname)[1].lower()
                with open(fpath, 'rb') as f: fb = f.read()
                if ext == ".xlsx":
                    import openpyxl, io
                    wb = openpyxl.load_workbook(io.BytesIO(fb))
                    for ws in wb.worksheets:
                        for row in ws.iter_rows(values_only=True):
                            for cell in row:
                                if cell and query.lower() in str(cell).lower():
                                    results.append(fname + ": " + str(cell)[:100])
                                    break
                elif ext == ".docx":
                    from docx import Document
                    import io
                    doc = Document(io.BytesIO(fb))
                    for p in doc.paragraphs:
                        if query.lower() in p.text.lower():
                            results.append(fname + ": " + p.text[:100])
                            break
                elif ext == ".pptx":
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(fb))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame and query.lower() in shape.text_frame.text.lower():
                                results.append(fname + ": " + shape.text_frame.text[:100])
                                break
            except Exception:
                pass
        if not results: return "No matches found for '" + str(query) + "'"
        return "**Search: " + str(query) + "** (" + str(len(results)) + " matches)\n" + "\n".join("- " + r for r in results[:20])

    # --- v3.4.0: Data Validation ---
    async def export_to_html(self, file_id: str, __user__=None, __request__=None) -> str:
        """Export any Office file to a styled HTML page."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": "File not found"})
        if ftype in ("xlsx","xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": "Export not supported for " + str(ftype)})
        base = os.path.splitext(filename)[0]
        html_c = content.replace('\n', '<br>\n')
        html_c = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_c, flags=re.MULTILINE)
        html_c = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_c, flags=re.MULTILINE)
        html_c = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html_c, flags=re.MULTILINE)
        html_c = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html_c)
        html_c = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html_c)
        html_c = re.sub(r'`(.+?)`', r'<code>\1</code>', html_c)
        html_c = re.sub(r'^- (.+)$', r'<li>\1</li>', html_c, flags=re.MULTILINE)
        html = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>""" + str(base) + """</title>
<style>
body { font-family: 'Segoe UI', system-ui, sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; line-height: 1.6; color: #333; background: #fafafa; }
h1 { color: #1a1a2e; border-bottom: 3px solid #e94560; padding-bottom: 10px; }
h2 { color: #16213e; margin-top: 30px; }
h3 { color: #0f3460; }
table { border-collapse: collapse; width: 100%; margin: 20px 0; }
th { background: #1a1a2e; color: white; padding: 10px; text-align: left; }
td { padding: 8px 10px; border-bottom: 1px solid #ddd; }
tr:nth-child(even) { background: #f5f5f5; }
code { background: #f0f0f0; padding: 2px 6px; border-radius: 3px; font-size: 0.9em; }
blockquote { border-left: 4px solid #e94560; margin: 20px 0; padding: 10px 20px; background: #fff5f5; }
</style>
</head>
<body>
""" + html_c + """
</body>
</html>"""
        html_bytes = html.encode('utf-8')
        url, fname = await self._save_and_link(html_bytes, base + ".html", __request__, __user__=__user__)
        if url: return "Exported to HTML: [" + str(fname) + "](" + str(url) + ")"
        return json.dumps({"error": "Could not save file"})


    # === v3.6.0: AI-Powered Features ===

    async def ai_analyze(self, file_id: str) -> str:
        """Extract document text for AI analysis. The LLM will analyze topics, sentiment, entities, and provide a summary."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": "File not found"})
        if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": "Unsupported format"})
        words = len(content.split())
        preview = content[:5000]
        return f"**Document: {filename}** ({words} words)\n\nAnalyze this document and provide:\n1. Main topics (3-5 bullet points)\n2. Sentiment (positive/negative/neutral)\n3. Key entities (people, companies, dates)\n4. Executive summary (2-3 sentences)\n\n```\n{preview}\n```" + ("\n\n... (truncated)" if len(content) > 5000 else "")





    async def smart_template(self, name: str, description: str, __user__=None, __request__=None) -> str:
        """Generate a document from a smart template that adapts to the conversation context."""
        templates = json.loads(self.valves.templates or "{}")
        if name in templates:
            content = templates[name]
            return await self.generate_document(content, name, __user__=__user__, __request__=__request__)
        return f"**Smart Template: {name}**\n\nDescription: {description}\n\nGenerate a professional document template for '{name}' with the following sections and {placeholders} for customization. Use markdown format with # headings, - bullets, and | tables."

    # === v3.6.0: Data Manipulation ===




    async def convert_data(self, file_id: str, target_format: str, __user__=None, __request__=None) -> str:
        """Convert between CSV, JSON, and XML formats. target_format: csv, json, xml."""
        import io, csv as _csv
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": "File not found"})
        base = os.path.splitext(filename)[0]
        text = file_bytes.decode('utf-8', errors='replace')
        if target_format == "json":
            try:
                if ftype == "csv" or filename.endswith('.csv'):
                    reader = _csv.DictReader(io.StringIO(text))
                    data = list(reader)
                elif filename.endswith('.xml'):
                    import xml.etree.ElementTree as ET
                    root = ET.fromstring(text)
                    data = [{child.tag: child.text for child in elem} for elem in root]
                else:
                    data = [{"content": text}]
                result = json.dumps(data, indent=2, ensure_ascii=False)
                ext = ".json"
            except Exception as e:
                return json.dumps({"error": f"Conversion failed: {str(e)}"})
        elif target_format == "csv":
            try:
                if ftype == "json" or filename.endswith('.json'):
                    data = json.loads(text)
                    if isinstance(data, list) and data:
                        out = io.StringIO()
                        writer = _csv.DictWriter(out, fieldnames=data[0].keys())
                        writer.writeheader(); writer.writerows(data)
                        result = out.getvalue()
                    else:
                        result = text
                else:
                    result = text
                ext = ".csv"
            except Exception as e:
                return json.dumps({"error": f"Conversion failed: {str(e)}"})
        elif target_format == "xml":
            try:
                if ftype == "json" or filename.endswith('.json'):
                    data = json.loads(text)
                    import xml.etree.ElementTree as ET
                    root = ET.Element("root")
                    for item in (data if isinstance(data, list) else [data]):
                        elem = ET.SubElement(root, "item")
                        for k, v in (item.items() if isinstance(item, dict) else {"value": str(item)}.items()):
                            child = ET.SubElement(elem, k)
                            child.text = str(v)
                    result = ET.tostring(root, encoding='unicode')
                else:
                    result = text
                ext = ".xml"
            except Exception as e:
                return json.dumps({"error": f"Conversion failed: {str(e)}"})
        else:
            return json.dumps({"error": f"Unsupported target format: {target_format}. Use: csv, json, xml"})
        result_bytes = result.encode('utf-8')
        url, fname = await self._save_and_link(result_bytes, f"{base}{ext}", __request__)
        if url: return f"Converted to {target_format.upper()}: [{fname}]({url})"
        return json.dumps({"error": "Could not save file"})

    # === v3.6.0: Enterprise Features ===

    async def compliance_check(self, file_id: str, standard: str = "gdpr") -> str:
        """Check document for compliance issues. standards: gdpr, accessibility, branding, all."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": "File not found"})
        issues = []
        if standard in ("gdpr", "all"):
            text = ""
            if ftype == "docx":
                from docx import Document; import io
                doc = Document(io.BytesIO(file_bytes))
                text = " ".join(p.text for p in doc.paragraphs)
            gdpr_keywords = ["email", "phone", "address", "name", "birth", "passport", "ssn", "tax id", "iban", "credit card", "ip address", "cookie"]
            found = [k for k in gdpr_keywords if k in text.lower()]
            if found: issues.append(f"GDPR: Personal data detected: {', '.join(found)}. Ensure consent and data processing agreement.")
        if standard in ("accessibility", "all"):
            if ftype == "docx":
                from docx import Document; import io
                doc = Document(io.BytesIO(file_bytes))
                headings = [p for p in doc.paragraphs if p.style.name.startswith('Heading')]
                if not headings: issues.append("Accessibility: No headings found. Add heading structure.")
                images = len([r for r in doc.part.rels.values() if "image" in r.reltype])
                if images > 0: issues.append(f"Accessibility: {images} image(s) found. Ensure alt text is provided.")
        if standard in ("branding", "all"):
            if ftype == "docx":
                from docx import Document; import io
                doc = Document(io.BytesIO(file_bytes))
                text = " ".join(p.text for p in doc.paragraphs)
                if "confidential" not in text.lower() and "draft" not in text.lower():
                    issues.append("Branding: No confidentiality marking found. Consider adding DRAFT or CONFIDENTIAL watermark.")
        if not issues: return f"Compliance check passed for {filename}. No issues found."
        return f"**Compliance Report: {filename}**\n\n" + "\n".join(f"- {i}" for i in issues)

    async def audit_log(self, action: str = "list", limit: int = 20) -> str:
        """View or manage the audit trail of document operations."""
        import time as _time
        conn2 = sqlite3.connect(_DB_PATH)
        try:
            conn2.row_factory = sqlite3.Row
            if action == "list":
                rows = conn2.execute("SELECT filename, created_at FROM file WHERE meta LIKE '%office-plugin%' ORDER BY created_at DESC LIMIT ?", (limit,)).fetchall()
                if not rows: return "No audit records found."
                result = f"**Audit Trail (last {len(rows)}):**\n"
                for r in rows:
                    ts = _time.strftime("%Y-%m-%d %H:%M", _time.localtime(r["created_at"])) if r["created_at"] else "unknown"
                    result += f"- {ts}: {r['filename']}\n"
                return result
            elif action == "stats":
                total = conn2.execute("SELECT COUNT(*) FROM file WHERE meta LIKE '%office-plugin%'").fetchone()[0]
                today = conn2.execute("SELECT COUNT(*) FROM file WHERE meta LIKE '%office-plugin%' AND date(created_at, 'unixepoch') = date('now')").fetchone()[0]
                return f"**Audit Stats:**\n- Total files: {total}\n- Today: {today}"
            return json.dumps({"error": f"Unknown action: {action}. Use: list, stats"})
        finally:
            conn2.close()

    async def retention_policy(self, policy: str = "view", days: int = 90, file_type: str = "all") -> str:
        """Manage document retention policies. policy: view, set, apply."""
        if policy == "view":
            sched = json.loads(self.valves.cleanup_schedule or "{}")
            if sched.get("enabled"):
                return f"Retention policy active: delete files older than {sched.get('days_old', 30)} days, every {sched.get('interval_hours', 24)} hours."
            return "No retention policy active. Use: retention_policy(policy='set', days=90)"
        elif policy == "set":
            self.valves.cleanup_schedule = json.dumps({"days_old": days, "interval_hours": 24, "enabled": True, "file_type": file_type})
            return f"Retention policy set: delete {file_type} files older than {days} days."
        elif policy == "apply":
            return await self.cleanup_files(days_old=days)
        return json.dumps({"error": f"Unknown policy: {policy}. Use: view, set, apply"})

    async def scheduled_report(self, action: str = "list", name: str = "", schedule: str = "", template: str = "", __user__=None, __request__=None) -> str:
        """Manage scheduled reports. action: list, create, delete, run."""
        reports = json.loads(self.valves.templates or "{}")
        scheduled = json.loads(self.valves.cleanup_schedule or "{}")
        if action == "list":
            sched_reports = {k: v for k, v in reports.items() if k.startswith("_scheduled_")}
            if not sched_reports: return "No scheduled reports."
            result = "**Scheduled Reports:**\n"
            for k, v in sched_reports.items():
                result += f"- {k.replace('_scheduled_', '')}: {v[:50]}...\n"
            return result
        elif action == "create":
            reports[f"_scheduled_{name}"] = json.dumps({"template": template, "schedule": schedule})
            self.valves.templates = json.dumps(reports)
            return f"Scheduled report '{name}' created."
        elif action == "delete":
            key = f"_scheduled_{name}"
            if key in reports:
                del reports[key]
                self.valves.templates = json.dumps(reports)
                return f"Scheduled report '{name}' deleted."
            return f"Report '{name}' not found."
        elif action == "run":
            key = f"_scheduled_{name}"
            if key not in reports: return f"Report '{name}' not found."
            cfg = json.loads(reports[key])
            return await self.generate_document(cfg.get("template", ""), name, __user__=__user__, __request__=__request__)
        return json.dumps({"error": f"Unknown action: {action}"})




    async def version_diff(self, file_id: str, version_label: str = "") -> str:
        """Show differences between current file and a previous version."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": "File not found"})
        base = os.path.splitext(filename)[0]
        ext = os.path.splitext(filename)[1]
        import glob as _glob, time as _time
        pattern = os.path.join(_UPLOAD_DIR, f"{base}_v*{ext}")
        versions = sorted(_glob.glob(pattern), reverse=True)
        if not versions: return f"No previous versions found for {filename}. Use version_file() to create versions."
        if version_label:
            versions = [v for v in versions if version_label in v]
            if not versions: return f"No version matching '{version_label}' found."
        latest = versions[0]
        vname = os.path.basename(latest)
        vtime = _time.strftime("%Y-%m-%d %H:%M", _time.localtime(os.path.getmtime(latest)))
        with open(latest, 'rb') as f: vbytes = f.read()
        if ftype in ("xlsx","xls"):
            curr = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
            prev = self._read_xlsx(vbytes, vname) if ftype == "xlsx" else self._read_xls(vbytes, vname)
        elif ftype == "docx":
            curr = self._read_docx(file_bytes, filename)
            prev = self._read_docx(vbytes, vname)
        elif ftype == "pptx":
            curr = self._read_pptx(file_bytes, filename)
            prev = self._read_pptx(vbytes, vname)
        else:
            return json.dumps({"error": f"Diff not supported for {ftype}"})
        cl = curr.split('\n'); pl = prev.split('\n')
        added = sum(1 for l in cl if l not in pl)
        removed = sum(1 for l in pl if l not in cl)
        return f"**Version Diff: {filename} vs {vname}** ({vtime})\n- Lines added: {added}\n- Lines removed: {removed}\n- Current: {len(cl)} lines\n- Previous: {len(pl)} lines"

    async def webhook_trigger(self, event: str = "test", url: str = "", file_id: str = "") -> str:
        """Trigger a webhook on document events. event: test, created, edited, deleted."""
        if not url:
            return json.dumps({"error": "url parameter required"})
        import urllib.request as _urllib
        payload = json.dumps({"event": event, "file_id": file_id, "timestamp": __import__('time').time()}).encode()
        try:
            req = _urllib.Request(url, data=payload, headers={"Content-Type": "application/json"})
            resp = _urllib.urlopen(req, timeout=10)
            return f"Webhook sent to {url}: HTTP {resp.getcode()}"
        except Exception as e:
            return json.dumps({"error": f"Webhook failed: {str(e)}"})

    async def import_from_api(self, url: str, data_path: str = "", output_filename: str = "api_data", __user__=None, __request__=None) -> str:
        """Import data from a REST API and export to Excel."""
        import urllib.request as _urllib, io
        try:
            req = _urllib.Request(url, headers={"Accept": "application/json", "User-Agent": "OpenWebUI/1.0"})
            resp = _urllib.urlopen(req, timeout=15)
            raw = json.loads(resp.read())
        except Exception as e:
            return json.dumps({"error": f"API request failed: {str(e)}"})
        data = raw
        if data_path:
            for key in data_path.split("."):
                if isinstance(data, dict):
                    data = data.get(key, [])
                elif isinstance(data, list) and key.isdigit():
                    data = data[int(key)]
        if not isinstance(data, list):
            data = [data] if isinstance(data, dict) else [{"value": str(data)}]
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        wb = Workbook(); ws = wb.active
        if data and isinstance(data[0], dict):
            headers = list(data[0].keys())
            for j, h in enumerate(headers, 1):
                cell = ws.cell(row=1, column=j, value=h)
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
            for i, item in enumerate(data, 2):
                for j, key in enumerate(headers, 1):
                    ws.cell(row=i, column=j, value=str(item.get(key, "")))
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        url_out, fname = await self._save_and_link(buf.getvalue(), f"{output_filename}.xlsx", __request__)
        if url_out: return f"Imported {len(data)} records from API: [{fname}]({url_out})"
        return json.dumps({"error": "Could not save file"})

# --- Method bindings from modular split ---
from .office_read import read_file, _read_xlsx, _read_xls, _read_docx, _read_pptx, _parse_csv_rows
Tools.read_file = read_file
Tools._read_xlsx = _read_xlsx
Tools._read_xls = _read_xls
Tools._read_docx = _read_docx
Tools._read_pptx = _read_pptx
Tools._parse_csv_rows = _parse_csv_rows

from .office_write import create_file, generate_document, generate_slides, generate_spreadsheet, add_content, add_chart, add_watermark, add_alt_text, add_speaker_notes, add_qr_code, add_data_validation, add_named_range, add_slide_transitions, add_pivot_table, sql_to_spreadsheet, fill_pdf_form, document_assembly, conditional_format, smart_fill
Tools.create_file = create_file
Tools.generate_document = generate_document
Tools.generate_slides = generate_slides
Tools.generate_spreadsheet = generate_spreadsheet
Tools.add_content = add_content
Tools.add_chart = add_chart
Tools.add_watermark = add_watermark
Tools.add_alt_text = add_alt_text
Tools.add_speaker_notes = add_speaker_notes
Tools.add_qr_code = add_qr_code
Tools.add_data_validation = add_data_validation
Tools.add_named_range = add_named_range
Tools.add_slide_transitions = add_slide_transitions
Tools.add_pivot_table = add_pivot_table
Tools.sql_to_spreadsheet = sql_to_spreadsheet
Tools.fill_pdf_form = fill_pdf_form
Tools.document_assembly = document_assembly
Tools.conditional_format = conditional_format
Tools.smart_fill = smart_fill

from .office_edit import replace_text, update_cells, modify_rows, protect_file, merge_sheets, batch_process, merge_pdfs, split_pdf, mail_merge, edit_metadata, check_accessibility, compare_documents, grammar_check, translate_document, classify_document, upload_to_drive, ocr_extract, translate_errors, document_stats
Tools.replace_text = replace_text
Tools.update_cells = update_cells
Tools.modify_rows = modify_rows
Tools.protect_file = protect_file
Tools.merge_sheets = merge_sheets
Tools.batch_process = batch_process
Tools.merge_pdfs = merge_pdfs
Tools.split_pdf = split_pdf
Tools.mail_merge = mail_merge
Tools.edit_metadata = edit_metadata
Tools.check_accessibility = check_accessibility
Tools.compare_documents = compare_documents
Tools.grammar_check = grammar_check
Tools.translate_document = translate_document
Tools.classify_document = classify_document
Tools.upload_to_drive = upload_to_drive
Tools.ocr_extract = ocr_extract
Tools.translate_errors = translate_errors
Tools.document_stats = document_stats

from .office_comments import add_comment, tracked_change, manage_revisions, create_odf
Tools.add_comment = add_comment
Tools.tracked_change = tracked_change
Tools.manage_revisions = manage_revisions
Tools.create_odf = create_odf
