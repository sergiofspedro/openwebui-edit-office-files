"""Edit Office Files — File reading (read_file, xlsx/xls/docx/pptx/csv readers)."""

import io
import json
import os
import platform
import re
import sqlite3
import sys
import traceback
from copy import copy
from typing import Any, Dict, Optional

from .constants import *
from .utils import *

def _read_xlsx(self, file_bytes: bytes, filename: str) -> str:
    """Extract text from an XLSX file."""
    from openpyxl import load_workbook
    import io
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    result = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            result.append("\t".join(cells))
    wb.close()
    return "\n\n".join(result)


def _read_xls(self, file_bytes: bytes, filename: str) -> str:
    """Extract text from an XLS file."""
    import xlrd
    import io
    wb = xlrd.open_workbook(file_contents=file_bytes)
    result = []
    for si in range(wb.nsheets):
        ws = wb.sheet_by_index(si)
        for r in range(ws.nrows):
            cells = [str(ws.cell_value(r, c)) for c in range(ws.ncols)]
            result.append("\t".join(cells))
    return "\n\n".join(result)


def _read_docx(self, file_bytes: bytes, filename: str) -> str:
    """Extract text from a DOCX file."""
    from docx import Document
    import io
    doc = Document(io.BytesIO(file_bytes))
    result = []
    for p in doc.paragraphs:
        if p.style.name.startswith("Heading"):
            level = p.style.name.replace("Heading ", "")
            result.append(f"{'#' * int(level)} {p.text}")
        else:
            result.append(p.text)
    for t in doc.tables:
        for row in t.rows:
            result.append(" | ".join(cell.text for cell in row.cells))
    return "\n\n".join(result)


def _read_pptx(self, file_bytes: bytes, filename: str) -> str:
    """Extract text from a PPTX file."""
    from pptx import Presentation
    import io
    prs = Presentation(io.BytesIO(file_bytes))
    result = []
    for si, slide in enumerate(prs.slides, 1):
        result.append(f"# Slide {si}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                result.append(shape.text)
    return "\n\n".join(result)


def _parse_csv_rows(self, content: str) -> list:
    """Parse CSV content into a list of rows with type conversion."""
    import csv as _csv_mod
    reader = _csv_mod.reader(io.StringIO(content))
    parsed_rows = []
    for csv_row in reader:
        converted = []
        for v in csv_row:
            v = v.strip()
            if v == '':
                converted.append(None)
            else:
                try:
                    converted.append(int(v))
                except ValueError:
                    try:
                        converted.append(float(v))
                    except ValueError:
                        if v.lower() == 'true':
                            converted.append(True)
                        elif v.lower() == 'false':
                            converted.append(False)
                        else:
                            converted.append(v)
        parsed_rows.append(converted)
    return parsed_rows

# -----------------------------------------------------------------
# Internal: save and return markdown link
# -----------------------------------------------------------------


async def read_file(
    self,
    file_id: str,
    max_rows: int = 500,
    sheet_name: str = "",
    row_start: int = 1,
    row_end: int = 0,
    __user__=None,
    __request__=None,
) -> str:
    """Read any Office file (.xlsx, .xls, .csv, .docx, .pptx) and return its contents as structured JSON.

        Auto-detects the file type from the file ID or filename.
        For xlsx/xls/csv: returns sheets with headers and rows.
        For docx: returns paragraphs with styles and tables.
        For pptx: returns slides with shapes and text.
        Legacy .doc and .ppt formats return a helpful error message.

        Args:
            file_id: The Open WebUI file ID (UUID) or filename
            max_rows: Maximum rows to return (default 500)
            sheet_name: Optional - read only this sheet (xlsx/xls only)
            row_start: Starting row (1-indexed, default 1)
            row_end: Ending row (0 = use max_rows limit)
        """
    try:
        file_data = _read_file_bytes(file_id)
        if file_data is None:
            return json.dumps({
                "error": (
                    f"Could not read file {file_id}. "
                    "Make sure the file was uploaded via the chat."
                )
            })

        # Detect type from filename in DB
        try:
            conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            row = conn.execute(
                "SELECT filename FROM file WHERE id = ?", (file_id,)
            ).fetchone()
            conn.close()
            filename = row[0] if row else file_id
        except Exception:
            filename = file_id

        file_type = _detect_type(filename)
        result: Dict[str, Any] = {
            "file_id": file_id,
            "filename": filename,
            "type": file_type,
        }

        if file_type == "xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True)
            result["sheets"] = []
            sheetnames = [sn for sn in wb.sheetnames] if not sheet_name else [sn for sn in wb.sheetnames if sn == sheet_name]
            if sheet_name and sheet_name not in wb.sheetnames:
                result["warning"] = f"Sheet '{sheet_name}' not found. Available sheets: {wb.sheetnames}"
            for sn in sheetnames:
                ws = wb[sn]
                sheet: Dict[str, Any] = {
                    "name": sn,
                    "headers": [],
                    "rows": [],
                    "total_rows": ws.max_row or 0,
                    "total_cols": ws.max_column or 0,
                }
                r_start = max(1, row_start)
                r_end = row_end if row_end > 0 else (ws.max_row or 0)
                r_end = min(r_end, r_start + max_rows - 1) if max_rows else r_end
                for ri, row in enumerate(ws.iter_rows(min_row=r_start, max_row=r_end), r_start):
                    rd = [_cell_value(c) for c in row]
                    if ri == 1:
                        sheet["headers"] = [str(v) if v is not None else "" for v in rd]
                    else:
                        sheet["rows"].append(rd)
                result["sheets"].append(sheet)
            wb.close()

        elif file_type == "xls":
            import xlrd
            xls_book = xlrd.open_workbook(file_contents=file_data)
            result["sheets"] = []
            sheet_indices = range(xls_book.nsheets)
            if sheet_name:
                sheet_indices = [i for i in range(xls_book.nsheets) if xls_book.sheet_by_index(i).name == sheet_name]
                if not sheet_indices:
                    all_names = [xls_book.sheet_by_index(i).name for i in range(xls_book.nsheets)]
                    result["warning"] = f"Sheet '{sheet_name}' not found. Available sheets: {all_names}"
            for sheet_idx in sheet_indices:
                xls_sheet = xls_book.sheet_by_index(sheet_idx)
                sheet: Dict[str, Any] = {
                    "name": xls_sheet.name,
                    "headers": [],
                    "rows": [],
                    "total_rows": xls_sheet.nrows,
                    "total_cols": xls_sheet.ncols,
                }
                r_start = max(0, row_start - 1)
                r_end = row_end if row_end > 0 else xls_sheet.nrows
                r_end = min(r_end, r_start + max_rows) if max_rows else r_end
                r_end = min(r_end, xls_sheet.nrows)
                for rx in range(r_start, r_end):
                    row_values = []
                    for cx in range(xls_sheet.ncols):
                        cell = xls_sheet.cell(rx, cx)
                        value = cell.value
                        if cell.ctype == xlrd.XL_CELL_DATE:
                            try:
                                import datetime as _dt
                                value = _dt.datetime(*xlrd.xldate_as_tuple(value, xls_book.datemode))
                            except Exception:
                                pass
                        elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                            value = bool(value)
                        row_values.append(value)
                    if rx == 0:
                        sheet["headers"] = [str(v) if v is not None else "" for v in row_values]
                    else:
                        sheet["rows"].append(row_values)
                result["sheets"].append(sheet)
            xls_book.release_resources()

        elif file_type == "csv":
            import csv as _csv_mod
            result["sheets"] = []
            d_bytes = file_data
            try:
                text = d_bytes.decode("utf-8-sig")
            except Exception:
                text = d_bytes.decode("utf-8", errors="replace")
            reader = _csv_mod.DictReader(io.StringIO(text))
            headers = reader.fieldnames or []
            rows = []
            for ri, row in enumerate(reader):
                if max_rows and ri >= max_rows:
                    break
                rows.append(dict(row))
            base_name = os.path.splitext(os.path.basename(filename))[0] if filename else "CSV"
            sheet = {
                "name": base_name,
                "headers": [str(h) for h in headers],
                "rows": [[r.get(h, "") for h in headers] for r in rows],
                "total_rows": len(rows),
                "total_cols": len(headers),
            }
            result["sheets"].append(sheet)

        elif file_type == "docx":
            from docx import Document
            from docx.shared import Pt
            doc = Document(io.BytesIO(file_data))
            paragraphs = []
            for p in doc.paragraphs:
                if p.text.strip():
                    style = p.style.name if p.style else "Normal"
                    runs_info = []
                    for run in p.runs:
                        run_data = {"text": run.text}
                        if run.font.highlight_color and run.font.highlight_color != WD_COLOR_INDEX.AUTO:
                            run_data["highlighted"] = True
                            run_data["highlight_color"] = str(run.font.highlight_color)
                        if run.font.bold:
                            run_data["bold"] = True
                        if run.font.italic:
                            run_data["italic"] = True
                        runs_info.append(run_data)
                    paragraphs.append({"style": style, "text": p.text, "runs": runs_info})
            tables = []
            for t in doc.tables:
                tbl = {"rows": []}
                for row in t.rows:
                    tbl["rows"].append([cell.text for cell in row.cells])
                tables.append(tbl)
            result["paragraphs"] = paragraphs
            result["tables"] = tables

        elif file_type == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_data))
            slides = []
            for si, slide in enumerate(prs.slides, 1):
                sdata: Dict[str, Any] = {"number": si, "shapes": []}
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        sdata["shapes"].append({
                            "type": str(shape.shape_type),
                            "name": shape.name,
                            "text": shape.text[:500],
                        })
                    if shape.has_table:
                        tbl = {"rows": []}
                        for row in shape.table.rows:
                            tbl["rows"].append([cell.text for cell in row.cells])
                        sdata["tables"] = tbl
                slides.append(sdata)
            result["slides"] = slides

        elif file_type == "doc":
            result["error"] = "Legacy .doc format is not supported. Please convert to .docx first."

        elif file_type == "ppt":
            result["error"] = "Legacy .ppt format is not supported. Please convert to .pptx first."

        else:
            result["error"] = f"Unsupported file type. Detected: {file_type}. Supported: xlsx, xls, docx, pptx"

        return json.dumps(result, indent=2, default=str, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

# -----------------------------------------------------------------
# ADD CONTENT
# -----------------------------------------------------------------


__all__ = [
    "_parse_csv_rows",
    "_read_docx",
    "_read_pptx",
    "_read_xls",
    "_read_xlsx",
    "read_file",
]
