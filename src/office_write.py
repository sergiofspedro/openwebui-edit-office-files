"""Edit Office Files — File creation and content generation (create, generate, add methods)."""

import io
import json
import os
import platform
import re
import sqlite3
import sys
import traceback
from copy import copy
from typing import Any, Dict, Optional

from .constants import *
from .utils import *

async def add_content(
    self,
    file_id: str,
    content: str,
    output_filename: str = "",
    __user__=None,
    __request__=None,
) -> str:
    """Add new content to an Office file while preserving original formatting.

        For spreadsheets (xlsx): content is CSV text with rows to add.
        For documents (docx): content is text to append at the end. Preserves original
            formatting and capitalization of the existing document.
        For presentations (pptx): each line defines a new slide. Use "---" as separator between slides.

        Args:
            file_id: File ID to edit
            content: Content to add (CSV for xlsx, text for docx/pptx)
            output_filename: Optional output filename
        """
    try:
        file_data = _read_file_bytes(file_id)
        if file_data is None:
            return json.dumps({"error": f"Could not read file {file_id}"})

        # Detect type
        try:
            conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            row = conn.execute(
                "SELECT filename FROM file WHERE id = ?", (file_id,)
            ).fetchone()
            conn.close()
            filename = row[0] if row else file_id
        except Exception:
            filename = file_id

        file_type = _detect_type(filename)
        out = io.BytesIO()
        out_name = output_filename

        if file_type == "xlsx":
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment
            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            ws = wb.active
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

            # Parse CSV content
            parsed_rows = self._parse_csv_rows(content)

            if not parsed_rows:
                return json.dumps({"error": "No rows provided in CSV content"})

            # Get reference styles from last row
            ref = {}
            if ws.max_row and ws.max_row >= 1:
                for cell in ws[ws.max_row]:
                    if cell.has_style:
                        ref[cell.column] = {
                            "font": copy(cell.font),
                            "fill": copy(cell.fill),
                            "border": copy(cell.border),
                            "alignment": copy(cell.alignment),
                            "number_format": cell.number_format,
                        }

            start = (ws.max_row or 0) + 1
            for i, rd in enumerate(parsed_rows):
                for j, v in enumerate(rd, 1):
                    cell = ws.cell(row=start + i, column=j)
                    if j in ref:
                        try:
                            cell.font = copy(ref[j]["font"])
                            cell.fill = copy(ref[j]["fill"])
                            cell.border = copy(ref[j]["border"])
                            cell.alignment = copy(ref[j]["alignment"])
                            cell.number_format = ref[j]["number_format"]
                        except Exception:
                            pass
                    cell.value = v

            wb.save(out)
            wb.close()

        elif file_type == "xls":
            # Convert .xls to .xlsx, then apply same add logic
            file_data = _xls_to_xlsx(file_data)
            file_type = "xlsx"
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment
            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            ws = wb.active

            # Parse CSV content
            parsed_rows = self._parse_csv_rows(content)

            if not parsed_rows:
                return json.dumps({"error": "No rows provided in CSV content"})

            ref = {}
            if ws.max_row and ws.max_row >= 1:
                for cell in ws[ws.max_row]:
                    if cell.has_style:
                        ref[cell.column] = {
                            "font": copy(cell.font),
                            "fill": copy(cell.fill),
                            "border": copy(cell.border),
                            "alignment": copy(cell.alignment),
                            "number_format": cell.number_format,
                        }

            start = (ws.max_row or 0) + 1
            for i, rd in enumerate(parsed_rows):
                for j, v in enumerate(rd, 1):
                    cell = ws.cell(row=start + i, column=j)
                    if j in ref:
                        try:
                            cell.font = copy(ref[j]["font"])
                            cell.fill = copy(ref[j]["fill"])
                            cell.border = copy(ref[j]["border"])
                            cell.alignment = copy(ref[j]["alignment"])
                            cell.number_format = ref[j]["number_format"]
                        except Exception:
                            pass
                    cell.value = v

            wb.save(out)
            wb.close()

        elif file_type == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_data))
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.docx"

            # Append paragraphs, preserving last paragraph's style
            last_style = "Normal"
            if doc.paragraphs:
                last_style = doc.paragraphs[-1].style.name if doc.paragraphs[-1].style else "Normal"

            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    continue
                if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
                    level = line.count('#')
                    text = line.lstrip('#').strip()
                    doc.add_heading(text, level=min(level, 3))
                elif line.startswith('- ') or line.startswith('* '):
                    text = line[2:].strip()
                    doc.add_paragraph(text, style='List Bullet')
                else:
                    segments = _parse_inline_md(line)
                    p = doc.add_paragraph(style=last_style)
                    for seg_text, fmt in segments:
                        if not seg_text:
                            continue
                        run = p.add_run(seg_text)
                        if fmt.get('code'):
                            run.font.name = "Consolas"
                            run.font.size = Pt(9)
                        if fmt.get('bold'):
                            run.font.bold = True
                        if fmt.get('italic'):
                            run.font.italic = True

            out = io.BytesIO()
            doc.save(out)

        elif file_type == "pptx":
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation(io.BytesIO(file_data))
            if not out_name:
                out_name = os.path.splitext(filename)[0] + "_edited.pptx"

            # Split content by "---" into slides
            slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
            blank_layout = prs.slide_layouts[6]  # Blank layout

            for spec in slide_specs:
                spec = spec.strip()
                if not spec:
                    continue
                lines = spec.split("\n")
                title = lines[0].strip()
                body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                slide = prs.slides.add_slide(blank_layout)
                # Add title
                txBox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                )
                tf = txBox.text_frame
                tf.text = _format_text(title, mode="format")
                p = tf.paragraphs[0]
                p.font.size = Inches(0.6)
                p.font.bold = True

                # Add body text
                if body:
                    txBox2 = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                    )
                    tf2 = txBox2.text_frame
                    tf2.text = _format_text(body, mode="format")
                    for para in tf2.paragraphs:
                        para.font.size = Inches(0.3)

            out = io.BytesIO()
            prs.save(out)

        else:
            return json.dumps({"error": f"Unsupported type: {file_type}"})

        out.seek(0)
        url, name = await self._save_and_link(out.read(), out_name, __request__)
        if url:
            return f"[{name}]({url})\n\nAdded content to {file_type.upper()} file, preserving original formatting."
        return json.dumps({"error": "Could not save file"})

    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

# -----------------------------------------------------------------
# REPLACE TEXT
# -----------------------------------------------------------------


async def create_file(
    self,
    file_type: str,
    content: str,
    output_filename: str = "",
    raw_text: bool = False,
    __user__=None,
    __request__=None,
) -> str:
    """Create a new Office file from scratch.

        For xlsx: content is CSV with headers on first line.
        For docx: supports markdown formatting — headings (#, ##, ###), bullets (-, *),
            and inline markdown (**bold**, *italic*, `code`).
        For pptx: each line defines a slide. Use "---" as separator between slides.

        Args:
            file_type: 'xlsx', 'docx', or 'pptx'
            content: Content specification
            output_filename: Output filename
            raw_text: If True, skip text formatting
        """
    fmt_mode = "preserve" if raw_text else "format"
    try:
        ftype = file_type.lower().replace(".", "")
        if ftype not in ("xlsx", "docx", "pptx"):
            return json.dumps({"error": f"Unsupported type: {file_type}. Use xlsx, docx, or pptx."})

        out_name = output_filename or f"document.{ftype}"
        out = io.BytesIO()

        if ftype == "xlsx":
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment
            wb = openpyxl.Workbook()
            ws = wb.active

            import csv as _csv_mod
            reader = _csv_mod.reader(io.StringIO(content))
            rows = list(reader)

            if rows:
                # First row = headers (styled)
                for j, h in enumerate(rows[0], 1):
                    c = ws.cell(row=1, column=j, value=h.strip())
                    c.font = Font(bold=True, color="FFFFFF", size=11)
                    c.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                    c.alignment = Alignment(horizontal="center")

                # Data rows
                for i, rd in enumerate(rows[1:], 2):
                    for j, v in enumerate(rd, 1):
                        v = v.strip()
                        c = ws.cell(row=i, column=j)
                        if v == '':
                            c.value = None
                        else:
                            try:
                                c.value = int(v)
                            except ValueError:
                                try:
                                    c.value = float(v)
                                except ValueError:
                                    if v.lower() == 'true':
                                        c.value = True
                                    elif v.lower() == 'false':
                                        c.value = False
                                    else:
                                        c.value = v
                        c.alignment = Alignment(
                            horizontal="center" if isinstance(c.value, (int, float)) else "left"
                        )

                # Auto-fit columns
                for col in ws.columns:
                    mx = max((len(str(c.value or "")) for c in col), default=5)
                    ws.column_dimensions[col[0].column_letter].width = min(mx + 3, 50)

            wb.save(out)
            wb.close()

        elif ftype == "docx":
            from docx import Document
            from docx.shared import Pt, RGBColor
            doc = Document()
            for raw_line in content.split("\n"):
                line = raw_line
                stripped = line.strip()

                if stripped == "":
                    # Empty line → paragraph break
                    doc.add_paragraph("")
                    continue

                # Heading detection (must start with # markers)
                if stripped.startswith("### "):
                    text = stripped[4:]
                    doc.add_heading(text, level=3)
                    continue
                if stripped.startswith("## "):
                    text = stripped[3:]
                    doc.add_heading(text, level=2)
                    continue
                if stripped.startswith("# "):
                    text = stripped[2:]
                    doc.add_heading(text, level=1)
                    continue

                # Bullet detection
                if stripped.startswith("- ") or stripped.startswith("* "):
                    text = _format_text(stripped[2:], mode=fmt_mode)
                    doc.add_paragraph(text, style='List Bullet')
                    continue

                # Body text — rich formatting via inline markdown
                text = _format_text(line, mode=fmt_mode)
                segments = _parse_inline_md(text)
                p = doc.add_paragraph()
                for seg_text, fmt in segments:
                    if not seg_text:
                        continue
                    run = p.add_run(seg_text)
                    if fmt.get('code'):
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)
                    else:
                        run.font.size = Pt(11)
                    if fmt.get('bold'):
                        run.font.bold = True
                    if fmt.get('italic'):
                        run.font.italic = True
                    if fmt.get('link'):
                        run.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
                        run.font.underline = True
            doc.save(out)

        elif ftype == "pptx":
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation()
            slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
            blank_layout = prs.slide_layouts[6]

            for spec in slide_specs:
                spec = spec.strip()
                if not spec:
                    continue
                lines = spec.split("\n")
                title = lines[0].strip()
                body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                slide = prs.slides.add_slide(blank_layout)
                txBox = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                )
                tf = txBox.text_frame
                tf.text = title
                p = tf.paragraphs[0]
                p.font.size = Inches(0.6)
                p.font.bold = True

                if body:
                    txBox2 = slide.shapes.add_textbox(
                        Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                    )
                    tf2 = txBox2.text_frame
                    tf2.text = body
                    for para in tf2.paragraphs:
                        para.font.size = Inches(0.3)

            prs.save(out)

        out.seek(0)
        url, name = await self._save_and_link(out.read(), out_name, __request__)
        if url:
            return f"[{name}]({url})\n\nCreated new {ftype.upper()} file."
        return json.dumps({"error": "Could not save file"})

    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


async def generate_document(self, content: str, title: str = "Document", theme: str = "professional", typography: str = "modern", raw_text: bool = False, __user__=None, __request__=None) -> str:
    """Generate a professional Word document with modern styling, emojis, cards, and visual elements.

        Headings preserve their original capitalization (no forced sentence case).
        Body text uses sentence case (first letter of each sentence capitalized, acronyms preserved).
        Supports inline markdown: **bold**, *italic*, `code`, [links](url), and ```code blocks```.

        Args:
            content: Markdown content to convert to a document
            title: Document title / filename
            theme: Visual theme - professional, modern, creative, corporate, minimal, elegant, ocean, sunset, forest, midnight
            typography: Font preset - modern, classic, serif, sans
            raw_text: If True, skip all text formatting — text is used exactly as provided.
        """
    from docx import Document
    from docx.shared import Inches, Pt, Cm, RGBColor, Emu
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml
    from lxml import etree
    import datetime, re as _re
    fmt_mode = "preserve" if raw_text else "format"
    
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Cm(2.0)
    section.bottom_margin = Cm(2.0)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)
    
    # --- Color Palettes (10+) ---
    palettes = {
        "professional": {"primary": "1F4E79", "accent": "2E75B6", "light": "D6E4F0", "text": "333333", "bg": "FFFFFF", "success": "27AE60", "warning": "F39C12", "danger": "E74C3C", "info": "3498DB"},
        "modern": {"primary": "2D3436", "accent": "6C5CE7", "light": "DFE6E9", "text": "2D3436", "bg": "FFFFFF", "success": "00B894", "warning": "FDCB6E", "danger": "FF7675", "info": "74B9FF"},
        "creative": {"primary": "E17055", "accent": "FDCB6E", "light": "FFF3E0", "text": "2D3436", "bg": "FFFAF5", "success": "00B894", "warning": "FDCB6E", "danger": "D63031", "info": "6C5CE7"},
        "corporate": {"primary": "003366", "accent": "CC0000", "light": "E8EEF4", "text": "1A1A1A", "bg": "FFFFFF", "success": "006633", "warning": "FF6600", "danger": "CC0000", "info": "0066CC"},
        "minimal": {"primary": "000000", "accent": "666666", "light": "F5F5F5", "text": "333333", "bg": "FAFAFA", "success": "333333", "warning": "666666", "danger": "000000", "info": "999999"},
        "elegant": {"primary": "4A235A", "accent": "8E44AD", "light": "F3E5F5", "text": "1A1A1A", "bg": "FDFBF7", "success": "27AE60", "warning": "D4AC0D", "danger": "C0392B", "info": "2980B9"},
        "ocean": {"primary": "0A3D62", "accent": "38ADA9", "light": "D1F2EB", "text": "1E272E", "bg": "F8FFFE", "success": "079992", "warning": "F6B93B", "danger": "E55039", "info": "3C6382"},
        "sunset": {"primary": "B33771", "accent": "FD7272", "light": "FFE4E4", "text": "2C3A47", "bg": "FFFBF5", "success": "58B19F", "warning": "F8B500", "danger": "E66767", "info": "786FA6"},
        "forest": {"primary": "1B4332", "accent": "40916C", "light": "D8F3DC", "text": "1A1A1A", "bg": "F7FFF7", "success": "2D6A4F", "warning": "B7B73F", "danger": "D00000", "info": "52B788"},
        "midnight": {"primary": "0F172A", "accent": "38BDF8", "light": "1E293B", "text": "F8FAFC", "bg": "0F172A", "success": "34D399", "warning": "FBBF24", "danger": "F87171", "info": "60A5FA"},
    }
    colors = palettes.get(theme, palettes["professional"])
    
    # --- Typography Presets ---
    fonts = {
        "modern": {"heading": "Calibri", "body": "Calibri"},
        "classic": {"heading": "Georgia", "body": "Calibri"},
        "serif": {"heading": "Cambria", "body": "Calibri"},
        "sans": {"heading": "Arial", "body": "Arial"},
    }
    font_pair = fonts.get(typography, fonts["modern"])
    
    def hex_to_rgb(h):
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    
    def add_colored_bar(doc, color_hex, height_pt=4):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(0)
        run = p.add_run(" ")
        run.font.size = Pt(height_pt)
        pPr = p._p.get_or_add_pPr()
        shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}" w:val="clear"/>')
        pPr.append(shd)
    
    def add_card_box(doc, lines, border_color, bg_color, icon="", title=""):
        table = doc.add_table(rows=1, cols=1)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        cell = table.rows[0].cells[0]
        cell.width = Inches(6.0)
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{bg_color}" w:val="clear"/>')
        tcPr.append(shd)
        borders = parse_xml(
            f'<w:tcBorders {nsdecls("w")}>'
            f'<w:left w:val="single" w:sz="24" w:space="8" w:color="{border_color}"/>'
            f'<w:top w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
            f'<w:bottom w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
            f'<w:right w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
            f'</w:tcBorders>'
        )
        tcPr.append(borders)
        p = cell.paragraphs[0]
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(4)
        if icon or title:
            header_text = f"{icon} {title}" if icon and title else (icon or title)
            run = p.add_run(_format_text(header_text, mode=fmt_mode))
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.name = font_pair["heading"]
            run.font.color.rgb = hex_to_rgb(border_color)
        for line in lines:
            p = cell.add_paragraph()
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after = Pt(2)
            text = _format_text(line, mode=fmt_mode)
            segments = _parse_inline_md(text)
            for seg_text, fmt in segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
                else:
                    run.font.size = Pt(10)
                    run.font.name = font_pair["body"]
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('italic'):
                    run.font.italic = True
                run.font.color.rgb = hex_to_rgb(colors["text"])
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
        doc.add_paragraph()
    
    def add_kpi_card(doc, value, label, color_hex, icon=""):
        table = doc.add_table(rows=1, cols=1)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        cell = table.rows[0].cells[0]
        cell.width = Inches(2.8)
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["light"]}" w:val="clear"/>')
        tcPr.append(shd)
        borders = parse_xml(f'<w:tcBorders {nsdecls("w")}><w:top w:val="single" w:sz="12" w:space="0" w:color="{color_hex}"/><w:left w:val="single" w:sz="4" w:space="0" w:color="DDDDDD"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="DDDDDD"/><w:right w:val="single" w:sz="4" w:space="0" w:color="DDDDDD"/></w:tcBorders>')
        tcPr.append(borders)
        p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(12)
        run = p.add_run(f"{icon} {value}" if icon else value)
        run.font.size = Pt(24); run.font.bold = True
        run.font.color.rgb = hex_to_rgb(color_hex); run.font.name = font_pair["heading"]
        p2 = cell.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p2.paragraph_format.space_after = Pt(12)
        label_text = _format_text(label, mode=fmt_mode)
        label_segments = _parse_inline_md(label_text)
        for seg_text, fmt in label_segments:
            if not seg_text:
                continue
            run2 = p2.add_run(seg_text)
            run2.font.size = Pt(9)
            run2.font.color.rgb = hex_to_rgb(colors["text"])
            run2.font.name = font_pair["body"]
            if fmt.get('bold'):
                run2.font.bold = True
            if fmt.get('italic'):
                run2.font.italic = True
            if fmt.get('code'):
                run2.font.name = "Consolas"
                run2.font.size = Pt(9)
    
    def add_progress_bar(doc, label, percentage, color_hex):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(4)
        p.paragraph_format.space_after = Pt(2)
        label_text = _format_text(label, mode=fmt_mode)
        label_segments = _parse_inline_md(label_text)
        for seg_text, fmt in label_segments:
            if not seg_text:
                continue
            run = p.add_run(seg_text)
            run.font.size = Pt(10); run.font.bold = True
            run.font.name = font_pair["body"]; run.font.color.rgb = hex_to_rgb(colors["text"])
            if fmt.get('italic'):
                run.font.italic = True
            if fmt.get('link'):
                run.font.color.rgb = hex_to_rgb("2563EB")
                run.font.underline = True
            if fmt.get('code'):
                run.font.name = "Consolas"
                run.font.size = Pt(9)
        run = p.add_run(f": {percentage}%")
        run.font.size = Pt(10); run.font.bold = True
        run.font.name = font_pair["body"]; run.font.color.rgb = hex_to_rgb(colors["text"])
        table = doc.add_table(rows=1, cols=2)
        table.alignment = WD_TABLE_ALIGNMENT.LEFT
        filled = table.rows[0].cells[0]; empty = table.rows[0].cells[1]
        filled.width = Inches(5.0 * percentage / 100)
        empty.width = Inches(5.0 * (100 - percentage) / 100)
        tcF = filled._tc; tcPrF = tcF.get_or_add_tcPr()
        shdF = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}" w:val="clear"/>')
        tcPrF.append(shdF)
        tcE = empty._tc; tcPrE = tcE.get_or_add_tcPr()
        shdE = parse_xml(f'<w:shd {nsdecls("w")} w:fill="E0E0E0" w:val="clear"/>')
        tcPrE.append(shdE)
        pF = filled.paragraphs[0]; runF = pF.add_run(" "); runF.font.size = Pt(6)
        pE = empty.paragraphs[0]; runE = pE.add_run(" "); runE.font.size = Pt(6)
    
    def add_step_guide(doc, steps):
        for i, step in enumerate(steps, 1):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(6)
            p.paragraph_format.space_after = Pt(2)
            run = p.add_run(f"  {i}  ")
            run.font.size = Pt(14); run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255); run.font.name = font_pair["heading"]
            rPr = run._r.get_or_add_rPr()
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["primary"]}" w:val="clear"/>')
            rPr.append(shd)
            text = _format_text(step, mode=fmt_mode)
            segments = _parse_inline_md(text)
            first = True
            for seg_text, fmt in segments:
                if not seg_text:
                    continue
                prefix = "  " if first else ""
                first = False
                run = p.add_run(prefix + seg_text)
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
                else:
                    run.font.size = Pt(11)
                    run.font.name = font_pair["body"]
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('italic'):
                    run.font.italic = True
                run.font.color.rgb = hex_to_rgb(colors["text"])
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
    
    def add_pull_quote(doc, text, author=""):
        table = doc.add_table(rows=1, cols=1)
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        cell = table.rows[0].cells[0]; cell.width = Inches(5.5)
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["light"]}" w:val="clear"/>')
        tcPr.append(shd)
        borders = parse_xml(f'<w:tcBorders {nsdecls("w")}><w:left w:val="single" w:sz="36" w:space="12" w:color="{colors["accent"]}"/><w:top w:val="none" w:sz="0" w:space="0" w:color="auto"/><w:bottom w:val="none" w:sz="0" w:space="0" w:color="auto"/><w:right w:val="none" w:sz="0" w:space="0" w:color="auto"/></w:tcBorders>')
        tcPr.append(borders)
        p = cell.paragraphs[0]; p.paragraph_format.space_before = Pt(12)
        p.paragraph_format.space_after = Pt(4)
        text_formatted = _format_text(text, mode=fmt_mode)
        segments = _parse_inline_md(text_formatted)
        # Opening quote
        run = p.add_run('"')
        run.font.size = Pt(14); run.font.italic = True
        run.font.name = font_pair["heading"]; run.font.color.rgb = hex_to_rgb(colors["text"])
        for seg_text, fmt in segments:
            if not seg_text:
                continue
            run = p.add_run(seg_text)
            run.font.size = Pt(14); run.font.italic = True
            run.font.name = font_pair["heading"]
            run.font.color.rgb = hex_to_rgb(colors["text"])
            if fmt.get('bold'):
                run.font.bold = True
            if fmt.get('link'):
                run.font.color.rgb = hex_to_rgb("2563EB")
                run.font.underline = True
            if fmt.get('code'):
                run.font.name = "Consolas"
                run.font.size = Pt(9)
        # Closing quote
        run = p.add_run('"')
        run.font.size = Pt(14); run.font.italic = True
        run.font.name = font_pair["heading"]; run.font.color.rgb = hex_to_rgb(colors["text"])
        if author:
            p2 = cell.add_paragraph(); p2.paragraph_format.space_after = Pt(12)
            p2.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            author_text = _format_text(author, mode=fmt_mode)
            author_segments = _parse_inline_md(author_text)
            first = True
            for seg_text, fmt in author_segments:
                if not seg_text:
                    continue
                prefix = "\u2014 " if first else ""
                first = False
                run = p2.add_run(prefix + seg_text)
                run.font.size = Pt(10); run.font.name = font_pair["body"]
                run.font.color.rgb = hex_to_rgb(colors["accent"])
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('italic'):
                    run.font.italic = True
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
        doc.add_paragraph()
    
    def add_comparison_table(doc, headers, rows):
        table = doc.add_table(rows=len(rows)+1, cols=len(headers))
        table.style = 'Colorful Grid Accent 1'
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        for j, h in enumerate(headers):
            cell = table.rows[0].cells[j]
            cell.text = ''
            p = cell.paragraphs[0]
            segments = _parse_inline_md(_format_text(h, mode=fmt_mode))
            for seg_text, fmt in segments:
                if not seg_text:
                    continue
                r = p.add_run(seg_text)
                r.font.size = Pt(10); r.font.bold = True; r.font.name = font_pair["heading"]
                if fmt.get('italic'):
                    r.font.italic = True
                if fmt.get('code'):
                    r.font.name = "Consolas"
                    r.font.size = Pt(9)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                if j < len(headers):
                    cell = table.rows[i+1].cells[j]
                    display = val
                    if val.lower() in ("yes","true","sim","si","oui","ja"): display = "\u2705 " + val
                    elif val.lower() in ("no","false","nao","non","nein"): display = "\u274c " + val
                    elif val.lower() in ("partial","maybe","talvez"): display = "\u26a0\ufe0f " + val
                    cell.text = ''
                    p = cell.paragraphs[0]
                    segments = _parse_inline_md(_format_text(display, mode=fmt_mode))
                    for seg_text, fmt in segments:
                        if not seg_text:
                            continue
                        r = p.add_run(seg_text)
                        r.font.size = Pt(10); r.font.name = font_pair["body"]
                        if fmt.get('bold'):
                            r.font.bold = True
                        if fmt.get('italic'):
                            r.font.italic = True
                        if fmt.get('code'):
                            r.font.name = "Consolas"
                            r.font.size = Pt(9)
        doc.add_paragraph()
    
    def add_timeline(doc, events):
        for i, (date, desc) in enumerate(events):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(4)
            p.paragraph_format.space_after = Pt(2)
            icon = "\u25cf" if i == 0 else "\u25cb"
            run = p.add_run(f"  {icon}  ")
            run.font.size = Pt(12); run.font.color.rgb = hex_to_rgb(colors["accent"])
            run.font.name = font_pair["heading"]
            date_text = _format_text(date, mode=fmt_mode)
            date_segments = _parse_inline_md(date_text)
            for seg_text, fmt in date_segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                run.font.size = Pt(10); run.font.bold = True
                run.font.name = font_pair["body"]; run.font.color.rgb = hex_to_rgb(colors["primary"])
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('italic'):
                    run.font.italic = True
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
            run_sep = p.add_run("  ")
            run_sep.font.size = Pt(10)
            run_sep.font.name = font_pair["body"]
            desc_text = _format_text(desc, mode=fmt_mode)
            desc_segments = _parse_inline_md(desc_text)
            for seg_text, fmt in desc_segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                run.font.size = Pt(10); run.font.name = font_pair["body"]
                run.font.color.rgb = hex_to_rgb(colors["text"])
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('italic'):
                    run.font.italic = True
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
    
    def add_status_badge(doc, text, status="info"):
        colors_map = {"success": colors["success"], "warning": colors["warning"], "danger": colors["danger"], "info": colors["info"]}
        icons = {"success": "\U0001f7e2", "warning": "\U0001f7e1", "danger": "\U0001f534", "info": "\U0001f535"}
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(2)
        p.paragraph_format.space_after = Pt(2)
        badge_color = hex_to_rgb(colors_map.get(status, colors["info"]))
        icon_text = icons.get(status, icons['info'])
        run = p.add_run(f" {icon_text} ")
        run.font.size = Pt(10); run.font.bold = True
        run.font.name = font_pair["body"]; run.font.color.rgb = badge_color
        badge_text = _format_text(text, mode=fmt_mode)
        badge_segments = _parse_inline_md(badge_text)
        for seg_text, fmt in badge_segments:
            if not seg_text:
                continue
            run = p.add_run(seg_text)
            run.font.size = Pt(10); run.font.bold = True
            run.font.name = font_pair["body"]; run.font.color.rgb = badge_color
            if fmt.get('italic'):
                run.font.italic = True
            if fmt.get('link'):
                run.font.color.rgb = hex_to_rgb("2563EB")
                run.font.underline = True
            if fmt.get('code'):
                run.font.name = "Consolas"
                run.font.size = Pt(9)
    
    def add_visual_separator(doc, style="dots"):
        separators = {"dots": "\u25cf \u25cf \u25cf", "line": "\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501", "dash": "\u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500"}
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(8)
        p.paragraph_format.space_after = Pt(8)
        run = p.add_run(separators.get(style, separators["dots"]))
        run.font.size = Pt(8); run.font.color.rgb = hex_to_rgb(colors["accent"])
    
    def _add_rich_paragraph(doc, text, style=None, font_name="Calibri", font_size=11, color=None):
        """Add a paragraph with rich formatting from inline markdown parsing."""
        segments = _parse_inline_md(text)
        if style and style.startswith('Heading'):
            try:
                level = int(style.split()[-1])
            except (ValueError, IndexError):
                level = 1
            p = doc.add_heading('', level=level)
            for run in p.runs:
                run.text = ''
        else:
            p = doc.add_paragraph(style=style) if style else doc.add_paragraph()
        for seg_text, fmt in segments:
            if not seg_text:
                continue
            run = p.add_run(seg_text)
            if fmt.get('code'):
                run.font.name = "Consolas"
                run.font.size = Pt(9)
            else:
                run.font.name = font_name
                run.font.size = Pt(font_size)
            if fmt.get('bold'):
                run.font.bold = True
            if fmt.get('italic'):
                run.font.italic = True
            if color:
                if isinstance(color, str):
                    run.font.color.rgb = hex_to_rgb(color)
                else:
                    run.font.color.rgb = color
            if fmt.get('link'):
                run.font.color.rgb = hex_to_rgb("2563EB")
                run.font.underline = True
        return p
    
    # --- Default Style ---
    style = doc.styles['Normal']
    font = style.font; font.name = font_pair["body"]; font.size = Pt(11)
    font.color.rgb = hex_to_rgb(colors["text"])
    style.paragraph_format.space_after = Pt(8)
    style.paragraph_format.line_spacing = 1.15
    
    for i in range(1, 4):
        hs = doc.styles[f'Heading {i}']
        hs.font.name = font_pair["heading"]
        hs.font.color.rgb = hex_to_rgb(colors["primary"])
        sizes = {1: 24, 2: 18, 3: 14}
        hs.font.size = Pt(sizes.get(i, 14))
        spaces = {1: (24, 12), 2: (18, 8), 3: (12, 6)}
        hs.paragraph_format.space_before = Pt(spaces[i][0])
        hs.paragraph_format.space_after = Pt(spaces[i][1])
    
    # --- Cover Page ---
    add_colored_bar(doc, colors["primary"], 4)
    cover_table = doc.add_table(rows=1, cols=1)
    cover_table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = cover_table.rows[0].cells[0]; cell.width = Inches(6.5)
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["primary"]}" w:val="clear"/>')
    tcPr.append(shd)
    p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(40); p.paragraph_format.space_after = Pt(8)
    run = p.add_run(title)
    run.font.size = Pt(28); run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255); run.font.name = font_pair["heading"]
    p2 = cell.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.paragraph_format.space_after = Pt(30)
    run2 = p2.add_run(datetime.datetime.now().strftime("%B %d, %Y"))
    run2.font.size = Pt(12); run2.font.color.rgb = RGBColor(200, 200, 200)
    run2.font.name = font_pair["body"]
    doc.add_paragraph()
    
    # --- Auto TOC ---
    toc_added = False
    
    # --- Process Content ---
    lines = content.split('\n')
    in_card = False; card_lines = []; card_icon = ""; card_title = ""
    in_kpi = False; kpi_data = []
    in_timeline = False; timeline_data = []
    in_steps = False; step_lines = []
    in_quote = False; quote_text = ""; quote_author = ""
    in_comparison = False; comp_headers = []; comp_rows = []
    in_progress = False; progress_data = []
    in_code_block = False; code_lines = []
    _timeline_re = _re.compile(r'^(\d{4}|\d{1,2}[/-]\d{1,2}|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*)\s*[-:]\s*(.+)$', _re.IGNORECASE)
    
    for line in lines:
        line = line.strip()
        if not line:
            if in_code_block:
                code_lines.append("")
                continue
            if in_card:
                add_card_box(doc, card_lines, colors["accent"], colors["light"], card_icon, card_title)
                card_lines = []; in_card = False
            if in_steps:
                if len(step_lines) >= 3:
                    add_step_guide(doc, step_lines)
                else:
                    for i, step_text in enumerate(step_lines, 1):
                        _add_rich_paragraph(doc, f"{i}. {_format_text(step_text, mode=fmt_mode)}", font_name=font_pair["body"], font_size=11, color=colors["text"])
                step_lines = []; in_steps = False
            if in_timeline:
                add_timeline(doc, timeline_data)
                timeline_data = []; in_timeline = False
            if in_comparison:
                add_comparison_table(doc, comp_headers, comp_rows)
                comp_headers = []; comp_rows = []; in_comparison = False
            if in_progress:
                for lbl, pct, clr in progress_data:
                    add_progress_bar(doc, lbl, pct, clr)
                progress_data = []; in_progress = False
            continue
        
        # Code block: ``` ... ```
        if line.startswith('```'):
            if not in_code_block:
                in_code_block = True
                code_lines = []
            else:
                for code_line in code_lines:
                    p = doc.add_paragraph()
                    p.paragraph_format.space_before = Pt(1)
                    p.paragraph_format.space_after = Pt(1)
                    run = p.add_run(code_line)
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
                    run.font.color.rgb = hex_to_rgb(colors["text"])
                in_code_block = False
                code_lines = []
            continue
        if in_code_block:
            code_lines.append(line)
            continue
        
        # Card: > 📊 **Title** or > content
        if line.startswith('> ') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_comparison and not in_progress:
            in_card = True
            text = line[2:]
            m = _re.match(r'^([\U0001F300-\U0001F9FF]\s*)?\*\*(.+?)\*\*', text)
            if m:
                card_icon = m.group(1).strip() if m.group(1) else ""
                card_title = m.group(2)
                remaining = text[m.end():].strip()
                if remaining: card_lines.append(remaining)
            else:
                card_lines.append(text)
            continue
        elif in_card and line.startswith('> '):
            card_lines.append(line[2:])
            continue
        
        # KPI: 📊 85% | Customer Satisfaction
        if '|' in line and '%' in line and not in_card and not in_timeline and not in_steps and not in_quote and not in_comparison and not in_progress:
            parts = [p.strip() for p in line.split('|')]
            if len(parts) >= 2 and any('%' in p for p in parts):
                in_kpi = True
                kpi_data.append(parts)
                continue
        elif in_kpi and '|' in line and '%' in line:
            parts = [p.strip() for p in line.split('|')]
            kpi_data.append(parts)
            continue
        elif in_kpi:
            # Render KPI cards
            kpi_table = doc.add_table(rows=1, cols=len(kpi_data))
            kpi_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            for j, kpi in enumerate(kpi_data):
                if j < len(kpi_data):
                    val = kpi[0] if len(kpi) > 0 else ""
                    lbl = kpi[1] if len(kpi) > 1 else ""
                    icon_match = _re.match(r'^([\U0001F300-\U0001F9FF]\s*)', val)
                    icon = icon_match.group(1).strip() if icon_match else ""
                    if icon: val = val[len(icon):].strip()
                    cell = kpi_table.rows[0].cells[j]
                    cell.width = Inches(2.8)
                    tcC = cell._tc; tcPrC = tcC.get_or_add_tcPr()
                    shdC = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["light"]}" w:val="clear"/>')
                    tcPrC.append(shdC)
                    pC = cell.paragraphs[0]; pC.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    pC.paragraph_format.space_before = Pt(12)
                    runC = pC.add_run(f"{icon} {val}" if icon else val)
                    runC.font.size = Pt(24); runC.font.bold = True
                    runC.font.color.rgb = hex_to_rgb(colors["primary"])
                    runC.font.name = font_pair["heading"]
                    pC2 = cell.add_paragraph(); pC2.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    pC2.paragraph_format.space_after = Pt(12)
                    runC2 = pC2.add_run(_format_text(lbl, mode=fmt_mode))
                    runC2.font.size = Pt(9); runC2.font.color.rgb = hex_to_rgb(colors["text"])
                    runC2.font.name = font_pair["body"]
            doc.add_paragraph()
            kpi_data = []; in_kpi = False
        
        # Timeline: 📅 2024 | Event description (pipe format)
        if '|' in line and _re.match(r'^[\U0001F300-\U0001F9FF\s]*\d{4}', line) and not in_card and not in_kpi and not in_steps and not in_quote and not in_comparison and not in_progress:
            in_timeline = True
            parts = [p.strip() for p in line.split('|', 1)]
            timeline_data.append((parts[0], parts[1] if len(parts) > 1 else ""))
            continue
        elif in_timeline and '|' in line and _re.match(r'^[\U0001F300-\U0001F9FF\s]*\d{4}', line):
            parts = [p.strip() for p in line.split('|', 1)]
            timeline_data.append((parts[0], parts[1] if len(parts) > 1 else ""))
            continue
        elif in_timeline:
            add_timeline(doc, timeline_data)
            timeline_data = []; in_timeline = False
        
        # Timeline (colon/dash format): 2024 - Event or Jan 15: Event
        if not in_timeline and not in_card and not in_kpi and not in_steps and not in_quote and not in_comparison and not in_progress:
            tl_m = _timeline_re.match(line)
            if tl_m:
                in_timeline = True
                timeline_data.append((tl_m.group(1), tl_m.group(2)))
                continue
        
        # Steps: 1. Step one / 2. Step two
        if _re.match(r'^\d+\.\s', line) and not in_card and not in_kpi and not in_timeline and not in_quote and not in_comparison and not in_progress:
            in_steps = True
            step_lines.append(_re.sub(r'^\d+\.\s', '', line))
            continue
        elif in_steps and _re.match(r'^\d+\.\s', line):
            step_lines.append(_re.sub(r'^\d+\.\s', '', line))
            continue
        elif in_steps:
            if len(step_lines) >= 3:
                add_step_guide(doc, step_lines)
            else:
                for i, step_text in enumerate(step_lines, 1):
                    _add_rich_paragraph(doc, f"{i}. {_format_text(step_text, mode=fmt_mode)}", font_name=font_pair["body"], font_size=11, color=colors["text"])
            step_lines = []; in_steps = False

        # Pull quote: "Quote text" — Author
        if line.startswith('"') and line.endswith('"') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_comparison and not in_progress:
            in_quote = True
            quote_text = line.strip('"')
            continue
        elif in_quote and (line.startswith('\u2014') or (line.strip() and not line.startswith('"') and not line.startswith('#') and not line.startswith('-') and not line.startswith('*') and not line.startswith('>') and not line.startswith('|') and not line.startswith('@') and not line.startswith('```'))):
            quote_author = line.lstrip('\u2014 ').strip()
            add_pull_quote(doc, quote_text, quote_author)
            in_quote = False; quote_text = ""; quote_author = ""
            continue
        
        # Comparison table: | Feature | A | B |
        if line.startswith('|') and line.endswith('|') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_progress:
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if all(c.startswith('---') for c in cells):
                continue
            if not in_comparison:
                in_comparison = True
                comp_headers = cells
            else:
                comp_rows.append(cells)
            continue
        elif in_comparison and line.startswith('|') and line.endswith('|'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if not all(c.startswith('---') for c in cells):
                comp_rows.append(cells)
            continue
        elif in_comparison:
            add_comparison_table(doc, comp_headers, comp_rows)
            comp_headers = []; comp_rows = []; in_comparison = False
        
        # Progress: Label: 75%
        if ':' in line and _re.search(r'(\d+)%', line) and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_comparison:
            in_progress = True
            parts = line.split(':', 1)
            lbl = parts[0].strip()
            pct_match = _re.search(r'(\d+)%', parts[1])
            pct = int(pct_match.group(1)) if pct_match else 0
            progress_data.append((lbl, pct, colors["accent"]))
            continue
        elif in_progress and ':' in line and _re.search(r'(\d+)%', line):
            parts = line.split(':', 1)
            lbl = parts[0].strip()
            pct_match = _re.search(r'(\d+)%', parts[1])
            pct = int(pct_match.group(1)) if pct_match else 0
            progress_data.append((lbl, pct, colors["accent"]))
            continue
        elif in_progress:
            for lbl, pct, clr in progress_data:
                add_progress_bar(doc, lbl, pct, clr)
            progress_data = []; in_progress = False
        
        # Status badge: @success Task complete
        if line.startswith('@') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_comparison and not in_progress:
            parts = line[1:].split(None, 1)
            status = parts[0].lower() if parts else "info"
            text = parts[1] if len(parts) > 1 else ""
            add_status_badge(doc, text, status)
            continue
        
        # Visual separator: --- or ***
        if line in ('---', '***', '...'):
            style_map = {'---': 'line', '***': 'dots', '...': 'dash'}
            add_visual_separator(doc, style_map.get(line, 'dots'))
            continue
        
        # Headings with emoji
        if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
            level = line.count('#')
            text = line.lstrip('#').strip()
            heading_level = min(level, 3)
            _add_rich_paragraph(doc, text, style=f'Heading {heading_level}', font_name=font_pair["heading"], font_size=sizes.get(heading_level, 14), color=colors["primary"])
            continue
        
        # Icon bullets
        if line.startswith('- ') or line.startswith('* '):
            bullet_text = line[2:].strip()
            text = _format_text(bullet_text, mode=fmt_mode)
            icon = "\u2022"
            if text.lower().startswith(('done','complete','yes','ok','success','conclu')):
                icon = "\u2705"
            elif text.lower().startswith(('no','not','fail','error','wrong','nao')):
                icon = "\u274c"
            elif text.lower().startswith(('warn','caution','careful','cuidado')):
                icon = "\u26a0\ufe0f"
            elif text.lower().startswith(('note','info','note','nota')):
                icon = "\U0001f4cc"
            elif text.lower().startswith(('idea','tip','suggestion','sugest')):
                icon = "\U0001f4a1"
            _add_rich_paragraph(doc, f"{icon} {text}", style='List Bullet', font_name=font_pair["body"], font_size=11, color=colors["text"])
            continue
        
        # Regular paragraph
        text = _format_text(line, mode=fmt_mode)
        _add_rich_paragraph(doc, text, font_name=font_pair["body"], font_size=11, color=colors["text"])
    
    # Render remaining
    if in_card:
        add_card_box(doc, card_lines, colors["accent"], colors["light"], card_icon, card_title)
    if in_steps:
        if len(step_lines) >= 3:
            add_step_guide(doc, step_lines)
        else:
            for i, step_text in enumerate(step_lines, 1):
                _add_rich_paragraph(doc, f"{i}. {_format_text(step_text, mode=fmt_mode)}", font_name=font_pair["body"], font_size=11, color=colors["text"])
    if in_timeline:
        add_timeline(doc, timeline_data)
    if in_comparison:
        add_comparison_table(doc, comp_headers, comp_rows)
    if in_progress:
        for lbl, pct, clr in progress_data:
            add_progress_bar(doc, lbl, pct, clr)
    if in_code_block:
        for code_line in code_lines:
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(1)
            p.paragraph_format.space_after = Pt(1)
            run = p.add_run(code_line)
            run.font.name = "Consolas"
            run.font.size = Pt(9)
            run.font.color.rgb = hex_to_rgb(colors["text"])
    
    # --- Branded Footer ---
    footer = section.footer
    footer.is_linked_to_previous = False
    ft = footer.add_table(rows=1, cols=3, width=Inches(6.5))
    ft.alignment = WD_TABLE_ALIGNMENT.CENTER
    c1 = ft.rows[0].cells[0]; c1.width = Inches(2)
    r1 = c1.paragraphs[0]; r1.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run1 = r1.add_run("Edit Office Files")
    run1.font.size = Pt(8); run1.font.color.rgb = hex_to_rgb(colors["accent"])
    run1.font.name = font_pair["body"]
    c2 = ft.rows[0].cells[1]; c2.width = Inches(2.5)
    r2 = c2.paragraphs[0]; r2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = r2.add_run("Page ")
    run2.font.size = Pt(8); run2.font.color.rgb = hex_to_rgb(colors["accent"])
    run2.font.name = font_pair["body"]
    fldChar1 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
    run2._r.append(fldChar1)
    instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
    run2._r.append(instrText)
    fldChar2 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
    run2._r.append(fldChar2)
    c3 = ft.rows[0].cells[2]; c3.width = Inches(2)
    r3 = c3.paragraphs[0]; r3.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run3 = r3.add_run(datetime.datetime.now().strftime("%Y-%m-%d"))
    run3.font.size = Pt(8); run3.font.color.rgb = hex_to_rgb(colors["accent"])
    run3.font.name = font_pair["body"]
    
    # --- Save ---
    file_bytes = io.BytesIO()
    doc.save(file_bytes)
    file_bytes.seek(0)
    url, fname = await self._save_and_link(file_bytes.getvalue(), f"{title}.docx", __request__)
    if url:
        return f"Document created: [{fname}]({url})"
    return json.dumps({"error": "Could not save file"})


async def generate_slides(self, content: str, title: str = "Presentation", theme: str = "modern", raw_text: bool = False, __user__=None, __request__=None) -> str:
    """Generate professional PowerPoint slides with modern design.

        Args:
            content: Markdown content (headings become slides, bullets become content)
            title: Presentation title for the first slide
            theme: Visual theme - modern, light, dark, corporate, creative, minimal
            raw_text: If True, skip text formatting
        Returns:
            Markdown link to the generated file
        """
    fmt_mode = "preserve" if raw_text else "format"
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu, Cm
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
        import datetime

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # --- Color Themes ---
        themes = {
            "modern": {"bg": "1A1A2E", "accent": "E94560", "text": "FFFFFF", "subtitle": "B0B0B0", "card": "16213E"},
            "light": {"bg": "FFFFFF", "accent": "2563EB", "text": "1E293B", "subtitle": "64748B", "card": "F1F5F9"},
            "dark": {"bg": "0F172A", "accent": "38BDF8", "text": "F8FAFC", "subtitle": "94A3B8", "card": "1E293B"},
            "corporate": {"bg": "FFFFFF", "accent": "003366", "text": "1A1A1A", "subtitle": "666666", "card": "F0F4F8"},
            "creative": {"bg": "FFF8F0", "accent": "FF6B6B", "text": "2D3436", "subtitle": "636E72", "card": "FFEAA7"},
            "minimal": {"bg": "FAFAFA", "accent": "000000", "text": "333333", "subtitle": "999999", "card": "F0F0F0"},
        }
        colors = themes.get(theme, themes["modern"])

        def hex_to_rgb(h):
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

        def set_slide_bg(slide, color_hex):
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(color_hex)

        def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = hex_to_rgb(color or colors["text"])
            p.font.name = 'Calibri'
            p.alignment = alignment
            return txBox

        def add_accent_bar(slide, left, top, width, height, color_hex):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
            shape.line.fill.background()
            return shape

        # --- Slide 1: Title Slide ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, colors["bg"])

        add_accent_bar(slide, 0, 3.0, 13.333, 0.06, colors["accent"])
        add_text_box(slide, 1.5, 1.5, 10, 1.5, title, font_size=44, bold=True, color=colors["text"])
        add_text_box(slide, 1.5, 3.3, 10, 0.8, "Generated %s" % datetime.datetime.now().strftime("%B %d, %Y"), font_size=18, color=colors["subtitle"])

        # --- Process content into slides ---
        lines = content.split('\n')
        current_slide_lines = []
        slide_count = 1

        for line in lines:
            line = line.strip()
            if not line:
                continue

            if line.startswith('# ') or line.startswith('## '):
                if current_slide_lines:
                    _render_content_slide(prs, current_slide_lines, colors, hex_to_rgb, add_text_box, add_accent_bar, set_slide_bg, slide_count, fmt_mode=fmt_mode)
                    slide_count += 1
                current_slide_lines = [line]
            else:
                current_slide_lines.append(line)

        if current_slide_lines:
            _render_content_slide(prs, current_slide_lines, colors, hex_to_rgb, add_text_box, add_accent_bar, set_slide_bg, slide_count, fmt_mode=fmt_mode)

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        url, fname = await self._save_and_link(out.getvalue(), "%s.pptx" % title, __request__)
        if url:
            return "Presentation created: [%s](%s)" % (fname, url)
        return json.dumps({"error": "Could not save file"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


async def generate_spreadsheet(self, content: str, title: str = "Spreadsheet", theme: str = "professional", raw_text: bool = False, __user__=None, __request__=None) -> str:
    """Generate a professional Excel spreadsheet with modern styling.

        Args:
            content: CSV or tab-delimited data (first row = headers, rest = data)
            title: Spreadsheet title / filename
            theme: Visual theme - professional, modern, corporate, minimal, colorful, pastel
            raw_text: If True, skip text formatting
        Returns:
            Markdown link to the generated file
        """
    fmt_mode = "preserve" if raw_text else "format"
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
        from openpyxl.worksheet.table import Table, TableStyleInfo
        import io

        # --- Color Themes ---
        themes = {
            "professional": {"header": "1F4E79", "header_font": "FFFFFF", "alt_row": "F0F5FA", "border": "CCCCCC", "accent": "2E75B6"},
            "modern": {"header": "2D3436", "header_font": "FFFFFF", "alt_row": "F5F5F5", "border": "DFE6E9", "accent": "6C5CE7"},
            "corporate": {"header": "003366", "header_font": "FFFFFF", "alt_row": "E8EEF4", "border": "B0C4DE", "accent": "CC0000"},
            "minimal": {"header": "333333", "header_font": "FFFFFF", "alt_row": "F8F8F8", "border": "DDDDDD", "accent": "666666"},
            "colorful": {"header": "E17055", "header_font": "FFFFFF", "alt_row": "FFF3E0", "border": "FABEAB", "accent": "FDCB6E"},
            "pastel": {"header": "6C5CE7", "header_font": "FFFFFF", "alt_row": "F3E5F5", "border": "CE93D8", "accent": "A29BFE"},
        }
        colors = themes.get(theme, themes["professional"])

        header_fill = PatternFill(start_color=colors["header"], end_color=colors["header"], fill_type="solid")
        header_font = Font(name='Calibri', size=11, bold=True, color=colors["header_font"])
        alt_fill = PatternFill(start_color=colors["alt_row"], end_color=colors["alt_row"], fill_type="solid")
        body_font = Font(name='Calibri', size=11)
        thin_border = Border(
            left=Side(style='thin', color=colors["border"]),
            right=Side(style='thin', color=colors["border"]),
            top=Side(style='thin', color=colors["border"]),
            bottom=Side(style='thin', color=colors["border"])
        )

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = title[:31]

        # Parse CSV or tab-delimited content
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        data = []
        for line in lines:
            sep = '\t' if '\t' in line else ','
            cells = [c.strip().strip('"') for c in line.split(sep)]
            data.append(cells)

        if not data:
            return json.dumps({"error": "No data provided"})

        # Write data
        for i, row in enumerate(data):
            if i == 0:
                ws.append(row)  # Headers: preserve original case
            else:
                ws.append([_format_text(c, mode=fmt_mode) if isinstance(c, str) else c for c in row])

        # Style header row
        for cell in ws[1]:
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center')
            cell.border = thin_border

        # Zebra striping for data rows
        for i, row in enumerate(ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=ws.max_column), start=2):
            if i % 2 == 0:
                for cell in row:
                    cell.fill = alt_fill
                    cell.border = thin_border
                    cell.font = body_font
            else:
                for cell in row:
                    cell.border = thin_border
                    cell.font = body_font

        # Auto-fit columns
        for col in ws.columns:
            max_len = 0
            col_letter = get_column_letter(col[0].column)
            for cell in col:
                try:
                    max_len = max(max_len, len(str(cell.value or '')))
                except:
                    pass
            ws.column_dimensions[col_letter].width = min(max_len + 3, 40)

        # Freeze top row
        ws.freeze_panes = 'A2'

        # Add auto-filter
        if ws.max_row > 1:
            ws.auto_filter.ref = "A1:%s%d" % (get_column_letter(ws.max_column), ws.max_row)

        # Add Excel Table if data exists
        if ws.max_row > 1 and ws.max_column > 0:
            try:
                tab = Table(
                    displayName="Table_" + ws.title.replace(' ', '')[:20],
                    ref="A1:%s%d" % (get_column_letter(ws.max_column), ws.max_row)
                )
                style = TableStyleInfo(
                    name="TableStyleMedium2",
                    showFirstColumn=False,
                    showLastColumn=False,
                    showRowStripes=True,
                    showColumnStripes=False
                )
                tab.tableStyleInfo = style
                ws.add_table(tab)
            except Exception:
                pass

        out = io.BytesIO()
        wb.save(out)
        out.seek(0)
        url, fname = await self._save_and_link(out.getvalue(), "%s.xlsx" % title, __request__)
        if url:
            return "Spreadsheet created: [%s](%s)" % (fname, url)
        return json.dumps({"error": "Could not save file"})
    except Exception as e:
        return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


async def add_chart(self, file_id: str, chart_type: str = "bar", title: str = "Chart", __user__=None, __request__=None) -> str:
    """Add a chart to an Excel spreadsheet. chart_type: bar, line, pie, scatter."""
    from openpyxl import load_workbook
    from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference
    from openpyxl.utils import get_column_letter
    import io
    
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes:
        return json.dumps({"error": f"File not found: {file_id}"})
    
    wb = load_workbook(io.BytesIO(file_bytes))
    ws = wb.active
    
    chart_types = {"bar": BarChart, "line": LineChart, "pie": PieChart, "scatter": ScatterChart}
    chart_class = chart_types.get(chart_type, BarChart)
    chart = chart_class()
    chart.title = title
    chart.style = 10
    
    if ws.max_row > 1:
        data = Reference(ws, min_col=1, min_row=1, max_row=ws.max_row, max_col=ws.max_column)
        chart.add_data(data, titles_from_data=True)
    
    chart_col = get_column_letter(ws.max_column + 2)
    ws.add_chart(chart, f"{chart_col}1")
    
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url:
        return f"Chart added: [{fname}]({url})"
    return json.dumps({"error": "Could not save file"})

# --- v3.2.0: Watermark ---


async def add_watermark(self, file_id: str, text: str = "DRAFT", __user__=None, __request__=None) -> str:
    """Add a diagonal watermark to a DOCX or PDF file."""
    import io
    
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes:
        return json.dumps({"error": f"File not found: {file_id}"})
    
    if ftype == "docx":
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document(io.BytesIO(file_bytes))
        for section in doc.sections:
            header = section.header
            header.is_linked_to_previous = False
            p = header.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(72)
            run.font.color.rgb = RGBColor(128, 128, 128)
        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
    elif ftype == "pdf":
        try:
            import fitz
            pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
            for page in pdf_doc:
                rect = page.rect
                page.insert_text((rect.width/2-100, rect.height/2), text, fontsize=72, color=(0.5,0.5,0.5), alpha=0.1, rotate=45)
            buf = io.BytesIO()
            pdf_doc.save(buf)
            pdf_doc.close()
            buf.seek(0)
        except ImportError:
            return json.dumps({"error": "PyMuPDF not installed. Install with: pip install PyMuPDF"})
    else:
        return json.dumps({"error": f"Watermark not supported for {ftype}. Use DOCX or PDF."})
    
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url:
        return f"Watermark added: [{fname}]({url})"
    return json.dumps({"error": "Could not save file"})

# --- v3.2.0: File Preview ---


async def add_alt_text(self, file_id: str, slide_num: int = 1, shape_index: int = 0, alt_text: str = "", __user__=None, __request__=None) -> str:
    """Add alt text to an image in a PowerPoint slide."""
    import io
    
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes:
        return json.dumps({"error": f"File not found: {file_id}"})
    
    if ftype != "pptx":
        return json.dumps({"error": "Alt text only supported for PPTX files"})
    
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    if slide_num > len(prs.slides):
        return json.dumps({"error": f"Slide {slide_num} not found. Has {len(prs.slides)} slides."})
    
    slide = prs.slides[slide_num - 1]
    pic_count = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:
            if pic_count == shape_index:
                shape.alt_text = alt_text
                buf = io.BytesIO()
                prs.save(buf)
                buf.seek(0)
                url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
                if url:
                    return f"Alt text added: [{fname}]({url})"
                return json.dumps({"error": "Could not save file"})
            pic_count += 1
    return json.dumps({"error": f"Image {shape_index} not found. Found {pic_count} images."})

# --- v3.2.0: Progress Indicators ---


async def add_speaker_notes(self, file_id: str, slide_num: int, notes: str, __user__=None, __request__=None) -> str:
    """Add speaker notes to a PowerPoint slide."""
    import io
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype != "pptx": return json.dumps({"error": "Speaker notes only supported for PPTX"})
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    if slide_num < 1 or slide_num > len(prs.slides):
        return json.dumps({"error": "Slide " + str(slide_num) + " not found. Has " + str(len(prs.slides)) + " slides."})
    slide = prs.slides[slide_num - 1]
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url: return "Speaker notes added to slide " + str(slide_num) + ": [" + str(fname) + "](" + str(url) + ")"
    return json.dumps({"error": "Could not save file"})

# --- v3.4.0: Document Stats ---


async def add_qr_code(self, file_id: str, data: str, __user__=None, __request__=None) -> str:
    """Add a QR code to a DOCX or PPTX file."""
    import io
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    try:
        import qrcode
        from PIL import Image
    except ImportError:
        return json.dumps({"error": "qrcode or Pillow not installed. pip install qrcode Pillow"})
    qr = qrcode.QRCode(version=1, box_size=10, border=2)
    qr.add_data(data); qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    qr_buf = io.BytesIO(); img.save(qr_buf, format='PNG'); qr_buf.seek(0)
    if ftype == "docx":
        from docx import Document
        from docx.shared import Inches
        doc = Document(io.BytesIO(file_bytes))
        doc.add_picture(qr_buf, width=Inches(2))
        buf = io.BytesIO(); doc.save(buf); buf.seek(0)
    elif ftype == "pptx":
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation(io.BytesIO(file_bytes))
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(qr_buf, Inches(4), Inches(2), Inches(2))
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    else: return json.dumps({"error": "QR codes only supported for DOCX and PPTX"})
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url: return "QR code added: [" + str(fname) + "](" + str(url) + ")"
    return json.dumps({"error": "Could not save file"})

# --- v3.4.0: Bulk Folder Ops ---


async def add_data_validation(self, file_id: str, sheet: str = "", col: str = "A", validation_type: str = "list", values: str = "", __user__=None, __request__=None) -> str:
    """Add data validation to an Excel column. Types: list, whole, decimal, date. values: comma-separated for list, or 'min,max' for numeric."""
    import io
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype not in ("xlsx","xls"): return json.dumps({"error": "Data validation only for Excel files"})
    from openpyxl import load_workbook
    from openpyxl.worksheet.datavalidation import DataValidation
    wb = load_workbook(io.BytesIO(file_bytes))
    ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
    col_letter = col.upper().replace("COL","").replace("COLUMN","").strip()
    if not col_letter.isalpha(): col_letter = "A"
    max_row = ws.max_row or 100
    cell_range = col_letter + "2:" + col_letter + str(max_row)
    if validation_type == "list":
        vals = [v.strip() for v in values.split(",") if v.strip()]
        formula = '"' + ",".join(vals) + '"'
        dv = DataValidation(type="list", formula1=formula, allow_blank=True)
    elif validation_type == "whole":
        parts = values.split(",")
        mn = int(parts[0]) if parts else 0
        mx = int(parts[1]) if len(parts) > 1 else 100
        dv = DataValidation(type="whole", operator="between", formula1=str(mn), formula2=str(mx))
    elif validation_type == "decimal":
        parts = values.split(",")
        mn = float(parts[0]) if parts else 0
        mx = float(parts[1]) if len(parts) > 1 else 100
        dv = DataValidation(type="decimal", operator="between", formula1=str(mn), formula2=str(mx))
    elif validation_type == "date":
        dv = DataValidation(type="date", operator="greaterThan", formula1="2000-01-01")
    else: return json.dumps({"error": "Unknown type: " + str(validation_type)})
    dv.error = "Invalid value"
    dv.errorTitle = "Validation Error"
    ws.add_data_validation(dv)
    dv.add(cell_range)
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url: return "Data validation added to column " + str(col_letter) + ": [" + str(fname) + "](" + str(url) + ")"
    return json.dumps({"error": "Could not save file"})

# --- v3.4.0: Named Ranges ---


async def add_named_range(self, file_id: str, name: str, range_str: str = "", __user__=None, __request__=None) -> str:
    """Define a named range in Excel. range_str: 'A1:B10' or auto-detected from active sheet."""
    import io
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype not in ("xlsx","xls"): return json.dumps({"error": "Named ranges only for Excel"})
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    from openpyxl.workbook.defined_name import DefinedName
    wb = load_workbook(io.BytesIO(file_bytes))
    ws = wb.active
    if not range_str:
        range_str = "A1:" + get_column_letter(ws.max_column) + str(ws.max_row)
    dn = DefinedName(name, attr_text=ws.title + "!$" + range_str.replace(":", ":$"))
    wb.defined_names.add(dn)
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url: return "Named range '" + str(name) + "' = " + str(range_str) + ": [" + str(fname) + "](" + str(url) + ")"
    return json.dumps({"error": "Could not save file"})

# --- v3.4.0: Slide Transitions ---


async def add_slide_transitions(self, file_id: str, transition_type: str = "fade", duration: float = 0.5, __user__=None, __request__=None) -> str:
    """Add transitions to all slides in a PPTX. Types: fade, push, wipe, split, random."""
    import io
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype != "pptx": return json.dumps({"error": "Transitions only for PPTX"})
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation(io.BytesIO(file_bytes))
    transitions = {"fade": 0, "push": 1, "wipe": 2, "split": 3, "random": 4}
    ttype = transitions.get(transition_type, 0)
    for slide in prs.slides:
        try:
            from pptx.oxml.ns import qn
            trans_elem = slide._element.find(qn('p:transition'))
            if trans_elem is None:
                from lxml import etree
                trans_elem = etree.SubElement(slide._element, qn('p:transition'))
            if transition_type == "fade":
                etree.SubElement(trans_elem, qn('p:fade'))
            elif transition_type == "push":
                etree.SubElement(trans_elem, qn('p:push'))
            elif transition_type == "wipe":
                etree.SubElement(trans_elem, qn('p:wipe'))
            elif transition_type == "split":
                etree.SubElement(trans_elem, qn('p:split'))
            trans_elem.set('advTm', str(int(duration * 1000)))
        except Exception:
            pass
    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
    if url: return "Transitions added (" + str(transition_type) + "): [" + str(fname) + "](" + str(url) + ")"
    return json.dumps({"error": "Could not save file"})

# --- v3.4.0: Export to HTML ---


async def smart_fill(self, file_id: str, section: str, instruction: str, __user__=None, __request__=None) -> str:
    """Fill a document section using AI based on instructions. The LLM will generate content for the specified section."""
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
    elif ftype == "docx": content = self._read_docx(file_bytes, filename)
    elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
    elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
    else: return json.dumps({"error": "Unsupported format"})
    return f"**Smart Fill: {filename}**\n\nSection to fill: **{section}**\nInstructions: {instruction}\n\nCurrent document content:\n```\n{content[:3000]}\n```\n\nPlease generate the content for the '{section}' section based on the instructions and existing document context."


async def add_pivot_table(self, file_id: str, rows_field: str = "", cols_field: str = "", data_field: str = "", aggregate: str = "sum", __user__=None, __request__=None) -> str:
    """Create a pivot table in Excel. aggregate: sum, count, average, min, max."""
    import io
    from collections import defaultdict
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype not in ("xlsx","xls"): return json.dumps({"error": "Pivot tables only for Excel"})
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    wb = load_workbook(io.BytesIO(file_bytes))
    ws = wb.active
    headers = [str(c.value or '') for c in ws[1]]
    if not rows_field:
        return f"**Available fields for pivot:** {', '.join(headers[:10])}\n\nUse: add_pivot_table(file_id, rows_field='FieldName', data_field='FieldName')"
    # Find column indices
    try:
        row_idx = headers.index(rows_field)
    except ValueError:
        return json.dumps({"error": f"Field '{rows_field}' not found in headers: {headers[:20]}"})
    data_idx = None
    if data_field:
        try:
            data_idx = headers.index(data_field)
        except ValueError:
            return json.dumps({"error": f"Field '{data_field}' not found in headers: {headers[:20]}"})
    # Read and aggregate data
    agg_map = defaultdict(list)
    for row in ws.iter_rows(min_row=2, values_only=True):
        row_key = str(row[row_idx]) if row[row_idx] is not None else "(blank)"
        if data_idx is not None and row[data_idx] is not None:
            try:
                agg_map[row_key].append(float(row[data_idx]))
            except (ValueError, TypeError):
                agg_map[row_key].append(0.0)
        elif data_idx is None:
            agg_map[row_key].append(1)  # count mode
    # Compute aggregate
    pivot_ws = wb.create_sheet("Pivot")
    pivot_ws.title = "Pivot"
    agg_name = aggregate.lower()
    if agg_name == "count":
        pivot_ws.cell(row=1, column=1, value=rows_field)
        pivot_ws.cell(row=1, column=2, value="Count")
        for r, (key, vals) in enumerate(sorted(agg_map.items()), 2):
            pivot_ws.cell(row=r, column=1, value=key)
            pivot_ws.cell(row=r, column=2, value=len(vals))
    else:
        pivot_ws.cell(row=1, column=1, value=rows_field)
        pivot_ws.cell(row=1, column=2, value=f"{agg_name} of {data_field}")
        for r, (key, vals) in enumerate(sorted(agg_map.items()), 2):
            pivot_ws.cell(row=r, column=1, value=key)
            if agg_name == "sum":
                pivot_ws.cell(row=r, column=2, value=sum(vals))
            elif agg_name == "average" or agg_name == "avg":
                pivot_ws.cell(row=r, column=2, value=sum(vals) / len(vals) if vals else 0)
            elif agg_name == "min":
                pivot_ws.cell(row=r, column=2, value=min(vals) if vals else 0)
            elif agg_name == "max":
                pivot_ws.cell(row=r, column=2, value=max(vals) if vals else 0)
            else:
                pivot_ws.cell(row=r, column=2, value=sum(vals))
    result = f"Pivot table created in sheet 'Pivot' ({len(agg_map)} rows). Fields: rows={rows_field}, data={data_field}, aggregate={aggregate}"
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__)
    if url: return f"{result}: [{fname}]({url})"
    return json.dumps({"error": "Could not save file"})


async def sql_to_spreadsheet(self, query: str, output_filename: str = "query_results", __user__=None, __request__=None) -> str:
    """Execute a SQL query on the local SQLite database and export results to Excel."""
    import io
    conn2 = sqlite3.connect(_DB_PATH)
    try:
        conn2.row_factory = sqlite3.Row
        rows = conn2.execute(query).fetchall()
    except Exception as e:
        return json.dumps({"error": f"SQL error: {str(e)}"})
    finally:
        conn2.close()
    if not rows: return "Query returned no results."
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    wb = Workbook(); ws = wb.active
    headers = list(rows[0].keys())
    for j, h in enumerate(headers, 1):
        cell = ws.cell(row=1, column=j, value=h)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    for i, row in enumerate(rows, 2):
        for j, key in enumerate(headers, 1):
            ws.cell(row=i, column=j, value=row[key])
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), f"{output_filename}.xlsx", __request__)
    if url: return f"Query results ({len(rows)} rows): [{fname}]({url})"
    return json.dumps({"error": "Could not save file"})


async def fill_pdf_form(self, file_id: str, field_values: str, __user__=None, __request__=None) -> str:
    """Fill a PDF form with values. field_values: 'field1=value1,field2=value2'."""
    import io
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype != "pdf": return json.dumps({"error": "PDF form filling only for PDF files"})
    try:
        import fitz
        pdf = fitz.open(stream=file_bytes, filetype="pdf")
        pairs = {}
        for pair in field_values.split(","):
            if "=" in pair:
                k, v = pair.split("=", 1)
                pairs[k.strip()] = v.strip()
        filled = 0
        for page in pdf:
            for widget in page.widgets():
                if widget.field_name in pairs:
                    widget.field_value = pairs[widget.field_name]
                    widget.update()
                    filled += 1
        buf = io.BytesIO(); pdf.save(buf); pdf.close(); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__)
        if url: return f"Filled {filled} field(s): [{fname}]({url})"
        return json.dumps({"error": "Could not save file"})
    except ImportError:
        return json.dumps({"error": "PyMuPDF not installed"})
    except Exception as e:
        return json.dumps({"error": str(e)})


async def document_assembly(self, template_name: str, data_file_id: str, output_prefix: str = "assembled", __user__=None, __request__=None) -> str:
    """Assemble multiple documents from a template and data source."""
    # Load template content
    templates = json.loads(self.valves.templates or "{}")
    if template_name not in templates:
        return f"Template '{template_name}' not found. Available: {', '.join(templates.keys())}"
    template_content = templates[template_name]
    
    # Load data
    d_bytes, d_name, d_type = self._resolve_file(data_file_id)
    if not d_bytes:
        return json.dumps({"error": f"Data file not found: {data_file_id}"})
    
    # Parse data rows
    import csv as _csv, io
    rows = []
    if d_type in ("xlsx", "xls"):
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(d_bytes))
        ws = wb.active
        headers = [str(c.value or '') for c in ws[1]]
        for row in ws.iter_rows(min_row=2, values_only=True):
            rows.append(dict(zip(headers, [str(v or '') for v in row])))
    else:
        reader = _csv.DictReader(io.StringIO(d_bytes.decode('utf-8')))
        rows = list(reader)
    
    if not rows:
        return json.dumps({"error": "No data rows found"})
    
    # Generate one document per data row
    results = []
    for i, row in enumerate(rows):
        content = template_content
        for key, value in row.items():
            content = content.replace(f"{{{key}}}", str(value))
        result = await self.generate_document(content, f"{output_prefix}_{i+1}", __user__=__user__, __request__=__request__)
        results.append(result)
    
    return f"Assembled {len(results)} documents from template '{template_name}'."


async def conditional_format(self, file_id: str, rules: str, __user__=None, __request__=None) -> str:
    """Apply conditional formatting rules to Excel. rules: 'col:A,op:>,val:100,color:27AE60;col:B,op:<,val:0,color:E74C3C'."""
    import io
    file_bytes, filename, ftype = self._resolve_file(file_id)
    if not file_bytes: return json.dumps({"error": "File not found"})
    if ftype not in ("xlsx","xls"): return json.dumps({"error": "Conditional formatting only for Excel"})
    from openpyxl import load_workbook
    from openpyxl.formatting.rule import CellIsRule
    from openpyxl.styles import PatternFill
    wb = load_workbook(io.BytesIO(file_bytes))
    ws = wb.active
    applied = 0
    for rule_str in rules.split(";"):
        parts = {}
        for p in rule_str.split(","):
            if ":" in p:
                k, v = p.split(":", 1)
                parts[k.strip()] = v.strip()
        if "col" not in parts: continue
        col = parts["col"].upper()
        op_map = {">": "greaterThan", "<": "lessThan", ">=": "greaterThanOrEqual", "<=": "lessThanOrEqual", "=": "equal", "!=": "notEqual"}
        op = op_map.get(parts.get("op", ">"), "greaterThan")
        val = parts.get("val", "0")
        color = parts.get("color", "27AE60")
        fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
        max_row = ws.max_row or 100
        cell_range = f"{col}2:{col}{max_row}"
        ws.conditional_formatting.add(cell_range, CellIsRule(operator=op, formula=[val], fill=fill))
        applied += 1
    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__)
    if url: return f"Applied {applied} conditional formatting rule(s): [{fname}]({url})"
    return json.dumps({"error": "Could not save file"})

# === v3.6.0: Collaboration Features ===


__all__ = [
    "add_alt_text",
    "add_chart",
    "add_content",
    "add_data_validation",
    "add_named_range",
    "add_pivot_table",
    "add_qr_code",
    "add_slide_transitions",
    "add_speaker_notes",
    "add_watermark",
    "conditional_format",
    "create_file",
    "document_assembly",
    "fill_pdf_form",
    "generate_document",
    "generate_slides",
    "generate_spreadsheet",
    "smart_fill",
    "sql_to_spreadsheet",
]
