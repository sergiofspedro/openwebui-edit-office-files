    async def add_content(
        self,
        file_id: str,
        content: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Add new content to an Office file while preserving original formatting.

        For spreadsheets (xlsx): content is CSV text with rows to add.
        For documents (docx): content is text to append at the end.
        For presentations (pptx): each line defines a new slide. Use "---" as separator between slides.

        Args:
            file_id: File ID to edit
            content: Content to add (CSV for xlsx, text for docx/pptx)
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            # Detect type
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = row[0] if row else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            out = io.BytesIO()
            out_name = output_filename

            if file_type == "xlsx":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                ws = wb.active
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

                # Parse CSV content
                import csv as _csv_mod
                reader = _csv_mod.reader(io.StringIO(content))
                parsed_rows = []
                for csv_row in reader:
                    converted = []
                    for v in csv_row:
                        v = v.strip()
                        if v == '':
                            converted.append(None)
                        else:
                            try:
                                converted.append(int(v))
                            except ValueError:
                                try:
                                    converted.append(float(v))
                                except ValueError:
                                    if v.lower() == 'true':
                                        converted.append(True)
                                    elif v.lower() == 'false':
                                        converted.append(False)
                                    else:
                                        converted.append(v)
                    parsed_rows.append(converted)

                if not parsed_rows:
                    return json.dumps({"error": "No rows provided in CSV content"})

                # Get reference styles from last row
                ref = {}
                if ws.max_row and ws.max_row >= 1:
                    for cell in ws[ws.max_row]:
                        if cell.has_style:
                            ref[cell.column] = {
                                "font": copy(cell.font),
                                "fill": copy(cell.fill),
                                "border": copy(cell.border),
                                "alignment": copy(cell.alignment),
                                "number_format": cell.number_format,
                            }

                start = (ws.max_row or 0) + 1
                for i, rd in enumerate(parsed_rows):
                    for j, v in enumerate(rd, 1):
                        cell = ws.cell(row=start + i, column=j)
                        if j in ref:
                            try:
                                cell.font = copy(ref[j]["font"])
                                cell.fill = copy(ref[j]["fill"])
                                cell.border = copy(ref[j]["border"])
                                cell.alignment = copy(ref[j]["alignment"])
                                cell.number_format = ref[j]["number_format"]
                            except Exception:
                                pass
                        cell.value = v

                wb.save(out)
                wb.close()

            elif file_type == "xls":
                # Convert .xls to .xlsx, then apply same add logic
                file_data = _xls_to_xlsx(file_data)
                file_type = "xlsx"
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                ws = wb.active

                # Parse CSV content
                import csv as _csv_mod
                reader = _csv_mod.reader(io.StringIO(content))
                parsed_rows = []
                for csv_row in reader:
                    converted = []
                    for v in csv_row:
                        v = v.strip()
                        if v == '':
                            converted.append(None)
                        else:
                            try:
                                converted.append(int(v))
                            except ValueError:
                                try:
                                    converted.append(float(v))
                                except ValueError:
                                    if v.lower() == 'true':
                                        converted.append(True)
                                    elif v.lower() == 'false':
                                        converted.append(False)
                                    else:
                                        converted.append(v)
                    parsed_rows.append(converted)

                if not parsed_rows:
                    return json.dumps({"error": "No rows provided in CSV content"})

                ref = {}
                if ws.max_row and ws.max_row >= 1:
                    for cell in ws[ws.max_row]:
                        if cell.has_style:
                            ref[cell.column] = {
                                "font": copy(cell.font),
                                "fill": copy(cell.fill),
                                "border": copy(cell.border),
                                "alignment": copy(cell.alignment),
                                "number_format": cell.number_format,
                            }

                start = (ws.max_row or 0) + 1
                for i, rd in enumerate(parsed_rows):
                    for j, v in enumerate(rd, 1):
                        cell = ws.cell(row=start + i, column=j)
                        if j in ref:
                            try:
                                cell.font = copy(ref[j]["font"])
                                cell.fill = copy(ref[j]["fill"])
                                cell.border = copy(ref[j]["border"])
                                cell.alignment = copy(ref[j]["alignment"])
                                cell.number_format = ref[j]["number_format"]
                            except Exception:
                                pass
                        cell.value = v

                wb.save(out)
                wb.close()

            elif file_type == "docx":
                from docx import Document
                doc = Document(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.docx"

                # Append paragraphs, preserving last paragraph's style
                last_style = "Normal"
                if doc.paragraphs:
                    last_style = doc.paragraphs[-1].style.name if doc.paragraphs[-1].style else "Normal"

                for line in content.split("\n"):
                    doc.add_paragraph(line, style=last_style)

                out = io.BytesIO()
                doc.save(out)

            elif file_type == "pptx":
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.pptx"

                # Split content by "---" into slides
                slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
                blank_layout = prs.slide_layouts[6]  # Blank layout

                for spec in slide_specs:
                    spec = spec.strip()
                    if not spec:
                        continue
                    lines = spec.split("\n")
                    title = lines[0].strip()
                    body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                    slide = prs.slides.add_slide(blank_layout)
                    # Add title
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                    )
                    tf = txBox.text_frame
                    tf.text = title
                    p = tf.paragraphs[0]
                    p.font.size = Inches(0.6)
                    p.font.bold = True

                    # Add body text
                    if body:
                        txBox2 = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.text = body
                        for para in tf2.paragraphs:
                            para.font.size = Inches(0.3)

                out = io.BytesIO()
                prs.save(out)

            else:
                return json.dumps({"error": f"Unsupported type: {file_type}"})

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__)
            if url:
                return f"[{name}]({url})\n\nAdded content to {file_type.upper()} file, preserving original formatting."
            return json.dumps({"error": "Could not save file"})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # REPLACE TEXT
    # -----------------------------------------------------------------
    async def replace_text(
        self,
        file_id: str,
        find_text: str,
        replace_with: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Find and replace text in any Office file while preserving original formatting.

        Works on cell values in xlsx, paragraph text in docx, and shape text in pptx.

        Args:
            file_id: File ID to edit
            find_text: Text to find
            replace_with: Text to replace with
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            # Detect type
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = row[0] if row else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            out = io.BytesIO()
            out_name = output_filename
            count = 0

            if file_type == "xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

                for sn in wb.sheetnames:
                    ws = wb[sn]
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.value is None:
                                continue
                            if isinstance(cell.value, str) and find_text in cell.value:
                                cell.value = cell.value.replace(find_text, replace_with)
                                count += 1
                            elif not isinstance(cell.value, str):
                                sval = str(cell.value)
                                if find_text in sval:
                                    cell.value = replace_with
                                    count += 1

                wb.save(out)
                wb.close()

            elif file_type == "xls":
                # Convert .xls to .xlsx, then apply same replace logic
                file_data = _xls_to_xlsx(file_data)
                file_type = "xlsx"
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_data))

                for sn in wb.sheetnames:
                    ws = wb[sn]
                    for row in ws.iter_rows():
                        for cell in row:
                            if cell.value is None:
                                continue
                            if isinstance(cell.value, str) and find_text in cell.value:
                                cell.value = cell.value.replace(find_text, replace_with)
                                count += 1
                            elif not isinstance(cell.value, str):
                                sval = str(cell.value)
                                if find_text in sval:
                                    cell.value = replace_with
                                    count += 1

                wb.save(out)
                wb.close()

            elif file_type == "docx":
                from docx import Document
                doc = Document(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.docx"

                for para in doc.paragraphs:
                    if find_text in para.text:
                        # Preserve formatting: replace in runs
                        full_text = para.text
                        if find_text in full_text:
                            new_text = full_text.replace(find_text, replace_with)
                            # Clear all runs and set new text in first run
                            if para.runs:
                                para.runs[0].text = new_text
                                for run in para.runs[1:]:
                                    run.text = ""
                                count += 1
                            else:
                                para.text = new_text
                                count += 1

                # Also replace in tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if find_text in cell.text:
                                for para in cell.paragraphs:
                                    if find_text in para.text:
                                        if para.runs:
                                            para.runs[0].text = para.text.replace(find_text, replace_with)
                                            for run in para.runs[1:]:
                                                run.text = ""
                                        else:
                                            para.text = para.text.replace(find_text, replace_with)
                                        count += 1

                doc.save(out)

            elif file_type == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.pptx"

                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            for para in shape.text_frame.paragraphs:
                                if find_text in para.text:
                                    if para.runs:
                                        para.runs[0].text = para.text.replace(find_text, replace_with)
                                        for run in para.runs[1:]:
                                            run.text = ""
                                    else:
                                        para.text = para.text.replace(find_text, replace_with)
                                    count += 1

                prs.save(out)

            else:
                return json.dumps({"error": f"Unsupported type: {file_type}"})

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__)
            if url:
                return f"[{name}]({url})\n\nReplaced '{find_text}' with '{replace_with}' in {count} place(s), preserving all formatting."
            return json.dumps({"error": "Could not save file"})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # UPDATE CELLS
    # -----------------------------------------------------------------
    async def update_cells(
        self,
        file_id: str,
        cells: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Update individual cells in an XLSX file. Pass a JSON array of cell updates.

        Each update: {"cell": "A1", "value": "new value"} or {"cell": "B2", "value": 42, "sheet": "Sheet1"}
        If sheet is omitted, uses the active sheet.

        Args:
            file_id: The file ID of the XLSX file to edit
            cells: JSON string - array of {cell, value[, sheet]} objects
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            try:
                conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row2 = conn2.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn2.close()
                filename = row2[0] if row2 else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            if file_type not in ("xlsx", "xls"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx/xls supported."})

            import openpyxl
            from openpyxl.utils import column_index_from_string
            import re as _re_cell

            try:
                updates = json.loads(cells)
                if not isinstance(updates, list):
                    updates = [updates]
            except json.JSONDecodeError:
                return json.dumps({"error": "cells must be valid JSON (array of {cell, value} objects)"})

            out_name = output_filename or os.path.splitext(filename)[0] + "_updated.xlsx"
            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            count = 0
            for upd in updates:
                cell_ref = upd.get("cell", "")
                value = upd.get("value", "")
                sname = upd.get("sheet", "")
                ws = wb[sname] if sname else wb.active
                m = _re_cell.match(r"([A-Za-z]+)(\d+)", cell_ref)
                if not m:
                    continue
                col_str, row_str = m.group(1), m.group(2)
                col_idx = column_index_from_string(col_str.upper())
                row_idx = int(row_str)
                ws.cell(row=row_idx, column=col_idx).value = value
                count += 1

            out = io.BytesIO()
            wb.save(out)
            wb.close()
            out.seek(0)

            url, name = await self._save_and_link(out.read(), out_name, __request__)
            if url:
                return f"[{name}]({url})\n\nUpdated {count} cell(s) in {file_type.upper()} file."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # MODIFY ROWS (insert/delete)
    # -----------------------------------------------------------------
    async def modify_rows(
        self,
        file_id: str,
        action: str,
        row_number: int = 1,
        count: int = 1,
        sheet_name: str = "",
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Insert or delete rows in an XLSX file.

        Args:
            file_id: The file ID of the XLSX file to edit
            action: 'insert' or 'delete'
            row_number: Row number to start at (1-indexed)
            count: Number of rows to insert/delete (default 1)
            sheet_name: Optional sheet name (default: active sheet)
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            try:
                conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row2 = conn2.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn2.close()
                filename = row2[0] if row2 else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            if file_type not in ("xlsx", "xls"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx/xls supported."})

            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            ws = wb[sheet_name] if sheet_name else wb.active
            out_name = output_filename or os.path.splitext(filename)[0] + "_modified.xlsx"

            action_lower = action.strip().lower()
            if action_lower == "insert":
                ws.insert_rows(idx=row_number, amount=count)
                msg = f"Inserted {count} row(s) at row {row_number}"
            elif action_lower == "delete":
                ws.delete_rows(idx=row_number, amount=count)
                msg = f"Deleted {count} row(s) starting at row {row_number}"
            else:
                return json.dumps({"error": f"Unknown action '{action}'. Use 'insert' or 'delete'."})

            out = io.BytesIO()
            wb.save(out)
            wb.close()
            out.seek(0)

            url, name = await self._save_and_link(out.read(), out_name, __request__)
            if url:
                return f"[{name}]({url})\n\n{msg} in {file_type.upper()} file."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # PASSWORD PROTECT FILE
    # -----------------------------------------------------------------
    async def create_file(
        self,
        file_type: str,
        content: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Create a new Office file from scratch.

        For xlsx: content is CSV with headers on first line.
        For docx: content is plain text (one paragraph per line).
        For pptx: each line defines a slide. Use "---" as separator between slides.

        Args:
            file_type: 'xlsx', 'docx', or 'pptx'
            content: Content specification
            output_filename: Output filename
        """
        try:
            ftype = file_type.lower().replace(".", "")
            if ftype not in ("xlsx", "docx", "pptx"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Use xlsx, docx, or pptx."})

            out_name = output_filename or f"document.{ftype}"
            out = io.BytesIO()

            if ftype == "xlsx":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.Workbook()
                ws = wb.active

                import csv as _csv_mod
                reader = _csv_mod.reader(io.StringIO(content))
                rows = list(reader)

                if rows:
                    # First row = headers (styled)
                    for j, h in enumerate(rows[0], 1):
                        c = ws.cell(row=1, column=j, value=h.strip())
                        c.font = Font(bold=True, color="FFFFFF", size=11)
                        c.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                        c.alignment = Alignment(horizontal="center")

                    # Data rows
                    for i, rd in enumerate(rows[1:], 2):
                        for j, v in enumerate(rd, 1):
                            v = v.strip()
                            c = ws.cell(row=i, column=j)
                            if v == '':
                                c.value = None
                            else:
                                try:
                                    c.value = int(v)
                                except ValueError:
                                    try:
                                        c.value = float(v)
                                    except ValueError:
                                        if v.lower() == 'true':
                                            c.value = True
                                        elif v.lower() == 'false':
                                            c.value = False
                                        else:
                                            c.value = v
                            c.alignment = Alignment(
                                horizontal="center" if isinstance(c.value, (int, float)) else "left"
                            )

                    # Auto-fit columns
                    for col in ws.columns:
                        mx = max((len(str(c.value or "")) for c in col), default=5)
                        ws.column_dimensions[col[0].column_letter].width = min(mx + 3, 50)

                wb.save(out)
                wb.close()

            elif ftype == "docx":
                from docx import Document
                from docx.shared import Pt
                doc = Document()
                for line in content.split("\n"):
                    doc.add_paragraph(line)
                doc.save(out)

            elif ftype == "pptx":
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
                blank_layout = prs.slide_layouts[6]

                for spec in slide_specs:
                    spec = spec.strip()
                    if not spec:
                        continue
                    lines = spec.split("\n")
                    title = lines[0].strip()
                    body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                    slide = prs.slides.add_slide(blank_layout)
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                    )
                    tf = txBox.text_frame
                    tf.text = title
                    p = tf.paragraphs[0]
                    p.font.size = Inches(0.6)
                    p.font.bold = True

                    if body:
                        txBox2 = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.text = body
                        for para in tf2.paragraphs:
                            para.font.size = Inches(0.3)

                prs.save(out)

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__)
            if url:
                return f"[{name}]({url})\n\nCreated new {ftype.upper()} file."
            return json.dumps({"error": "Could not save file"})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})



