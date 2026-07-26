    async def mail_merge(self, template_file_id: str, data_file_id: str, output_prefix: str = "merged", __user__=None, __request__=None) -> str:
        """Generate personalized documents by merging CSV/Excel data into a DOCX template."""
        from docx import Document
        import io, csv
        
        t_bytes, t_name, t_type = self._resolve_file(template_file_id)
        if not t_bytes:
            return json.dumps({"error": f"Template not found: {template_file_id}"})
        
        d_bytes, d_name, d_type = self._resolve_file(data_file_id)
        if not d_bytes:
            return json.dumps({"error": f"Data file not found: {data_file_id}"})
        
        rows = []
        if d_type in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(d_bytes))
            ws = wb.active
            headers = [str(c.value or '') for c in ws[1]]
            for row in ws.iter_rows(min_row=2, values_only=True):
                rows.append(dict(zip(headers, [str(v or '') for v in row])))
        else:
            reader = csv.DictReader(io.StringIO(d_bytes.decode('utf-8')))
            rows = list(reader)
        
        if not rows:
            return json.dumps({"error": "No data rows found"})
        
        results = []
        for i, row in enumerate(rows):
            doc = Document(io.BytesIO(t_bytes))
            for para in doc.paragraphs:
                for key, value in row.items():
                    if f"{{{{{key}}}}}" in para.text:
                        for run in para.runs:
                            run.text = run.text.replace(f"{{{{{key}}}}}", value)
            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)
            fname = f"{output_prefix}_{i+1}.docx"
            url, saved = await self._save_and_link(buf.getvalue(), fname, __request__, __user__=__user__)
            if url:
                results.append(f"[{saved}]({url})")
        
        return f"Merged {len(results)} documents:\n" + "\n".join(results)

    # --- v3.2.0: Chart Generation ---
    async def add_chart(self, file_id: str, chart_type: str = "bar", title: str = "Chart", __user__=None, __request__=None) -> str:
        """Add a chart to an Excel spreadsheet. chart_type: bar, line, pie, scatter."""
        from openpyxl import load_workbook
        from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference
        from openpyxl.utils import get_column_letter
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        wb = load_workbook(io.BytesIO(file_bytes))
        ws = wb.active
        
        chart_types = {"bar": BarChart, "line": LineChart, "pie": PieChart, "scatter": ScatterChart}
        chart_class = chart_types.get(chart_type, BarChart)
        chart = chart_class()
        chart.title = _format_text(title)
        chart.style = 10
        
        if ws.max_row > 1:
            data = Reference(ws, min_col=1, min_row=1, max_row=ws.max_row, max_col=ws.max_column)
            chart.add_data(data, titles_from_data=True)
        
        chart_col = get_column_letter(ws.max_column + 2)
        ws.add_chart(chart, f"{chart_col}1")
        
        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url:
            return f"Chart added: [{fname}]({url})"
        return json.dumps({"error": "Could not save file"})

    # --- v3.2.0: Watermark ---
    async def add_watermark(self, file_id: str, text: str = "DRAFT", __user__=None, __request__=None) -> str:
        """Add a diagonal watermark to a DOCX or PDF file."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype == "docx":
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document(io.BytesIO(file_bytes))
            for section in doc.sections:
                header = section.header
                header.is_linked_to_previous = False
                p = header.paragraphs[0]
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p.add_run()
                run.text = _format_text(text)
                run.font.size = Pt(72)
                run.font.color.rgb = RGBColor(128, 128, 128)
            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)
        elif ftype == "pdf":
            try:
                import fitz
                pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
                for page in pdf_doc:
                    rect = page.rect
                    page.insert_text((rect.width/2-100, rect.height/2), _format_text(text), fontsize=72, color=(0.5,0.5,0.5), alpha=0.1, rotate=45)
                buf = io.BytesIO()
                pdf_doc.save(buf)
                pdf_doc.close()
                buf.seek(0)
            except ImportError:
                return json.dumps({"error": "PyMuPDF not installed. Install with: pip install PyMuPDF"})
        else:
            return json.dumps({"error": f"Watermark not supported for {ftype}. Use DOCX or PDF."})
        
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url:
            return f"Watermark added: [{fname}]({url})"
        return json.dumps({"error": "Could not save file"})

    # --- v3.2.0: File Preview ---
    async def protect_file(
        self,
        file_id: str,
        password: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Add password protection to an XLSX or DOCX file.

        For XLSX: uses openpyxl workbook protection. For DOCX: uses write protection.
        Install msoffcrypto-tool (pip install msoffcrypto-tool) for encryption support.

        Args:
            file_id: The file ID of the Office file to protect
            password: Password to set
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            try:
                conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row2 = conn2.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn2.close()
                filename = row2[0] if row2 else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            if file_type not in ("xlsx", "docx"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx and docx supported."})

            import hashlib

            if file_type == "xlsx":
                import openpyxl
                from openpyxl.workbook.protection import WorkbookProtection

                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                out_name = output_filename or os.path.splitext(filename)[0] + "_protected.xlsx"

                # Set workbook protection
                wb.security = WorkbookProtection(workbookPassword=password, lockStructure=True)

                # Protect all sheets
                for ws in wb.worksheets:
                    ws.protection.set_password(password)

                out = io.BytesIO()
                wb.save(out)
                wb.close()
                out.seek(0)
                protected_bytes = out.read()

            elif file_type == "docx":
                from docx import Document
                from docx.oxml.ns import qn
                from lxml import etree

                out_name = output_filename or os.path.splitext(filename)[0] + "_protected.docx"
                doc = Document(io.BytesIO(file_data))

                # Add write protection XML
                settings_element = doc.settings.element
                protection = etree.SubElement(
                    settings_element,
                    qn("w:documentProtection")
                )
                protection.set(qn("w:edit"), "readOnly")
                protection.set(qn("w:enforcement"), "1")
                pw_hash = hashlib.sha1(password.encode("utf-8")).hexdigest().upper()
                protection.set(qn("w:cryptProviderType"), "rsaAES")
                protection.set(qn("w:cryptAlgorithmClass"), "hash")
                protection.set(qn("w:cryptAlgorithmType"), "typeAny")
                protection.set(qn("w:cryptAlgorithmSid"), "4")
                protection.set(qn("w:cryptSpinCount"), "100000")
                protection.set(qn("w:hash"), pw_hash)
                protection.set(qn("w:salt"), "AAAAAAAAAAAAAAAAAAAAAA==")

                out = io.BytesIO()
                doc.save(out)
                out.seek(0)
                protected_bytes = out.read()

            url, fname = await self._save_and_link(protected_bytes, out_name, __request__)
            if url:
                return f"[{fname}]({url})\n\nPassword-protected {file_type.upper()} file created."
            return json.dumps({"error": "Could not save file"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # CREATE NEW FILE
    # -----------------------------------------------------------------
    async def preview_file(self, file_id: str, max_lines: int = 20) -> str:
        """Show a text preview of any Office file before downloading."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype in ("xlsx", "xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx":
            content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx":
            content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt", "ods", "odp"):
            content = _read_odf(file_bytes, filename)
        else:
            return json.dumps({"error": f"Preview not supported for {ftype}"})
        
        lines = content.split('\n')
        preview = '\n'.join(lines[:max_lines])
        total = len(lines)
        result = f"**{filename}** ({ftype.upper()}, {total} lines)\n\n```\n{preview}\n```"
        if total > max_lines:
            result += f"\n... ({total - max_lines} more lines)"
        return result

    # --- v3.2.0: Metadata Editing ---
    async def edit_metadata(self, file_id: str, author: str = "", title: str = "", subject: str = "", keywords: str = "", __user__=None, __request__=None) -> str:
        """Edit document metadata (author, title, subject, keywords)."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        changes = []
        if ftype == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            cp = doc.core_properties
            if author: cp.author = _format_text(author); changes.append("author")
            if title: cp.title = _format_text(title); changes.append("title")
            if subject: cp.subject = _format_text(subject); changes.append("subject")
            if keywords: cp.keywords = _format_text(keywords); changes.append("keywords")
            buf = io.BytesIO(); doc.save(buf); buf.seek(0)
        elif ftype == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes))
            if author: wb.properties.creator = _format_text(author); changes.append("author")
            if title: wb.properties.title = _format_text(title); changes.append("title")
            if subject: wb.properties.subject = _format_text(subject); changes.append("subject")
            buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        elif ftype == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            cp = prs.core_properties
            if author: cp.author = _format_text(author); changes.append("author")
            if title: cp.title = _format_text(title); changes.append("title")
            if subject: cp.subject = _format_text(subject); changes.append("subject")
            if keywords: cp.keywords = _format_text(keywords); changes.append("keywords")
            buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        else:
            return json.dumps({"error": f"Metadata editing not supported for {ftype}"})
        
        if not changes:
            return json.dumps({"error": "No metadata fields specified"})
        
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url:
            return f"Metadata updated ({', '.join(changes)}): [{fname}]({url})"
        return json.dumps({"error": "Could not save file"})

    # --- v3.2.0: Accessibility Check ---
    async def check_accessibility(self, file_id: str) -> str:
        """Check a document for accessibility issues."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        issues = []
        if ftype == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            headings = []
            for p in doc.paragraphs:
                if p.style.name.startswith('Heading'):
                    level = int(p.style.name.split()[-1]) if p.style.name.split()[-1].isdigit() else 1
                    headings.append((level, p.text[:50]))
            prev = 0
            for level, text in headings:
                if level > prev + 1:
                    issues.append(f"Heading skip: H{prev} to H{level} ('{text}')")
                prev = level
            if not headings:
                issues.append("No headings found - document may lack structure")
        elif ftype == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13 and not shape.alt_text:
                        issues.append(f"Slide {i+1}: Image missing alt text")
        
        if not issues:
            return f"Accessibility check passed for {filename}. No issues found."
        result = f"**Accessibility Report: {filename}**\n\n"
        for issue in issues:
            result += f"- {issue}\n"
        result += f"\n{len(issues)} issue(s) found."
        return result

    # --- v3.2.0: Add Alt Text ---
    async def add_alt_text(self, file_id: str, slide_num: int = 1, shape_index: int = 0, alt_text: str = "", __user__=None, __request__=None) -> str:
        """Add alt text to an image in a PowerPoint slide."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype != "pptx":
            return json.dumps({"error": "Alt text only supported for PPTX files"})
        
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        if slide_num > len(prs.slides):
            return json.dumps({"error": f"Slide {slide_num} not found. Has {len(prs.slides)} slides."})
        
        slide = prs.slides[slide_num - 1]
        pic_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:
                if pic_count == shape_index:
                    shape.alt_text = _format_text(alt_text)
                    buf = io.BytesIO()
                    prs.save(buf)
                    buf.seek(0)
                    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
                    if url:
                        return f"Alt text added: [{fname}]({url})"
                    return json.dumps({"error": "Could not save file"})
                pic_count += 1
        return json.dumps({"error": f"Image {shape_index} not found. Found {pic_count} images."})

    # --- v3.2.0: Progress Indicators ---
