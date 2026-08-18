"""
title: Edit Office Files
author: giofsp
author_url: https://github.com/sergiofspedro
description: Unified tool to read, edit, and create Office files (.xlsx, .xls, .docx, .pptx) preserving original formatting and styles. Supports markdown rendering in DOCX (headings, bold, italic, code, links). Detects highlights, bold, italic formatting. Detects legacy .doc and .ppt. Note: Track changes are not supported. For 2+ comments on one file, use add_comments (not repeated add_comment calls).
version: 4.0.2
requirements: openpyxl, python-docx, python-pptx, xlrd, odfpy, docx-revisions, lxml, PyMuPDF, Pillow, pytesseract, qrcode, google-api-python-client, google-auth
"""

import io
import json
import os
import platform
import re
import sqlite3
import sys
import traceback
from copy import copy
from typing import Any, Dict, Optional

from pydantic import BaseModel, Field

import base64 as _b64_mod


def _validate_outbound_url(url: str) -> Optional[str]:
    """Reject a URL that would let this plugin read local files or reach internal/cloud
    network targets (SSRF). Used by every function that fetches a caller-supplied URL
    (import_from_url, import_from_api, webhook_trigger) -- these are reachable via prompt
    injection in an uploaded document, not just direct user requests, so `file://` and
    loopback/link-local/private-network targets must be blocked unconditionally.

    Returns an error message string if the URL should be rejected, or None if it's OK.
    """
    import ipaddress
    import socket
    from urllib.parse import urlparse

    parsed = urlparse(url)
    if parsed.scheme not in ("http", "https"):
        return f"URL scheme {parsed.scheme!r} is not allowed. Only http/https URLs are supported."
    host = parsed.hostname
    if not host:
        return "URL has no hostname."
    try:
        infos = socket.getaddrinfo(host, None)
    except socket.gaierror as e:
        return f"Could not resolve hostname {host!r}: {e}"
    for info in infos:
        addr = info[4][0]
        try:
            ip = ipaddress.ip_address(addr)
        except ValueError:
            continue
        if ip.is_loopback or ip.is_link_local or ip.is_private or ip.is_reserved or ip.is_multicast:
            return f"URL resolves to a non-public address ({addr}) -- refusing to fetch internal/local network targets."
    return None


def _get_owui_data_dir() -> str:
    """Return the Open WebUI data directory for the current OS."""
    data_dir = os.environ.get("OPEN_WEBUI_DATA_DIR", "")
    if data_dir:
        return data_dir
    home = os.path.expanduser("~")
    system = platform.system()
    if system == "Windows":
        return os.path.join(os.environ.get("APPDATA", home), "open-webui", "data")
    if system == "Darwin":
        return os.path.join(home, "Library", "Application Support", "open-webui", "data")
    # Linux
    xdg_data_home = os.environ.get("XDG_DATA_HOME", "")
    if xdg_data_home:
        return os.path.join(xdg_data_home, "open-webui", "data")
    return os.path.join(home, ".open-webui", "data")

def _get_owui_uploads_dir() -> str:
    """Return the Open WebUI uploads directory for the current OS."""
    data_dir = os.environ.get("OPEN_WEBUI_DATA_DIR", "")
    if data_dir:
        return os.path.join(data_dir, "data", "uploads")
    home = os.path.expanduser("~")
    system = platform.system()
    if system == "Windows":
        return os.path.join(os.environ.get("APPDATA", home), "open-webui", "data", "uploads")
    if system == "Darwin":
        return os.path.join(home, "Library", "Application Support", "open-webui", "data", "uploads")
    # Linux
    xdg_data_home = os.environ.get("XDG_DATA_HOME", "")
    if xdg_data_home:
        return os.path.join(xdg_data_home, "open-webui", "data", "uploads")
    return os.path.join(home, ".open-webui", "data", "uploads")

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
_data_dir = os.environ.get("OPEN_WEBUI_DATA_DIR", "")
if _data_dir:
    _DB_PATH = os.path.join(_data_dir, "data", "webui.db")
else:
    _DB_PATH = os.path.join(_get_owui_data_dir(), "webui.db")

_data_dir = os.environ.get("OPEN_WEBUI_DATA_DIR", "")
if _data_dir:
    _UPLOAD_DIR = os.path.join(_data_dir, "data", "uploads")
else:
    _UPLOAD_DIR = _get_owui_uploads_dir()

_EXPORT_DIR = os.environ.get("OWUI_EXPORTS_DIR", os.path.join(os.path.expanduser("~"), "open-webui", "exports"))

# ---------------------------------------------------------------------------
# PPTX namespace constants (used by add_comment)
# ---------------------------------------------------------------------------
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
_CM_REL_TYPE = "http://schemas.microsoft.com/office/2016/09/relationships/commentsModern"
_CT_MODERN = "application/vnd.ms-office.presentation.commentsModern"
_CT_AUTHORS = "application/vnd.ms-office.presentation.commentsAuthors"

__all__ = [
    "_get_owui_data_dir", "_get_owui_uploads_dir",
    "_data_dir", "_DB_PATH", "_UPLOAD_DIR", "_EXPORT_DIR",
    "_P_NS", "_P14_NS", "_R_NS", "_PKG_REL_NS", "_CT_NS",
    "_CM_REL_TYPE", "_CT_MODERN", "_CT_AUTHORS",
]

def _resolve_file_path(file_id: str) -> Optional[str]:
    """Resolve an Open WebUI file UUID to an absolute disk path."""
    try:
        conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
        try:
            row = conn.execute(
                "SELECT path FROM file WHERE id = ?", (file_id,)
            ).fetchone()
        finally:
            conn.close()
    except Exception as exc:
        print(f"[office] DB lookup failed for {file_id}: {exc}", file=sys.stderr)
        return None

    if not row or not row[0]:
        # Fallback: try by filename
        try:
            conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
            try:
                row = conn.execute(
                    "SELECT path FROM file WHERE filename LIKE ?",
                    (f"%{file_id}%",),
                ).fetchone()
            finally:
                conn.close()
        except Exception:
            pass

    if not row or not row[0]:
        print(f"[office] No path for file_id {file_id}", file=sys.stderr)
        return None

    path = row[0]
    if os.path.isfile(path):
        return path

    # Fallback: uploads directory
    candidate = os.path.join(_UPLOAD_DIR, os.path.basename(path))
    if os.path.isfile(candidate):
        return candidate

    # Last resort: UUID prefix match in uploads
    prefix = file_id.split("-")[0] if "-" in file_id else file_id[:8]
    if os.path.isdir(_UPLOAD_DIR):
        for name in os.listdir(_UPLOAD_DIR):
            if name.startswith(prefix):
                candidate = os.path.join(_UPLOAD_DIR, name)
                if os.path.isfile(candidate):
                    return candidate

    print(f"[office] File not found on disk: {path}", file=sys.stderr)
    return None


def _read_file_bytes(file_id: str) -> Optional[bytes]:
    """Return raw bytes for an Open WebUI file."""
    path = _resolve_file_path(file_id)
    if path is None:
        return None
    # Path traversal guard: ensure resolved path stays inside allowed directories
    try:
        _abs = os.path.realpath(path)
        _allowed = False
        for _base in (_UPLOAD_DIR, _EXPORT_DIR,
                       os.path.join(os.environ.get("OPEN_WEBUI_DATA_DIR", ""), "data"),
                       _get_owui_data_dir()):
            if _base and _abs.startswith(os.path.realpath(_base) + os.sep):
                _allowed = True
                break
        if not _allowed:
            print(f"[office] Path traversal blocked: {path}", file=sys.stderr)
            return None
    except Exception as exc:
        print(f"[office] Path traversal check failed for {path}: {exc}", file=sys.stderr)
        return None
    try:
        # A byte-size cap before reading, not after -- without this, an oversized file
        # (accidental or malicious upload) gets fully buffered into memory here, then
        # copied again into io.BytesIO() by every caller, then expanded further by
        # openpyxl/python-docx's object graph -- multiple multiples of the file's size,
        # for a file nobody meant to process this way.
        _MAX_FILE_BYTES = 200 * 1024 * 1024  # 200 MB
        size = os.path.getsize(path)
        if size > _MAX_FILE_BYTES:
            print(f"[office] File too large ({size} bytes > {_MAX_FILE_BYTES}): {path}", file=sys.stderr)
            return None
        with open(path, "rb") as fh:
            return fh.read()
    except Exception as exc:
        print(f"[office] Read failed for {path}: {exc}", file=sys.stderr)
        return None


_ERROR_TRANSLATIONS = {
    "en": {"file_not_found": "File not found", "could_not_save": "Could not save file", "unsupported": "Unsupported format"},
    "pt": {"file_not_found": "Ficheiro nao encontrado", "could_not_save": "Nao foi possivel guardar", "unsupported": "Formato nao suportado"},
    "es": {"file_not_found": "Archivo no encontrado", "could_not_save": "No se pudo guardar", "unsupported": "Formato no soportado"},
    "fr": {"file_not_found": "Fichier introuvable", "could_not_save": "Impossible d'enregistrer", "unsupported": "Format non pris en charge"},
    "de": {"file_not_found": "Datei nicht gefunden", "could_not_save": "Konnte nicht gespeichert werden", "unsupported": "Nicht unterstutztes Format"},
}


def _ensure_ext(name: str, ext: str) -> str:
    """Ensure `name` ends with `.ext`, without doubling it if already present."""
    ext = ext.lower().lstrip(".")
    return name if name.lower().endswith(f".{ext}") else f"{name}.{ext}"


def _detect_type(filename: str) -> str:
    """Detect file type from extension."""
    ext = os.path.splitext(filename)[1].lower()
    if ext in (".xlsx",):
        return "xlsx"
    if ext in (".xls",):
        return "xls"
    if ext in (".docx",):
        return "docx"
    if ext in (".pptx",):
        return "pptx"
    if ext in (".ppt",):
        return "ppt"
    if ext in (".odt",):
        return "odt"
    if ext in (".ods",):
        return "ods"
    if ext in (".odp",):
        return "odp"
    if ext in (".csv",):
        return "csv"
    if ext in (".pdf",):
        return "pdf"
    return "unknown"


def _cell_value(cell):
    """Extract a JSON-safe value from an openpyxl cell."""
    import datetime

    v = cell.value
    if v is None:
        return None
    if isinstance(v, (datetime.datetime, datetime.date)):
        return v.isoformat()
    if isinstance(v, (int, float, str, bool)):
        return v
    return str(v)


def _xls_to_xlsx(xls_data: bytes) -> bytes:
    """Convert .xls bytes to .xlsx bytes using xlrd + openpyxl.

    Returns an in-memory .xlsx workbook as bytes.
    """
    import xlrd
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment

    xls_book = xlrd.open_workbook(file_contents=xls_data)
    xlsx_wb = openpyxl.Workbook()
    # Remove the default sheet; we'll add one per xls sheet
    xlsx_wb.remove(xlsx_wb.active)

    for sheet_idx in range(xls_book.nsheets):
        xls_sheet = xls_book.sheet_by_index(sheet_idx)
        ws = xlsx_wb.create_sheet(title=xls_sheet.name[:31])  # Excel 31-char limit

        for rx in range(xls_sheet.nrows):
            for cx in range(xls_sheet.ncols):
                cell = xls_sheet.cell(rx, cx)
                value = cell.value

                # xlrd date handling: if cell type is XL_CELL_DATE, convert
                if cell.ctype == xlrd.XL_CELL_DATE:
                    try:
                        import datetime as _dt
                        value = _dt.datetime(*xlrd.xldate_as_tuple(value, xls_book.datemode))
                    except Exception:
                        pass
                elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                    value = bool(value)

                ws.cell(row=rx + 1, column=cx + 1, value=value)

        # Auto-fit column widths (rough estimate)
        for col_cells in ws.columns:
            max_len = 0
            for cell in col_cells:
                try:
                    max_len = max(max_len, len(str(cell.value or "")))
                except Exception:
                    pass
            letter = openpyxl.utils.get_column_letter(col_cells[0].column)
            ws.column_dimensions[letter].width = min(max_len + 2, 50)

    out = io.BytesIO()
    xlsx_wb.save(out)
    xlsx_wb.close()
    xls_book.release_resources()
    out.seek(0)
    return out.read()


def _format_text(text: str, mode: str = "format") -> str:
    """Apply consistent text formatting: normalize dashes, sentence case, preserve acronyms.

    Args:
        text: The text to format.
        mode: "format" (default) applies sentence case + acronym preservation + dash normalization.
              "preserve" returns text unchanged.
    """
    if mode == "preserve":
        return text
    if not isinstance(text, str) or not text or text.startswith('='):
        return text
    
    # 1. Normalize dashes
    text = text.replace('\u2014', '-').replace('\u2015', '-').replace('\u2013', '-')
    
    # 2. Split into sentences and apply sentence case
    acronyms = {'API', 'PDF', 'HTML', 'CSS', 'JSON', 'SQL', 'AI', 'UI', 'UX', 'ID', 'URL', 'HTTP', 'HTTPS', 'FTP', 'SSH', 'DNS', 'IP', 'TCP', 'UDP', 'SSL', 'TLS', 'REST', 'CRUD', 'YAML', 'XML', 'SVG', 'PNG', 'JPG', 'GIF', 'CSV', 'DOCX', 'PPTX', 'XLSX', 'DOC', 'PPT', 'XLS', 'ISO', 'RGB', 'CMYK', 'PHP', 'NPM', 'YARN', 'CLI', 'GUI', 'IDE', 'SDK', 'JDK', 'JRE', 'JVM', 'DB', 'SQLite', 'MySQL', 'NoSQL', 'OAuth', 'JWT', 'CORS', 'MVC', 'MVP', 'MVVM', 'SPA', 'PWA', 'SSR', 'CSR', 'SSG', 'ISR', 'CDN', 'DHCP', 'NAT', 'VPN', 'LAN', 'WAN', 'MAC', 'BIOS', 'UEFI', 'GPT', 'MBR', 'RAID', 'NAS', 'SAN', 'AWS', 'GCP', 'Azure', 'SaaS', 'PaaS', 'IaaS', 'FaaS', 'CaaS', 'K8s', 'Docker', 'Kubernetes'}
    
    import re as _re
    
    # Protect acronyms with placeholders (longest first to avoid partial matches)
    acronym_map = {}
    for i, acro in enumerate(sorted(acronyms, key=len, reverse=True)):
        placeholder = f'\x00{i:04d}\x00'
        pattern = _re.compile(r'(?<![a-zA-Z])' + _re.escape(acro) + r'(?![a-zA-Z])', _re.IGNORECASE)
        text = pattern.sub(placeholder, text)
        acronym_map[placeholder] = acro
    
    # Smart sentence splitting: avoid splitting on abbreviations
    # Heuristic 1: A word containing an internal period (e.g., "U.S.A.", "e.g.", "i.e.", "Ph.D.")
    #              is an abbreviation — don't split after it (handles U.S.A., U.K., E.U., a.k.a., w.r.t., etc.)
    # Mini-list: short common abbreviations without internal periods (Dr., Mr., St., etc.)
    # Note: month abbreviations (Jan, Feb, etc.) are included since they're commonly followed by
    # a number (e.g., "Jan 15") which would pass the sentence-split check.
    _MINI_ABBREV = {'Dr', 'Mr', 'Mrs', 'Ms', 'Prof', 'Sr', 'Jr', 'St', 'vs', 'etc',
                    'Jan', 'Feb', 'Mar', 'Apr', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec'}

    # Python's re requires fixed-width lookbehind, so a variable-length alternation of
    # abbreviations (Dr|Mrs|Prof|...) can't live inside (?<!...). Split on the fixed-width
    # boundary first, then merge a split back together whenever the preceding fragment ends
    # with a mini-abbreviation or a word containing an internal period (U.S.A., e.g., i.e.).
    _boundary_re = _re.compile(r'(?<=[.!?])\s+')
    _last_token_re = _re.compile(r'(\S+)$')
    _internal_period_re = _re.compile(r'[A-Za-z]\.[A-Za-z]')

    def _ends_with_abbrev(fragment):
        m = _last_token_re.search(fragment)
        if not m:
            return False
        token = m.group(1)
        if token.rstrip('.!?') in _MINI_ABBREV:
            return True
        return bool(_internal_period_re.search(token))

    raw_parts = _boundary_re.split(text)
    sentences = []
    buf = ""
    for part in raw_parts:
        if buf and _ends_with_abbrev(buf):
            buf += " " + part
        else:
            if buf:
                sentences.append(buf)
            buf = part
    if buf:
        sentences.append(buf)

    # Sentence case with intentional capitalization preservation
    # Rule: words with ANY uppercase letter are preserved as-is (proper nouns,
    # acronyms not in the list, product names, the pronoun "I", etc.).
    # Words that are entirely lowercase get sentence-cased.
    formatted_sentences = []
    for s in sentences:
        s = s.strip()
        if not s:
            continue
        
        # Split into words while preserving spacing
        words = s.split(' ')
        formatted_words = []
        
        for word in words:
            if not word:
                formatted_words.append(word)
                continue
            
            # Skip placeholder-protected acronyms (they start with \x00)
            if '\x00' in word:
                formatted_words.append(word)
                continue
            
            # Check if word has intentional capitalization
            # (any uppercase letter means the user intentionally capitalized it)
            has_uppercase = any(c.isupper() for c in word)
            
            if has_uppercase:
                # Preserve original capitalization
                formatted_words.append(word)
            else:
                # Safe to lowercase
                formatted_words.append(word.lower())
        
        s = ' '.join(formatted_words)
        
        # Capitalize first letter of the first word in the sentence
        if s and s[0].isalpha():
            s = s[0].upper() + s[1:]
        
        formatted_sentences.append(s)
    
    text = ' '.join(formatted_sentences)
    
    # Restore acronyms
    for placeholder, acro in acronym_map.items():
        text = text.replace(placeholder, acro)
    
    return text


def _parse_inline_md(text: str) -> list[tuple[str, dict]]:
    """Parse inline markdown in *text* and return segments with formatting metadata.

    Returns a list of ``(plain_text, formatting_dict)`` tuples. The formatting dict
    may contain the keys ``bold``, ``italic``, ``code``, or ``link``.

    Parsing order (priority, first matched wins):
        1. Links — ``[text](url)``
        2. Bold — ``**text**``
        3. Italic — ``*text*``  (conflicts with ``* `` bullet lists are avoided)
        4. Inline code — `` `code` ``
    """
    import re as _re

    if not isinstance(text, str):
        text = str(text) if text is not None else ""

    if text == "":
        return [("", {})]

    # Regexes — link first to prevent its brackets from being consumed by bold/italic
    _link_re = _re.compile(r'\[(.+?)\]\(([^()]*(?:\([^()]*\)[^()]*)*)\)')
    _bold_re = _re.compile(r'\*\*(.+?)\*\*')
    _italic_re = _re.compile(r'(?<!\*)\*(?!\s)(.+?)(?<!\s)\*(?!\*)')
    _code_re = _re.compile(r'`(.+?)`')

    segments: list[tuple[str, dict]] = []

    # Work on a mutable list of (start, end, kind, data) markers.
    # Kind: 'link', 'bold', 'italic', 'code'
    markers: list[tuple[int, int, str, str]] = []

    # Links
    for m in _link_re.finditer(text):
        markers.append((m.start(), m.end(), 'link', (m.group(1), m.group(2))))  # (label, url)

    # Bold
    for m in _bold_re.finditer(text):
        # Only add if this range doesn't overlap with a higher-priority marker
        markers.append((m.start(), m.end(), 'bold', m.group(1)))

    # Italic
    for m in _italic_re.finditer(text):
        markers.append((m.start(), m.end(), 'italic', m.group(1)))

    # Code
    for m in _code_re.finditer(text):
        markers.append((m.start(), m.end(), 'code', m.group(1)))

    # Sort by start position, then by end position (longer span first wins ties)
    markers.sort(key=lambda x: (x[0], -x[1]))

    # Remove overlapping markers: keep the earliest-starting, and if tied,
    # the longer one; discard any later marker whose start < retained end.
    cleaned: list[tuple[int, int, str, str]] = []
    for mkr in markers:
        if not cleaned or mkr[0] >= cleaned[-1][1]:
            cleaned.append(mkr)
        else:
            # Overlap — keep only if this marker starts at the same position
            # as the last kept one but is longer (shouldn't happen with our
            # sort, but be safe) or is a link that overlaps with something else.
            pass

    markers = cleaned

    pos = 0
    for start, end, kind, data in markers:
        # Plain text before this marker
        if start > pos:
            segments.append((text[pos:start], {}))
        # The formatted segment
        fmt: dict = {}
        if kind == 'link':
            label, url = data
            fmt['link'] = url
            text_part = label
        else:
            if kind == 'bold':
                fmt['bold'] = True
            elif kind == 'italic':
                fmt['italic'] = True
            elif kind == 'code':
                fmt['code'] = True
            text_part = data
        segments.append((text_part, fmt))
        pos = end

    # Trailing plain text
    if pos < len(text):
        segments.append((text[pos:], {}))

    return segments


# ---------------------------------------------------------------------------
# Office Plugin Registry — allows external plugins to register document processors
# ---------------------------------------------------------------------------
_office_plugins: dict = {}

def register_office_plugin(name: str):
    """Decorator to register an office document processor plugin."""
    def decorator(func):
        _office_plugins[name] = func
        return func
    return decorator

def _call_office_plugins(plugin_type: str, *args, **kwargs) -> dict:
    """Call all registered plugins of a given type and return results dict."""
    results = {}
    for pname, plugin in _office_plugins.items():
        try:
            if callable(plugin):
                results[pname] = plugin(plugin_type, *args, **kwargs)
        except Exception as e:
            print(f"[office-plugin] '{pname}' failed on '{plugin_type}': {e}", file=sys.stderr)
    return results

def _encode_filename(filename: str) -> str:
    """Encode filename to base64 for safe SQLite3 storage (handles international chars, special chars, path traversal)."""
    if not filename:
        return filename
    safe = filename.encode('utf-8')
    encoded = _b64_mod.urlsafe_b64encode(safe).decode('ascii').rstrip('=')
    _, ext = os.path.splitext(filename)
    return encoded + ext

def _decode_filename(encoded_name: str) -> str:
    """Decode a base64-encoded filename back to original. Returns as-is if decoding fails."""
    try:
        # _encode_filename base64-encodes the FULL original name (extension included) then
        # appends the real extension again as a plain suffix so path-based extension checks
        # keep working on the encoded form. So the decoded payload already carries the
        # extension -- don't re-append os.path.splitext(encoded_name)[1] on top of it.
        base = os.path.splitext(encoded_name)[0]
        padding = 4 - len(base) % 4
        if padding != 4:
            base += '=' * padding
        return _b64_mod.urlsafe_b64decode(base).decode('utf-8')
    except Exception:
        return encoded_name



def _read_odf(file_bytes: bytes, filename: str) -> str:
    """Read ODF files (.odt, .ods, .odp) and return structured text."""
    ext = os.path.splitext(filename)[1].lower()
    
    try:
        if ext == ".ods":
            from odf.opendocument import load
            from odf.table import Table, TableRow, TableCell
            from odf.text import P

            doc = load(io.BytesIO(file_bytes))
            result = []
            _MAX_REPEAT = 20  # cap on how many times one repeated cell/row is echoed
            for table in doc.getElementsByType(Table):
                for row in table.getElementsByType(TableRow):
                    # LibreOffice compacts runs of identical/blank cells into ONE
                    # <table:table-cell table:number-columns-repeated="N"> element --
                    # ignoring that attribute (the old code did) misaligns every column
                    # after the first repeated run. Blank cells are still expanded (for
                    # column alignment) but capped, since real files routinely encode
                    # "rest of the row is empty" as one cell repeated thousands of times.
                    row_repeat_raw = row.getAttribute('numberrowsrepeated')
                    try:
                        row_repeat = int(row_repeat_raw) if row_repeat_raw else 1
                    except (TypeError, ValueError):
                        row_repeat = 1

                    cells = []
                    for cell in row.getElementsByType(TableCell):
                        text_parts = []
                        for p in cell.getElementsByType(P):
                            try:
                                text_parts.append(str(p))
                            except Exception:
                                pass
                        cell_text = " ".join(text_parts).strip()
                        repeat_raw = cell.getAttribute('numbercolumnsrepeated')
                        try:
                            repeat = int(repeat_raw) if repeat_raw else 1
                        except (TypeError, ValueError):
                            repeat = 1
                        repeat = min(repeat, _MAX_REPEAT)
                        cells.extend([cell_text] * repeat)

                    line = " | ".join(cells)
                    if row_repeat > 1 and not any(c.strip() for c in cells):
                        # An entirely blank row repeated many times is filler
                        # ("rest of the sheet is empty") -- skip it rather than
                        # emitting hundreds/thousands of blank lines.
                        continue
                    emit = min(row_repeat, _MAX_REPEAT) if row_repeat > 1 else 1
                    for _ in range(emit):
                        result.append(line)
                    if row_repeat > emit:
                        result.append(f"(+{row_repeat - emit} more identical row(s))")
            return "\n".join(result) if result else "(empty spreadsheet)"
        
        elif ext == ".odt":
            from odf.opendocument import load
            from odf.text import P, H
            
            doc = load(io.BytesIO(file_bytes))
            result = []
            for elem in doc.getElementsByType(H):
                try:
                    result.append(f"## {str(elem)}")
                except Exception:
                    pass
            for elem in doc.getElementsByType(P):
                try:
                    text = str(elem).strip()
                    if text:
                        result.append(text)
                except Exception:
                    pass
            return "\n\n".join(result) if result else "(empty document)"
        
        elif ext == ".odp":
            from odf.opendocument import load
            from odf.text import P
            from odf.draw import Page

            doc = load(io.BytesIO(file_bytes))
            result = []
            # Iterate real slide boundaries (draw:page, same API already used correctly
            # in create_odf's odp-writing branch) instead of counting non-empty
            # paragraphs -- the old code incremented slide_num per PARAGRAPH, so a
            # single 5-bullet slide reported as "Slide 1" through "Slide 5".
            for i, page in enumerate(doc.getElementsByType(Page), 1):
                texts = []
                for elem in page.getElementsByType(P):
                    try:
                        text = str(elem).strip()
                        if text:
                            texts.append(text)
                    except Exception:
                        pass
                if texts:
                    result.append(f"Slide {i}:\n" + "\n".join(texts))
                else:
                    result.append(f"Slide {i}: (no text)")
            return "\n\n".join(result) if result else "(empty presentation)"
        
        return f"Unsupported ODF format: {ext}"
    except ImportError:
        return "Error: odfpy library not installed. Install with: pip install odfpy"
    except Exception as e:
        return f"Error reading {ext} file: {str(e)}"



# ---------------------------------------------------------------------------
# Professional document helpers
# ---------------------------------------------------------------------------
def _add_callout_box(doc, lines, colors, hex_to_rgb, fmt_mode="format"):
    """Add a professional callout box (note/tip/warning)."""
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    first_line = lines[0].lower() if lines else ""
    if first_line.startswith("**warning") or first_line.startswith("**alert"):
        border_color = "E74C3C"
        bg_color = "FDEDEC"
        icon = "\u26a0\ufe0f"
    elif first_line.startswith("**tip") or first_line.startswith("**pro tip"):
        border_color = "27AE60"
        bg_color = "E8F8F5"
        icon = "\U0001f4a1"
    elif first_line.startswith("**note") or first_line.startswith("**info"):
        border_color = colors.get("accent", "2E75B6")
        bg_color = colors.get("light", "D6E4F0")
        icon = "\U0001f4cc"

    table = doc.add_table(rows=1, cols=1)
    table.alignment = 1
    cell = table.rows[0].cells[0]
    cell.width = Inches(6.0)
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{bg_color}" w:val="clear"/>')
    tcPr.append(shading)

    borders = parse_xml(
        f'<w:tcBorders {nsdecls("w")}>'
        f'<w:left w:val="single" w:sz="24" w:space="8" w:color="{border_color}"/>'
        f'<w:top w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
        f'<w:bottom w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
        f'<w:right w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
        f'</w:tcBorders>'
    )
    tcPr.append(borders)

    for i, line in enumerate(lines):
        text = _format_text(line, mode=fmt_mode)
        segments = _parse_inline_md(text)
        if i == 0:
            p = cell.paragraphs[0]
            p.paragraph_format.space_before = Pt(6)
            run = p.add_run(f"{icon} ")
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.name = 'Calibri'
            for seg_text, fmt in segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                run.font.size = Pt(11)
                run.font.bold = True
                run.font.name = 'Calibri'
                if fmt.get('italic'):
                    run.font.italic = True
                if fmt.get('code'):
                    run.font.name = 'Consolas'
                    run.font.size = Pt(9)
        else:
            p = cell.add_paragraph()
            p.paragraph_format.space_before = Pt(2)
            for seg_text, fmt in segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                run.font.size = Pt(10)
                run.font.name = 'Calibri'
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('italic'):
                    run.font.italic = True
                if fmt.get('code'):
                    run.font.name = 'Consolas'
                    run.font.size = Pt(9)

    doc.add_paragraph()


def _add_professional_table(doc, rows, colors, hex_to_rgb, fmt_mode="format"):
    """Add a professionally styled table."""
    from docx.shared import Pt, RGBColor, Inches
    from docx.oxml.ns import qn, nsdecls
    from docx.oxml import parse_xml

    if not rows:
        return

    num_cols = max(len(r) for r in rows)
    table = doc.add_table(rows=len(rows), cols=num_cols)
    table.style = 'Colorful Grid Accent 1'
    table.alignment = 1

    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            if j < num_cols:
                cell = table.rows[i].cells[j]
                cell.text = ''
                p = cell.paragraphs[0]
                p.paragraph_format.space_before = Pt(2)
                p.paragraph_format.space_after = Pt(2)
                segments = _parse_inline_md(_format_text(cell_text, mode=fmt_mode))
                for seg_text, fmt in segments:
                    if not seg_text:
                        continue
                    run = p.add_run(seg_text)
                    run.font.size = Pt(10)
                    run.font.name = 'Calibri'
                    if i == 0:
                        run.font.bold = True
                    if fmt.get('bold'):
                        run.font.bold = True
                    if fmt.get('italic'):
                        run.font.italic = True
                    if fmt.get('code'):
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)

    doc.add_paragraph()


def _render_content_slide(prs, lines, colors, hex_to_rgb, add_text_box, add_accent_bar, set_slide_bg, slide_num, fmt_mode="format"):
    """Render a content slide with professional layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, colors["bg"])

    y_pos = 0.5
    for line in lines:
        line = line.strip()
        if not line:
            continue

        if line.startswith('# ') or line.startswith('## '):
            text = line.lstrip('#').strip()
            add_accent_bar(slide, 0.8, y_pos + 0.15, 0.06, 0.5, colors["accent"])
            add_text_box(slide, 1.1, y_pos, 11, 0.7, text, font_size=28, bold=True, color=colors["text"])
            y_pos += 0.9
        elif line.startswith('- ') or line.startswith('* '):
            text = _format_text(line[2:].strip(), mode=fmt_mode)
            add_text_box(slide, 1.5, y_pos, 10.5, 0.5, "\u2022 " + text, font_size=16, color=colors["text"])
            y_pos += 0.5
        elif line.startswith('|'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if any(c.startswith('---') for c in cells):
                continue
            add_text_box(slide, 1.0, y_pos, 11, 0.5, "  |  ".join(cells), font_size=14, color=colors["text"])
            y_pos += 0.4
        else:
            text = _format_text(line, mode=fmt_mode)
            add_text_box(slide, 1.0, y_pos, 11.5, 0.5, text, font_size=16, color=colors["text"])
            y_pos += 0.5

        if y_pos > 6.5:
            break


# ---------------------------------------------------------------------------
# Excerpt matching helpers (used by add_comment/add_comments for PDF + DOCX)
# ---------------------------------------------------------------------------
_QUOTE_DASH_REPLACEMENTS = {
    "‘": "'", "’": "'", "“": '"', "”": '"',
    "–": "-", "—": "-", "―": "-",
}


def _normalize_chars(text: str) -> str:
    """Character-for-character normalization (curly quotes/dashes -> plain).

    A 1:1 character mapping, so it never changes string length: offsets computed
    on the result map exactly back to offsets in the original text. Used where we
    need to know precisely where a match starts/ends (DOCX run splitting), as
    opposed to just whether it matches at all.
    """
    if not text:
        return ""
    for src, dst in _QUOTE_DASH_REPLACEMENTS.items():
        text = text.replace(src, dst)
    return text


def _normalize_match_text(text: str) -> str:
    """Normalize text for excerpt matching: curly quotes/dashes -> plain, collapse whitespace."""
    return re.sub(r"\s+", " ", _normalize_chars(text)).strip()


def _split_on_ellipsis(excerpt: str) -> list:
    """Split an excerpt on literal '...'/'…' markers (denoting omitted text).

    Returns a list of non-empty, normalized segments in order. A single-element
    list means the excerpt has no ellipsis. Callers should try the full excerpt
    (segments joined) first, and fall back to segments[0] if that fails to match --
    since PDFs/DOCX never contain the literal omission marker itself.
    """
    normalized = _normalize_match_text(excerpt)
    if not normalized:
        return []
    parts = re.split(r"\.\.\.|…", normalized)
    segments = [p.strip() for p in parts if p.strip()]
    return segments or [normalized]


def _loose_text_pattern(candidate: str) -> str:
    """Build a regex from a whitespace-collapsed candidate, treating any run of
    whitespace in the target text as flexible (matches single/double spaces, line
    wraps, tabs -- whatever the source document happens to use)."""
    parts = [p for p in candidate.split(" ") if p]
    if not parts:
        return ""
    return r"\s+".join(re.escape(p) for p in parts)


def _quote_style_variants(straight_text: str) -> list:
    """Return quote-style variants of an already whitespace-normalized, straight-quote
    string, for matching against text we can't normalize ourselves (PDF page content
    via PyMuPDF's search_for). Order: as given, then a curly-quote guess (naive: every
    ' becomes a right single quote, and "..." pairs become left/right double quotes --
    good enough for the common case of at most one quoted span per excerpt)."""
    variants = [straight_text]
    if "'" in straight_text or '"' in straight_text:
        curly = straight_text.replace("'", "’")
        curly = re.sub(r'"([^"]*)"', "“\\1”", curly)
        if curly != straight_text:
            variants.append(curly)
    return variants


__all__ = [
    "_office_plugins", "register_office_plugin", "_call_office_plugins",
    "_resolve_file_path", "_read_file_bytes", "_detect_type", "_cell_value",
    "_xls_to_xlsx", "_format_text", "_parse_inline_md", "_encode_filename",
    "_decode_filename", "_read_odf", "_add_callout_box", "_add_professional_table",
    "_render_content_slide", "_normalize_match_text", "_split_on_ellipsis",
    "_normalize_chars", "_loose_text_pattern", "_quote_style_variants",
]

# =========================================================================
# =========================================================================
class Tools:
    class Valves(BaseModel):
        base_url: Optional[str] = Field(
            default=None,
            description="Override the base URL for download links. Auto-detected from X-Original-Host header or WEBUI_URL env var if unset.",
        )
        file_url_pattern: Optional[str] = Field(
            default=None,
            description="Custom file download URL pattern. Use {file_id} placeholder. Examples: /api/v1/files/{file_id}/content (default), /api/files/{file_id}. If unset, uses standard Open WebUI URL.",
        )
        templates: Optional[str] = Field(default="{}", description="JSON map of template names to content strings.")
        cleanup_schedule: Optional[str] = Field(default="{}", description="JSON schedule for auto-cleanup.")
        language: Optional[str] = Field(default="en", description="Language for error messages: en, pt, es, fr, de.")
        pass

    def __init__(self):
        self.valves = self.Valves()

    def _resolve_file(self, file_id: str):
        """Resolve a file ID to (bytes, filename, file_type)."""
        path = _resolve_file_path(file_id)
        if not path:
            return None, None, None
        filename = os.path.basename(path)
        ftype = _detect_type(filename)
        if ftype == "unknown":
            # Files this tool itself creates via _save_and_link are stored on disk under
            # their bare file_id (no extension), so the path's basename never carries a
            # usable extension. Fall back to the DB's filename column, which always does.
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                try:
                    row = conn.execute("SELECT filename FROM file WHERE id = ?", (file_id,)).fetchone()
                finally:
                    conn.close()
                if row and row[0]:
                    # The DB column stores the base64-encoded name (_encode_filename) --
                    # decode it back before handing it to callers, otherwise every
                    # chained call (add_comment -> add_comment on the returned file_id)
                    # re-encodes an already-encoded name, growing it on each hop.
                    filename = _decode_filename(row[0])
                    ftype = _detect_type(filename)
            except Exception:
                pass
        file_bytes = _read_file_bytes(file_id)
        return file_bytes, filename, ftype

    def _read_xlsx(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from an XLSX file."""
        from openpyxl import load_workbook
        import io
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        result = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                result.append("\t".join(cells))
        wb.close()
        return "\n\n".join(result)

    def _read_xls(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from an XLS file."""
        import xlrd
        import io
        wb = xlrd.open_workbook(file_contents=file_bytes)
        result = []
        for si in range(wb.nsheets):
            ws = wb.sheet_by_index(si)
            for r in range(ws.nrows):
                cells = [str(ws.cell_value(r, c)) for c in range(ws.ncols)]
                result.append("\t".join(cells))
        return "\n\n".join(result)

    def _read_docx(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from a DOCX file."""
        from docx import Document
        import io
        doc = Document(io.BytesIO(file_bytes))
        result = []
        for p in doc.paragraphs:
            m = re.match(r"Heading (\d+)$", p.style.name) if p.style and p.style.name else None
            if m:
                result.append(f"{'#' * int(m.group(1))} {p.text}")
            else:
                result.append(p.text)
        for t in doc.tables:
            for row in t.rows:
                result.append(" | ".join(cell.text for cell in row.cells))
        return "\n\n".join(result)

    def _read_pptx(self, file_bytes: bytes, filename: str) -> str:
        """Extract text from a PPTX file."""
        from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(file_bytes))
        result = []
        for si, slide in enumerate(prs.slides, 1):
            result.append(f"# Slide {si}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    result.append(shape.text)
        return "\n\n".join(result)

    def _parse_csv_rows(self, content: str) -> list:
        """Parse CSV content into a list of rows with type conversion."""
        import csv as _csv_mod
        reader = _csv_mod.reader(io.StringIO(content))
        parsed_rows = []
        for csv_row in reader:
            converted = []
            for v in csv_row:
                v = v.strip()
                if v == '':
                    converted.append(None)
                # A leading zero followed by another digit ("02134", "00123") means the
                # value is an identifier (zip code, account number), not a number -- int()
                # would strip the leading zero. Also reject underscores: Python 3 accepts
                # "1_000" as a numeric literal, so int("1_000") silently returns 1000 for
                # what was meant to be the literal text "1_000".
                elif re.match(r'^-?0\d', v) or '_' in v:
                    converted.append(v)
                else:
                    try:
                        converted.append(int(v))
                    except ValueError:
                        try:
                            converted.append(float(v))
                        except ValueError:
                            if v.lower() == 'true':
                                converted.append(True)
                            elif v.lower() == 'false':
                                converted.append(False)
                            else:
                                converted.append(v)
            parsed_rows.append(converted)
        return parsed_rows

    def _last_populated_row(self, ws) -> int:
        """Return the last row index that actually has a non-empty cell value.

        ws.max_row includes trailing rows that only carry formatting/blank
        cells, so appending after it can leave a large gap of empty rows
        before new content when the sheet has such trailing rows.
        """
        for row in range(ws.max_row or 0, 0, -1):
            for cell in ws[row]:
                if cell.value is not None and str(cell.value).strip() != "":
                    return row
        return 0

    # -----------------------------------------------------------------
    # Internal: save and return markdown link
    # -----------------------------------------------------------------
    async def _save_and_link(self, file_bytes: bytes, filename: str, __request__=None, __user__=None) -> tuple:
        """Save file via Open WebUI's own Storage/Files layer and return a working download URL.

        Previously this wrote bytes directly into a locally-computed uploads directory
        (guessed from a CUSTOM env var, `OPEN_WEBUI_DATA_DIR`, that Open WebUI itself does
        not use or recognize -- it only "worked" on deployments that happen to set that exact
        variable) and inserted a row via raw SQL. On any deployment without that variable
        (i.e. a stock install) the file was written to a path Open WebUI's own file-serving
        code never looks at, and the returned link 401'd/404'd. It was also silently broken
        on any S3/GCS/Azure-backed deployment (`STORAGE_PROVIDER` env var), since those
        require `file.path` to be a Storage URI, not a local filesystem path.

        This now calls the exact same sequence Open WebUI's own upload endpoint uses
        (`routers/files.py:upload_file_handler`): `Storage.upload_file()` (blocking, so run
        in a thread) writes the bytes through whatever backend is actually configured, then
        `Files.insert_new_file()` registers it the same way a real upload would. The
        resulting `/api/v1/files/{id}/content` link is guaranteed to be servable, on any
        storage backend, with no environment-variable guessing required.
        """
        import base64 as _b64
        import hashlib
        import io as _io
        import uuid as _uuid

        ext = os.path.splitext(filename)[1].lower()
        mt = {
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".xls": "application/vnd.ms-excel",
            ".doc": "application/msword",
            ".ppt": "application/vnd.ms-powerpoint",
            ".pdf": "application/pdf",
            ".odt": "application/vnd.oasis.opendocument.text",
            ".ods": "application/vnd.oasis.opendocument.spreadsheet",
            ".odp": "application/vnd.oasis.opendocument.presentation",
        }
        content_type = mt.get(ext, "application/octet-stream")

        try:
            import asyncio as _asyncio
            from open_webui.storage.provider import Storage
            from open_webui.models.files import Files, FileForm

            file_id = str(_uuid.uuid4())
            user_id = __user__.get("id", "") if __user__ and isinstance(__user__, dict) else ""

            tags = {"OpenWebUI-User-Id": user_id, "OpenWebUI-File-Id": file_id}
            # Storage.upload_file is a blocking call (local disk write, or a network call for
            # S3/GCS/Azure) -- must run off the event loop, same as the real upload endpoint.
            contents, file_path = await _asyncio.to_thread(
                Storage.upload_file, _io.BytesIO(file_bytes), f"{file_id}_{filename}", tags
            )
            file_hash = hashlib.sha256(contents).hexdigest()

            file_item = await Files.insert_new_file(
                user_id,
                FileForm(
                    id=file_id,
                    hash=file_hash,
                    filename=filename,
                    path=file_path,
                    data={},
                    meta={
                        "name": filename,
                        "content_type": content_type,
                        "size": len(contents),
                        "source": "office-plugin",
                        "generated": True,
                    },
                ),
            )
            if file_item is None:
                raise RuntimeError("Files.insert_new_file returned None")

            base_url = self.valves.base_url
            if not base_url and __request__:
                try:
                    host = __request__.headers.get("x-original-host")
                    if host:
                        base_url = f"https://{host}"
                except Exception:
                    pass
            if not base_url:
                base_url = os.environ.get("WEBUI_URL", "http://localhost:3000")
            base_url = base_url.rstrip("/")

            # Standard Open WebUI file download URL (works for most versions)
            # Customize via file_url_pattern valve if needed
            if self.valves.file_url_pattern:
                url = f"{base_url}{self.valves.file_url_pattern.replace('{file_id}', file_id)}"
            else:
                url = f"{base_url}/api/v1/files/{file_id}/content"
            return (url, filename)

        except Exception as e:
            print(f"[office] Save failed: {e}", file=sys.stderr)
            try:
                if len(file_bytes) > 1_000_000:  # 1MB limit
                    return (None, None)
                data = _b64.b64encode(file_bytes).decode("ascii")
                return (f"data:{content_type};base64,{data}", filename)
            except Exception:
                return (None, None)

    # -----------------------------------------------------------------
    # READ
    # -----------------------------------------------------------------
    async def read_file(
        self,
        file_id: str,
        max_rows: int = 500,
        sheet_name: str = "",
        row_start: int = 1,
        row_end: int = 0,
        page_start: int = 1,
        page_end: int = 0,
        __user__=None,
        __request__=None,
    ) -> str:
        """Read any Office file (.xlsx, .xls, .csv, .docx, .pptx, .pdf) and return its
        contents as structured JSON.

        Auto-detects the file type from the file ID or filename.
        For xlsx/xls/csv: returns sheets with headers and rows.
        For docx: returns paragraphs with styles and tables.
        For pptx: returns slides with shapes and text.
        For pdf: returns real embedded text per page (fast, no OCR) -- see page_start/
            page_end. Pages with little/no extractable text are flagged as likely scanned
            images; use ocr_extract for those specific pages.
        Legacy .doc and .ppt formats return a helpful error message.

        Args:
            file_id: The Open WebUI file ID (UUID) or filename
            max_rows: Maximum rows to return (default 500)
            sheet_name: Optional - read only this sheet (xlsx/xls only)
            row_start: Starting row (1-indexed, default 1) -- xlsx/xls only, ignored for
                csv/docx/pptx (csv is limited only by max_rows; docx/pptx have no row concept)
            row_end: Ending row (0 = use max_rows limit) -- xlsx/xls only, same as row_start
            page_start: Starting page (1-indexed, default 1) -- pdf only, ignored otherwise
            page_end: Ending page (0 = up to 20 pages from page_start) -- pdf only. Capped
                at 20 pages per call regardless of the value given; call again with a higher
                page_start to continue through a long PDF.
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({
                    "error": (
                        f"Could not read file {file_id}. "
                        "Make sure the file was uploaded via the chat."
                    )
                })

            # Detect type from filename in DB
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = _decode_filename(row[0]) if row and row[0] else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            result: Dict[str, Any] = {
                "file_id": file_id,
                "filename": filename,
                "type": file_type,
            }

            if file_type == "xlsx":
                import openpyxl
                # read_only=True avoids building openpyxl's full in-memory cell-object
                # graph for a read-only operation -- _read_xlsx already does this;
                # read_file (the more commonly reached-for function) didn't, so a large
                # sheet cost far more memory/time here than it needed to.
                wb = openpyxl.load_workbook(io.BytesIO(file_data), data_only=True, read_only=True)
                result["sheets"] = []
                sheetnames = [sn for sn in wb.sheetnames] if not sheet_name else [sn for sn in wb.sheetnames if sn == sheet_name]
                if sheet_name and sheet_name not in wb.sheetnames:
                    result["warning"] = f"Sheet '{sheet_name}' not found. Available sheets: {wb.sheetnames}"
                for sn in sheetnames:
                    ws = wb[sn]
                    sheet: Dict[str, Any] = {
                        "name": sn,
                        "headers": [],
                        "rows": [],
                        "total_rows": ws.max_row or 0,
                        "total_cols": ws.max_column or 0,
                    }
                    r_start = max(1, row_start)
                    r_end = row_end if row_end > 0 else (ws.max_row or 0)
                    r_end = min(r_end, r_start + max_rows - 1) if max_rows else r_end
                    # Headers always come from row 1, fetched independently of the data
                    # window -- the old code only populated "headers" when ri == 1, which
                    # never happened once row_start > 1 (e.g. paging through a large sheet
                    # with row_start=501), silently returning empty headers on every page
                    # after the first. This also fixes an off-by-one: the header row no
                    # longer eats one of max_rows's data-row slots.
                    header_row = next(ws.iter_rows(min_row=1, max_row=1), None)
                    if header_row:
                        sheet["headers"] = [str(_cell_value(c)) if _cell_value(c) is not None else "" for c in header_row]
                    data_start = max(r_start, 2)
                    for row in ws.iter_rows(min_row=data_start, max_row=r_end):
                        sheet["rows"].append([_cell_value(c) for c in row])
                    result["sheets"].append(sheet)
                wb.close()

            elif file_type == "xls":
                import xlrd
                xls_book = xlrd.open_workbook(file_contents=file_data)
                result["sheets"] = []
                sheet_indices = range(xls_book.nsheets)
                if sheet_name:
                    sheet_indices = [i for i in range(xls_book.nsheets) if xls_book.sheet_by_index(i).name == sheet_name]
                    if not sheet_indices:
                        all_names = [xls_book.sheet_by_index(i).name for i in range(xls_book.nsheets)]
                        result["warning"] = f"Sheet '{sheet_name}' not found. Available sheets: {all_names}"
                for sheet_idx in sheet_indices:
                    xls_sheet = xls_book.sheet_by_index(sheet_idx)
                    sheet: Dict[str, Any] = {
                        "name": xls_sheet.name,
                        "headers": [],
                        "rows": [],
                        "total_rows": xls_sheet.nrows,
                        "total_cols": xls_sheet.ncols,
                    }
                    r_start = max(0, row_start - 1)
                    r_end = row_end if row_end > 0 else xls_sheet.nrows
                    r_end = min(r_end, r_start + max_rows) if max_rows else r_end
                    r_end = min(r_end, xls_sheet.nrows)
                    for rx in range(r_start, r_end):
                        row_values = []
                        for cx in range(xls_sheet.ncols):
                            cell = xls_sheet.cell(rx, cx)
                            value = cell.value
                            if cell.ctype == xlrd.XL_CELL_DATE:
                                try:
                                    import datetime as _dt
                                    value = _dt.datetime(*xlrd.xldate_as_tuple(value, xls_book.datemode))
                                except Exception:
                                    pass
                            elif cell.ctype == xlrd.XL_CELL_BOOLEAN:
                                value = bool(value)
                            row_values.append(value)
                        if rx == 0:
                            sheet["headers"] = [str(v) if v is not None else "" for v in row_values]
                        else:
                            sheet["rows"].append(row_values)
                    result["sheets"].append(sheet)
                xls_book.release_resources()

            elif file_type == "csv":
                import csv as _csv_mod
                result["sheets"] = []
                d_bytes = file_data
                try:
                    text = d_bytes.decode("utf-8-sig")
                except Exception:
                    text = d_bytes.decode("utf-8", errors="replace")
                reader = _csv_mod.DictReader(io.StringIO(text))
                headers = reader.fieldnames or []
                rows = []
                for ri, row in enumerate(reader):
                    if max_rows and ri >= max_rows:
                        break
                    rows.append(dict(row))
                base_name = os.path.splitext(os.path.basename(filename))[0] if filename else "CSV"
                sheet = {
                    "name": base_name,
                    "headers": [str(h) for h in headers],
                    "rows": [[r.get(h, "") for h in headers] for r in rows],
                    "total_rows": len(rows),
                    "total_cols": len(headers),
                }
                result["sheets"].append(sheet)

            elif file_type == "docx":
                from docx import Document
                from docx.shared import Pt
                from docx.enum.text import WD_COLOR_INDEX
                doc = Document(io.BytesIO(file_data))
                paragraphs = []
                for p in doc.paragraphs:
                    if p.text.strip():
                        style = p.style.name if p.style else "Normal"
                        runs_info = []
                        for run in p.runs:
                            run_data = {"text": run.text}
                            if run.font.highlight_color and run.font.highlight_color != WD_COLOR_INDEX.AUTO:
                                run_data["highlighted"] = True
                                run_data["highlight_color"] = str(run.font.highlight_color)
                            if run.font.bold:
                                run_data["bold"] = True
                            if run.font.italic:
                                run_data["italic"] = True
                            runs_info.append(run_data)
                        paragraphs.append({"style": style, "text": p.text, "runs": runs_info})
                tables = []
                for t in doc.tables:
                    tbl = {"rows": []}
                    for row in t.rows:
                        tbl["rows"].append([cell.text for cell in row.cells])
                    tables.append(tbl)
                result["paragraphs"] = paragraphs
                result["tables"] = tables

            elif file_type == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_data))
                slides = []
                for si, slide in enumerate(prs.slides, 1):
                    sdata: Dict[str, Any] = {"number": si, "shapes": []}
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            sdata["shapes"].append({
                                "type": str(shape.shape_type),
                                "name": shape.name,
                                "text": shape.text[:500],
                            })
                        if shape.has_table:
                            tbl = {"rows": []}
                            for row in shape.table.rows:
                                tbl["rows"].append([cell.text for cell in row.cells])
                            # A slide with 2+ tables previously kept only the last one --
                            # "tables" (plural key) held a single dict, overwritten each
                            # time. Append to a list instead.
                            sdata.setdefault("tables", []).append(tbl)
                    slides.append(sdata)
                result["slides"] = slides

            elif file_type == "pdf":
                import fitz
                pdf = fitz.open(stream=file_data, filetype="pdf")
                total_pages = pdf.page_count
                p_start = max(1, page_start)
                p_end = page_end if page_end > 0 else total_pages
                p_end = min(p_end, total_pages)
                _MAX_PDF_PAGES_PER_CALL = 20
                if p_end - p_start + 1 > _MAX_PDF_PAGES_PER_CALL:
                    p_end = p_start + _MAX_PDF_PAGES_PER_CALL - 1
                pages = []
                sparse_pages = []
                for i in range(p_start - 1, p_end):
                    text = pdf.load_page(i).get_text().strip()
                    if len(text) < 20:
                        sparse_pages.append(i + 1)
                    pages.append({"page": i + 1, "text": text})
                pdf.close()
                result["total_pages"] = total_pages
                result["pages"] = pages
                if p_end < total_pages:
                    result["note"] = (
                        f"Returned pages {p_start}-{p_end} of {total_pages}. "
                        f"Call again with page_start={p_end + 1} to continue."
                    )
                if sparse_pages:
                    result["sparse_pages"] = sparse_pages
                    result["sparse_pages_note"] = (
                        "Little/no extractable text on these pages (likely scanned images) -- "
                        "use ocr_extract for these specific pages."
                    )

            elif file_type == "doc":
                result["error"] = "Legacy .doc format is not supported. Please convert to .docx first."

            elif file_type == "ppt":
                result["error"] = "Legacy .ppt format is not supported. Please convert to .pptx first."

            else:
                result["error"] = f"Unsupported file type. Detected: {file_type}. Supported: xlsx, xls, docx, pptx, pdf"

            return json.dumps(result, indent=2, default=str, ensure_ascii=False)

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # ADD CONTENT
    # -----------------------------------------------------------------
    async def add_content(
        self,
        file_id: str,
        content: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Add new content to an Office file while preserving original formatting.

        For spreadsheets (xlsx): content is CSV text with rows to add.
        For documents (docx): content is text to append at the end. Preserves original
            formatting and capitalization of the existing document.
        For presentations (pptx): content is split into slides on "---" lines. Within each
            "---"-separated block, the FIRST line becomes that slide's title and the remaining
            lines become its body text -- it is not one slide per line.

        Args:
            file_id: File ID to edit
            content: Content to add (CSV for xlsx, text for docx/pptx)
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            # Detect type
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = _decode_filename(row[0]) if row and row[0] else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            out = io.BytesIO()
            out_name = output_filename

            if file_type == "xlsx":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                ws = wb.active
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

                # Parse CSV content
                parsed_rows = self._parse_csv_rows(content)

                if not parsed_rows:
                    return json.dumps({"error": "No rows provided in CSV content"})

                # Get reference styles from last row
                last_data_row = self._last_populated_row(ws)
                ref = {}
                if last_data_row >= 1:
                    for cell in ws[last_data_row]:
                        if cell.has_style:
                            ref[cell.column] = {
                                "font": copy(cell.font),
                                "fill": copy(cell.fill),
                                "border": copy(cell.border),
                                "alignment": copy(cell.alignment),
                                "number_format": cell.number_format,
                            }

                start = last_data_row + 1
                for i, rd in enumerate(parsed_rows):
                    for j, v in enumerate(rd, 1):
                        cell = ws.cell(row=start + i, column=j)
                        if j in ref:
                            try:
                                cell.font = copy(ref[j]["font"])
                                cell.fill = copy(ref[j]["fill"])
                                cell.border = copy(ref[j]["border"])
                                cell.alignment = copy(ref[j]["alignment"])
                                cell.number_format = ref[j]["number_format"]
                            except Exception:
                                pass
                        cell.value = v

                wb.save(out)
                wb.close()

            elif file_type == "xls":
                # Convert .xls to .xlsx, then apply same add logic
                file_data = _xls_to_xlsx(file_data)
                file_type = "xlsx"
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                ws = wb.active

                # Parse CSV content
                parsed_rows = self._parse_csv_rows(content)

                if not parsed_rows:
                    return json.dumps({"error": "No rows provided in CSV content"})

                last_data_row = self._last_populated_row(ws)
                ref = {}
                if last_data_row >= 1:
                    for cell in ws[last_data_row]:
                        if cell.has_style:
                            ref[cell.column] = {
                                "font": copy(cell.font),
                                "fill": copy(cell.fill),
                                "border": copy(cell.border),
                                "alignment": copy(cell.alignment),
                                "number_format": cell.number_format,
                            }

                start = last_data_row + 1
                for i, rd in enumerate(parsed_rows):
                    for j, v in enumerate(rd, 1):
                        cell = ws.cell(row=start + i, column=j)
                        if j in ref:
                            try:
                                cell.font = copy(ref[j]["font"])
                                cell.fill = copy(ref[j]["fill"])
                                cell.border = copy(ref[j]["border"])
                                cell.alignment = copy(ref[j]["alignment"])
                                cell.number_format = ref[j]["number_format"]
                            except Exception:
                                pass
                        cell.value = v

                wb.save(out)
                wb.close()

            elif file_type == "docx":
                from docx import Document
                from docx.shared import Pt
                doc = Document(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.docx"

                # Append paragraphs, preserving last paragraph's style
                last_style = "Normal"
                if doc.paragraphs:
                    last_style = doc.paragraphs[-1].style.name if doc.paragraphs[-1].style else "Normal"

                for line in content.split("\n"):
                    line = line.strip()
                    if not line:
                        continue
                    if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
                        level = len(line) - len(line.lstrip('#'))
                        text = line.lstrip('#').strip()
                        doc.add_heading(text, level=min(level, 3))
                    elif line.startswith('- ') or line.startswith('* '):
                        text = line[2:].strip()
                        doc.add_paragraph(text, style='List Bullet')
                    else:
                        segments = _parse_inline_md(line)
                        p = doc.add_paragraph(style=last_style)
                        for seg_text, fmt in segments:
                            if not seg_text:
                                continue
                            run = p.add_run(seg_text)
                            if fmt.get('code'):
                                run.font.name = "Consolas"
                                run.font.size = Pt(9)
                            if fmt.get('bold'):
                                run.font.bold = True
                            if fmt.get('italic'):
                                run.font.italic = True

                out = io.BytesIO()
                doc.save(out)

            elif file_type == "pptx":
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.pptx"

                # Split content by "---" into slides
                slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
                blank_layout = prs.slide_layouts[6]  # Blank layout

                for spec in slide_specs:
                    spec = spec.strip()
                    if not spec:
                        continue
                    lines = spec.split("\n")
                    title = lines[0].strip()
                    body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                    slide = prs.slides.add_slide(blank_layout)
                    # Add title
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                    )
                    tf = txBox.text_frame
                    tf.text = _format_text(title, mode="format")
                    p = tf.paragraphs[0]
                    p.font.size = Inches(0.6)
                    p.font.bold = True

                    # Add body text
                    if body:
                        txBox2 = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.text = _format_text(body, mode="format")
                        for para in tf2.paragraphs:
                            para.font.size = Inches(0.3)

                out = io.BytesIO()
                prs.save(out)

            else:
                return json.dumps({"error": f"Unsupported type: {file_type}"})

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__, __user__=__user__)
            if url:
                return f"[{name}]({url})\n\nAdded content to {file_type.upper()} file, preserving original formatting."
            return json.dumps({"error": self._err("could_not_save")})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # REPLACE TEXT
    # -----------------------------------------------------------------
    async def replace_text(
        self,
        file_id: str,
        find_text: str,
        replace_with: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Find and replace text in any Office file while preserving original formatting
        and capitalization of the existing document.

        Works on cell values in xlsx, paragraph text in docx, and shape text in pptx.

        Args:
            file_id: File ID to edit
            find_text: Text to find
            replace_with: Text to replace with
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            # Detect type
            try:
                conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row = conn.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn.close()
                filename = _decode_filename(row[0]) if row and row[0] else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            out = io.BytesIO()
            out_name = output_filename
            count = 0

            if file_type == "xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"

                for sn in wb.sheetnames:
                    ws = wb[sn]
                    for row in ws.iter_rows():
                        for cell in row:
                            # Only string cells are text-replaced. A non-string cell (number,
                            # date, bool) previously matched by substring-testing its str()
                            # form but then OVERWRITING THE WHOLE CELL with replace_with --
                            # e.g. replace_text(id, "0", "O") turned every number/date cell
                            # containing a "0" into the literal string "O". Skipping non-string
                            # cells entirely avoids destroying numeric/date data.
                            if isinstance(cell.value, str) and find_text in cell.value:
                                cell.value = cell.value.replace(find_text, replace_with)
                                count += 1

                wb.save(out)
                wb.close()

            elif file_type == "xls":
                # Convert .xls to .xlsx, then apply same replace logic
                file_data = _xls_to_xlsx(file_data)
                file_type = "xlsx"
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.xlsx"
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_data))

                for sn in wb.sheetnames:
                    ws = wb[sn]
                    for row in ws.iter_rows():
                        for cell in row:
                            # Only string cells are text-replaced. A non-string cell (number,
                            # date, bool) previously matched by substring-testing its str()
                            # form but then OVERWRITING THE WHOLE CELL with replace_with --
                            # e.g. replace_text(id, "0", "O") turned every number/date cell
                            # containing a "0" into the literal string "O". Skipping non-string
                            # cells entirely avoids destroying numeric/date data.
                            if isinstance(cell.value, str) and find_text in cell.value:
                                cell.value = cell.value.replace(find_text, replace_with)
                                count += 1

                wb.save(out)
                wb.close()

            elif file_type == "docx":
                from docx import Document
                doc = Document(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.docx"

                cross_run_misses = 0

                def _replace_in_paragraph(para):
                    nonlocal count, cross_run_misses
                    if find_text not in para.text:
                        return
                    if not para.runs:
                        para.text = para.text.replace(find_text, replace_with)
                        count += 1
                        return
                    # Only count a real replacement -- a paragraph can contain find_text
                    # while NO SINGLE RUN does (Word splits phrases across runs constantly:
                    # spell-check marks, bold spans, language tags). The old code counted
                    # the paragraph as replaced regardless, reporting success on a file that
                    # was never actually changed.
                    changed = False
                    for run in para.runs:
                        if find_text in run.text:
                            run.text = run.text.replace(find_text, replace_with)
                            changed = True
                    if changed:
                        count += 1
                    else:
                        cross_run_misses += 1

                for para in doc.paragraphs:
                    _replace_in_paragraph(para)

                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                _replace_in_paragraph(para)

                doc.save(out)

            elif file_type == "pptx":
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_data))
                if not out_name:
                    out_name = os.path.splitext(filename)[0] + "_edited.pptx"

                cross_run_misses = 0
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame"):
                            for para in shape.text_frame.paragraphs:
                                if find_text not in para.text:
                                    continue
                                if not para.runs:
                                    para.text = para.text.replace(find_text, replace_with)
                                    count += 1
                                    continue
                                changed = False
                                for run in para.runs:
                                    if find_text in run.text:
                                        run.text = run.text.replace(find_text, replace_with)
                                        changed = True
                                if changed:
                                    count += 1
                                else:
                                    cross_run_misses += 1

                prs.save(out)

            else:
                return json.dumps({"error": f"Unsupported type: {file_type}"})

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__, __user__=__user__)
            if not url:
                return json.dumps({"error": self._err("could_not_save")})
            result = f"[{name}]({url})\n\nReplaced '{find_text}' with '{replace_with}' in {count} place(s), preserving all formatting."
            if count == 0:
                result = f"No replacements made -- '{find_text}' was not found (or only spans a formatting boundary Word split into separate runs). No file was changed: [{name}]({url}) is identical to the original."
            elif locals().get("cross_run_misses"):
                result += f"\n{cross_run_misses} additional occurrence(s) were seen but NOT replaced because '{find_text}' spans a formatting boundary (split across separate runs)."
            return result

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # UPDATE CELLS
    # -----------------------------------------------------------------
    async def update_cells(
        self,
        file_id: str,
        cells: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Update individual cells in an XLSX file. Pass a JSON array of cell updates.

        Each update: {"cell": "A1", "value": "new value"} or {"cell": "B2", "value": 42, "sheet": "Sheet1"}
        If sheet is omitted, uses the active sheet.

        Args:
            file_id: The file ID of the XLSX file to edit
            cells: JSON string - array of {cell, value[, sheet]} objects
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            try:
                conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row2 = conn2.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn2.close()
                filename = _decode_filename(row2[0]) if row2 and row2[0] else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            if file_type not in ("xlsx", "xls"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx/xls supported."})

            import openpyxl
            from openpyxl.utils import column_index_from_string
            import re as _re_cell

            try:
                updates = json.loads(cells)
                if not isinstance(updates, list):
                    updates = [updates]
            except json.JSONDecodeError:
                return json.dumps({"error": "cells must be valid JSON (array of {cell, value} objects)"})

            out_name = output_filename or os.path.splitext(filename)[0] + "_updated.xlsx"
            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            count = 0
            skipped = []
            for i, upd in enumerate(updates):
                cell_ref = upd.get("cell", "")
                if "value" not in upd:
                    skipped.append(f"entry {i}: missing 'value' key")
                    continue
                value = upd.get("value")
                sname = upd.get("sheet", "")
                if sname and sname not in wb.sheetnames:
                    skipped.append(f"entry {i}: sheet {sname!r} not found (available: {wb.sheetnames})")
                    continue
                ws = wb[sname] if sname else wb.active
                # fullmatch, not match: "A1:B10" would previously match at the start
                # (A1) and silently write one cell while the caller believed a range
                # was updated.
                m = _re_cell.fullmatch(r"([A-Za-z]+)(\d+)", cell_ref)
                if not m:
                    skipped.append(f"entry {i}: invalid cell reference {cell_ref!r}")
                    continue
                col_str, row_str = m.group(1), m.group(2)
                col_idx = column_index_from_string(col_str.upper())
                row_idx = int(row_str)
                ws.cell(row=row_idx, column=col_idx).value = value
                count += 1

            out = io.BytesIO()
            wb.save(out)
            wb.close()
            out.seek(0)

            url, name = await self._save_and_link(out.read(), out_name, __request__, __user__=__user__)
            if not url:
                return json.dumps({"error": self._err("could_not_save")})
            result = f"[{name}]({url})\n\nUpdated {count} cell(s) in {file_type.upper()} file."
            if skipped:
                result += f"\n{len(skipped)} entries skipped: " + "; ".join(skipped)
            return result
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # MODIFY ROWS (insert/delete)
    # -----------------------------------------------------------------
    def _shift_merged_ranges(self, ws, row_number: int, count: int, is_insert: bool):
        """Unmerge every merged range on `ws`, and return the (min_col, min_row, max_col,
        max_row) bounds each should be re-merged at after an insert/delete of `count` rows
        at `row_number` -- openpyxl's insert_rows()/delete_rows() move cell VALUES only,
        they do not touch merged-cell ranges at all, so a merge spanning the shifted rows
        would otherwise silently end up covering the wrong cells (or straddling a
        now-nonexistent boundary, which Excel treats as corrupt).

        Must be called (and the returned ranges re-merged via `_remerge_ranges`) around the
        insert_rows()/delete_rows() call, not instead of it.
        """
        old_ranges = list(ws.merged_cells.ranges)
        new_ranges = []
        del_start, del_end = row_number, row_number + count - 1

        def _map_row(r):
            if is_insert:
                return r if r < row_number else r + count
            # delete: rows before the deleted block keep their index; rows after shift up;
            # rows inside the deleted block collapse to the block's start (callers clamp).
            if r < del_start:
                return r
            if r > del_end:
                return r - count
            return del_start

        for mcr in old_ranges:
            min_col, min_row, max_col, max_row = mcr.bounds
            new_min_row = _map_row(min_row)
            new_max_row = _map_row(max_row) if is_insert else (
                del_start - 1 if del_start <= max_row <= del_end else _map_row(max_row)
            )
            if new_min_row <= new_max_row:
                new_ranges.append((min_col, new_min_row, max_col, new_max_row))
            # else: the range was entirely consumed by the deleted rows -- dropped.

        for mcr in old_ranges:
            ws.unmerge_cells(str(mcr))
        return new_ranges

    def _remerge_ranges(self, ws, ranges):
        from openpyxl.utils import get_column_letter
        for min_col, min_row, max_col, max_row in ranges:
            try:
                ws.merge_cells(f"{get_column_letter(min_col)}{min_row}:{get_column_letter(max_col)}{max_row}")
            except Exception:
                pass  # skip a range that's no longer valid rather than corrupt the file

    async def modify_rows(
        self,
        file_id: str,
        action: str,
        row_number: int = 1,
        count: int = 1,
        sheet_name: str = "",
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Insert or delete rows in an XLSX file. Merged cell ranges are repositioned to
        match the shift. Formulas are NOT adjusted (openpyxl has no equivalent of Excel's
        own reference-rewriting on insert/delete) -- if the sheet contains formulas, the
        response includes a warning so you know to check them manually.

        Args:
            file_id: The file ID of the XLSX file to edit
            action: 'insert' or 'delete'
            row_number: Row number to start at (1-indexed)
            count: Number of rows to insert/delete (default 1)
            sheet_name: Optional sheet name (default: active sheet)
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            try:
                conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row2 = conn2.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn2.close()
                filename = _decode_filename(row2[0]) if row2 and row2[0] else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            if file_type not in ("xlsx", "xls"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx/xls supported."})

            import openpyxl

            wb = openpyxl.load_workbook(io.BytesIO(file_data))
            ws = wb[sheet_name] if sheet_name else wb.active
            out_name = output_filename or os.path.splitext(filename)[0] + "_modified.xlsx"

            action_lower = action.strip().lower()
            if action_lower == "insert":
                new_ranges = self._shift_merged_ranges(ws, row_number, count, is_insert=True)
                ws.insert_rows(idx=row_number, amount=count)
                self._remerge_ranges(ws, new_ranges)
                msg = f"Inserted {count} row(s) at row {row_number}"
            elif action_lower == "delete":
                new_ranges = self._shift_merged_ranges(ws, row_number, count, is_insert=False)
                ws.delete_rows(idx=row_number, amount=count)
                self._remerge_ranges(ws, new_ranges)
                msg = f"Deleted {count} row(s) starting at row {row_number}"
            else:
                return json.dumps({"error": f"Unknown action '{action}'. Use 'insert' or 'delete'."})

            formula_count = sum(
                1 for row in ws.iter_rows() for cell in row
                if getattr(cell, "data_type", None) == "f"
            )

            out = io.BytesIO()
            wb.save(out)
            wb.close()
            out.seek(0)

            url, name = await self._save_and_link(out.read(), out_name, __request__, __user__=__user__)
            if not url:
                return json.dumps({"error": self._err("could_not_save")})
            result = f"[{name}]({url})\n\n{msg} in {file_type.upper()} file."
            if formula_count:
                result += (
                    f"\nWarning: this sheet contains {formula_count} formula(s). Formulas are "
                    f"NOT adjusted for the row shift and may now reference the wrong cells -- "
                    f"verify them manually."
                )
            return result
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # PASSWORD PROTECT FILE
    # -----------------------------------------------------------------
    async def protect_file(
        self,
        file_id: str,
        password: str,
        output_filename: str = "",
        __user__=None,
        __request__=None,
    ) -> str:
        """Add a structural edit-lock to an XLSX or DOCX file. This does NOT encrypt the file
        or restrict opening it -- anyone can still open and read it in Excel/Word without the
        password. It only sets a workbook/worksheet protection flag (xlsx, via openpyxl) or a
        document-protection flag (docx) that restricts editing within the app, which is a much
        weaker guarantee than real password encryption and can be removed by re-saving the file
        or by tools that ignore the flag. Do not use this for anything that needs a genuine
        opening password.

        Args:
            file_id: The file ID of the Office file to protect
            password: Password to set
            output_filename: Optional output filename
        """
        try:
            file_data = _read_file_bytes(file_id)
            if file_data is None:
                return json.dumps({"error": f"Could not read file {file_id}"})

            try:
                conn2 = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
                row2 = conn2.execute(
                    "SELECT filename FROM file WHERE id = ?", (file_id,)
                ).fetchone()
                conn2.close()
                filename = _decode_filename(row2[0]) if row2 and row2[0] else file_id
            except Exception:
                filename = file_id

            file_type = _detect_type(filename)
            if file_type not in ("xlsx", "docx"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Only xlsx and docx supported."})

            import hashlib

            if file_type == "xlsx":
                import openpyxl
                from openpyxl.workbook.protection import WorkbookProtection

                wb = openpyxl.load_workbook(io.BytesIO(file_data))
                out_name = output_filename or os.path.splitext(filename)[0] + "_protected.xlsx"

                # Set workbook protection
                wb.security = WorkbookProtection(workbookPassword=password, lockStructure=True)

                # Protect all sheets. set_password() only assigns the password -- it does
                # NOT flip SheetProtection.sheet to True (its default is False), so without
                # this the saved file has a <sheetProtection> element that's actually
                # disabled and every sheet opens fully editable.
                for ws in wb.worksheets:
                    ws.protection.set_password(password)
                    ws.protection.sheet = True

                out = io.BytesIO()
                wb.save(out)
                wb.close()
                out.seek(0)
                protected_bytes = out.read()

            elif file_type == "docx":
                from docx import Document
                from docx.oxml.ns import qn
                from lxml import etree

                out_name = output_filename or os.path.splitext(filename)[0] + "_protected.docx"
                doc = Document(io.BytesIO(file_data))

                # Add write protection XML
                settings_element = doc.settings.element
                protection = etree.SubElement(
                    settings_element,
                    qn("w:documentProtection")
                )
                protection.set(qn("w:edit"), "readOnly")
                protection.set(qn("w:enforcement"), "1")
                pw_hash = hashlib.sha256(password.encode("utf-8")).hexdigest().upper()
                protection.set(qn("w:cryptProviderType"), "rsaAES")
                protection.set(qn("w:cryptAlgorithmClass"), "hash")
                protection.set(qn("w:cryptAlgorithmType"), "typeAny")
                protection.set(qn("w:cryptAlgorithmSid"), "4")
                protection.set(qn("w:cryptSpinCount"), "100000")
                protection.set(qn("w:hash"), pw_hash)
                protection.set(qn("w:salt"), "AAAAAAAAAAAAAAAAAAAAAA==")

                out = io.BytesIO()
                doc.save(out)
                out.seek(0)
                protected_bytes = out.read()

            url, fname = await self._save_and_link(protected_bytes, out_name, __request__, __user__=__user__)
            if url:
                return f"[{fname}]({url})\n\nPassword-protected {file_type.upper()} file created."
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    # -----------------------------------------------------------------
    # CREATE NEW FILE
    # -----------------------------------------------------------------
    async def create_file(
        self,
        file_type: str,
        content: str,
        output_filename: str = "",
        raw_text: bool = False,
        __user__=None,
        __request__=None,
    ) -> str:
        """Create a new Office file from scratch.

        For xlsx: content is CSV with headers on first line.
        For docx: supports markdown formatting — headings (#, ##, ###), bullets (-, *),
            and inline markdown (**bold**, *italic*, `code`).
        For pptx: content is split into slides on "---" lines. Within each "---"-separated
            block, the FIRST line becomes that slide's title and the remaining lines become its
            body text -- it is not one slide per line.

        Args:
            file_type: 'xlsx', 'docx', or 'pptx'
            content: Content specification
            output_filename: Output filename (the correct extension is added automatically if omitted)
            raw_text: If True, skip text formatting
        """
        fmt_mode = "preserve" if raw_text else "format"
        try:
            ftype = file_type.lower().replace(".", "")
            if ftype not in ("xlsx", "docx", "pptx"):
                return json.dumps({"error": f"Unsupported type: {file_type}. Use xlsx, docx, or pptx."})

            out_name = output_filename or f"document.{ftype}"
            if not out_name.lower().endswith(f".{ftype}"):
                out_name = f"{out_name}.{ftype}"
            out = io.BytesIO()

            if ftype == "xlsx":
                import openpyxl
                from openpyxl.styles import Font, PatternFill, Alignment
                wb = openpyxl.Workbook()
                ws = wb.active

                import csv as _csv_mod
                reader = _csv_mod.reader(io.StringIO(content))
                rows = list(reader)

                if rows:
                    # First row = headers (styled)
                    for j, h in enumerate(rows[0], 1):
                        c = ws.cell(row=1, column=j, value=h.strip())
                        c.font = Font(bold=True, color="FFFFFF", size=11)
                        c.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
                        c.alignment = Alignment(horizontal="center")

                    # Data rows
                    for i, rd in enumerate(rows[1:], 2):
                        for j, v in enumerate(rd, 1):
                            v = v.strip()
                            c = ws.cell(row=i, column=j)
                            if v == '':
                                c.value = None
                            else:
                                try:
                                    c.value = int(v)
                                except ValueError:
                                    try:
                                        c.value = float(v)
                                    except ValueError:
                                        if v.lower() == 'true':
                                            c.value = True
                                        elif v.lower() == 'false':
                                            c.value = False
                                        else:
                                            c.value = v
                            c.alignment = Alignment(
                                horizontal="center" if isinstance(c.value, (int, float)) else "left"
                            )

                    # Auto-fit columns
                    for col in ws.columns:
                        mx = max((len(str(c.value or "")) for c in col), default=5)
                        ws.column_dimensions[col[0].column_letter].width = min(mx + 3, 50)

                wb.save(out)
                wb.close()

            elif ftype == "docx":
                from docx import Document
                from docx.shared import Pt, RGBColor
                doc = Document()
                for raw_line in content.split("\n"):
                    line = raw_line
                    stripped = line.strip()

                    if stripped == "":
                        # Empty line → paragraph break
                        doc.add_paragraph("")
                        continue

                    # Heading detection (must start with # markers)
                    if stripped.startswith("### "):
                        text = stripped[4:]
                        doc.add_heading(text, level=3)
                        continue
                    if stripped.startswith("## "):
                        text = stripped[3:]
                        doc.add_heading(text, level=2)
                        continue
                    if stripped.startswith("# "):
                        text = stripped[2:]
                        doc.add_heading(text, level=1)
                        continue

                    # Bullet detection
                    if stripped.startswith("- ") or stripped.startswith("* "):
                        text = _format_text(stripped[2:], mode=fmt_mode)
                        doc.add_paragraph(text, style='List Bullet')
                        continue

                    # Body text — rich formatting via inline markdown
                    text = _format_text(line, mode=fmt_mode)
                    segments = _parse_inline_md(text)
                    p = doc.add_paragraph()
                    for seg_text, fmt in segments:
                        if not seg_text:
                            continue
                        run = p.add_run(seg_text)
                        if fmt.get('code'):
                            run.font.name = "Consolas"
                            run.font.size = Pt(9)
                        else:
                            run.font.size = Pt(11)
                        if fmt.get('bold'):
                            run.font.bold = True
                        if fmt.get('italic'):
                            run.font.italic = True
                        if fmt.get('link'):
                            run.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
                            run.font.underline = True
                doc.save(out)

            elif ftype == "pptx":
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                slide_specs = re.split(r'\n---\n|\r\n---\r\n|\n---\n', content)
                blank_layout = prs.slide_layouts[6]

                for spec in slide_specs:
                    spec = spec.strip()
                    if not spec:
                        continue
                    lines = spec.split("\n")
                    title = lines[0].strip()
                    body = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""

                    slide = prs.slides.add_slide(blank_layout)
                    txBox = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
                    )
                    tf = txBox.text_frame
                    tf.text = title
                    p = tf.paragraphs[0]
                    p.font.size = Inches(0.6)
                    p.font.bold = True

                    if body:
                        txBox2 = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
                        )
                        tf2 = txBox2.text_frame
                        tf2.text = body
                        for para in tf2.paragraphs:
                            para.font.size = Inches(0.3)

                prs.save(out)

            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__, __user__=__user__)
            if url:
                return f"[{name}]({url})\n\nCreated new {ftype.upper()} file."
            return json.dumps({"error": self._err("could_not_save")})

        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})



    async def generate_document(self, content: str, title: str = "Document", theme: str = "professional", typography: str = "modern", raw_text: bool = False, __user__=None, __request__=None) -> str:
        """Generate a professional Word document with modern styling, emojis, cards, and visual
        elements. Prefer this over add_content/create_file for docx when you want rich visual
        formatting -- those two only support plain markdown (headings/bold/italic/lists/tables).

        Headings preserve their original capitalization (no forced sentence case).
        Body text uses sentence case (first letter of each sentence capitalized, acronyms preserved).
        Supports inline markdown: **bold**, *italic*, `code`, [links](url), and ```code blocks```.

        Additional line-level patterns are auto-detected and rendered as rich visual elements
        (each must be on its own line; consecutive matching lines are grouped into one element):
        - Card box: a line starting with "> " (optionally "> 🎯 **Title**" on the first line)
        - KPI cards: a line containing both "|" and "%", e.g. "85% | Customer Satisfaction"
        - Timeline: pipe-delimited lines starting with an emoji/4-digit year, e.g.
          "📅 2024 | Event", or colon/dash date lines like "Jan 15: Event" / "2024 - Event"
        - Step guide: 3+ consecutive numbered lines ("1. Step one") get step-badge styling;
          fewer than 3 render as plain numbered paragraphs instead
        - Pull quote: a line wrapped in straight double quotes, followed by an "— Author" line
        - Comparison table: markdown pipe-table syntax "| Feature | A | B |" (yes/no/partial
          values get auto ✅/❌/⚠️ icons)
        - Progress bar: a "Label: NN%" line
        - Status badge: a line starting with "@success", "@warning", "@danger", or "@info"
        - Visual separator: a bare "---", "***", or "..." line
        - Bullet lines ("- "/"* ") get an auto icon based on leading keywords (done/fail/warn/note/idea)

        Args:
            content: Markdown content to convert to a document
            title: Document title / filename
            theme: Visual theme - professional, modern, creative, corporate, minimal, elegant, ocean, sunset, forest, midnight
            typography: Font preset - modern, classic, serif, sans
            raw_text: If True, skip all text formatting — text is used exactly as provided.
        """
        from docx import Document
        from docx.shared import Inches, Pt, Cm, RGBColor, Emu
        from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml.ns import qn, nsdecls
        from docx.oxml import parse_xml
        from lxml import etree
        import datetime, re as _re
        fmt_mode = "preserve" if raw_text else "format"
        
        doc = Document()
        section = doc.sections[0]
        section.top_margin = Cm(2.0)
        section.bottom_margin = Cm(2.0)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
        
        # --- Color Palettes (10+) ---
        palettes = {
            "professional": {"primary": "1F4E79", "accent": "2E75B6", "light": "D6E4F0", "text": "333333", "bg": "FFFFFF", "success": "27AE60", "warning": "F39C12", "danger": "E74C3C", "info": "3498DB"},
            "modern": {"primary": "2D3436", "accent": "6C5CE7", "light": "DFE6E9", "text": "2D3436", "bg": "FFFFFF", "success": "00B894", "warning": "FDCB6E", "danger": "FF7675", "info": "74B9FF"},
            "creative": {"primary": "E17055", "accent": "FDCB6E", "light": "FFF3E0", "text": "2D3436", "bg": "FFFAF5", "success": "00B894", "warning": "FDCB6E", "danger": "D63031", "info": "6C5CE7"},
            "corporate": {"primary": "003366", "accent": "CC0000", "light": "E8EEF4", "text": "1A1A1A", "bg": "FFFFFF", "success": "006633", "warning": "FF6600", "danger": "CC0000", "info": "0066CC"},
            "minimal": {"primary": "000000", "accent": "666666", "light": "F5F5F5", "text": "333333", "bg": "FAFAFA", "success": "333333", "warning": "666666", "danger": "000000", "info": "999999"},
            "elegant": {"primary": "4A235A", "accent": "8E44AD", "light": "F3E5F5", "text": "1A1A1A", "bg": "FDFBF7", "success": "27AE60", "warning": "D4AC0D", "danger": "C0392B", "info": "2980B9"},
            "ocean": {"primary": "0A3D62", "accent": "38ADA9", "light": "D1F2EB", "text": "1E272E", "bg": "F8FFFE", "success": "079992", "warning": "F6B93B", "danger": "E55039", "info": "3C6382"},
            "sunset": {"primary": "B33771", "accent": "FD7272", "light": "FFE4E4", "text": "2C3A47", "bg": "FFFBF5", "success": "58B19F", "warning": "F8B500", "danger": "E66767", "info": "786FA6"},
            "forest": {"primary": "1B4332", "accent": "40916C", "light": "D8F3DC", "text": "1A1A1A", "bg": "F7FFF7", "success": "2D6A4F", "warning": "B7B73F", "danger": "D00000", "info": "52B788"},
            "midnight": {"primary": "0F172A", "accent": "38BDF8", "light": "1E293B", "text": "F8FAFC", "bg": "0F172A", "success": "34D399", "warning": "FBBF24", "danger": "F87171", "info": "60A5FA"},
        }
        colors = palettes.get(theme, palettes["professional"])
        
        # --- Typography Presets ---
        fonts = {
            "modern": {"heading": "Calibri", "body": "Calibri"},
            "classic": {"heading": "Georgia", "body": "Calibri"},
            "serif": {"heading": "Cambria", "body": "Calibri"},
            "sans": {"heading": "Arial", "body": "Arial"},
        }
        font_pair = fonts.get(typography, fonts["modern"])
        
        def hex_to_rgb(h):
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        
        def add_colored_bar(doc, color_hex, height_pt=4):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(0)
            p.paragraph_format.space_after = Pt(0)
            run = p.add_run(" ")
            run.font.size = Pt(height_pt)
            pPr = p._p.get_or_add_pPr()
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}" w:val="clear"/>')
            pPr.append(shd)
        
        def add_card_box(doc, lines, border_color, bg_color, icon="", title=""):
            table = doc.add_table(rows=1, cols=1)
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            cell = table.rows[0].cells[0]
            cell.width = Inches(6.0)
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{bg_color}" w:val="clear"/>')
            tcPr.append(shd)
            borders = parse_xml(
                f'<w:tcBorders {nsdecls("w")}>'
                f'<w:left w:val="single" w:sz="24" w:space="8" w:color="{border_color}"/>'
                f'<w:top w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
                f'<w:bottom w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
                f'<w:right w:val="single" w:sz="4" w:space="0" w:color="{border_color}"/>'
                f'</w:tcBorders>'
            )
            tcPr.append(borders)
            p = cell.paragraphs[0]
            p.paragraph_format.space_before = Pt(8)
            p.paragraph_format.space_after = Pt(4)
            if icon or title:
                header_text = f"{icon} {title}" if icon and title else (icon or title)
                run = p.add_run(_format_text(header_text, mode=fmt_mode))
                run.font.size = Pt(13)
                run.font.bold = True
                run.font.name = font_pair["heading"]
                run.font.color.rgb = hex_to_rgb(border_color)
            for line in lines:
                p = cell.add_paragraph()
                p.paragraph_format.space_before = Pt(2)
                p.paragraph_format.space_after = Pt(2)
                text = _format_text(line, mode=fmt_mode)
                segments = _parse_inline_md(text)
                for seg_text, fmt in segments:
                    if not seg_text:
                        continue
                    run = p.add_run(seg_text)
                    if fmt.get('code'):
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)
                    else:
                        run.font.size = Pt(10)
                        run.font.name = font_pair["body"]
                    if fmt.get('bold'):
                        run.font.bold = True
                    if fmt.get('italic'):
                        run.font.italic = True
                    run.font.color.rgb = hex_to_rgb(colors["text"])
                    if fmt.get('link'):
                        run.font.color.rgb = hex_to_rgb("2563EB")
                        run.font.underline = True
            doc.add_paragraph()
        
        def add_kpi_card(doc, value, label, color_hex, icon=""):
            table = doc.add_table(rows=1, cols=1)
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            cell = table.rows[0].cells[0]
            cell.width = Inches(2.8)
            tc = cell._tc; tcPr = tc.get_or_add_tcPr()
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["light"]}" w:val="clear"/>')
            tcPr.append(shd)
            borders = parse_xml(f'<w:tcBorders {nsdecls("w")}><w:top w:val="single" w:sz="12" w:space="0" w:color="{color_hex}"/><w:left w:val="single" w:sz="4" w:space="0" w:color="DDDDDD"/><w:bottom w:val="single" w:sz="4" w:space="0" w:color="DDDDDD"/><w:right w:val="single" w:sz="4" w:space="0" w:color="DDDDDD"/></w:tcBorders>')
            tcPr.append(borders)
            p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_before = Pt(12)
            run = p.add_run(f"{icon} {value}" if icon else value)
            run.font.size = Pt(24); run.font.bold = True
            run.font.color.rgb = hex_to_rgb(color_hex); run.font.name = font_pair["heading"]
            p2 = cell.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p2.paragraph_format.space_after = Pt(12)
            label_text = _format_text(label, mode=fmt_mode)
            label_segments = _parse_inline_md(label_text)
            for seg_text, fmt in label_segments:
                if not seg_text:
                    continue
                run2 = p2.add_run(seg_text)
                run2.font.size = Pt(9)
                run2.font.color.rgb = hex_to_rgb(colors["text"])
                run2.font.name = font_pair["body"]
                if fmt.get('bold'):
                    run2.font.bold = True
                if fmt.get('italic'):
                    run2.font.italic = True
                if fmt.get('code'):
                    run2.font.name = "Consolas"
                    run2.font.size = Pt(9)
        
        def add_progress_bar(doc, label, percentage, color_hex):
            # Clamp: a value outside 0-100 (e.g. "Growth: 150%") would otherwise produce a
            # negative table-cell width, which Word treats as corrupt content.
            percentage = max(0, min(100, percentage))
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(4)
            p.paragraph_format.space_after = Pt(2)
            label_text = _format_text(label, mode=fmt_mode)
            label_segments = _parse_inline_md(label_text)
            for seg_text, fmt in label_segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                run.font.size = Pt(10); run.font.bold = True
                run.font.name = font_pair["body"]; run.font.color.rgb = hex_to_rgb(colors["text"])
                if fmt.get('italic'):
                    run.font.italic = True
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
            run = p.add_run(f": {percentage}%")
            run.font.size = Pt(10); run.font.bold = True
            run.font.name = font_pair["body"]; run.font.color.rgb = hex_to_rgb(colors["text"])
            table = doc.add_table(rows=1, cols=2)
            table.alignment = WD_TABLE_ALIGNMENT.LEFT
            filled = table.rows[0].cells[0]; empty = table.rows[0].cells[1]
            filled.width = Inches(5.0 * percentage / 100)
            empty.width = Inches(5.0 * (100 - percentage) / 100)
            tcF = filled._tc; tcPrF = tcF.get_or_add_tcPr()
            shdF = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{color_hex}" w:val="clear"/>')
            tcPrF.append(shdF)
            tcE = empty._tc; tcPrE = tcE.get_or_add_tcPr()
            shdE = parse_xml(f'<w:shd {nsdecls("w")} w:fill="E0E0E0" w:val="clear"/>')
            tcPrE.append(shdE)
            pF = filled.paragraphs[0]; runF = pF.add_run(" "); runF.font.size = Pt(6)
            pE = empty.paragraphs[0]; runE = pE.add_run(" "); runE.font.size = Pt(6)
        
        def add_step_guide(doc, steps):
            for i, step in enumerate(steps, 1):
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(2)
                run = p.add_run(f"  {i}  ")
                run.font.size = Pt(14); run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255); run.font.name = font_pair["heading"]
                rPr = run._r.get_or_add_rPr()
                shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["primary"]}" w:val="clear"/>')
                rPr.append(shd)
                text = _format_text(step, mode=fmt_mode)
                segments = _parse_inline_md(text)
                first = True
                for seg_text, fmt in segments:
                    if not seg_text:
                        continue
                    prefix = "  " if first else ""
                    first = False
                    run = p.add_run(prefix + seg_text)
                    if fmt.get('code'):
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)
                    else:
                        run.font.size = Pt(11)
                        run.font.name = font_pair["body"]
                    if fmt.get('bold'):
                        run.font.bold = True
                    if fmt.get('italic'):
                        run.font.italic = True
                    run.font.color.rgb = hex_to_rgb(colors["text"])
                    if fmt.get('link'):
                        run.font.color.rgb = hex_to_rgb("2563EB")
                        run.font.underline = True
        
        def add_pull_quote(doc, text, author=""):
            table = doc.add_table(rows=1, cols=1)
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            cell = table.rows[0].cells[0]; cell.width = Inches(5.5)
            tc = cell._tc; tcPr = tc.get_or_add_tcPr()
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["light"]}" w:val="clear"/>')
            tcPr.append(shd)
            borders = parse_xml(f'<w:tcBorders {nsdecls("w")}><w:left w:val="single" w:sz="36" w:space="12" w:color="{colors["accent"]}"/><w:top w:val="none" w:sz="0" w:space="0" w:color="auto"/><w:bottom w:val="none" w:sz="0" w:space="0" w:color="auto"/><w:right w:val="none" w:sz="0" w:space="0" w:color="auto"/></w:tcBorders>')
            tcPr.append(borders)
            p = cell.paragraphs[0]; p.paragraph_format.space_before = Pt(12)
            p.paragraph_format.space_after = Pt(4)
            text_formatted = _format_text(text, mode=fmt_mode)
            segments = _parse_inline_md(text_formatted)
            # Opening quote
            run = p.add_run('"')
            run.font.size = Pt(14); run.font.italic = True
            run.font.name = font_pair["heading"]; run.font.color.rgb = hex_to_rgb(colors["text"])
            for seg_text, fmt in segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                run.font.size = Pt(14); run.font.italic = True
                run.font.name = font_pair["heading"]
                run.font.color.rgb = hex_to_rgb(colors["text"])
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
            # Closing quote
            run = p.add_run('"')
            run.font.size = Pt(14); run.font.italic = True
            run.font.name = font_pair["heading"]; run.font.color.rgb = hex_to_rgb(colors["text"])
            if author:
                p2 = cell.add_paragraph(); p2.paragraph_format.space_after = Pt(12)
                p2.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                author_text = _format_text(author, mode=fmt_mode)
                author_segments = _parse_inline_md(author_text)
                first = True
                for seg_text, fmt in author_segments:
                    if not seg_text:
                        continue
                    prefix = "\u2014 " if first else ""
                    first = False
                    run = p2.add_run(prefix + seg_text)
                    run.font.size = Pt(10); run.font.name = font_pair["body"]
                    run.font.color.rgb = hex_to_rgb(colors["accent"])
                    if fmt.get('bold'):
                        run.font.bold = True
                    if fmt.get('italic'):
                        run.font.italic = True
                    if fmt.get('link'):
                        run.font.color.rgb = hex_to_rgb("2563EB")
                        run.font.underline = True
                    if fmt.get('code'):
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)
            doc.add_paragraph()
        
        def add_comparison_table(doc, headers, rows):
            table = doc.add_table(rows=len(rows)+1, cols=len(headers))
            table.style = 'Colorful Grid Accent 1'
            table.alignment = WD_TABLE_ALIGNMENT.CENTER
            for j, h in enumerate(headers):
                cell = table.rows[0].cells[j]
                cell.text = ''
                p = cell.paragraphs[0]
                segments = _parse_inline_md(_format_text(h, mode=fmt_mode))
                for seg_text, fmt in segments:
                    if not seg_text:
                        continue
                    r = p.add_run(seg_text)
                    r.font.size = Pt(10); r.font.bold = True; r.font.name = font_pair["heading"]
                    if fmt.get('italic'):
                        r.font.italic = True
                    if fmt.get('code'):
                        r.font.name = "Consolas"
                        r.font.size = Pt(9)
            for i, row in enumerate(rows):
                for j, val in enumerate(row):
                    if j < len(headers):
                        cell = table.rows[i+1].cells[j]
                        display = val
                        if val.lower() in ("yes","true","sim","si","oui","ja"): display = "\u2705 " + val
                        elif val.lower() in ("no","false","nao","non","nein"): display = "\u274c " + val
                        elif val.lower() in ("partial","maybe","talvez"): display = "\u26a0\ufe0f " + val
                        cell.text = ''
                        p = cell.paragraphs[0]
                        segments = _parse_inline_md(_format_text(display, mode=fmt_mode))
                        for seg_text, fmt in segments:
                            if not seg_text:
                                continue
                            r = p.add_run(seg_text)
                            r.font.size = Pt(10); r.font.name = font_pair["body"]
                            if fmt.get('bold'):
                                r.font.bold = True
                            if fmt.get('italic'):
                                r.font.italic = True
                            if fmt.get('code'):
                                r.font.name = "Consolas"
                                r.font.size = Pt(9)
            doc.add_paragraph()
        
        def add_timeline(doc, events):
            for i, (date, desc) in enumerate(events):
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(4)
                p.paragraph_format.space_after = Pt(2)
                icon = "\u25cf" if i == 0 else "\u25cb"
                run = p.add_run(f"  {icon}  ")
                run.font.size = Pt(12); run.font.color.rgb = hex_to_rgb(colors["accent"])
                run.font.name = font_pair["heading"]
                date_text = _format_text(date, mode=fmt_mode)
                date_segments = _parse_inline_md(date_text)
                for seg_text, fmt in date_segments:
                    if not seg_text:
                        continue
                    run = p.add_run(seg_text)
                    run.font.size = Pt(10); run.font.bold = True
                    run.font.name = font_pair["body"]; run.font.color.rgb = hex_to_rgb(colors["primary"])
                    if fmt.get('bold'):
                        run.font.bold = True
                    if fmt.get('italic'):
                        run.font.italic = True
                    if fmt.get('code'):
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)
                    if fmt.get('link'):
                        run.font.color.rgb = hex_to_rgb("2563EB")
                        run.font.underline = True
                run_sep = p.add_run("  ")
                run_sep.font.size = Pt(10)
                run_sep.font.name = font_pair["body"]
                desc_text = _format_text(desc, mode=fmt_mode)
                desc_segments = _parse_inline_md(desc_text)
                for seg_text, fmt in desc_segments:
                    if not seg_text:
                        continue
                    run = p.add_run(seg_text)
                    run.font.size = Pt(10); run.font.name = font_pair["body"]
                    run.font.color.rgb = hex_to_rgb(colors["text"])
                    if fmt.get('bold'):
                        run.font.bold = True
                    if fmt.get('italic'):
                        run.font.italic = True
                    if fmt.get('code'):
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)
                    if fmt.get('link'):
                        run.font.color.rgb = hex_to_rgb("2563EB")
                        run.font.underline = True
        
        def add_status_badge(doc, text, status="info"):
            colors_map = {"success": colors["success"], "warning": colors["warning"], "danger": colors["danger"], "info": colors["info"]}
            icons = {"success": "\U0001f7e2", "warning": "\U0001f7e1", "danger": "\U0001f534", "info": "\U0001f535"}
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(2)
            p.paragraph_format.space_after = Pt(2)
            badge_color = hex_to_rgb(colors_map.get(status, colors["info"]))
            icon_text = icons.get(status, icons['info'])
            run = p.add_run(f" {icon_text} ")
            run.font.size = Pt(10); run.font.bold = True
            run.font.name = font_pair["body"]; run.font.color.rgb = badge_color
            badge_text = _format_text(text, mode=fmt_mode)
            badge_segments = _parse_inline_md(badge_text)
            for seg_text, fmt in badge_segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                run.font.size = Pt(10); run.font.bold = True
                run.font.name = font_pair["body"]; run.font.color.rgb = badge_color
                if fmt.get('italic'):
                    run.font.italic = True
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
        
        def add_visual_separator(doc, style="dots"):
            separators = {"dots": "\u25cf \u25cf \u25cf", "line": "\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501\u2501", "dash": "\u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500 \u2500"}
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_before = Pt(8)
            p.paragraph_format.space_after = Pt(8)
            run = p.add_run(separators.get(style, separators["dots"]))
            run.font.size = Pt(8); run.font.color.rgb = hex_to_rgb(colors["accent"])
        
        def _add_rich_paragraph(doc, text, style=None, font_name="Calibri", font_size=11, color=None):
            """Add a paragraph with rich formatting from inline markdown parsing."""
            segments = _parse_inline_md(text)
            if style and style.startswith('Heading'):
                try:
                    level = int(style.split()[-1])
                except (ValueError, IndexError):
                    level = 1
                p = doc.add_heading('', level=level)
                for run in p.runs:
                    run.text = ''
            else:
                p = doc.add_paragraph(style=style) if style else doc.add_paragraph()
            for seg_text, fmt in segments:
                if not seg_text:
                    continue
                run = p.add_run(seg_text)
                if fmt.get('code'):
                    run.font.name = "Consolas"
                    run.font.size = Pt(9)
                else:
                    run.font.name = font_name
                    run.font.size = Pt(font_size)
                if fmt.get('bold'):
                    run.font.bold = True
                if fmt.get('italic'):
                    run.font.italic = True
                if color:
                    if isinstance(color, str):
                        run.font.color.rgb = hex_to_rgb(color)
                    else:
                        run.font.color.rgb = color
                if fmt.get('link'):
                    run.font.color.rgb = hex_to_rgb("2563EB")
                    run.font.underline = True
            return p
        
        # --- Default Style ---
        style = doc.styles['Normal']
        font = style.font; font.name = font_pair["body"]; font.size = Pt(11)
        font.color.rgb = hex_to_rgb(colors["text"])
        style.paragraph_format.space_after = Pt(8)
        style.paragraph_format.line_spacing = 1.15
        
        for i in range(1, 4):
            hs = doc.styles[f'Heading {i}']
            hs.font.name = font_pair["heading"]
            hs.font.color.rgb = hex_to_rgb(colors["primary"])
            sizes = {1: 24, 2: 18, 3: 14}
            hs.font.size = Pt(sizes.get(i, 14))
            spaces = {1: (24, 12), 2: (18, 8), 3: (12, 6)}
            hs.paragraph_format.space_before = Pt(spaces[i][0])
            hs.paragraph_format.space_after = Pt(spaces[i][1])
        
        # --- Cover Page ---
        add_colored_bar(doc, colors["primary"], 4)
        cover_table = doc.add_table(rows=1, cols=1)
        cover_table.alignment = WD_TABLE_ALIGNMENT.CENTER
        cell = cover_table.rows[0].cells[0]; cell.width = Inches(6.5)
        tc = cell._tc; tcPr = tc.get_or_add_tcPr()
        shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["primary"]}" w:val="clear"/>')
        tcPr.append(shd)
        p = cell.paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(40); p.paragraph_format.space_after = Pt(8)
        run = p.add_run(title)
        run.font.size = Pt(28); run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255); run.font.name = font_pair["heading"]
        p2 = cell.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p2.paragraph_format.space_after = Pt(30)
        run2 = p2.add_run(datetime.datetime.now().strftime("%B %d, %Y"))
        run2.font.size = Pt(12); run2.font.color.rgb = RGBColor(200, 200, 200)
        run2.font.name = font_pair["body"]
        doc.add_paragraph()
        
        # --- Auto TOC ---
        toc_added = False
        
        # --- Process Content ---
        lines = content.split('\n')
        in_card = False; card_lines = []; card_icon = ""; card_title = ""
        in_kpi = False; kpi_data = []
        in_timeline = False; timeline_data = []
        in_steps = False; step_lines = []
        in_quote = False; quote_text = ""; quote_author = ""
        in_comparison = False; comp_headers = []; comp_rows = []
        in_progress = False; progress_data = []
        in_code_block = False; code_lines = []
        # Group 2 is capped at 100 chars (real timeline/agenda entries are short; a match
        # longer than that is far more likely an ordinary sentence that happens to start
        # with a date, e.g. "2024: was a great year for the company because..."). A
        # separate numeric-range check below further rejects "3-4 people attended"-style
        # false positives from the M/D date alternative.
        _timeline_re = _re.compile(r'^(\d{4}|\d{1,2}[/-]\d{1,2}|(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*)\s*[-:]\s*(.{1,100})$', _re.IGNORECASE)
        _md_date_re = _re.compile(r'^(\d{1,2})[/-](\d{1,2})$')

        def _is_plausible_timeline_entry(token, description):
            # Bare years and month names are already unambiguous; only the M/D-style
            # alternative needs range validation to rule out ordinary number ranges
            # ("3-4 people", "10-15 minutes") being mistaken for a date.
            md = _md_date_re.match(token)
            if md:
                a, b = int(md.group(1)), int(md.group(2))
                if not (1 <= a <= 12 and 1 <= b <= 31):
                    return False
            # The 100-char cap on the regex alone doesn't catch short-but-still-a-sentence
            # cases ("2024: was a great year for the company because it grew fast" is only
            # 55 chars). Real timeline/agenda entries are short phrases (a handful of
            # words), not full sentences -- cap at 8 words as the more reliable signal.
            if len(description.split()) > 8:
                return False
            return True

        def flush_kpi():
            # KPI blocks had no flush handler at blank-line or end-of-content boundaries
            # (unlike card/steps/timeline/comparison/progress, which all do) -- a document
            # ending in a KPI block, or with a blank line after one, silently lost it.
            nonlocal kpi_data, in_kpi
            if not kpi_data:
                in_kpi = False
                return
            kpi_table = doc.add_table(rows=1, cols=len(kpi_data))
            kpi_table.alignment = WD_TABLE_ALIGNMENT.CENTER
            for j, kpi in enumerate(kpi_data):
                val = kpi[0] if len(kpi) > 0 else ""
                lbl = kpi[1] if len(kpi) > 1 else ""
                icon_match = _re.match(r'^([\U0001F300-\U0001F9FF]\s*)', val)
                icon = icon_match.group(1).strip() if icon_match else ""
                if icon: val = val[len(icon):].strip()
                cell = kpi_table.rows[0].cells[j]
                cell.width = Inches(2.8)
                tcC = cell._tc; tcPrC = tcC.get_or_add_tcPr()
                shdC = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["light"]}" w:val="clear"/>')
                tcPrC.append(shdC)
                pC = cell.paragraphs[0]; pC.alignment = WD_ALIGN_PARAGRAPH.CENTER
                pC.paragraph_format.space_before = Pt(12)
                runC = pC.add_run(f"{icon} {val}" if icon else val)
                runC.font.size = Pt(24); runC.font.bold = True
                runC.font.color.rgb = hex_to_rgb(colors["primary"])
                runC.font.name = font_pair["heading"]
                pC2 = cell.add_paragraph(); pC2.alignment = WD_ALIGN_PARAGRAPH.CENTER
                pC2.paragraph_format.space_after = Pt(12)
                runC2 = pC2.add_run(_format_text(lbl, mode=fmt_mode))
                runC2.font.size = Pt(9); runC2.font.color.rgb = hex_to_rgb(colors["text"])
                runC2.font.name = font_pair["body"]
            doc.add_paragraph()
            kpi_data = []; in_kpi = False

        def flush_quote():
            # Same gap as KPI -- a document ending in (or with a blank line after) a pull
            # quote lost it entirely, since only the "next line starts the attribution"
            # branch ever rendered it.
            nonlocal quote_text, quote_author, in_quote
            if not quote_text:
                in_quote = False
                return
            add_pull_quote(doc, quote_text, quote_author)
            in_quote = False; quote_text = ""; quote_author = ""
        
        for line in lines:
            line = line.strip()
            if not line:
                if in_code_block:
                    code_lines.append("")
                    continue
                if in_card:
                    add_card_box(doc, card_lines, colors["accent"], colors["light"], card_icon, card_title)
                    card_lines = []; in_card = False
                if in_steps:
                    if len(step_lines) >= 3:
                        add_step_guide(doc, step_lines)
                    else:
                        for i, step_text in enumerate(step_lines, 1):
                            _add_rich_paragraph(doc, f"{i}. {_format_text(step_text, mode=fmt_mode)}", font_name=font_pair["body"], font_size=11, color=colors["text"])
                    step_lines = []; in_steps = False
                if in_timeline:
                    add_timeline(doc, timeline_data)
                    timeline_data = []; in_timeline = False
                if in_comparison:
                    add_comparison_table(doc, comp_headers, comp_rows)
                    comp_headers = []; comp_rows = []; in_comparison = False
                if in_progress:
                    for lbl, pct, clr in progress_data:
                        add_progress_bar(doc, lbl, pct, clr)
                    progress_data = []; in_progress = False
                if in_kpi:
                    flush_kpi()
                if in_quote:
                    flush_quote()
                continue

            # Code block: ``` ... ```
            if line.startswith('```'):
                if not in_code_block:
                    in_code_block = True
                    code_lines = []
                else:
                    for code_line in code_lines:
                        p = doc.add_paragraph()
                        p.paragraph_format.space_before = Pt(1)
                        p.paragraph_format.space_after = Pt(1)
                        run = p.add_run(code_line)
                        run.font.name = "Consolas"
                        run.font.size = Pt(9)
                        run.font.color.rgb = hex_to_rgb(colors["text"])
                    in_code_block = False
                    code_lines = []
                continue
            if in_code_block:
                code_lines.append(line)
                continue
            
            # Card: > 📊 **Title** or > content
            if line.startswith('> ') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_comparison and not in_progress:
                in_card = True
                text = line[2:]
                m = _re.match(r'^([\U0001F300-\U0001F9FF]\s*)?\*\*(.+?)\*\*', text)
                if m:
                    card_icon = m.group(1).strip() if m.group(1) else ""
                    card_title = m.group(2)
                    remaining = text[m.end():].strip()
                    if remaining: card_lines.append(remaining)
                else:
                    card_lines.append(text)
                continue
            elif in_card and line.startswith('> '):
                card_lines.append(line[2:])
                continue
            
            # KPI: 📊 85% | Customer Satisfaction
            if '|' in line and '%' in line and not in_card and not in_timeline and not in_steps and not in_quote and not in_comparison and not in_progress:
                parts = [p.strip() for p in line.split('|')]
                if len(parts) >= 2 and any('%' in p for p in parts):
                    in_kpi = True
                    kpi_data.append(parts)
                    continue
            elif in_kpi and '|' in line and '%' in line:
                parts = [p.strip() for p in line.split('|')]
                kpi_data.append(parts)
                continue
            elif in_kpi:
                flush_kpi()
            
            # Timeline: 📅 2024 | Event description (pipe format)
            if '|' in line and _re.match(r'^[\U0001F300-\U0001F9FF\s]*\d{4}', line) and not in_card and not in_kpi and not in_steps and not in_quote and not in_comparison and not in_progress:
                in_timeline = True
                parts = [p.strip() for p in line.split('|', 1)]
                timeline_data.append((parts[0], parts[1] if len(parts) > 1 else ""))
                continue
            elif in_timeline and '|' in line and _re.match(r'^[\U0001F300-\U0001F9FF\s]*\d{4}', line):
                parts = [p.strip() for p in line.split('|', 1)]
                timeline_data.append((parts[0], parts[1] if len(parts) > 1 else ""))
                continue
            elif in_timeline:
                add_timeline(doc, timeline_data)
                timeline_data = []; in_timeline = False
            
            # Timeline (colon/dash format): 2024 - Event or Jan 15: Event
            if not in_timeline and not in_card and not in_kpi and not in_steps and not in_quote and not in_comparison and not in_progress:
                tl_m = _timeline_re.match(line)
                if tl_m and _is_plausible_timeline_entry(tl_m.group(1), tl_m.group(2)):
                    in_timeline = True
                    timeline_data.append((tl_m.group(1), tl_m.group(2)))
                    continue
            
            # Steps: 1. Step one / 2. Step two
            if _re.match(r'^\d+\.\s', line) and not in_card and not in_kpi and not in_timeline and not in_quote and not in_comparison and not in_progress:
                in_steps = True
                step_lines.append(_re.sub(r'^\d+\.\s', '', line))
                continue
            elif in_steps and _re.match(r'^\d+\.\s', line):
                step_lines.append(_re.sub(r'^\d+\.\s', '', line))
                continue
            elif in_steps:
                if len(step_lines) >= 3:
                    add_step_guide(doc, step_lines)
                else:
                    for i, step_text in enumerate(step_lines, 1):
                        _add_rich_paragraph(doc, f"{i}. {_format_text(step_text, mode=fmt_mode)}", font_name=font_pair["body"], font_size=11, color=colors["text"])
                step_lines = []; in_steps = False

            # Pull quote: "Quote text" — Author
            if line.startswith('"') and line.endswith('"') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_comparison and not in_progress:
                in_quote = True
                quote_text = line.strip('"')
                continue
            elif in_quote and (line.startswith('\u2014') or (line.strip() and not line.startswith('"') and not line.startswith('#') and not line.startswith('-') and not line.startswith('*') and not line.startswith('>') and not line.startswith('|') and not line.startswith('@') and not line.startswith('```'))):
                quote_author = line.lstrip('\u2014 ').strip()
                flush_quote()
                continue
            
            # Comparison table: | Feature | A | B |
            if line.startswith('|') and line.endswith('|') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_progress:
                cells = [c.strip() for c in line.split('|')[1:-1]]
                if all(_re.fullmatch(r':?-{2,}:?', c) for c in cells):
                    continue
                if not in_comparison:
                    in_comparison = True
                    comp_headers = cells
                else:
                    comp_rows.append(cells)
                continue
            elif in_comparison and line.startswith('|') and line.endswith('|'):
                cells = [c.strip() for c in line.split('|')[1:-1]]
                if not all(_re.fullmatch(r':?-{2,}:?', c) for c in cells):
                    comp_rows.append(cells)
                continue
            elif in_comparison:
                add_comparison_table(doc, comp_headers, comp_rows)
                comp_headers = []; comp_rows = []; in_comparison = False
            
            # Progress: Label: 75%
            # Anchored to the WHOLE line (not just "contains a ':' and a '%' somewhere"),
            # and requires the percentage to be the entire remainder after the label --
            # otherwise ordinary prose like "Note: sales grew by 12% last quarter" gets
            # misread as a progress bar labeled "Note" at 12%.
            _progress_m = _re.match(r'^(.+?):\s*(\d{1,3})%\s*$', line)
            if _progress_m and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_comparison:
                in_progress = True
                progress_data.append((_progress_m.group(1).strip(), int(_progress_m.group(2)), colors["accent"]))
                continue
            elif in_progress and _progress_m:
                progress_data.append((_progress_m.group(1).strip(), int(_progress_m.group(2)), colors["accent"]))
                continue
            elif in_progress:
                for lbl, pct, clr in progress_data:
                    add_progress_bar(doc, lbl, pct, clr)
                progress_data = []; in_progress = False
            
            # Status badge: @success Task complete
            if line.startswith('@') and not in_card and not in_kpi and not in_timeline and not in_steps and not in_quote and not in_comparison and not in_progress:
                parts = line[1:].split(None, 1)
                status = parts[0].lower() if parts else "info"
                text = parts[1] if len(parts) > 1 else ""
                add_status_badge(doc, text, status)
                continue
            
            # Visual separator: --- or ***
            if line in ('---', '***', '...'):
                style_map = {'---': 'line', '***': 'dots', '...': 'dash'}
                add_visual_separator(doc, style_map.get(line, 'dots'))
                continue
            
            # Headings with emoji
            if line.startswith('# ') or line.startswith('## ') or line.startswith('### '):
                level = len(line) - len(line.lstrip('#'))
                text = line.lstrip('#').strip()
                heading_level = min(level, 3)
                _add_rich_paragraph(doc, text, style=f'Heading {heading_level}', font_name=font_pair["heading"], font_size=sizes.get(heading_level, 14), color=colors["primary"])
                continue
            
            # Icon bullets
            if line.startswith('- ') or line.startswith('* '):
                bullet_text = line[2:].strip()
                text = _format_text(bullet_text, mode=fmt_mode)
                icon = "\u2022"
                if text.lower().startswith(('done','complete','yes','ok','success','conclu')):
                    icon = "\u2705"
                elif text.lower().startswith(('no','not','fail','error','wrong','nao')):
                    icon = "\u274c"
                elif text.lower().startswith(('warn','caution','careful','cuidado')):
                    icon = "\u26a0\ufe0f"
                elif text.lower().startswith(('note','info','note','nota')):
                    icon = "\U0001f4cc"
                elif text.lower().startswith(('idea','tip','suggestion','sugest')):
                    icon = "\U0001f4a1"
                _add_rich_paragraph(doc, f"{icon} {text}", style='List Bullet', font_name=font_pair["body"], font_size=11, color=colors["text"])
                continue
            
            # Regular paragraph
            text = _format_text(line, mode=fmt_mode)
            _add_rich_paragraph(doc, text, font_name=font_pair["body"], font_size=11, color=colors["text"])
        
        # Render remaining
        if in_card:
            add_card_box(doc, card_lines, colors["accent"], colors["light"], card_icon, card_title)
        if in_steps:
            if len(step_lines) >= 3:
                add_step_guide(doc, step_lines)
            else:
                for i, step_text in enumerate(step_lines, 1):
                    _add_rich_paragraph(doc, f"{i}. {_format_text(step_text, mode=fmt_mode)}", font_name=font_pair["body"], font_size=11, color=colors["text"])
        if in_timeline:
            add_timeline(doc, timeline_data)
        if in_comparison:
            add_comparison_table(doc, comp_headers, comp_rows)
        if in_progress:
            for lbl, pct, clr in progress_data:
                add_progress_bar(doc, lbl, pct, clr)
        if in_kpi:
            flush_kpi()
        if in_quote:
            flush_quote()
        if in_code_block:
            for code_line in code_lines:
                p = doc.add_paragraph()
                p.paragraph_format.space_before = Pt(1)
                p.paragraph_format.space_after = Pt(1)
                run = p.add_run(code_line)
                run.font.name = "Consolas"
                run.font.size = Pt(9)
                run.font.color.rgb = hex_to_rgb(colors["text"])
        
        # --- Branded Footer ---
        footer = section.footer
        footer.is_linked_to_previous = False
        ft = footer.add_table(rows=1, cols=3, width=Inches(6.5))
        ft.alignment = WD_TABLE_ALIGNMENT.CENTER
        c1 = ft.rows[0].cells[0]; c1.width = Inches(2)
        r1 = c1.paragraphs[0]; r1.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run1 = r1.add_run("Edit Office Files")
        run1.font.size = Pt(8); run1.font.color.rgb = hex_to_rgb(colors["accent"])
        run1.font.name = font_pair["body"]
        c2 = ft.rows[0].cells[1]; c2.width = Inches(2.5)
        r2 = c2.paragraphs[0]; r2.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run2 = r2.add_run("Page ")
        run2.font.size = Pt(8); run2.font.color.rgb = hex_to_rgb(colors["accent"])
        run2.font.name = font_pair["body"]
        fldChar1 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="begin"/>')
        run2._r.append(fldChar1)
        instrText = parse_xml(f'<w:instrText {nsdecls("w")} xml:space="preserve"> PAGE </w:instrText>')
        run2._r.append(instrText)
        fldChar2 = parse_xml(f'<w:fldChar {nsdecls("w")} w:fldCharType="end"/>')
        run2._r.append(fldChar2)
        c3 = ft.rows[0].cells[2]; c3.width = Inches(2)
        r3 = c3.paragraphs[0]; r3.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run3 = r3.add_run(datetime.datetime.now().strftime("%Y-%m-%d"))
        run3.font.size = Pt(8); run3.font.color.rgb = hex_to_rgb(colors["accent"])
        run3.font.name = font_pair["body"]
        
        # --- Save ---
        file_bytes = io.BytesIO()
        doc.save(file_bytes)
        file_bytes.seek(0)
        url, fname = await self._save_and_link(file_bytes.getvalue(), _ensure_ext(title, "docx"), __request__, __user__=__user__)
        if url:
            return f"Document created: [{fname}]({url})"
        return json.dumps({"error": self._err("could_not_save")})


    async def generate_slides(self, content: str, title: str = "Presentation", theme: str = "modern", raw_text: bool = False, __user__=None, __request__=None) -> str:
        """Generate professional PowerPoint slides with modern design.

        Args:
            content: Markdown content (headings become slides, bullets become content)
            title: Presentation title for the first slide
            theme: Visual theme - modern, light, dark, corporate, creative, minimal
            raw_text: If True, skip text formatting
        Returns:
            Markdown link to the generated file
        """
        fmt_mode = "preserve" if raw_text else "format"
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu, Cm
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE
            import datetime

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # --- Color Themes ---
            themes = {
                "modern": {"bg": "1A1A2E", "accent": "E94560", "text": "FFFFFF", "subtitle": "B0B0B0", "card": "16213E"},
                "light": {"bg": "FFFFFF", "accent": "2563EB", "text": "1E293B", "subtitle": "64748B", "card": "F1F5F9"},
                "dark": {"bg": "0F172A", "accent": "38BDF8", "text": "F8FAFC", "subtitle": "94A3B8", "card": "1E293B"},
                "corporate": {"bg": "FFFFFF", "accent": "003366", "text": "1A1A1A", "subtitle": "666666", "card": "F0F4F8"},
                "creative": {"bg": "FFF8F0", "accent": "FF6B6B", "text": "2D3436", "subtitle": "636E72", "card": "FFEAA7"},
                "minimal": {"bg": "FAFAFA", "accent": "000000", "text": "333333", "subtitle": "999999", "card": "F0F0F0"},
            }
            colors = themes.get(theme, themes["modern"])

            def hex_to_rgb(h):
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

            def set_slide_bg(slide, color_hex):
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = hex_to_rgb(color_hex)

            def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
                txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(font_size)
                p.font.bold = bold
                p.font.color.rgb = hex_to_rgb(color or colors["text"])
                p.font.name = 'Calibri'
                p.alignment = alignment
                return txBox

            def add_accent_bar(slide, left, top, width, height, color_hex):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
                shape.line.fill.background()
                return shape

            # --- Slide 1: Title Slide ---
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide, colors["bg"])

            add_accent_bar(slide, 0, 3.0, 13.333, 0.06, colors["accent"])
            add_text_box(slide, 1.5, 1.5, 10, 1.5, title, font_size=44, bold=True, color=colors["text"])
            add_text_box(slide, 1.5, 3.3, 10, 0.8, "Generated %s" % datetime.datetime.now().strftime("%B %d, %Y"), font_size=18, color=colors["subtitle"])

            # --- Process content into slides ---
            lines = content.split('\n')
            current_slide_lines = []
            slide_count = 1

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if line.startswith('# ') or line.startswith('## '):
                    if current_slide_lines:
                        _render_content_slide(prs, current_slide_lines, colors, hex_to_rgb, add_text_box, add_accent_bar, set_slide_bg, slide_count, fmt_mode=fmt_mode)
                        slide_count += 1
                    current_slide_lines = [line]
                else:
                    current_slide_lines.append(line)

            if current_slide_lines:
                _render_content_slide(prs, current_slide_lines, colors, hex_to_rgb, add_text_box, add_accent_bar, set_slide_bg, slide_count, fmt_mode=fmt_mode)

            out = io.BytesIO()
            prs.save(out)
            out.seek(0)
            url, fname = await self._save_and_link(out.getvalue(), _ensure_ext(title, "pptx"), __request__, __user__=__user__)
            if url:
                return "Presentation created: [%s](%s)" % (fname, url)
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def generate_spreadsheet(self, content: str, title: str = "Spreadsheet", theme: str = "professional", raw_text: bool = False, __user__=None, __request__=None) -> str:
        """Generate a professional Excel spreadsheet with modern styling.

        Args:
            content: CSV or tab-delimited data (first row = headers, rest = data)
            title: Spreadsheet title / filename
            theme: Visual theme - professional, modern, corporate, minimal, colorful, pastel
            raw_text: If True, skip text formatting
        Returns:
            Markdown link to the generated file
        """
        fmt_mode = "preserve" if raw_text else "format"
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
            from openpyxl.worksheet.table import Table, TableStyleInfo
            import io

            # --- Color Themes ---
            themes = {
                "professional": {"header": "1F4E79", "header_font": "FFFFFF", "alt_row": "F0F5FA", "border": "CCCCCC", "accent": "2E75B6"},
                "modern": {"header": "2D3436", "header_font": "FFFFFF", "alt_row": "F5F5F5", "border": "DFE6E9", "accent": "6C5CE7"},
                "corporate": {"header": "003366", "header_font": "FFFFFF", "alt_row": "E8EEF4", "border": "B0C4DE", "accent": "CC0000"},
                "minimal": {"header": "333333", "header_font": "FFFFFF", "alt_row": "F8F8F8", "border": "DDDDDD", "accent": "666666"},
                "colorful": {"header": "E17055", "header_font": "FFFFFF", "alt_row": "FFF3E0", "border": "FABEAB", "accent": "FDCB6E"},
                "pastel": {"header": "6C5CE7", "header_font": "FFFFFF", "alt_row": "F3E5F5", "border": "CE93D8", "accent": "A29BFE"},
            }
            colors = themes.get(theme, themes["professional"])

            header_fill = PatternFill(start_color=colors["header"], end_color=colors["header"], fill_type="solid")
            header_font = Font(name='Calibri', size=11, bold=True, color=colors["header_font"])
            alt_fill = PatternFill(start_color=colors["alt_row"], end_color=colors["alt_row"], fill_type="solid")
            body_font = Font(name='Calibri', size=11)
            thin_border = Border(
                left=Side(style='thin', color=colors["border"]),
                right=Side(style='thin', color=colors["border"]),
                top=Side(style='thin', color=colors["border"]),
                bottom=Side(style='thin', color=colors["border"])
            )

            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = title[:31]

            # Parse CSV or tab-delimited content
            lines = [l.strip() for l in content.split('\n') if l.strip()]
            data = []
            for line in lines:
                sep = '\t' if '\t' in line else ','
                cells = [c.strip().strip('"') for c in line.split(sep)]
                data.append(cells)

            if not data:
                return json.dumps({"error": "No data provided"})

            # Write data
            for i, row in enumerate(data):
                if i == 0:
                    ws.append(row)  # Headers: preserve original case
                else:
                    ws.append([_format_text(c, mode=fmt_mode) if isinstance(c, str) else c for c in row])

            # Style header row
            for cell in ws[1]:
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border

            # Zebra striping for data rows
            for i, row in enumerate(ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=ws.max_column), start=2):
                if i % 2 == 0:
                    for cell in row:
                        cell.fill = alt_fill
                        cell.border = thin_border
                        cell.font = body_font
                else:
                    for cell in row:
                        cell.border = thin_border
                        cell.font = body_font

            # Auto-fit columns
            for col in ws.columns:
                max_len = 0
                col_letter = get_column_letter(col[0].column)
                for cell in col:
                    try:
                        max_len = max(max_len, len(str(cell.value or '')))
                    except:
                        pass
                ws.column_dimensions[col_letter].width = min(max_len + 3, 40)

            # Freeze top row
            ws.freeze_panes = 'A2'

            # Add auto-filter
            if ws.max_row > 1:
                ws.auto_filter.ref = "A1:%s%d" % (get_column_letter(ws.max_column), ws.max_row)

            # Add Excel Table if data exists
            if ws.max_row > 1 and ws.max_column > 0:
                try:
                    tab = Table(
                        displayName="Table_" + ws.title.replace(' ', '')[:20],
                        ref="A1:%s%d" % (get_column_letter(ws.max_column), ws.max_row)
                    )
                    style = TableStyleInfo(
                        name="TableStyleMedium2",
                        showFirstColumn=False,
                        showLastColumn=False,
                        showRowStripes=True,
                        showColumnStripes=False
                    )
                    tab.tableStyleInfo = style
                    ws.add_table(tab)
                except Exception:
                    pass

            out = io.BytesIO()
            wb.save(out)
            out.seek(0)
            url, fname = await self._save_and_link(out.getvalue(), _ensure_ext(title, "xlsx"), __request__, __user__=__user__)
            if url:
                return "Spreadsheet created: [%s](%s)" % (fname, url)
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def tracked_change(self, file_id: str, change_type: str, content: str, author: str = "Reviewer", paragraph_index: int = -1, output_filename: str = "", __user__=None, __request__=None) -> str:
        """Apply tracked changes (redlines) to a Word document with custom author name.
        Preserves original formatting and capitalization of the existing document.
        DOCX only -- returns an error for any other format.

        change_type: replace (use old_text|||new_text), insert (append text with redline; note
        this always appends the new paragraph at the END of the document -- paragraph_index is
        ignored for insert), delete (mark paragraph at paragraph_index for deletion)
        author: Name shown in Word's Track Changes (e.g., "Sergio Pedro")
        """
        try:
            import sqlite3 as s3
            conn2 = s3.connect(_DB_PATH)
            try:
                row = conn2.execute("SELECT filename, meta FROM file WHERE id=?", (file_id,)).fetchone()
                if not row:
                    row = conn2.execute("SELECT filename, meta FROM file WHERE filename LIKE ?", (f"%{file_id}%",)).fetchone()
                if not row:
                    return json.dumps({"error": self._err("file_not_found")})
                filename = _decode_filename(row[0]) if row[0] else row[0]
                meta = json.loads(row[1]) if row[1] else {}
                data = _read_file_bytes(meta.get("path", file_id))
                if not data:
                    return json.dumps({"error": "File not found on disk"})
            finally:
                conn2.close()

            # Check file type — track changes only supported for DOCX
            ftype = _detect_type(filename)
            if ftype != "docx":
                return json.dumps({"error": f"Track changes are only supported for DOCX files. {ftype.upper()} format does not support revision marks."})

            from docx import Document
            from docx_revisions import RevisionParagraph
            doc = Document(io.BytesIO(data))
            out_name = output_filename or filename
            results = []
    
            if change_type == "replace":
                parts = content.split("|||", 1)
                if len(parts) != 2:
                    return json.dumps({"error": "Format: old_text|||new_text"})
                find_t, replace_t = parts
                for i, p in enumerate(doc.paragraphs):
                    if paragraph_index >= 0 and i != paragraph_index:
                        continue
                    if find_t in p.text:
                        rp = RevisionParagraph.from_paragraph(p)
                        cnt = rp.replace_tracked(find_t, replace_t, author=author)
                        results.append(f"Para {i}: {cnt} replacements")
            elif change_type == "insert":
                p = doc.add_paragraph()
                rp = RevisionParagraph.from_paragraph(p)
                rp.add_tracked_insertion(content, author=author)
                results.append("Inserted tracked text")
            elif change_type == "delete":
                idx = int(content) if content.isdigit() else paragraph_index
                if idx >= 0 and idx < len(doc.paragraphs):
                    p = doc.paragraphs[idx]
                    rp = RevisionParagraph.from_paragraph(p)
                    rp.add_tracked_deletion(0, len(p.text), author=author)
                    results.append(f"Marked para {idx} for deletion")
                else:
                    return json.dumps({"error": f"Paragraph index {idx} out of range (document has {len(doc.paragraphs)} paragraphs)."})
            else:
                return json.dumps({"error": f"Unknown change_type: {change_type!r}. Use: replace, insert, delete."})

            if not results:
                # No branch above actually changed anything (e.g. "replace" whose old_text
                # wasn't found in any paragraph) -- saving here would produce an unchanged
                # file while still reporting success.
                return json.dumps({"error": "No tracked changes were applied -- the target text/paragraph was not found. No file was changed."})

            out = io.BytesIO()
            doc.save(out)
            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__, __user__=__user__)
            if url:
                return f"[{name}]({url})\n\nTracked changes by '{author}':\n" + "\n".join(results)
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})


    async def merge_sheets(self, file_ids: str, output_filename: str = "", __user__=None, __request__=None) -> str:
        """Merge every sheet from multiple Excel (.xlsx) files into one new workbook, preserving
        cell styles. Each source sheet is renamed "<source_filename>_<original_sheet_name>" in
        the output so sheets from different files never collide.

        Args:
            file_ids: Comma-separated list of file IDs to merge, e.g. "id1,id2,id3"
            output_filename: Optional output filename (default "merged_workbook.xlsx")
        """
        try:
            import sqlite3 as s3, openpyxl, io, os
            from copy import copy
            conn2 = s3.connect(_DB_PATH)
            try:
                ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
                wb_out = openpyxl.Workbook()
                wb_out.remove(wb_out.active)
                merged = 0
                for fid in ids:
                    row = conn2.execute("SELECT filename, meta FROM file WHERE id=?", (fid,)).fetchone()
                    if not row:
                        row = conn2.execute("SELECT filename, meta FROM file WHERE filename LIKE ?", ("%"+fid+"%",)).fetchone()
                    if not row:
                        continue
                    filename = _decode_filename(row[0]) if row[0] else row[0]
                    meta = json.loads(row[1]) if row[1] else {}
                    data = _read_file_bytes(meta.get("path", fid))
                    if not data:
                        continue
                    wb_src = openpyxl.load_workbook(io.BytesIO(data))
                    base_name = os.path.splitext(os.path.basename(filename))[0][:15]
                    for sn in wb_src.sheetnames:
                        ws_src = wb_src[sn]
                        sheet_name = (base_name + "_" + sn)[:31]
                        if sheet_name in wb_out.sheetnames:
                            # Two source files truncating to the same 15-char base_name
                            # (e.g. "Quarterly_Report_2023.xlsx" / "...2024.xlsx" both
                            # become "Quarterly_Repor") would otherwise collide -- Excel
                            # sheet names must be unique, so disambiguate with a suffix
                            # instead of letting a later sheet silently overwrite an
                            # earlier one under the same name.
                            suffix = 2
                            while f"{sheet_name[:29]}_{suffix}" in wb_out.sheetnames:
                                suffix += 1
                            sheet_name = f"{sheet_name[:29]}_{suffix}"
                        ws_out = wb_out.create_sheet(title=sheet_name)
                        for ri, row_data in enumerate(ws_src.iter_rows(), 1):
                            for ci, cell in enumerate(row_data, 1):
                                out_cell = ws_out.cell(row=ri, column=ci, value=cell.value)
                                if cell.has_style:
                                    out_cell.font = copy(cell.font)
                                    out_cell.fill = copy(cell.fill)
                                    out_cell.border = copy(cell.border)
                                    out_cell.alignment = copy(cell.alignment)
                                    out_cell.number_format = cell.number_format
                        merged += 1
                    wb_src.close()
            finally:
                conn2.close()
            if merged == 0:
                return json.dumps({"error": "No files could be merged"})
            out = io.BytesIO()
            wb_out.save(out)
            out.seek(0)
            fname = output_filename or "merged_workbook.xlsx"
            url, name = await self._save_and_link(out.read(), fname, __request__, __user__=__user__)
            if url:
                return f"[{name}]({url})\n\nMerged {merged} sheets from {len(ids)} files."
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def batch_process(self, file_ids: str, operation: str, params: str = "", output_filename: str = "", __user__=None, __request__=None) -> str:
        """Apply the same operation to multiple files in one call.

        Args:
            file_ids: Comma-separated list of file IDs, e.g. "id1,id2,id3"
            operation: Only two operations are supported:
                "replace" -- find-and-replace text. params must be "old_text|||new_text".
                "add_rows" -- append CSV rows (xlsx) via add_content(). params is the CSV content.
                Any other value returns an "unsupported operation" result for every file.
            params: Operation-specific payload, format depends on `operation` (see above).
            output_filename: Not currently used by either operation (each file keeps its own
                per-function output naming).
        """
        try:
            ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
            results = []
            failed = 0
            for fid in ids:
                if operation == "replace":
                    parts = params.split("|||", 1)
                    if len(parts) != 2:
                        results.append(f"  {fid}: FAILED -- params must be 'old_text|||new_text'")
                        failed += 1
                        continue
                    # Capture and inspect each call's actual return value -- the old code
                    # discarded it and unconditionally printed "replaced"/"rows added" even
                    # when the underlying call errored or made zero changes.
                    r = await self.replace_text(fid, parts[0], parts[1], "", __user__, __request__)
                    if r.strip().startswith("{") or r.startswith("No replacements made"):
                        results.append(f"  {fid}: FAILED -- {r[:150]}")
                        failed += 1
                    else:
                        results.append(f"  {fid}: replaced")
                elif operation == "add_rows":
                    r = await self.add_content(fid, params, "", __user__, __request__)
                    if r.strip().startswith("{"):
                        results.append(f"  {fid}: FAILED -- {r[:150]}")
                        failed += 1
                    else:
                        results.append(f"  {fid}: rows added")
                else:
                    results.append(f"  {fid}: unsupported operation '{operation}'")
                    failed += 1
            if results:
                summary = f"Batch processed {len(ids)} files ({len(ids) - failed} succeeded, {failed} failed):\n" + "\n".join(results)
                return summary
            return json.dumps({"error": f"No files processed. Unsupported operation: {operation}"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def auto_backup(self, __user__=None, __request__=None) -> str:
        """Create a timestamped backup copy of the entire Open WebUI database (webui.db) --
        not just this plugin's files. Useful to call before a risky bulk operation
        (batch_process, bulk_folder_ops with delete_old, etc.). Takes no parameters; each call
        writes a new "webui_backup_<timestamp>.db" file to a local backups/ directory (not an
        uploaded file, no download link is returned)."""
        try:
            import datetime
            db_path = _DB_PATH
            backup_dir = os.path.join(_get_owui_data_dir(), "backups")
            os.makedirs(backup_dir, exist_ok=True)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_name = f"webui_backup_{timestamp}.db"
            backup_path = os.path.join(backup_dir, backup_name)
            # sqlite3's own Connection.backup() (not shutil.copy2) -- a plain file copy of
            # a LIVE database can land mid-write, copying a page that's only half-flushed
            # to disk, which produces a torn/unrecoverable backup file. backup() uses
            # SQLite's online backup API, which is safe to run against a database another
            # process is actively writing to.
            src = sqlite3.connect(f"file:{db_path}?mode=ro", uri=True)
            try:
                dst = sqlite3.connect(backup_path)
                try:
                    src.backup(dst)
                finally:
                    dst.close()
            finally:
                src.close()
            size_kb = os.path.getsize(backup_path) / 1024
            return json.dumps({"success": True, "backup_path": backup_path, "size_kb": round(size_kb,1), "message": f"Backup: {backup_name} ({size_kb:.1f} KB)"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})



    async def merge_pdfs(self, file_ids: str, output_filename: str = "", __user__=None, __request__=None) -> str:
        """Concatenate multiple PDFs into one file, in the order given.

        Args:
            file_ids: Comma-separated list of PDF file IDs, e.g. "id1,id2,id3" -- pages appear
                in the output in this exact order.
            output_filename: Optional output filename (default "merged.pdf")
        """
        try:
            import fitz, sqlite3 as s3, io, os
            conn2 = s3.connect(_DB_PATH)
            try:
                ids = [fid.strip() for fid in file_ids.split(",") if fid.strip()]
                merger = fitz.open()
                count = 0
                for fid in ids:
                    row = conn2.execute("SELECT meta FROM file WHERE id=?", (fid,)).fetchone()
                    if not row:
                        row = conn2.execute("SELECT meta FROM file WHERE filename LIKE ?", ("%"+fid+"%",)).fetchone()
                    if not row:
                        continue
                    meta = json.loads(row[0]) if row[0] else {}
                    data = _read_file_bytes(meta.get("path", fid))
                    if not data:
                        continue
                    src = fitz.open(stream=data, filetype="pdf")
                    merger.insert_pdf(src)
                    src.close()
                    count += 1
            finally:
                conn2.close()
            if count == 0:
                merger.close()
                return json.dumps({"error": "No PDFs could be merged"})
            out = io.BytesIO()
            merger.save(out)
            merger.close()
            out.seek(0)
            fname = output_filename or "merged.pdf"
            url, name = await self._save_and_link(out.read(), fname, __request__, __user__=__user__)
            if url:
                return f"[{name}]({url})\n\nMerged {count} PDFs into one file."
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def split_pdf(self, file_id: str, pages_per_file: int = 1, output_filename: str = "", __user__=None, __request__=None) -> str:
        """Split a PDF into multiple smaller files, each with `pages_per_file` consecutive pages
        (the last file may have fewer). Output files are named "part_<start>_<end>.pdf" using
        1-based page numbers; `output_filename` is not currently used.

        Complement of merge_pdfs() -- use that to combine files back together.
        """
        try:
            import fitz, sqlite3 as s3, io, os
            conn2 = s3.connect(_DB_PATH)
            try:
                row = conn2.execute("SELECT meta FROM file WHERE id=?", (file_id,)).fetchone()
                if not row:
                    row = conn2.execute("SELECT meta FROM file WHERE filename LIKE ?", ("%"+file_id+"%",)).fetchone()
                if not row:
                    return json.dumps({"error": self._err("file_not_found")})
                meta = json.loads(row[0]) if row[0] else {}
                data = _read_file_bytes(meta.get("path", file_id))
                if not data:
                    return json.dumps({"error": "File not found on disk"})
            finally:
                conn2.close()
            src = fitz.open(stream=data, filetype="pdf")
            total_pages = src.page_count
            urls = []
            for start in range(0, total_pages, pages_per_file):
                end = min(start + pages_per_file, total_pages)
                sub = fitz.open()
                sub.insert_pdf(src, from_page=start, to_page=end-1)
                out = io.BytesIO()
                sub.save(out)
                sub.close()
                out.seek(0)
                part_name = f"part_{start+1}_{end}.pdf"
                url, name = await self._save_and_link(out.read(), part_name, __request__, __user__=__user__)
                if url:
                    urls.append(f"[{name}]({url})")
            src.close()
            if urls:
                return "Split into " + str(len(urls)) + " files:\n" + "\n".join(urls)
            return json.dumps({"error": "Could not split PDF"})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    async def tool_stats(self, __user__=None, __request__=None) -> str:
        """Diagnostics/introspection call, not document-related: reports counts of active
        Open WebUI tools and functions, how many files this plugin has created, and the
        database file size. Takes no parameters."""
        try:
            import sqlite3 as s3
            conn2 = s3.connect(_DB_PATH)
            try:
                try:
                    tool_count = conn2.execute("SELECT COUNT(*) FROM tool WHERE is_active=1").fetchone()[0]
                except Exception:
                    tool_count = 0
                try:
                    func_count = conn2.execute("SELECT COUNT(*) FROM function WHERE is_active=1").fetchone()[0]
                except Exception:
                    func_count = 0
                try:
                    model_count = conn2.execute("SELECT COUNT(*) FROM model WHERE is_active=1").fetchone()[0]
                except Exception:
                    model_count = 0
                exports_dir = os.path.join(_get_owui_data_dir(), "exports")
                export_count = len([f for f in os.listdir(exports_dir) if os.path.isfile(os.path.join(exports_dir, f))]) if os.path.exists(exports_dir) else 0
                db_size_kb = os.path.getsize(_DB_PATH) / 1024
            finally:
                conn2.close()
            return json.dumps({
                "tools": tool_count,
                "functions": func_count,
                "models": model_count,
                "exported_files": export_count,
                "db_size_kb": round(db_size_kb, 1)
            }, indent=2)
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

    

    async def cleanup_files(self, days_old: int = 30) -> str:
        """Remove generated Office files older than N days. Default 30 days."""
        import time as _time
        
        cutoff = int(_time.time()) - (days_old * 86400)
        
        conn = sqlite3.connect(_DB_PATH)
        try:
            conn.row_factory = sqlite3.Row
            rows = conn.execute(
                "SELECT id, filename, created_at FROM file WHERE meta LIKE '%office-plugin%' AND created_at < ?",
                (cutoff,)
            ).fetchall()
            
            if not rows:
                return f"No office-generated files older than {days_old} days found."
            
            deleted = []
            errors = []
            for row in rows:
                try:
                    fpath = os.path.join(_UPLOAD_DIR, row["id"])
                    if os.path.exists(fpath):
                        os.remove(fpath)
                    conn.execute("DELETE FROM file WHERE id = ?", (row["id"],))
                    deleted.append(row["filename"])
                except Exception as e:
                    errors.append(f"{row['filename']}: {e}")
            
            conn.commit()
        finally:
            conn.close()
        
        result = f"Cleaned up {len(deleted)} file(s) older than {days_old} days:\n"
        for f in deleted:
            result += f"- {f}\n"
        if errors:
            result += f"\nErrors ({len(errors)}):\n"
            for e in errors:
                result += f"- {e}\n"
        return result

    async def manage_revisions(self, file_id: str, action: str, output_filename: str = "", __user__=None, __request__=None) -> str:
        """List, accept_all or reject_all tracked changes in a Word document."""
        try:
            import sqlite3 as s3
            conn2 = s3.connect(_DB_PATH)
            try:
                row = conn2.execute("SELECT filename, meta FROM file WHERE id=?", (file_id,)).fetchone()
                if not row:
                    row = conn2.execute("SELECT filename, meta FROM file WHERE filename LIKE ?", (f"%{file_id}%",)).fetchone()
                if not row:
                    return json.dumps({"error": self._err("file_not_found")})
                filename = _decode_filename(row[0]) if row[0] else row[0]
                meta = json.loads(row[1]) if row[1] else {}
                data = _read_file_bytes(meta.get("path", file_id))
                if not data:
                    return json.dumps({"error": "File not found on disk"})
            finally:
                conn2.close()

            # Check file type — revision management only supported for DOCX
            ftype = _detect_type(filename)
            if ftype != "docx":
                return json.dumps({"error": f"Revision management is only supported for DOCX files. {ftype.upper()} format does not support revision marks."})

            from docx_revisions import RevisionDocument

            if action == "list":
                rdoc = RevisionDocument(io.BytesIO(data))
                revs = []
                failed = 0
                # `rdoc.paragraphs` explicitly EXCLUDES table paragraphs by design (per the
                # docx_revisions docstring) -- tracked changes inside a table cell were
                # silently missing from this list. `all_paragraphs` walks the body and
                # recurses into all tables (including nested tables) and already returns
                # ready-to-use RevisionParagraph objects, so the old code's extra
                # RevisionParagraph.from_paragraph(para) re-wrap is also removed.
                for rp in rdoc.all_paragraphs:
                    try:
                        if rp.has_track_changes:
                            for ins in rp.insertions:
                                revs.append({"type": "insertion", "author": ins.author, "text": ins.text[:100]})
                            for d in rp.deletions:
                                revs.append({"type": "deletion", "author": d.author, "text": d.text[:100]})
                    except Exception:
                        failed += 1
                result = {"revisions": revs, "count": len(revs)}
                if failed:
                    result["paragraphs_with_errors"] = failed
                return json.dumps(result, indent=2)
    
            out_name = output_filename or filename
            rdoc = RevisionDocument(io.BytesIO(data))
            if action == "accept_all":
                rdoc.accept_all()
                msg = "All track changes accepted"
            elif action == "reject_all":
                rdoc.reject_all()
                msg = "All track changes rejected"
            else:
                return json.dumps({"error": f"Unknown action: {action}"})
    
            out = io.BytesIO()
            rdoc.save(out)
            out.seek(0)
            url, name = await self._save_and_link(out.read(), out_name, __request__, __user__=__user__)
            if url:
                return f"[{name}]({url})\n\n{msg}."
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": str(e), "traceback": traceback.format_exc()})



    # --- v3.2.0: ODF Write ---
    async def create_odf(self, content: str, filename: str = "document", format: str = "odt", raw_text: bool = False, __user__=None, __request__=None) -> str:
        """Create a new ODF file (.odt, .ods, .odp).

        Args:
            content: Content for the file (CSV for .ods, text for .odt, markdown for .odp)
            filename: Output filename
            format: 'odt', 'ods', or 'odp'
            raw_text: If True, skip text formatting
        """
        fmt_mode = "preserve" if raw_text else "format"

        try:
            from odf.opendocument import OpenDocumentText, OpenDocumentSpreadsheet, OpenDocumentPresentation
            from odf.text import P, H
            from odf.table import Table, TableRow, TableCell
            import io

            if format == "ods":
                doc = OpenDocumentSpreadsheet()
                lines = content.split('\n')
                table = Table(name="Sheet1")
                for line in lines:
                    line = line.strip()
                    if not line: continue
                    cells = [c.strip() for c in line.split(',')]
                    row = TableRow()
                    for c in cells:
                        cell = TableCell()
                        cell.addElement(P(text=_format_text(c, mode=fmt_mode)))
                        row.addElement(cell)
                    table.addElement(row)
                doc.spreadsheet.addElement(table)
            elif format == "odp":
                # ODP requires real slides: a draw:page per slide (with a master page and a
                # text frame) -- adding text:h/text:p directly to office:presentation (the old
                # code) is not valid ODP structure and produces an empty/broken presentation
                # when opened in LibreOffice/PowerPoint, even though it "saves" without error.
                from odf.style import Style, MasterPage, PageLayout, PageLayoutProperties
                from odf.draw import Page, Frame, TextBox

                doc = OpenDocumentPresentation()
                layout = PageLayout(name="PL1")
                layout.addElement(PageLayoutProperties(pagewidth="28cm", pageheight="15.75cm", printorientation="landscape"))
                doc.automaticstyles.addElement(layout)
                master = MasterPage(name="Default", pagelayoutname=layout)
                doc.masterstyles.addElement(master)

                frame_style = Style(name="FrameStyle", family="presentation")
                doc.automaticstyles.addElement(frame_style)

                slides = []
                current = []
                for line in content.split('\n'):
                    line = line.strip()
                    if not line:
                        continue
                    if line.startswith('# ') and current:
                        slides.append(current)
                        current = [line]
                    else:
                        current.append(line)
                if current:
                    slides.append(current)
                if not slides:
                    slides = [["(empty slide)"]]

                for slide_lines in slides:
                    page = Page(masterpagename=master)
                    doc.presentation.addElement(page)
                    frame = Frame(stylename=frame_style, width="25cm", height="13.75cm", x="1.5cm", y="1cm")
                    page.addElement(frame)
                    textbox = TextBox()
                    frame.addElement(textbox)
                    for line in slide_lines:
                        text = line[2:].strip() if line.startswith('# ') else line
                        textbox.addElement(P(text=_format_text(text, mode=fmt_mode)))
            else:
                doc = OpenDocumentText()
                for line in content.split('\n'):
                    line = line.strip()
                    if not line: continue
                    if line.startswith('# '):
                        doc.text.addElement(H(outlinelevel=1, text=line[2:]))
                    elif line.startswith('## '):
                        doc.text.addElement(H(outlinelevel=2, text=line[3:]))
                    else:
                        doc.text.addElement(P(text=_format_text(line, mode=fmt_mode)))
            
            buf = io.BytesIO()
            doc.write(buf)
            buf.seek(0)
            url, fname = await self._save_and_link(buf.getvalue(), f"{filename}.{format}", __request__, __user__=__user__)
            if url:
                return f"ODF file created: [{fname}]({url})"
            return json.dumps({"error": self._err("could_not_save")})
        except ImportError:
            return json.dumps({"error": "odfpy not installed. Install with: pip install odfpy"})
        except Exception as e:
            return json.dumps({"error": str(e)})

    # --- v3.2.0: Format Conversion ---
    async def convert_format(self, file_id: str, target_format: str, __user__=None, __request__=None) -> str:
        """Convert between Office formats. Same-category conversions are clean: docx<->odt,
        xlsx<->ods, pptx<->odp. Cross-category conversions (e.g. xlsx->pptx, docx->xlsx) are
        technically accepted -- the source content is read then written via the target format's
        generator -- but produce garbled/lossy output since content isn't restructured for the
        target type. Prefer same-category conversions."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        base_name = os.path.splitext(filename)[0]
        
        if ftype in ("xlsx", "xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx":
            content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx":
            content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt", "ods", "odp"):
            content = _read_odf(file_bytes, filename)
        else:
            return json.dumps({"error": f"Unsupported source format: {ftype}"})
        
        if target_format in ("odt", "ods", "odp"):
            return await self.create_odf(content, base_name, target_format, __user__=__user__, __request__=__request__)
        elif target_format == "docx":
            return await self.generate_document(content, base_name, __user__=__user__, __request__=__request__)
        elif target_format == "pptx":
            return await self.generate_slides(content, base_name, __user__=__user__, __request__=__request__)
        elif target_format == "xlsx":
            return await self.generate_spreadsheet(content, base_name, __user__=__user__, __request__=__request__)
        return json.dumps({"error": f"Unsupported target format: {target_format}"})

    # --- v3.2.0: Template System ---
    async def save_template(self, name: str, content: str) -> str:
        """Save a document template for reuse."""
        templates = json.loads(self.valves.templates or "{}")
        templates[name] = content
        self.valves.templates = json.dumps(templates)
        return f"Template '{name}' saved."

    async def use_template(self, name: str, __user__=None, __request__=None, **kwargs) -> str:
        """Generate a document from a template saved via save_template(), replacing
        `{placeholder}` markers with values. Pass placeholder values as extra keyword arguments
        matching the placeholder names, e.g. use_template(name='x', client_name='Acme')
        replaces `{client_name}` in the template with 'Acme'."""
        templates = json.loads(self.valves.templates or "{}")
        if name not in templates:
            return f"Template '{name}' not found. Available: {', '.join(templates.keys())}"
        content = templates[name]
        for key, value in kwargs.items():
            content = content.replace(f"{{{key}}}", str(value))
        return await self.generate_document(content, name, __user__=__user__, __request__=__request__)

    async def list_templates(self) -> str:
        """List all saved templates."""
        templates = json.loads(self.valves.templates or "{}")
        if not templates:
            return "No templates saved."
        result = "Available templates:\n"
        for name in templates:
            preview = templates[name][:50].replace('\n', ' ')
            result += f"- {name}: {preview}...\n"
        return result

    # --- v3.2.0: Scheduled Cleanup ---
    async def schedule_cleanup(self, days_old: int = 30, interval_hours: int = 24) -> str:
        """Store a cleanup policy (days_old, interval_hours) for later use. This tool has no
        background scheduler -- the policy only takes effect when cleanup_files() or
        retention_policy(policy="apply") is called explicitly."""
        schedule = {"days_old": days_old, "interval_hours": interval_hours, "enabled": interval_hours > 0}
        self.valves.cleanup_schedule = json.dumps(schedule)
        if interval_hours > 0:
            return (
                f"Cleanup policy stored: remove files older than {days_old} days. "
                f"This does not run automatically -- call cleanup_files() or "
                f"retention_policy(policy=\"apply\") whenever you want it applied."
            )
        return "Stored cleanup policy disabled."

    # --- v3.2.0: Mail Merge ---
    async def mail_merge(self, template_file_id: str, data_file_id: str, output_prefix: str = "merged", __user__=None, __request__=None) -> str:
        """Generate personalized documents by merging CSV/Excel data into a DOCX template.
        In the template, use `{{field_name}}` placeholders matching the data file's column
        headers (for xlsx/xls, the first row is treated as headers) -- e.g. `{{first_name}}`.
        Produces one output document per data row."""
        from docx import Document
        import io, csv
        
        t_bytes, t_name, t_type = self._resolve_file(template_file_id)
        if not t_bytes:
            return json.dumps({"error": f"Template not found: {template_file_id}"})
        
        d_bytes, d_name, d_type = self._resolve_file(data_file_id)
        if not d_bytes:
            return json.dumps({"error": f"Data file not found: {data_file_id}"})
        
        rows = []
        if d_type in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(d_bytes))
            ws = wb.active
            headers = [str(c.value or '') for c in ws[1]]
            for row in ws.iter_rows(min_row=2, values_only=True):
                # `str(v or '')` treats 0, 0.0, and False as falsy -> silently blanked to
                # "" instead of "0"/"0.0"/"False". Only a real None should become "".
                rows.append(dict(zip(headers, ['' if v is None else str(v) for v in row])))
        else:
            reader = csv.DictReader(io.StringIO(d_bytes.decode('utf-8-sig')))
            rows = list(reader)

        if not rows:
            return json.dumps({"error": "No data rows found"})

        _MAX_ROWS = 200
        mail_merge_truncated_note = ""
        if len(rows) > _MAX_ROWS:
            mail_merge_truncated_note = f"\n\nData had {len(rows)} rows; only the first {_MAX_ROWS} were merged."
            rows = rows[:_MAX_ROWS]

        def merge_paragraph(para, row):
            for key, value in row.items():
                token = f"{{{{{key}}}}}"
                if token not in para.text:
                    continue
                # Prefer a per-run replace (preserves formatting exactly) when the token
                # fits entirely inside one run. Word often splits a typed "{{field}}"
                # across separate runs (spell-check state, rsid boundaries) -- when no
                # single run contains the whole token, the old code silently left the
                # literal "{{field}}" in the output. Fall back to a paragraph-level
                # rebuild in that case so the substitution actually happens.
                if any(token in run.text for run in para.runs):
                    for run in para.runs:
                        run.text = run.text.replace(token, value)
                elif para.runs:
                    merged = "".join(r.text for r in para.runs).replace(token, value)
                    para.runs[0].text = merged
                    for r in para.runs[1:]:
                        r.text = ""

        results = []
        errors = []
        for i, row in enumerate(rows):
            try:
                doc = Document(io.BytesIO(t_bytes))
                for para in doc.paragraphs:
                    merge_paragraph(para, row)
                for table in doc.tables:
                    for trow in table.rows:
                        for cell in trow.cells:
                            for para in cell.paragraphs:
                                merge_paragraph(para, row)
                buf = io.BytesIO()
                doc.save(buf)
                buf.seek(0)
                fname = f"{output_prefix}_{i+1}.docx"
                url, saved = await self._save_and_link(buf.getvalue(), fname, __request__, __user__=__user__)
                if url:
                    results.append(f"[{saved}]({url})")
                else:
                    errors.append(f"row {i+1}: could not save")
            except Exception as e:
                errors.append(f"row {i+1}: {str(e)}")

        result = f"Merged {len(results)} documents:\n" + "\n".join(results)
        if errors:
            result += f"\n{len(errors)} row(s) failed: " + "; ".join(errors)
        result += mail_merge_truncated_note
        return result

    # --- v3.2.0: Chart Generation ---
    async def add_chart(self, file_id: str, chart_type: str = "bar", title: str = "Chart", __user__=None, __request__=None) -> str:
        """Add a chart to an Excel spreadsheet. chart_type: bar, line, pie, scatter."""
        from openpyxl import load_workbook
        from openpyxl.chart import BarChart, LineChart, PieChart, ScatterChart, Reference
        from openpyxl.utils import get_column_letter
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        if ftype not in ("xlsx", "xls"):
            return json.dumps({"error": f"Charts only supported for Excel files. Got: {ftype}"})

        try:
            wb = load_workbook(io.BytesIO(file_bytes))
            ws = wb.active
            if ws.max_column < 2:
                return json.dumps({"error": "Sheet needs at least 2 columns (labels + data) to chart."})

            chart_types = {"bar": BarChart, "line": LineChart, "pie": PieChart, "scatter": ScatterChart}
            chart_class = chart_types.get(chart_type, BarChart)
            chart = chart_class()
            chart.title = title
            chart.style = 10

            if ws.max_row > 1:
                # Column 1 = category labels (e.g. month names), columns 2+ = actual data
                # series. Including column 1 in the data range (the old code) makes openpyxl
                # plot the labels themselves as a bogus numeric-looking series, and never
                # calls set_categories, so the category axis shows 1,2,3... instead of the
                # real labels.
                data = Reference(ws, min_col=2, min_row=1, max_row=ws.max_row, max_col=ws.max_column)
                cats = Reference(ws, min_col=1, min_row=2, max_row=ws.max_row)
                chart.add_data(data, titles_from_data=True)
                chart.set_categories(cats)

            chart_col = get_column_letter(ws.max_column + 2)
            ws.add_chart(chart, f"{chart_col}1")

            buf = io.BytesIO()
            wb.save(buf)
            buf.seek(0)
            url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
            if url:
                return f"Chart added: [{fname}]({url})"
            return json.dumps({"error": self._err("could_not_save")})
        except Exception as e:
            return json.dumps({"error": f"Chart creation failed: {str(e)}"})

    # --- v3.2.0: Watermark ---
    async def add_watermark(self, file_id: str, text: str = "DRAFT", __user__=None, __request__=None) -> str:
        """Add a watermark to a DOCX or PDF file. PDF: a real diagonal, rotated 45-degree
        watermark across the page. DOCX: NOT diagonal -- a large centered gray text banner
        added to the page header instead, since python-docx has no rotated-text support."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype == "docx":
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document(io.BytesIO(file_bytes))
            for section in doc.sections:
                header = section.header
                header.is_linked_to_previous = False
                p = header.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                run = p.add_run()
                run.text = text
                run.font.size = Pt(72)
                run.font.color.rgb = RGBColor(128, 128, 128)
            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)
        elif ftype == "pdf":
            try:
                import fitz
                pdf_doc = fitz.open(stream=file_bytes, filetype="pdf")
                for page in pdf_doc:
                    rect = page.rect
                    # insert_text's `rotate` only accepts multiples of 90 -- an arbitrary
                    # 45-degree diagonal needs `morph` (pivot point + rotation matrix)
                    # instead. `alpha` isn't a real insert_text parameter; use fill_opacity.
                    center = fitz.Point(rect.width / 2, rect.height / 2)
                    mat = fitz.Matrix(1, 1).prerotate(45)
                    page.insert_text(
                        (rect.width / 2 - 100, rect.height / 2), text, fontsize=72,
                        color=(0.5, 0.5, 0.5), fill_opacity=0.3, morph=(center, mat),
                    )
                buf = io.BytesIO()
                pdf_doc.save(buf)
                pdf_doc.close()
                buf.seek(0)
            except ImportError:
                return json.dumps({"error": "PyMuPDF not installed. Install with: pip install PyMuPDF"})
            except Exception as e:
                return json.dumps({"error": f"Watermark failed: {str(e)}"})
        else:
            return json.dumps({"error": f"Watermark not supported for {ftype}. Use DOCX or PDF."})
        
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url:
            return f"Watermark added: [{fname}]({url})"
        return json.dumps({"error": self._err("could_not_save")})

    # --- v3.2.0: File Preview ---
    async def preview_file(self, file_id: str, max_lines: int = 20) -> str:
        """Show a text preview of any Office file before downloading."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype in ("xlsx", "xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx":
            content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx":
            content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt", "ods", "odp"):
            content = _read_odf(file_bytes, filename)
        else:
            return json.dumps({"error": f"Preview not supported for {ftype}"})
        
        lines = content.split('\n')
        preview = '\n'.join(lines[:max_lines])
        total = len(lines)
        result = f"**{filename}** ({ftype.upper()}, {total} lines)\n\n```\n{preview}\n```"
        if total > max_lines:
            result += f"\n... ({total - max_lines} more lines)"
        return result

    # --- v3.2.0: Metadata Editing ---
    async def edit_metadata(self, file_id: str, author: str = "", title: str = "", subject: str = "", keywords: str = "", __user__=None, __request__=None) -> str:
        """Edit document metadata (author, title, subject, keywords)."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        changes = []
        if ftype == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            cp = doc.core_properties
            if author: cp.author = author; changes.append("author")
            if title: cp.title = title; changes.append("title")
            if subject: cp.subject = subject; changes.append("subject")
            if keywords: cp.keywords = keywords; changes.append("keywords")
            buf = io.BytesIO(); doc.save(buf); buf.seek(0)
        elif ftype == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(file_bytes))
            if author: wb.properties.creator = author; changes.append("author")
            if title: wb.properties.title = title; changes.append("title")
            if subject: wb.properties.subject = subject; changes.append("subject")
            buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        elif ftype == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            cp = prs.core_properties
            if author: cp.author = author; changes.append("author")
            if title: cp.title = title; changes.append("title")
            if subject: cp.subject = subject; changes.append("subject")
            if keywords: cp.keywords = keywords; changes.append("keywords")
            buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        else:
            return json.dumps({"error": f"Metadata editing not supported for {ftype}"})
        
        if not changes:
            return json.dumps({"error": "No metadata fields specified"})
        
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url:
            return f"Metadata updated ({', '.join(changes)}): [{fname}]({url})"
        return json.dumps({"error": self._err("could_not_save")})

    # --- v3.2.0: Accessibility Check ---
    def _pptx_shape_cNvPr(self, shape):
        """Return the shape's `p:cNvPr` XML element, or None. python-pptx has no `alt_text`
        property on Shape/Picture (despite it being a natural-looking attribute name to
        reach for) -- reading OR writing `shape.alt_text` silently creates/reads a throwaway
        Python instance attribute that is never serialized to the file. Real alt text lives
        in the `descr` attribute of this element."""
        from pptx.oxml.ns import qn
        el = shape._element
        for tag in ("p:nvPicPr", "p:nvSpPr", "p:nvGrpSpPr", "p:nvCxnSpPr", "p:nvGraphicFramePr"):
            nv = el.find(qn(tag))
            if nv is not None:
                cNvPr = nv.find(qn("p:cNvPr"))
                if cNvPr is not None:
                    return cNvPr
        return None

    async def check_accessibility(self, file_id: str) -> str:
        """Check a DOCX or PPTX document for accessibility issues (heading structure, image alt
        text). Other formats (xlsx, pdf, etc.) return an explicit "not supported" result rather
        than a false "passed, no issues" (no checks are run against them).
        """
        import io

        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})

        if ftype not in ("docx", "pptx"):
            return json.dumps({"error": f"check_accessibility only supports DOCX and PPTX. Got: {ftype}. No checks were run."})

        issues = []
        if ftype == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            headings = []
            for p in doc.paragraphs:
                if p.style.name.startswith('Heading'):
                    level = int(p.style.name.split()[-1]) if p.style.name.split()[-1].isdigit() else 1
                    headings.append((level, p.text[:50]))
            prev = 0
            for level, text in headings:
                if level > prev + 1:
                    issues.append(f"Heading skip: H{prev} to H{level} ('{text}')")
                prev = level
            if not headings:
                issues.append("No headings found - document may lack structure")
        elif ftype == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:
                        cNvPr = self._pptx_shape_cNvPr(shape)
                        if not (cNvPr is not None and cNvPr.get("descr")):
                            issues.append(f"Slide {i+1}: Image missing alt text")
        
        if not issues:
            return f"Accessibility check passed for {filename}. No issues found."
        result = f"**Accessibility Report: {filename}**\n\n"
        for issue in issues:
            result += f"- {issue}\n"
        result += f"\n{len(issues)} issue(s) found."
        return result

    # --- v3.2.0: Add Alt Text ---
    async def add_alt_text(self, file_id: str, slide_num: int = 1, shape_index: int = 0, alt_text: str = "", __user__=None, __request__=None) -> str:
        """Add alt text to an image in a PowerPoint slide."""
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype != "pptx":
            return json.dumps({"error": "Alt text only supported for PPTX files"})
        
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        if slide_num > len(prs.slides):
            return json.dumps({"error": f"Slide {slide_num} not found. Has {len(prs.slides)} slides."})
        
        slide = prs.slides[slide_num - 1]
        pic_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:
                if pic_count == shape_index:
                    cNvPr = self._pptx_shape_cNvPr(shape)
                    if cNvPr is None:
                        return json.dumps({"error": "Could not locate this image's XML element to set alt text."})
                    cNvPr.set("descr", alt_text)
                    buf = io.BytesIO()
                    prs.save(buf)
                    buf.seek(0)
                    url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
                    if url:
                        return f"Alt text added: [{fname}]({url})"
                    return json.dumps({"error": self._err("could_not_save")})
                pic_count += 1
        return json.dumps({"error": f"Image {shape_index} not found. Found {pic_count} images."})

    # --- v3.2.0: Progress Indicators ---
    async def _progress(self, current: int, total: int, operation: str = "Processing", __event_emitter__=None) -> None:
        """Emit progress via __event_emitter__ if available."""
        if __event_emitter__ is None:
            return
        try:
            await __event_emitter__({"type": "status", "data": {"description": f"{operation}: {current}/{total}", "done": current >= total}})
        except Exception:
            pass


    # --- v3.3.0: Document Comparison ---
    async def compare_documents(self, file_id_a: str, file_id_b: str) -> str:
        """Compare two documents and show differences. Supports xlsx, xls, docx, pptx, odt,
        ods, odp, pdf. If either file is an unsupported format, returns an explicit error
        instead of silently comparing empty text (which would falsely report "0 differences"
        for two completely different files)."""
        a_bytes, a_name, a_type = self._resolve_file(file_id_a)
        b_bytes, b_name, b_type = self._resolve_file(file_id_b)
        if not a_bytes or not b_bytes:
            return json.dumps({"error": "One or both files not found"})

        _SUPPORTED = ("xlsx", "xls", "docx", "pptx", "odt", "ods", "odp", "pdf")
        if a_type not in _SUPPORTED or b_type not in _SUPPORTED:
            return json.dumps({
                "error": f"compare_documents does not support this format pair ({a_type}, {b_type}). "
                         f"Supported: {', '.join(_SUPPORTED)}."
            })

        def get_text(ftype, fb, fn):
            if ftype in ("xlsx","xls"):
                return self._read_xlsx(fb, fn) if ftype == "xlsx" else self._read_xls(fb, fn)
            elif ftype == "docx": return self._read_docx(fb, fn)
            elif ftype == "pptx": return self._read_pptx(fb, fn)
            elif ftype in ("odt","ods","odp"): return _read_odf(fb, fn)
            elif ftype == "pdf":
                import fitz
                pdf = fitz.open(stream=fb, filetype="pdf")
                text = "\n".join(pdf.load_page(i).get_text() for i in range(pdf.page_count))
                pdf.close()
                return text
            return ""

        text_a = get_text(a_type, a_bytes, a_name)
        text_b = get_text(b_type, b_bytes, b_name)
        
        lines_a = text_a.split('\n')
        lines_b = text_b.split('\n')
        
        result = f"**Comparison: {a_name} vs {b_name}**\n\n"
        added = 0; removed = 0; changed = 0
        max_len = max(len(lines_a), len(lines_b))
        
        for i in range(max_len):
            la = lines_a[i] if i < len(lines_a) else None
            lb = lines_b[i] if i < len(lines_b) else None
            if la is None:
                result += f"+ L{i+1}: {lb}\n"; added += 1
            elif lb is None:
                result += f"- L{i+1}: {la}\n"; removed += 1
            elif la != lb:
                result += f"~ L{i+1}:\n  - {la}\n  + {lb}\n"; changed += 1
        
        result += f"\n**Summary:** {added} added, {removed} removed, {changed} changed"
        return result

    # --- v3.3.0: Export to Markdown ---
    async def export_to_markdown(self, file_id: str, __user__=None, __request__=None) -> str:
        """Export any Office file to Markdown format -- returns the extracted text INLINE
        (capped at 8000 characters, with a truncation note) as well as saving it as a
        downloadable .md file. For a long PDF where you need more than the first 8000
        characters, use read_file(file_id, page_start=..., page_end=...) instead, which
        pages through PDF text without a single-call size limit.

        For PDF specifically: this extracts the real embedded text directly (fast, one call,
        no OCR) with each page clearly labeled "--- Page N ---" -- use this first whenever you
        need a PDF's text with real page numbers (e.g. to confirm where an excerpt sits before
        calling add_comment/add_comments), instead of ocr_extract. Pages with little or no
        extractable text (likely scanned images) are flagged in the output; only use ocr_extract
        for those specific flagged pages, not the whole document.
        """
        import io
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        if ftype in ("xlsx","xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx":
            content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx":
            content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"):
            content = _read_odf(file_bytes, filename)
        elif ftype == "pdf":
            try:
                import fitz
            except ImportError:
                return json.dumps({"error": "PyMuPDF not installed"})
            pdf = fitz.open(stream=file_bytes, filetype="pdf")
            parts = []
            sparse_pages = []
            for i in range(pdf.page_count):
                text = pdf.load_page(i).get_text().strip()
                # A page with almost no extractable text is likely a scanned image,
                # not a real gap in an otherwise text-based PDF -- flag it instead of
                # silently omitting it, so the caller knows ocr_extract may be needed.
                if len(text) < 20:
                    sparse_pages.append(i + 1)
                parts.append(f"--- Page {i + 1} ---\n{text}" if text else f"--- Page {i + 1} ---\n[no extractable text]")
            pdf.close()
            content = "\n\n".join(parts)
            if sparse_pages:
                content += (
                    f"\n\n[Pages with little or no extractable text (likely scanned images): "
                    f"{', '.join(str(p) for p in sparse_pages)}. Use ocr_extract for these pages.]"
                )
        else:
            return json.dumps({"error": f"Export not supported for {ftype}"})

        base = os.path.splitext(filename)[0]
        md_content = f"# {base}\n\n{content}"
        md_bytes = md_content.encode('utf-8')

        url, fname = await self._save_and_link(md_bytes, f"{base}.md", __request__, __user__=__user__)
        if not url:
            return json.dumps({"error": self._err("could_not_save")})

        _INLINE_CAP = 8000
        preview = content[:_INLINE_CAP]
        if len(content) > _INLINE_CAP:
            preview += f"\n\n...[truncated, {len(content) - _INLINE_CAP} more characters -- see the full file at the link below, or use read_file(page_start=...) for PDF]"
        return f"Exported to Markdown: [{fname}]({url})\n\n{preview}"

    async def find_text(self, file_id: str, query: str, max_results: int = 20, context_chars: int = 200) -> str:
        """Find where a piece of text occurs in a document -- page (pdf/pptx), paragraph
        index (docx), or cell (xlsx/xls) -- with a short surrounding snippet per hit.

        Use this BEFORE add_comment/add_comments when you need to confirm where an excerpt
        is (or whether it's unique) without reading the whole document. For PDF, this uses
        the same word-level matching as add_comment's excerpt search, so a page reported
        here is exactly the page add_comment will find that excerpt on.

        Args:
            file_id: File to search.
            query: Exact text to find. May contain "..." to mark text omitted between two
                quoted spans, same as add_comment's excerpt parameter.
            max_results: Cap on the number of hits returned (default 20).
            context_chars: Characters of surrounding context per hit (default 200; docx/xlsx
                only -- pdf/pptx snippets are the matched shape/page text, already short).
        """
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": self._err("file_not_found")})

        hits = []
        if ftype == "pdf":
            import fitz
            pdf = fitz.open(stream=file_bytes, filetype="pdf")
            qnorm = _normalize_match_text(_split_on_ellipsis(query)[0]) if _split_on_ellipsis(query) else ""
            for pidx in range(pdf.page_count):
                if len(hits) >= max_results:
                    break
                page = pdf.load_page(pidx)
                occurrences, partial = self._pdf_find_all_excerpt_occurrences(page, query)
                if not occurrences:
                    continue
                page_text_norm = _normalize_match_text(page.get_text())
                snippet = page_text_norm
                if qnorm:
                    pos = page_text_norm.lower().find(qnorm.lower())
                    if pos >= 0:
                        start = max(0, pos - context_chars // 2)
                        end = min(len(page_text_norm), pos + len(qnorm) + context_chars // 2)
                        snippet = ("..." if start > 0 else "") + page_text_norm[start:end] + ("..." if end < len(page_text_norm) else "")
                for _ in occurrences:
                    hits.append({"page": pidx + 1, "snippet": snippet[:context_chars * 2]})
                    if len(hits) >= max_results:
                        break
            pdf.close()

        elif ftype == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = list(doc.paragraphs)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        paragraphs.extend(cell.paragraphs)
            i = 1
            while len(hits) < max_results:
                para, runs, _offsets, _partial = self._docx_find_excerpt_runs(doc, query, i)
                if not para:
                    break
                try:
                    idx = list(doc.paragraphs).index(para)
                except ValueError:
                    idx = None
                text = para.text
                snippet = text[:context_chars] + ("..." if len(text) > context_chars else "")
                hits.append({"paragraph_index": idx, "in_table": idx is None, "snippet": snippet})
                i += 1

        elif ftype == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            qnorm = _normalize_match_text(query).lower()
            for si, slide in enumerate(prs.slides, 1):
                if len(hits) >= max_results:
                    break
                for shape in slide.shapes:
                    if not (hasattr(shape, "text") and shape.text):
                        continue
                    t = shape.text
                    if qnorm and qnorm in _normalize_match_text(t).lower():
                        snippet = t[:context_chars] + ("..." if len(t) > context_chars else "")
                        hits.append({"slide": si, "snippet": snippet})
                        if len(hits) >= max_results:
                            break

        elif ftype in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True) if ftype == "xlsx" else None
            if ftype == "xls":
                return json.dumps({"error": "find_text does not support legacy .xls -- convert to .xlsx first."})
            qnorm = _normalize_match_text(query).lower()
            for ws in wb.worksheets:
                if len(hits) >= max_results:
                    break
                for row in ws.iter_rows():
                    if len(hits) >= max_results:
                        break
                    for cell in row:
                        if cell.value and qnorm in _normalize_match_text(str(cell.value)).lower():
                            hits.append({"sheet": ws.title, "cell": cell.coordinate, "snippet": str(cell.value)[:context_chars]})
                            if len(hits) >= max_results:
                                break
        else:
            return json.dumps({"error": f"find_text not supported for {ftype}. Supported: pdf, docx, pptx, xlsx."})

        return json.dumps({"query": query, "matches": len(hits), "hits": hits}, indent=2, ensure_ascii=False)

    # --- v3.3.0: Import from URL ---
    async def import_from_url(self, url: str, title: str = "Web Document", __user__=None, __request__=None) -> str:
        """Fetch a web page and convert it to a Word document. Extraction is crude: HTML tags
        are stripped with regex (no readability/main-content detection), so nav bars, footers,
        and script/style text may leak into the output alongside the real article text.
        Content is truncated to the first 50,000 characters."""
        url_err = _validate_outbound_url(url)
        if url_err:
            return json.dumps({"error": url_err})
        try:
            import urllib.request as _urllib
            req = _urllib.Request(url, headers={"User-Agent": "Mozilla/5.0"})
            resp = _urllib.urlopen(req, timeout=15)
            html = resp.read().decode('utf-8', errors='replace')
        except Exception as e:
            return json.dumps({"error": f"Failed to fetch URL: {str(e)}"})
        
        # Simple HTML to text
        import re as _re
        text = _re.sub(r'<script[^>]*>.*?</script>', '', html, flags=_re.DOTALL|_re.IGNORECASE)
        text = _re.sub(r'<style[^>]*>.*?</style>', '', text, flags=_re.DOTALL|_re.IGNORECASE)
        text = _re.sub(r'<[^>]+>', '\n', text)
        text = _re.sub(r'\n\s*\n', '\n\n', text)
        text = _re.sub(r'[ \t]+', ' ', text)
        text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
        
        if not text:
            return json.dumps({"error": "No text content extracted from URL"})

        doc_result = await self.generate_document(text[:50000], title, __user__=__user__, __request__=__request__)
        _INLINE_CAP = 4000
        preview = text[:_INLINE_CAP]
        if len(text) > _INLINE_CAP:
            preview += f"\n...[truncated, {len(text) - _INLINE_CAP} more characters -- see the full file at the link above]"
        return f"{doc_result}\n\n{preview}"

    # --- v3.3.0: File Versioning ---
    async def version_file(self, file_id: str, label: str = "", __user__=None, __request__=None) -> str:
        """Save a versioned copy of a file before editing."""
        import time as _time
        
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": f"File not found: {file_id}"})
        
        base, ext = os.path.splitext(filename)
        ts = _time.strftime("%Y%m%d_%H%M%S")
        label_str = f"_{label}" if label else ""
        version_name = f"{base}_v{ts}{label_str}{ext}"
        
        url, fname = await self._save_and_link(file_bytes, version_name, __request__, __user__=__user__)
        if url:
            return f"Version saved: [{fname}]({url})"
        return json.dumps({"error": "Could not save version"})

    # --- v3.3.0: Cloud Storage (Google Drive) ---
    async def upload_to_drive(self, file_id: str, folder_id: str = "root") -> str:
        """Upload a file to Google Drive. Requires google-api-python-client and credentials."""
        try:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
            from googleapiclient.http import MediaIoBaseUpload
            import io as _io
            
            file_bytes, filename, ftype = self._resolve_file(file_id)
            if not file_bytes:
                return json.dumps({"error": f"File not found: {file_id}"})
            
            creds_path = os.environ.get("GOOGLE_CREDENTIALS", "")
            if not creds_path or not os.path.exists(creds_path):
                return json.dumps({"error": "GOOGLE_CREDENTIALS env var not set or file not found"})
            
            mime_map = {"xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                       "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                       "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                       "pdf": "application/pdf"}
            mime = mime_map.get(ftype, "application/octet-stream")
            
            creds = service_account.Credentials.from_service_account_file(creds_path, scopes=["https://www.googleapis.com/auth/drive.file"])
            service = build("drive", "v3", credentials=creds)
            
            media = MediaIoBaseUpload(_io.BytesIO(file_bytes), mimetype=mime, resumable=True)
            file_metadata = {"name": filename, "parents": [folder_id]}
            drive_file = service.files().create(body=file_metadata, media_body=media, fields="id,webViewLink").execute()
            
            return f"Uploaded to Google Drive: {drive_file.get('webViewLink', drive_file.get('id'))}"
        except ImportError:
            return json.dumps({"error": "google-api-python-client not installed. pip install google-api-python-client google-auth"})
        except Exception as e:
            return json.dumps({"error": f"Drive upload failed: {str(e)}"})

    # --- v3.3.0: OCR ---
    async def ocr_extract(self, file_id: str, language: str = "eng", max_pages: int = 25) -> str:
        """Extract text from images in a document using OCR. PDF only -- does not OCR images
        embedded inside DOCX/PPTX files. Requires pytesseract and Pillow.

        Args:
            file_id: File ID to OCR
            language: Tesseract language code(s), e.g. "eng", "por", "eng+por"
            max_pages: Maximum number of PDF pages to OCR (default 25). Tesseract runs as a
                subprocess per page and is slow -- large PDFs are capped to keep this from
                running for many minutes. Pages beyond the cap are skipped, not silently dropped;
                the response says how many were processed vs. skipped.
        """
        try:
            import pytesseract
            from PIL import Image
            import asyncio as _asyncio

            file_bytes, filename, ftype = self._resolve_file(file_id)
            if not file_bytes:
                return json.dumps({"error": f"File not found: {file_id}"})

            if ftype == "pdf":
                try:
                    import fitz
                    pdf = fitz.open(stream=file_bytes, filetype="pdf")
                    total_pages = pdf.page_count
                    pages_to_process = total_pages if max_pages <= 0 else min(total_pages, max_pages)
                    loop = _asyncio.get_running_loop()

                    def _ocr_page(page_index, lang):
                        # Runs in a worker thread: PDF page RENDER (get_pixmap, ~100-300ms
                        # of GIL-blocking C code) and the Tesseract subprocess call both
                        # block, so both now happen here -- previously get_pixmap ran on
                        # the main event loop before handing only the OCR step off, so a
                        # 25-page document still froze the server for several seconds
                        # during rendering alone.
                        page = pdf.load_page(page_index)
                        pix = page.get_pixmap(dpi=200)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        return pytesseract.image_to_string(img, lang=lang)

                    results = []
                    for i in range(pages_to_process):
                        try:
                            text = await _asyncio.wait_for(
                                loop.run_in_executor(None, _ocr_page, i, language),
                                timeout=60,
                            )
                        except _asyncio.TimeoutError:
                            results.append(f"--- Page {i+1} ---\n[OCR timed out on this page after 60s, skipped]")
                            continue
                        if text.strip():
                            results.append(f"--- Page {i+1} ---\n{text.strip()}")
                    pdf.close()

                    output = "\n\n".join(results) if results else "No text found in PDF images"
                    if total_pages > pages_to_process:
                        output += (
                            f"\n\n[Processed {pages_to_process} of {total_pages} pages -- "
                            f"stopped at max_pages={max_pages}. Call again with a higher "
                            f"max_pages to continue.]"
                        )
                    return output
                except ImportError:
                    return json.dumps({"error": "PyMuPDF not installed"})
            else:
                return json.dumps({"error": f"OCR only supported for PDF files. Got: {ftype}"})
        except ImportError:
            return json.dumps({"error": "pytesseract or Pillow not installed. pip install pytesseract Pillow"})
        except Exception as e:
            return json.dumps({"error": f"OCR failed: {str(e)}"})

    def _err(self, key: str) -> str:
        """Look up a shared error message in the language set via translate_errors()."""
        lang = self.valves.language or "en"
        table = _ERROR_TRANSLATIONS.get(lang, _ERROR_TRANSLATIONS["en"])
        return table.get(key, _ERROR_TRANSLATIONS["en"].get(key, key))

    # --- v3.3.0: i18n Error Messages ---
    async def translate_errors(self, language: str = "en") -> str:
        """Set the language for error messages. Supported: en, pt, es, fr, de.

        Only applies to the shared "File not found" / "Could not save file" / "Unsupported
        format" messages used across most functions -- function-specific error text (e.g.
        messages that include a filename or reason) stays in English.
        """
        translations = _ERROR_TRANSLATIONS
        if language not in translations:
            return json.dumps({"error": f"Language '{language}' not supported. Available: {', '.join(translations.keys())}"})
        self.valves.language = language
        return f"Language set to {language}. Error messages will now appear in {language}."


    # --- v3.4.0: AI Summarize ---
    async def ai_summarize(self, file_id: str) -> str:
        """Extract document text for LLM summarization -- returns raw text (first 3000 chars,
        with a truncation note if longer), it does not generate a summary itself; the calling
        LLM reads the returned text and writes the summary. Supports xlsx, xls, docx, pptx,
        odt, ods, odp. Not supported: pdf."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": "File not found: " + str(file_id)})
        if ftype in ("xlsx","xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": "Unsupported format: " + str(ftype)})
        words = len(content.split())
        preview = content[:3000]
        return "**" + str(filename) + "** (" + str(words) + " words)\n\n" + preview + ("\n\n... (truncated, " + str(words) + " total words)" if len(content) > 3000 else "")

    # --- v3.4.0: Speaker Notes ---
    async def add_speaker_notes(self, file_id: str, slide_num: int, notes: str, __user__=None, __request__=None) -> str:
        """Add speaker notes to a PowerPoint slide."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype != "pptx": return json.dumps({"error": "Speaker notes only supported for PPTX"})
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        if slide_num < 1 or slide_num > len(prs.slides):
            return json.dumps({"error": "Slide " + str(slide_num) + " not found. Has " + str(len(prs.slides)) + " slides."})
        slide = prs.slides[slide_num - 1]
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url: return "Speaker notes added to slide " + str(slide_num) + ": [" + str(fname) + "](" + str(url) + ")"
        return json.dumps({"error": self._err("could_not_save")})

    # --- v3.4.0: Document Stats ---
    async def document_stats(self, file_id: str) -> str:
        """Show document statistics: word count, reading time, complexity."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype in ("xlsx","xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": self._err("unsupported")})
        words = len(content.split())
        chars = len(content)
        lines = content.count('\n') + 1
        sentences = content.count('.') + content.count('!') + content.count('?')
        reading_time = max(1, words // 200)
        avg_word_len = sum(len(w) for w in content.split()) / max(1, words)
        complexity = "Easy" if avg_word_len < 4 else ("Medium" if avg_word_len < 5.5 else "Complex")
        return "**Stats: " + str(filename) + "**\n- Words: " + str(words) + "\n- Characters: " + str(chars) + "\n- Lines: " + str(lines) + "\n- Sentences: " + str(sentences) + "\n- Reading time: ~" + str(reading_time) + " min\n- Avg word length: " + str(round(avg_word_len, 1)) + "\n- Complexity: " + str(complexity)

    # --- v3.4.0: QR Codes ---
    async def add_qr_code(self, file_id: str, data: str, __user__=None, __request__=None) -> str:
        """Add a QR code to a DOCX or PPTX file."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        try:
            import qrcode
            from PIL import Image
        except ImportError:
            return json.dumps({"error": "qrcode or Pillow not installed. pip install qrcode Pillow"})
        qr = qrcode.QRCode(version=1, box_size=10, border=2)
        qr.add_data(data); qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        qr_buf = io.BytesIO(); img.save(qr_buf, format='PNG'); qr_buf.seek(0)
        if ftype == "docx":
            from docx import Document
            from docx.shared import Inches
            doc = Document(io.BytesIO(file_bytes))
            doc.add_picture(qr_buf, width=Inches(2))
            buf = io.BytesIO(); doc.save(buf); buf.seek(0)
        elif ftype == "pptx":
            from pptx import Presentation
            from pptx.util import Inches
            prs = Presentation(io.BytesIO(file_bytes))
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(qr_buf, Inches(4), Inches(2), Inches(2))
            buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        else: return json.dumps({"error": "QR codes only supported for DOCX and PPTX"})
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url: return "QR code added: [" + str(fname) + "](" + str(url) + ")"
        return json.dumps({"error": self._err("could_not_save")})

    # --- v3.4.0: Bulk Folder Ops ---
    async def bulk_folder_ops(self, operation: str = "list", pattern: str = "*", __user__=None, __request__=None) -> str:
        """Apply an operation to all files in the uploads folder. Operations: list, delete_old
        (PERMANENT, irreversible deletion of every file matching `pattern` whose modified time
        is older than a fixed 30 days -- not configurable, no confirmation step), stats."""
        import glob as _glob, time as _time
        uploads = _UPLOAD_DIR
        # `pattern` is caller-supplied and fed straight into glob() -- ".."  or a path
        # separator would let it walk outside the uploads directory entirely (e.g.
        # pattern="../../*" reaching arbitrary files on the host), and delete_old would
        # then PERMANENTLY DELETE whatever matched. Reject it outright rather than trying
        # to sanitize a glob pattern.
        if ".." in pattern or "/" in pattern or "\\" in pattern or os.path.isabs(pattern):
            return json.dumps({"error": "Invalid pattern -- must be a simple glob within the uploads folder (no '..', '/', or '\\\\')."})
        uploads_real = os.path.realpath(uploads)
        if operation == "list":
            files = sorted(_glob.glob(os.path.join(uploads, pattern)))
            if not files: return "No files found matching: " + str(pattern)
            result = "**Files in uploads (" + str(len(files)) + "):**\n"
            for f in files[:50]:
                size = os.path.getsize(f)
                mtime = _time.strftime("%Y-%m-%d %H:%M", _time.localtime(os.path.getmtime(f)))
                result += "- " + os.path.basename(f) + " (" + str(round(size/1024,1)) + " KB, " + str(mtime) + ")\n"
            if len(files) > 50: result += "... and " + str(len(files)-50) + " more"
            return result
        elif operation == "delete_old":
            cutoff = _time.time() - (30 * 86400)
            deleted = 0
            for f in _glob.glob(os.path.join(uploads, pattern)):
                # Belt-and-suspenders: even with the pattern check above, confirm the
                # resolved path is still inside the uploads directory before removing it.
                if not os.path.realpath(f).startswith(uploads_real + os.sep):
                    continue
                if os.path.getmtime(f) < cutoff:
                    os.remove(f); deleted += 1
            return "Deleted " + str(deleted) + " files older than 30 days."
        elif operation == "stats":
            files = _glob.glob(os.path.join(uploads, pattern))
            total_size = sum(os.path.getsize(f) for f in files)
            return "**Uploads Stats:**\n- Files: " + str(len(files)) + "\n- Total size: " + str(round(total_size/1024/1024,1)) + " MB"
        return json.dumps({"error": "Unknown operation: " + str(operation) + ". Use: list, delete_old, stats"})

    # --- v3.4.0: File Search ---
    async def file_search(self, query: str, file_type: str = "all") -> str:
        """Search for text across files in the uploads folder. file_type: all, xlsx, docx,
        pptx, pdf. Note: only matches files that keep a real extension on disk (typically
        user uploads) -- files this plugin itself creates are saved without a filename/
        extension on disk (the human-readable name only lives in Open WebUI's file record),
        so freshly generated documents will not show up here."""
        import glob as _glob
        results = []
        patterns = {"all": "*.*", "xlsx": "*.xlsx", "docx": "*.docx", "pptx": "*.pptx", "pdf": "*.pdf"}
        pattern = patterns.get(file_type, "*.*")
        for fpath in _glob.glob(os.path.join(_UPLOAD_DIR, pattern)):
            try:
                fname = os.path.basename(fpath)
                ext = os.path.splitext(fname)[1].lower()
                with open(fpath, 'rb') as f: fb = f.read()
                if ext == ".pdf":
                    import fitz
                    pdf = fitz.open(stream=fb, filetype="pdf")
                    for pidx in range(pdf.page_count):
                        page_text = pdf.load_page(pidx).get_text()
                        if query.lower() in page_text.lower():
                            pos = page_text.lower().find(query.lower())
                            snippet = page_text[max(0, pos - 40):pos + len(query) + 40].replace("\n", " ")
                            results.append(f"{fname} (page {pidx + 1}): {snippet[:100]}")
                            break
                    pdf.close()
                elif ext == ".xlsx":
                    import openpyxl, io
                    wb = openpyxl.load_workbook(io.BytesIO(fb))
                    for ws in wb.worksheets:
                        for row in ws.iter_rows(values_only=True):
                            for cell in row:
                                if cell and query.lower() in str(cell).lower():
                                    results.append(fname + ": " + str(cell)[:100])
                                    break
                elif ext == ".docx":
                    from docx import Document
                    import io
                    doc = Document(io.BytesIO(fb))
                    for p in doc.paragraphs:
                        if query.lower() in p.text.lower():
                            results.append(fname + ": " + p.text[:100])
                            break
                elif ext == ".pptx":
                    from pptx import Presentation
                    import io
                    prs = Presentation(io.BytesIO(fb))
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame and query.lower() in shape.text_frame.text.lower():
                                results.append(fname + ": " + shape.text_frame.text[:100])
                                break
            except Exception:
                pass
        if not results: return "No matches found for '" + str(query) + "'"
        return "**Search: " + str(query) + "** (" + str(len(results)) + " matches)\n" + "\n".join("- " + r for r in results[:20])

    # --- v3.4.0: Data Validation ---
    async def add_data_validation(self, file_id: str, sheet: str = "", col: str = "A", validation_type: str = "list", values: str = "", __user__=None, __request__=None) -> str:
        """Add data validation to an Excel column. Types: list, whole, decimal, date. values: comma-separated for list, or 'min,max' for numeric."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype not in ("xlsx","xls"): return json.dumps({"error": "Data validation only for Excel files"})
        from openpyxl import load_workbook
        from openpyxl.worksheet.datavalidation import DataValidation
        wb = load_workbook(io.BytesIO(file_bytes))
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
        col_letter = col.upper().replace("COL","").replace("COLUMN","").strip()
        if not col_letter.isalpha(): col_letter = "A"
        max_row = ws.max_row or 100
        cell_range = col_letter + "2:" + col_letter + str(max_row)
        if validation_type == "list":
            vals = [v.strip() for v in values.split(",") if v.strip()]
            formula = '"' + ",".join(vals) + '"'
            dv = DataValidation(type="list", formula1=formula, allow_blank=True)
        elif validation_type == "whole":
            parts = [p.strip() for p in values.split(",") if p.strip()]
            mn = int(parts[0]) if parts else 0
            mx = int(parts[1]) if len(parts) > 1 else 100
            dv = DataValidation(type="whole", operator="between", formula1=str(mn), formula2=str(mx))
        elif validation_type == "decimal":
            parts = [p.strip() for p in values.split(",") if p.strip()]
            mn = float(parts[0]) if parts else 0
            mx = float(parts[1]) if len(parts) > 1 else 100
            dv = DataValidation(type="decimal", operator="between", formula1=str(mn), formula2=str(mx))
        elif validation_type == "date":
            dv = DataValidation(type="date", operator="greaterThan", formula1="2000-01-01")
        else: return json.dumps({"error": "Unknown type: " + str(validation_type)})
        dv.error = "Invalid value"
        dv.errorTitle = "Validation Error"
        ws.add_data_validation(dv)
        dv.add(cell_range)
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url: return "Data validation added to column " + str(col_letter) + ": [" + str(fname) + "](" + str(url) + ")"
        return json.dumps({"error": self._err("could_not_save")})

    # --- v3.4.0: Named Ranges ---
    async def add_named_range(self, file_id: str, name: str, range_str: str = "", __user__=None, __request__=None) -> str:
        """Define a named range in Excel. range_str: 'A1:B10' or auto-detected from active sheet."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype not in ("xlsx","xls"): return json.dumps({"error": "Named ranges only for Excel"})
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
        from openpyxl.workbook.defined_name import DefinedName
        wb = load_workbook(io.BytesIO(file_bytes))
        ws = wb.active
        if not range_str:
            range_str = "A1:" + get_column_letter(ws.max_column) + str(ws.max_row)
        if not re.fullmatch(r"[A-Za-z_][A-Za-z0-9_.]*", name) or re.fullmatch(r"[A-Za-z]{1,3}\d+", name):
            return json.dumps({"error": f"Invalid name {name!r} -- must start with a letter/underscore, contain no spaces, and not look like a cell reference (e.g. 'Q1')."})
        # Quote the sheet name if it needs it (contains a space or other character not
        # legal in a bare reference) -- e.g. "Sales Data!$A$1" is not a valid Excel
        # reference, it must be "'Sales Data'!$A$1". Also make the range fully absolute
        # ($A$1:$B$10, not the old code's column-only "$A1:$B10") so it doesn't shift if
        # rows/columns are later inserted.
        sheet_title = ws.title
        if not re.fullmatch(r"[A-Za-z_][A-Za-z0-9_.]*", sheet_title):
            sheet_ref = "'" + sheet_title.replace("'", "''") + "'"
        else:
            sheet_ref = sheet_title
        abs_range = re.sub(r"([A-Za-z]+)(\d+)", r"$\1$\2", range_str)
        dn = DefinedName(name, attr_text=f"{sheet_ref}!{abs_range}")
        wb.defined_names.add(dn)
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url: return "Named range '" + str(name) + "' = " + str(range_str) + ": [" + str(fname) + "](" + str(url) + ")"
        return json.dumps({"error": self._err("could_not_save")})

    # --- v3.4.0: Slide Transitions ---
    async def add_slide_transitions(self, file_id: str, transition_type: str = "fade", duration: float = 0.5, __user__=None, __request__=None) -> str:
        """Add transitions to all slides in a PPTX. Types: fade, push, wipe, split, random
        (picks a different one of fade/push/wipe/split for each slide)."""
        import io, random as _random
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype != "pptx": return json.dumps({"error": "Transitions only for PPTX"})
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        prs = Presentation(io.BytesIO(file_bytes))
        real_types = ["fade", "push", "wipe", "split"]
        failures = 0
        for slide in prs.slides:
            try:
                trans_elem = slide._element.find(qn('p:transition'))
                if trans_elem is not None:
                    # Remove and rebuild rather than appending duplicate child elements
                    # (e.g. two <p:fade/> under one <p:transition/>) on a second call.
                    slide._element.remove(trans_elem)
                trans_elem = etree.SubElement(slide._element, qn('p:transition'))
                # CT_Slide requires p:transition to precede p:timing -- appending after an
                # existing p:timing (the old code always appended last) produces
                # out-of-schema-order XML that PowerPoint may flag for repair.
                timing_elem = slide._element.find(qn('p:timing'))
                if timing_elem is not None:
                    timing_elem.addprevious(trans_elem)
                this_type = _random.choice(real_types) if transition_type == "random" else transition_type
                if this_type == "fade":
                    etree.SubElement(trans_elem, qn('p:fade'))
                elif this_type == "push":
                    etree.SubElement(trans_elem, qn('p:push'))
                elif this_type == "wipe":
                    etree.SubElement(trans_elem, qn('p:wipe'))
                elif this_type == "split":
                    etree.SubElement(trans_elem, qn('p:split'))
                trans_elem.set('advTm', str(int(duration * 1000)))
            except Exception:
                failures += 1
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if not url:
            return json.dumps({"error": self._err("could_not_save")})
        result = "Transitions added (" + str(transition_type) + "): [" + str(fname) + "](" + str(url) + ")"
        if failures:
            result += f"\n{failures} of {len(prs.slides)} slide(s) failed to get a transition."
        return result

    # --- v3.4.0: Export to HTML ---
    async def export_to_html(self, file_id: str, __user__=None, __request__=None) -> str:
        """Export any Office file to a styled HTML page."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype in ("xlsx","xls"):
            content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": "Export not supported for " + str(ftype)})
        base = os.path.splitext(filename)[0]
        html_c = content.replace('\n', '<br>\n')
        html_c = re.sub(r'^### (.+)$', r'<h3>\1</h3>', html_c, flags=re.MULTILINE)
        html_c = re.sub(r'^## (.+)$', r'<h2>\1</h2>', html_c, flags=re.MULTILINE)
        html_c = re.sub(r'^# (.+)$', r'<h1>\1</h1>', html_c, flags=re.MULTILINE)
        html_c = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', html_c)
        html_c = re.sub(r'\*(.+?)\*', r'<em>\1</em>', html_c)
        html_c = re.sub(r'`(.+?)`', r'<code>\1</code>', html_c)
        html_c = re.sub(r'^- (.+)$', r'<li>\1</li>', html_c, flags=re.MULTILINE)
        html = """<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>""" + str(base) + """</title>
<style>
body { font-family: 'Segoe UI', system-ui, sans-serif; max-width: 800px; margin: 40px auto; padding: 0 20px; line-height: 1.6; color: #333; background: #fafafa; }
h1 { color: #1a1a2e; border-bottom: 3px solid #e94560; padding-bottom: 10px; }
h2 { color: #16213e; margin-top: 30px; }
h3 { color: #0f3460; }
table { border-collapse: collapse; width: 100%; margin: 20px 0; }
th { background: #1a1a2e; color: white; padding: 10px; text-align: left; }
td { padding: 8px 10px; border-bottom: 1px solid #ddd; }
tr:nth-child(even) { background: #f5f5f5; }
code { background: #f0f0f0; padding: 2px 6px; border-radius: 3px; font-size: 0.9em; }
blockquote { border-left: 4px solid #e94560; margin: 20px 0; padding: 10px 20px; background: #fff5f5; }
</style>
</head>
<body>
""" + html_c + """
</body>
</html>"""
        html_bytes = html.encode('utf-8')
        url, fname = await self._save_and_link(html_bytes, base + ".html", __request__, __user__=__user__)
        if url: return "Exported to HTML: [" + str(fname) + "](" + str(url) + ")"
        return json.dumps({"error": self._err("could_not_save")})


    # === v3.6.0: AI-Powered Features ===

    async def ai_analyze(self, file_id: str) -> str:
        """Extract document text (first 5000 chars) plus an analysis prompt for the calling LLM
        to act on -- returns raw text and instructions, not an analysis result itself. Supports
        xlsx, xls, docx, pptx, odt, ods, odp. Not supported: pdf."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": self._err("unsupported")})
        words = len(content.split())
        preview = content[:5000]
        return f"**Document: {filename}** ({words} words)\n\nAnalyze this document and provide:\n1. Main topics (3-5 bullet points)\n2. Sentiment (positive/negative/neutral)\n3. Key entities (people, companies, dates)\n4. Executive summary (2-3 sentences)\n\n```\n{preview}\n```" + ("\n\n... (truncated)" if len(content) > 5000 else "")

    async def smart_fill(self, file_id: str, section: str, instruction: str, __user__=None, __request__=None) -> str:
        """Fill a document section using AI based on instructions. The LLM will generate content for the specified section."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": self._err("unsupported")})
        return f"**Smart Fill: {filename}**\n\nSection to fill: **{section}**\nInstructions: {instruction}\n\nCurrent document content:\n```\n{content[:3000]}\n```\n\nPlease generate the content for the '{section}' section based on the instructions and existing document context."

    async def grammar_check(self, file_id: str) -> str:
        """Extract document text (first 4000 chars) plus a grammar/style review prompt for the
        calling LLM to act on -- returns raw text and instructions, not corrections itself.
        Supports xlsx, xls, docx, pptx, odt, ods, odp. Not supported: pdf."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": self._err("unsupported")})
        return f"**Grammar Check: {filename}**\n\nReview this document for:\n1. Grammar errors\n2. Spelling mistakes\n3. Style inconsistencies\n4. Passive voice overuse\n5. Readability issues\n\nProvide corrections with line references:\n\n```\n{content[:4000]}\n```"

    async def translate_document(self, file_id: str, target_language: str, __user__=None, __request__=None) -> str:
        """Extract document text (first 4000 chars) plus a translation prompt for the calling
        LLM to act on -- returns raw text and instructions, not a translated document itself.
        Supports xlsx, xls, docx, pptx, odt, ods, odp. Not supported: pdf."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": self._err("unsupported")})
        return f"**Translate to {target_language}: {filename}**\n\nTranslate the following document to {target_language}. Preserve all formatting markers (# for headings, | for tables, - for bullets). Keep numbers, dates, and proper names unchanged.\n\n```\n{content[:4000]}\n```"

    async def classify_document(self, file_id: str) -> str:
        """Extract document text (first 2000 chars) plus a classification prompt for the calling
        LLM to act on -- returns raw text and instructions, not a classification result itself.
        Supports xlsx, xls, docx, pptx, odt, ods, odp. Not supported: pdf."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype in ("xlsx","xls"): content = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
        elif ftype == "docx": content = self._read_docx(file_bytes, filename)
        elif ftype == "pptx": content = self._read_pptx(file_bytes, filename)
        elif ftype in ("odt","ods","odp"): content = _read_odf(file_bytes, filename)
        else: return json.dumps({"error": self._err("unsupported")})
        return f"**Classify: {filename}**\n\nAnalyze this document and provide:\n1. Document type (report, proposal, invoice, contract, presentation, spreadsheet, letter, memo, manual, other)\n2. Primary theme/topic\n3. Department (finance, HR, marketing, engineering, sales, legal, operations, other)\n4. Confidentiality level (public, internal, confidential, restricted)\n5. Suggested tags (3-5 keywords)\n\n```\n{content[:2000]}\n```"

    async def smart_template(self, name: str, description: str, __user__=None, __request__=None) -> str:
        """Generate a document from a smart template that adapts to the conversation context.

        If `name` matches a template already saved via save_template(), generates and saves a
        real document from it (same as use_template()). Otherwise, no document is created --
        this returns guidance text describing the requested template so the calling model can
        draft the content itself, then save it with save_template() for reuse.
        """
        templates = json.loads(self.valves.templates or "{}")
        if name in templates:
            content = templates[name]
            return await self.generate_document(content, name, __user__=__user__, __request__=__request__)
        return f"**Smart Template: {name}**\n\nDescription: {description}\n\nNo saved template named '{name}' yet. Draft a professional document template for '{name}' with the following sections and use `{{placeholder}}` markers for customization, then save it with save_template(name='{name}', content=...). Use markdown format with # headings, - bullets, and | tables."

    # === v3.6.0: Data Manipulation ===

    async def add_pivot_table(self, file_id: str, rows_field: str = "", cols_field: str = "", data_field: str = "", aggregate: str = "sum", __user__=None, __request__=None) -> str:
        """Create a pivot table in Excel (xlsx/xls only), grouping by rows_field and aggregating
        data_field. If rows_field is omitted, this is a discovery call: instead of creating
        anything, it returns the first 10 column headers so you can pick valid field names for
        rows_field/data_field. aggregate: sum, count, average, min, max -- omitting data_field
        forces count mode regardless of the aggregate value."""
        import io
        from collections import defaultdict
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype not in ("xlsx","xls"): return json.dumps({"error": "Pivot tables only for Excel"})
        from openpyxl import load_workbook
        from openpyxl.utils import get_column_letter
        wb = load_workbook(io.BytesIO(file_bytes))
        ws = wb.active
        headers = [str(c.value or '') for c in ws[1]]
        if not rows_field:
            return f"**Available fields for pivot:** {', '.join(headers[:10])}\n\nUse: add_pivot_table(file_id, rows_field='FieldName', data_field='FieldName')"
        # Find column indices
        try:
            row_idx = headers.index(rows_field)
        except ValueError:
            return json.dumps({"error": f"Field '{rows_field}' not found in headers: {headers[:20]}"})
        data_idx = None
        if data_field:
            try:
                data_idx = headers.index(data_field)
            except ValueError:
                return json.dumps({"error": f"Field '{data_field}' not found in headers: {headers[:20]}"})
        # Read and aggregate data
        agg_map = defaultdict(list)
        skipped_non_numeric = 0
        for row in ws.iter_rows(min_row=2, values_only=True):
            row_key = str(row[row_idx]) if row[row_idx] is not None else "(blank)"
            if data_idx is not None and row[data_idx] is not None:
                try:
                    agg_map[row_key].append(float(row[data_idx]))
                except (ValueError, TypeError):
                    # A text-formatted number ("1,234.00", "$500") would previously coerce
                    # to 0.0 and silently corrupt the sum/average -- skip it instead and
                    # report the count, so a value doesn't get counted as zero.
                    skipped_non_numeric += 1
            elif data_idx is None:
                agg_map[row_key].append(1)  # count mode
        # Compute aggregate. Reuse the "Pivot" sheet name if it already exists -- openpyxl's
        # create_sheet()/title setter silently DE-DUPLICATES a colliding name into "Pivot1",
        # "Pivot2", etc. instead of raising, so a second run would write to a new sheet while
        # this function's own success message kept saying "Pivot", leaving the user reading
        # stale numbers in the sheet actually named "Pivot".
        if "Pivot" in wb.sheetnames:
            del wb["Pivot"]
        pivot_ws = wb.create_sheet("Pivot")
        agg_name = aggregate.lower()
        if agg_name == "count":
            pivot_ws.cell(row=1, column=1, value=rows_field)
            pivot_ws.cell(row=1, column=2, value="Count")
            for r, (key, vals) in enumerate(sorted(agg_map.items()), 2):
                pivot_ws.cell(row=r, column=1, value=key)
                pivot_ws.cell(row=r, column=2, value=len(vals))
        else:
            pivot_ws.cell(row=1, column=1, value=rows_field)
            pivot_ws.cell(row=1, column=2, value=f"{agg_name} of {data_field}")
            for r, (key, vals) in enumerate(sorted(agg_map.items()), 2):
                pivot_ws.cell(row=r, column=1, value=key)
                if agg_name == "sum":
                    pivot_ws.cell(row=r, column=2, value=sum(vals))
                elif agg_name == "average" or agg_name == "avg":
                    pivot_ws.cell(row=r, column=2, value=sum(vals) / len(vals) if vals else 0)
                elif agg_name == "min":
                    pivot_ws.cell(row=r, column=2, value=min(vals) if vals else 0)
                elif agg_name == "max":
                    pivot_ws.cell(row=r, column=2, value=max(vals) if vals else 0)
                else:
                    pivot_ws.cell(row=r, column=2, value=sum(vals))
        col2_header = "Count" if agg_name == "count" else f"{agg_name} of {data_field}"
        preview_rows = []
        for r in range(2, pivot_ws.max_row + 1):
            preview_rows.append(f"| {pivot_ws.cell(row=r, column=1).value} | {pivot_ws.cell(row=r, column=2).value} |")
        _PREVIEW_CAP = 50
        preview = f"| {rows_field} | {col2_header} |\n|---|---|\n" + "\n".join(preview_rows[:_PREVIEW_CAP])
        if len(preview_rows) > _PREVIEW_CAP:
            preview += f"\n...({len(preview_rows) - _PREVIEW_CAP} more rows in the saved file)"
        result = f"Pivot table created in sheet 'Pivot' ({len(agg_map)} rows). Fields: rows={rows_field}, data={data_field}, aggregate={aggregate}"
        if skipped_non_numeric:
            result += f"\n{skipped_non_numeric} row(s) had a non-numeric {data_field!r} value and were excluded from the aggregate (not counted as 0)."
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url: return f"{result}: [{fname}]({url})\n\n{preview}"
        return json.dumps({"error": self._err("could_not_save")})

    async def sql_to_spreadsheet(self, query: str, output_filename: str = "query_results", __user__=None, __request__=None) -> str:
        """Execute a SQL query on the local SQLite database and export results to Excel.
        WARNING: `query` runs verbatim against the app's own internal database -- this executes
        arbitrary SQL (not restricted to SELECT), including INSERT/UPDATE/DELETE/DROP against
        live application data. Use with caution and prefer read-only SELECT queries."""
        import io
        conn2 = sqlite3.connect(_DB_PATH)
        try:
            conn2.row_factory = sqlite3.Row
            cursor = conn2.execute(query)
            rows = cursor.fetchall()
            # sqlite3's default isolation_level opens an implicit transaction for any DML
            # statement; without an explicit commit(), conn2.close() below silently ROLLS
            # IT BACK -- so a query like "UPDATE ..." would report "no results" while
            # leaving the database completely unchanged, contradicting the docstring's
            # promise that non-SELECT statements actually run.
            conn2.commit()
        except Exception as e:
            conn2.rollback()
            return json.dumps({"error": f"SQL error: {str(e)}"})
        finally:
            conn2.close()
        if not rows:
            return f"Query executed (0 rows returned). If this was a write statement (INSERT/UPDATE/DELETE), it was committed -- {cursor.rowcount if cursor.rowcount >= 0 else 'unknown'} row(s) affected."
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        wb = Workbook(); ws = wb.active
        headers = list(rows[0].keys())
        for j, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=j, value=h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        for i, row in enumerate(rows, 2):
            for j, key in enumerate(headers, 1):
                ws.cell(row=i, column=j, value=row[key])
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), _ensure_ext(output_filename, "xlsx"), __request__, __user__=__user__)
        if not url:
            return json.dumps({"error": self._err("could_not_save")})
        _PREVIEW_CAP = 50
        preview_lines = ["| " + " | ".join(headers) + " |", "|" + "---|" * len(headers)]
        for row in rows[:_PREVIEW_CAP]:
            preview_lines.append("| " + " | ".join(str(row[h]) for h in headers) + " |")
        preview = "\n".join(preview_lines)
        if len(rows) > _PREVIEW_CAP:
            preview += f"\n...({len(rows) - _PREVIEW_CAP} more rows in the saved file)"
        return f"Query results ({len(rows)} rows): [{fname}]({url})\n\n{preview}"

    async def fill_pdf_form(self, file_id: str, field_values: str, __user__=None, __request__=None) -> str:
        """Fill a PDF form with values. field_values: 'field1=value1,field2=value2'."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype != "pdf": return json.dumps({"error": "PDF form filling only for PDF files"})
        try:
            import fitz
            pdf = fitz.open(stream=file_bytes, filetype="pdf")
            pairs = {}
            for pair in field_values.split(","):
                if "=" in pair:
                    k, v = pair.split("=", 1)
                    pairs[k.strip()] = v.strip()
            filled = 0
            matched_keys = set()
            all_field_names = []
            for page in pdf:
                for widget in page.widgets():
                    all_field_names.append(widget.field_name)
                    if widget.field_name in pairs:
                        widget.field_value = pairs[widget.field_name]
                        widget.update()
                        filled += 1
                        matched_keys.add(widget.field_name)
            if filled == 0:
                pdf.close()
                unmatched = list(pairs.keys())
                return json.dumps({
                    "error": "No fields were filled -- none of the requested field names matched this PDF's form fields.",
                    "requested_fields": unmatched,
                    "actual_field_names": all_field_names,
                })
            buf = io.BytesIO(); pdf.save(buf); pdf.close(); buf.seek(0)
            url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
            if not url:
                return json.dumps({"error": self._err("could_not_save")})
            result = f"Filled {filled} field(s): [{fname}]({url})"
            unmatched = [k for k in pairs if k not in matched_keys]
            if unmatched:
                result += f"\n{len(unmatched)} requested field(s) not found on this form: {unmatched}"
            return result
        except ImportError:
            return json.dumps({"error": "PyMuPDF not installed"})
        except Exception as e:
            return json.dumps({"error": str(e)})

    async def convert_data(self, file_id: str, target_format: str, __user__=None, __request__=None) -> str:
        """Convert between CSV, JSON, and XML formats. target_format: csv, json, xml.

        Only these source->target pairs do a real conversion: csv->json, xml->json, json->csv,
        json->xml. Any other pair (e.g. csv->xml, xml->csv) is not supported and returns an
        explicit error -- it does NOT silently write the original unconverted content under the
        new extension.
        """
        import io, csv as _csv
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        base = os.path.splitext(filename)[0]
        text = file_bytes.decode('utf-8', errors='replace')
        if target_format == "json":
            try:
                if ftype == "csv" or filename.endswith('.csv'):
                    reader = _csv.DictReader(io.StringIO(text))
                    data = list(reader)
                elif filename.endswith('.xml'):
                    import xml.etree.ElementTree as ET
                    root = ET.fromstring(text)
                    data = [{child.tag: child.text for child in elem} for elem in root]
                else:
                    data = [{"content": text}]
                result = json.dumps(data, indent=2, ensure_ascii=False)
                ext = ".json"
            except Exception as e:
                return json.dumps({"error": f"Conversion failed: {str(e)}"})
        elif target_format == "csv":
            try:
                if ftype == "json" or filename.endswith('.json'):
                    data = json.loads(text)
                    if isinstance(data, list) and data:
                        out = io.StringIO()
                        writer = _csv.DictWriter(out, fieldnames=data[0].keys())
                        writer.writeheader(); writer.writerows(data)
                        result = out.getvalue()
                    else:
                        return json.dumps({"error": "Conversion failed: JSON source is not a non-empty list of records"})
                else:
                    return json.dumps({"error": f"convert_data to csv only supports a JSON source. Got: {ftype}. No file was written."})
                ext = ".csv"
            except Exception as e:
                return json.dumps({"error": f"Conversion failed: {str(e)}"})
        elif target_format == "xml":
            try:
                if ftype == "json" or filename.endswith('.json'):
                    data = json.loads(text)
                    import xml.etree.ElementTree as ET
                    root = ET.Element("root")
                    for item in (data if isinstance(data, list) else [data]):
                        elem = ET.SubElement(root, "item")
                        for k, v in (item.items() if isinstance(item, dict) else {"value": str(item)}.items()):
                            child = ET.SubElement(elem, k)
                            child.text = str(v)
                    result = ET.tostring(root, encoding='unicode')
                else:
                    return json.dumps({"error": f"convert_data to xml only supports a JSON source. Got: {ftype}. No file was written."})
                ext = ".xml"
            except Exception as e:
                return json.dumps({"error": f"Conversion failed: {str(e)}"})
        else:
            return json.dumps({"error": f"Unsupported target format: {target_format}. Use: csv, json, xml"})
        result_bytes = result.encode('utf-8')
        url, fname = await self._save_and_link(result_bytes, f"{base}{ext}", __request__, __user__=__user__)
        if not url:
            return json.dumps({"error": self._err("could_not_save")})
        _INLINE_CAP = 4000
        preview = result[:_INLINE_CAP]
        if len(result) > _INLINE_CAP:
            preview += f"\n...[truncated, {len(result) - _INLINE_CAP} more characters -- see the full file at the link below]"
        return f"Converted to {target_format.upper()}: [{fname}]({url})\n\n{preview}"

    # === v3.6.0: Enterprise Features ===

    async def compliance_check(self, file_id: str, standard: str = "gdpr") -> str:
        """Check a DOCX document for compliance issues (GDPR keyword scan, heading/alt-text
        accessibility, confidentiality marking). standard: gdpr, accessibility, branding, all.

        Only DOCX is supported -- other formats (xlsx, pptx, pdf, etc.) return an explicit
        "not supported" result rather than a false "passed, no issues" (no checks are run
        against them).
        """
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype != "docx":
            return json.dumps({"error": f"compliance_check only supports DOCX. Got: {ftype}. No checks were run."})
        issues = []
        if standard in ("gdpr", "all"):
            from docx import Document; import io
            doc = Document(io.BytesIO(file_bytes))
            text = " ".join(p.text for p in doc.paragraphs)
            gdpr_keywords = ["email", "phone", "address", "name", "birth", "passport", "ssn", "tax id", "iban", "credit card", "ip address", "cookie"]
            found = [k for k in gdpr_keywords if k in text.lower()]
            if found: issues.append(f"GDPR: Personal data detected: {', '.join(found)}. Ensure consent and data processing agreement.")
        if standard in ("accessibility", "all"):
            from docx import Document; import io
            doc = Document(io.BytesIO(file_bytes))
            headings = [p for p in doc.paragraphs if p.style.name.startswith('Heading')]
            if not headings: issues.append("Accessibility: No headings found. Add heading structure.")
            images = len([r for r in doc.part.rels.values() if "image" in r.reltype])
            if images > 0: issues.append(f"Accessibility: {images} image(s) found. Ensure alt text is provided.")
        if standard in ("branding", "all"):
            from docx import Document; import io
            doc = Document(io.BytesIO(file_bytes))
            text = " ".join(p.text for p in doc.paragraphs)
            if "confidential" not in text.lower() and "draft" not in text.lower():
                issues.append("Branding: No confidentiality marking found. Consider adding DRAFT or CONFIDENTIAL watermark.")
        if not issues: return f"Compliance check passed for {filename}. No issues found."
        return f"**Compliance Report: {filename}**\n\n" + "\n".join(f"- {i}" for i in issues)

    async def audit_log(self, action: str = "list", limit: int = 20) -> str:
        """View or manage the audit trail of document operations."""
        import time as _time
        conn2 = sqlite3.connect(_DB_PATH)
        try:
            conn2.row_factory = sqlite3.Row
            if action == "list":
                rows = conn2.execute("SELECT filename, created_at FROM file WHERE meta LIKE '%office-plugin%' ORDER BY created_at DESC LIMIT ?", (limit,)).fetchall()
                if not rows: return "No audit records found."
                result = f"**Audit Trail (last {len(rows)}):**\n"
                for r in rows:
                    ts = _time.strftime("%Y-%m-%d %H:%M", _time.localtime(r["created_at"])) if r["created_at"] else "unknown"
                    result += f"- {ts}: {r['filename']}\n"
                return result
            elif action == "stats":
                total = conn2.execute("SELECT COUNT(*) FROM file WHERE meta LIKE '%office-plugin%'").fetchone()[0]
                today = conn2.execute("SELECT COUNT(*) FROM file WHERE meta LIKE '%office-plugin%' AND date(created_at, 'unixepoch') = date('now')").fetchone()[0]
                return f"**Audit Stats:**\n- Total files: {total}\n- Today: {today}"
            return json.dumps({"error": f"Unknown action: {action}. Use: list, stats"})
        finally:
            conn2.close()

    async def retention_policy(self, policy: str = "view", days: int = 90, file_type: str = "all") -> str:
        """Manage document retention policies. policy: view, set, apply."""
        if policy == "view":
            sched = json.loads(self.valves.cleanup_schedule or "{}")
            if sched.get("enabled"):
                return f"Retention policy active: delete files older than {sched.get('days_old', 30)} days, every {sched.get('interval_hours', 24)} hours."
            return "No retention policy active. Use: retention_policy(policy='set', days=90)"
        elif policy == "set":
            self.valves.cleanup_schedule = json.dumps({"days_old": days, "interval_hours": 24, "enabled": True, "file_type": file_type})
            return (
                f"Retention policy stored: {file_type} files older than {days} days. "
                f"This does not run automatically -- call retention_policy(policy=\"apply\") "
                f"whenever you want it applied."
            )
        elif policy == "apply":
            return await self.cleanup_files(days_old=days)
        return json.dumps({"error": f"Unknown policy: {policy}. Use: view, set, apply"})

    async def scheduled_report(self, action: str = "list", name: str = "", schedule: str = "", template: str = "", __user__=None, __request__=None) -> str:
        """Manage scheduled reports. action: list, create, delete, run."""
        reports = json.loads(self.valves.templates or "{}")
        scheduled = json.loads(self.valves.cleanup_schedule or "{}")
        if action == "list":
            sched_reports = {k: v for k, v in reports.items() if k.startswith("_scheduled_")}
            if not sched_reports: return "No scheduled reports."
            result = "**Scheduled Reports:**\n"
            for k, v in sched_reports.items():
                result += f"- {k.replace('_scheduled_', '')}: {v[:50]}...\n"
            return result
        elif action == "create":
            reports[f"_scheduled_{name}"] = json.dumps({"template": template, "schedule": schedule})
            self.valves.templates = json.dumps(reports)
            return f"Scheduled report '{name}' created."
        elif action == "delete":
            key = f"_scheduled_{name}"
            if key in reports:
                del reports[key]
                self.valves.templates = json.dumps(reports)
                return f"Scheduled report '{name}' deleted."
            return f"Report '{name}' not found."
        elif action == "run":
            key = f"_scheduled_{name}"
            if key not in reports: return f"Report '{name}' not found."
            cfg = json.loads(reports[key])
            return await self.generate_document(cfg.get("template", ""), name, __user__=__user__, __request__=__request__)
        return json.dumps({"error": f"Unknown action: {action}"})

    async def document_assembly(self, template_name: str, data_file_id: str, output_prefix: str = "assembled", __user__=None, __request__=None) -> str:
        """Assemble multiple documents by merging each row of data_file_id into template_name
        (a template previously saved via save_template()) -- produces one document per row.

        Args:
            template_name: Name of an existing saved template (see save_template/list_templates).
            data_file_id: xlsx/xls get real structured parsing (first row = headers). Any other
                format falls back to a raw CSV-decode attempt, which will fail or produce wrong
                results for non-CSV files -- prefer xlsx/xls or genuine CSV for data_file_id.
            output_prefix: Prefix for each generated document's filename.
        """
        # Load template content
        templates = json.loads(self.valves.templates or "{}")
        if template_name not in templates:
            return f"Template '{template_name}' not found. Available: {', '.join(templates.keys())}"
        template_content = templates[template_name]
        
        # Load data
        d_bytes, d_name, d_type = self._resolve_file(data_file_id)
        if not d_bytes:
            return json.dumps({"error": f"Data file not found: {data_file_id}"})
        
        # Parse data rows
        import csv as _csv, io
        rows = []
        if d_type in ("xlsx", "xls"):
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(d_bytes))
            ws = wb.active
            headers = [str(c.value or '') for c in ws[1]]
            for row in ws.iter_rows(min_row=2, values_only=True):
                rows.append(dict(zip(headers, [str(v or '') for v in row])))
        else:
            reader = _csv.DictReader(io.StringIO(d_bytes.decode('utf-8')))
            rows = list(reader)
        
        if not rows:
            return json.dumps({"error": "No data rows found"})

        _MAX_ROWS = 200
        truncated_note = ""
        if len(rows) > _MAX_ROWS:
            truncated_note = f"\n\nData had {len(rows)} rows; only the first {_MAX_ROWS} were assembled."
            rows = rows[:_MAX_ROWS]

        # Generate one document per data row
        results = []
        for i, row in enumerate(rows):
            content = template_content
            for key, value in row.items():
                content = content.replace(f"{{{{{key}}}}}", str(value)).replace(f"{{{key}}}", str(value))
            result = await self.generate_document(content, f"{output_prefix}_{i+1}", __user__=__user__, __request__=__request__)
            results.append(result)

        return f"Assembled {len(results)} documents from template '{template_name}':\n\n" + "\n".join(results) + truncated_note

    async def conditional_format(self, file_id: str, rules: str, __user__=None, __request__=None) -> str:
        """Apply conditional formatting rules to Excel. rules: 'col:A,op:>,val:100,color:27AE60;col:B,op:<,val:0,color:E74C3C'."""
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        if ftype not in ("xlsx","xls"): return json.dumps({"error": "Conditional formatting only for Excel"})
        from openpyxl import load_workbook
        from openpyxl.formatting.rule import CellIsRule
        from openpyxl.styles import PatternFill
        wb = load_workbook(io.BytesIO(file_bytes))
        ws = wb.active
        applied = 0
        for rule_str in rules.split(";"):
            parts = {}
            for p in rule_str.split(","):
                if ":" in p:
                    k, v = p.split(":", 1)
                    parts[k.strip()] = v.strip()
            if "col" not in parts: continue
            col = parts["col"].upper()
            op_map = {">": "greaterThan", "<": "lessThan", ">=": "greaterThanOrEqual", "<=": "lessThanOrEqual", "=": "equal", "!=": "notEqual"}
            op = op_map.get(parts.get("op", ">"), "greaterThan")
            val = parts.get("val", "0")
            color = parts.get("color", "27AE60")
            fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
            max_row = ws.max_row or 100
            cell_range = f"{col}2:{col}{max_row}"
            ws.conditional_formatting.add(cell_range, CellIsRule(operator=op, formula=[val], fill=fill))
            applied += 1
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if url: return f"Applied {applied} conditional formatting rule(s): [{fname}]({url})"
        return json.dumps({"error": self._err("could_not_save")})

    # === v3.6.0: Collaboration Features ===

    async def add_comment(self, file_id: str, text: str, author: str = "Reviewer", paragraph_index: int = 0, cell_ref: str = "A1", slide_num: int = 1, page_num: Optional[int] = None, excerpt: str = "", match_index: int = 1, __user__=None, __request__=None) -> str:
        """Add ONE review comment to a Word, Excel, PowerPoint, or PDF file -- for 2+ comments,
        call add_comments() instead (one call, one output file); calling add_comment() in a
        loop creates a separate file per comment, not one file with all the comments.

        Each call starts from file_id's original content and saves a new, independent file --
        it does not accumulate onto a previous add_comment output. To add another comment on top
        of one you just added, pass the file_id returned by that call. For adding many comments
        to the same PDF or DOCX in one go, use add_comments() instead.

        Args:
            file_id: File ID to comment on
            text: Comment text
            author: Name shown in the comment (e.g., "Sergio Pedro")
            paragraph_index: For DOCX: paragraph index to attach comment to. Used only if
                excerpt is empty or not found (default 0)
            cell_ref: For XLSX: cell reference (default "A1")
            slide_num: For PPTX: slide number (default 1)
            page_num: For PDF: page number to search on. Optional when `excerpt` is given --
                every page is searched in order and the comment lands wherever the excerpt is
                found (no need to know the page in advance). Required when `excerpt` is empty
                (nowhere else to place the sticky note), and used as the fallback location if
                a given excerpt isn't found anywhere.
            excerpt: For DOCX/PDF: exact quoted text to locate and anchor the comment to.
                Recommended over paragraph_index/fixed positioning -- highlights the matched
                text (PDF) or the matched run(s) (DOCX) instead of the whole paragraph/page.
                May contain "..." to mark text omitted between two quoted spans; the omitted
                middle is never searched for literally. For PDF, use find_text() first if
                you're unsure an excerpt is unique or want to confirm which page it's on.
            match_index: If excerpt appears more than once, which occurrence to use (1-based,
                default 1). For PDF, this counts occurrences across the WHOLE document when
                page_num is omitted, or on just that page when page_num is given.
        """
        import io
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})

        if ftype == "docx":
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            err, warning = self._docx_add_comment(doc, text, author, excerpt=excerpt, match_index=match_index, paragraph_index=paragraph_index)
            if err:
                return json.dumps({"error": err})
            buf = io.BytesIO(); doc.save(buf); buf.seek(0)
            url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
            if url:
                msg = f"Comment added by {author}: [{fname}]({url})"
                if warning:
                    msg += f"\nWarning: {warning}"
                return msg
            return json.dumps({"error": self._err("could_not_save")})

        elif ftype == "xlsx":
            from openpyxl import load_workbook
            from openpyxl.comments import Comment
            wb = load_workbook(io.BytesIO(file_bytes))
            ws = wb.active
            comment = Comment(text, author)
            ws[cell_ref].comment = comment
            buf = io.BytesIO(); wb.save(buf); buf.seek(0)
            url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
            if url: return f"Comment added by {author} on cell {cell_ref}: [{fname}]({url})"
            return json.dumps({"error": self._err("could_not_save")})

        elif ftype == "pptx":
            try:
                import zipfile
                import uuid as _uuid
                from datetime import datetime, timezone
                from lxml import etree


                slide_name = "ppt/slides/slide%d.xml" % slide_num
                rels_name = "ppt/slides/_rels/slide%d.xml.rels" % slide_num
                modern_name = "ppt/comments/commentModern%d.xml" % slide_num
                authors_name = "ppt/commentsAuthors.xml"

                with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                    entries = {name: z.read(name) for name in z.namelist()}

                if slide_name not in entries:
                    return json.dumps({"error": "Slide %d not found in presentation" % slide_num})

                etree.register_namespace("p", _P_NS)
                etree.register_namespace("p14", _P14_NS)
                etree.register_namespace("r", _R_NS)

                utc_now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%S.000")
                para_guid = "{" + str(_uuid.uuid4()).upper() + "}"

                # --- modern comment part: compute next comment idx ---
                if modern_name in entries:
                    mroot = etree.fromstring(entries[modern_name])
                    max_idx = 0
                    for el in mroot:
                        if el.tag == "{%s}cm" % _P_NS:
                            try:
                                max_idx = max(max_idx, int(el.get("idx") or 0))
                            except (TypeError, ValueError):
                                pass
                    next_idx = max_idx + 1
                else:
                    mroot = etree.Element("{%s}cmLst" % _P_NS)
                    next_idx = 1

                # --- commentsAuthors part: resolve or create author id ---
                if authors_name in entries:
                    aroot = etree.fromstring(entries[authors_name])
                else:
                    aroot = etree.Element("{%s}cmAuthorLst" % _P14_NS)

                author_id = None
                existing_ids = []
                author_els = []
                for el in aroot:
                    if el.tag == "{%s}cmAuthor" % _P14_NS:
                        author_els.append(el)
                        try:
                            existing_ids.append(int(el.get("id") or 0))
                        except (TypeError, ValueError):
                            existing_ids.append(0)
                        if el.get("name") == author:
                            author_id = el.get("id")

                if author_id is not None:
                    for el in author_els:
                        if el.get("name") == author:
                            el.set("lastIdx", str(next_idx))
                            break
                else:
                    new_id = (max(existing_ids) + 1) if existing_ids else 0
                    initials = "".join(w[0] for w in re.split(r"[\s._\-]+", author.strip()) if w)[:4].upper() or "AU"
                    new_author = etree.SubElement(aroot, "{%s}cmAuthor" % _P14_NS)
                    new_author.set("id", str(new_id))
                    new_author.set("name", author)
                    new_author.set("initials", initials)
                    new_author.set("lastIdx", "1")
                    new_author.set("clrIndex", str(len(author_els) % 14))
                    author_id = str(new_id)

                # --- append modern comment ---
                cm = etree.SubElement(mroot, "{%s}cm" % _P_NS)
                cm.set("authorId", author_id)
                cm.set("dt", utc_now)
                cm.set("idx", str(next_idx))
                pos = etree.SubElement(cm, "{%s}pos" % _P_NS)
                pos.set("x", "100")
                pos.set("y", "100")
                cm_txt = etree.SubElement(cm, "{%s}text" % _P_NS)
                cm_txt.text = text
                cm_extLst = etree.SubElement(cm, "{%s}extLst" % _P_NS)
                cm_ext = etree.SubElement(cm_extLst, "{%s}ext" % _P_NS)
                cm_ext.set("uri", "{D6B160D4-4F5E-48AF-9E34-8E18A86A1F0A}")
                comment_ex = etree.SubElement(cm_ext, "{%s}commentEx" % _P14_NS)
                comment_ex.set("paraId", para_guid)
                comment_ex.set("dt", utc_now)
                comment_ex.set("parentIdx", "0")
                ce_txt = etree.SubElement(comment_ex, "{%s}text" % _P14_NS)
                ce_txt.text = text
                ce_author = etree.SubElement(comment_ex, "{%s}authorId" % _P14_NS)
                ce_author.text = author_id

                # --- [Content_Types].xml: add overrides if missing ---
                ct_root = etree.fromstring(entries["[Content_Types].xml"])
                has_authors_ct = any(
                    el.tag == "{%s}Override" % _CT_NS and el.get("PartName") == "/ppt/commentsAuthors.xml"
                    for el in ct_root
                )
                if not has_authors_ct:
                    o = etree.SubElement(ct_root, "{%s}Override" % _CT_NS)
                    o.set("PartName", "/ppt/commentsAuthors.xml")
                    o.set("ContentType", _CT_AUTHORS)
                has_modern_ct = any(
                    el.tag == "{%s}Override" % _CT_NS and el.get("PartName") == "/ppt/comments/commentModern%d.xml" % slide_num
                    for el in ct_root
                )
                if not has_modern_ct:
                    o = etree.SubElement(ct_root, "{%s}Override" % _CT_NS)
                    o.set("PartName", "/ppt/comments/commentModern%d.xml" % slide_num)
                    o.set("ContentType", _CT_MODERN)

                # --- slide rels: add commentsModern relationship ---
                if rels_name in entries:
                    rroot = etree.fromstring(entries[rels_name])
                else:
                    rroot = etree.Element("{%s}Relationships" % _PKG_REL_NS, nsmap={None: _PKG_REL_NS})
                max_rid = 0
                for el in rroot:
                    rid = el.get("Id") or ""
                    m = re.match(r"rId(\d+)$", rid)
                    if m:
                        max_rid = max(max_rid, int(m.group(1)))
                new_rid = max_rid + 1
                target = "../comments/commentModern%d.xml" % slide_num
                # Deduplicate: reuse existing Relationship to same target (avoid duplicate on repeat calls)
                existing_rel = None
                for child in rroot:
                    if child.tag == "{%s}Relationship" % _PKG_REL_NS and child.get("Type") == _CM_REL_TYPE and child.get("Target") == target:
                        existing_rel = child
                        break
                if existing_rel is not None:
                    existing_rel.set("Id", "rId%d" % new_rid)
                else:
                    rel_el = etree.SubElement(rroot, "{%s}Relationship" % _PKG_REL_NS)
                    rel_el.set("Id", "rId%d" % new_rid)
                    rel_el.set("Type", _CM_REL_TYPE)
                    rel_el.set("Target", target)

                # --- slide XML: attach commentRel extension ---
                sroot = etree.fromstring(entries[slide_name])
                extLst = None
                for child in sroot:
                    if child.tag == "{%s}extLst" % _P_NS:
                        extLst = child
                        break
                if extLst is None:
                    extLst = etree.SubElement(sroot, "{%s}extLst" % _P_NS)
                # Deduplicate: reuse existing commentRel ext if present (avoid duplicate URI on repeat calls)
                slide_ext = None
                for child in extLst:
                    if child.tag == "{%s}ext" % _P_NS and child.get("uri") == "{6950BFC3-D8DA-4A85-94F7-54DA5524770B}":
                        slide_ext = child
                        break
                if slide_ext is None:
                    slide_ext = etree.SubElement(extLst, "{%s}ext" % _P_NS)
                    slide_ext.set("uri", "{6950BFC3-D8DA-4A85-94F7-54DA5524770B}")
                # Update or create commentRel child
                comment_rel = None
                for child in slide_ext:
                    if child.tag == "{%s}commentRel" % _P14_NS:
                        comment_rel = child
                        break
                if comment_rel is None:
                    comment_rel = etree.SubElement(slide_ext, "{%s}commentRel" % _P14_NS)
                comment_rel.set("{%s}id" % _R_NS, "rId%d" % new_rid)

                # --- serialize modified parts ---
                def _ser(root):
                    return etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)

                entries[modern_name] = _ser(mroot)
                entries[authors_name] = _ser(aroot)
                entries[rels_name] = _ser(rroot)
                entries[slide_name] = _ser(sroot)
                entries["[Content_Types].xml"] = _ser(ct_root)

                # --- rebuild zip (OPC: [Content_Types].xml first) ---
                buf = io.BytesIO()
                with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zout:
                    zout.writestr("[Content_Types].xml", entries["[Content_Types].xml"])
                    for name in entries:
                        if name != "[Content_Types].xml":
                            zout.writestr(name, entries[name])
                buf.seek(0)

                url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
                if url:
                    return f"Comment added by {author} on slide {slide_num}: [{fname}]({url})"
                return json.dumps({"error": self._err("could_not_save")})
            except Exception as e:
                return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

        elif ftype == "pdf":
            try:
                import fitz
                pdf = fitz.open(stream=file_bytes, filetype="pdf")
                err, warning, page_used = self._pdf_add_comment_annot(pdf, page_num, text, author, excerpt=excerpt, match_index=match_index)
                if err:
                    pdf.close()
                    return json.dumps({"error": err})
                buf = io.BytesIO()
                pdf.save(buf)
                pdf.close()
                buf.seek(0)
                url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
                if url:
                    msg = f"Comment added by {author} on page {page_used}: [{fname}]({url})"
                    if warning:
                        msg += f"\nWarning: {warning}"
                    return msg
                return json.dumps({"error": self._err("could_not_save")})
            except Exception as e:
                return json.dumps({"error": str(e), "traceback": traceback.format_exc()})

        else:
            return json.dumps({"error": f"Comments not supported for {ftype}. Supported: DOCX, XLSX, PPTX, PDF."})

    def _pdf_find_all_excerpt_occurrences(self, page, excerpt: str):
        """Find every occurrence of an excerpt on a fitz Page, word-by-word.

        `page.search_for` returns one rect PER LINE FRAGMENT, not per occurrence -- a
        quote that wraps a line comes back as 2+ separate rects for what is really one
        match, which breaks any attempt to count/index "occurrences" on top of it. This
        instead tokenizes the page via `page.get_text("words")` (each word carries its
        own rect) and slides the excerpt's word sequence over the page's word sequence,
        comparing punctuation-stripped, case-folded tokens. Each match's rects (one per
        word in the match, correctly spanning line wraps) are grouped as one occurrence.

        Handles excerpts containing '...'/'…' the same way as before: tries the full
        normalized excerpt first, then falls back to just the text before the first
        ellipsis.

        Returns (occurrences: list[list[fitz.Rect]], partial: bool) in reading order.
        `occurrences` is [] if nothing matched either candidate.
        """
        import fitz
        import string as _string_mod

        segments = _split_on_ellipsis(excerpt)
        if not segments:
            return [], False

        candidates = [(_normalize_match_text(excerpt), False)]
        if len(segments) > 1 and segments[0] != candidates[0][0]:
            candidates.append((segments[0], True))

        words = page.get_text("words")
        # sort into reading order: block, then line, then word position
        words = sorted(words, key=lambda w: (w[5], w[6], w[7]))
        page_tokens = [_normalize_chars(w[4]).lower().strip(_string_mod.punctuation) for w in words]

        for candidate_text, partial in candidates:
            if not candidate_text:
                continue
            q_tokens = [
                t.strip(_string_mod.punctuation)
                for t in re.findall(r"\S+", _normalize_chars(candidate_text).lower())
            ]
            q_tokens = [t for t in q_tokens if t]
            if not q_tokens:
                continue
            n, m = len(page_tokens), len(q_tokens)
            occurrences = []
            for i in range(n - m + 1):
                if page_tokens[i:i + m] == q_tokens:
                    occurrences.append([fitz.Rect(words[k][:4]) for k in range(i, i + m)])
            if occurrences:
                return occurrences, partial
        return [], False

    def _pdf_find_excerpt_quads(self, page, excerpt: str, match_index: int = 1):
        """Find the match_index'th occurrence of an excerpt on a fitz Page.

        Thin wrapper around `_pdf_find_all_excerpt_occurrences` for callers that want a
        single indexed occurrence (kept for compatibility with existing call sites).

        Returns (rects: list[fitz.Rect] or None, partial: bool, occurrence_count: int).
        `rects` is None both when nothing matched (occurrence_count == 0) and when
        match_index is out of range for the matches found (occurrence_count > 0) --
        callers must check occurrence_count to tell these apart.
        """
        occurrences, partial = self._pdf_find_all_excerpt_occurrences(page, excerpt)
        if not occurrences:
            return None, False, 0
        idx = match_index - 1
        if not (0 <= idx < len(occurrences)):
            return None, partial, len(occurrences)
        return occurrences[idx], partial, len(occurrences)

    def _pdf_add_comment_annot(self, pdf, page_num, text: str, author: str, excerpt: str = "", match_index: int = 1):
        """Add one comment annotation to an already-open fitz PDF.

        `page_num` (1-based) may be None when `excerpt` is given -- in that case every
        page is searched in order and `match_index` counts occurrences across the WHOLE
        document (not per-page), so the caller doesn't need to know the page in advance.
        `page_num` is still required when `excerpt` is empty (nowhere else to place the
        sticky note) or as the fallback location if the excerpt isn't found anywhere.

        Without `excerpt`: sticky-note icon at a fixed default position on `page_num`
        (legacy behavior; `page_num` required).
        With `excerpt`: highlights the matched text (all rects of the matched occurrence,
        so a quote spanning a line wrap is highlighted in full) and attaches the comment
        to that highlight -- clicking it in a PDF reader shows author + text, same as a
        manual "select text, add comment" in Acrobat/Preview. If the excerpt can't be
        found anywhere (or `match_index` exceeds the occurrences found), falls back to a
        sticky note on `page_num` (if given) and reports a warning -- never silently
        mis-places the comment on the wrong occurrence.

        Returns (error: Optional[str], warning: Optional[str], page_used: Optional[int]).
        """
        import fitz

        if page_num is not None and (page_num < 1 or page_num > pdf.page_count):
            return f"Page {page_num} not found in PDF (has {pdf.page_count} pages).", None, None

        if excerpt:
            pages_to_search = [page_num - 1] if page_num is not None else list(range(pdf.page_count))
            remaining = match_index
            total_found = 0
            partial_any = False
            for pidx in pages_to_search:
                page = pdf.load_page(pidx)
                occurrences, partial = self._pdf_find_all_excerpt_occurrences(page, excerpt)
                if partial:
                    partial_any = True
                total_found += len(occurrences)
                if len(occurrences) >= remaining and remaining >= 1:
                    rects = occurrences[remaining - 1]
                    annot = page.add_highlight_annot(rects)
                    annot.set_info(title=author, content=text, subject="Comment")
                    warning = (
                        "Matched only the text before the '...' in the excerpt -- verify placement."
                        if partial else None
                    )
                    return None, warning, pidx + 1
                remaining -= len(occurrences)

            # Not found (or match_index exceeded total occurrences across all pages searched).
            if page_num is None:
                if total_found > 0:
                    return (
                        f"Excerpt found {total_found} time(s) across the document, but "
                        f"match_index={match_index} is out of range and no page_num was given "
                        f"to fall back to.",
                        None,
                        None,
                    )
                return f"Excerpt not found anywhere in the {pdf.page_count}-page document.", None, None

            page = pdf.load_page(page_num - 1)
            annot = page.add_text_annot(fitz.Point(72, 72), text)
            annot.set_info(title=author, content=text, subject="Comment")
            if total_found > 0:
                warning = (
                    f"Excerpt found {total_found} time(s), but match_index={match_index} is out "
                    f"of range -- placed a default sticky note on page {page_num} instead."
                )
            else:
                warning = f"Excerpt not found on page {page_num} -- placed a default sticky note instead of highlighting the text."
            return None, warning, page_num

        if page_num is None:
            return "page_num is required when excerpt is not given.", None, None
        page = pdf.load_page(page_num - 1)
        annot = page.add_text_annot(fitz.Point(72, 72), text)
        annot.set_info(title=author, content=text, subject="Comment")
        return None, None, page_num

    def _docx_find_excerpt_runs(self, doc, excerpt: str, match_index: int = 1):
        """Locate an excerpt across a DOCX's body paragraphs and table cells.

        Returns (paragraph, runs, single_run_offsets, partial):
            paragraph: the matched Paragraph, or None if nothing matched
            runs: list of Run objects overlapping the match
            single_run_offsets: (run, local_start, local_end) when the excerpt fits
                entirely inside one run (enabling a precise split into just that
                span); None when the match spans a run boundary
            partial: True if only the text before the excerpt's '...' was matched

        Returns (None, None, None, False) if the excerpt wasn't found anywhere.
        """
        segments = _split_on_ellipsis(excerpt)
        if not segments:
            return None, None, None, False

        full_candidate = _normalize_match_text(excerpt)
        candidates = [(full_candidate, False)]
        if len(segments) > 1 and segments[0] != full_candidate:
            candidates.append((segments[0], True))

        paragraphs = list(doc.paragraphs)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    paragraphs.extend(cell.paragraphs)

        for candidate_text, partial in candidates:
            pattern = _loose_text_pattern(candidate_text)
            if not pattern:
                continue
            occurrence = 0
            for para in paragraphs:
                if not para.runs:
                    continue
                run_spans = []
                normalized_parts = []
                pos = 0
                for run in para.runs:
                    norm = _normalize_chars(run.text)
                    run_spans.append((run, pos, pos + len(norm)))
                    normalized_parts.append(norm)
                    pos += len(norm)
                para_text = "".join(normalized_parts)

                for m in re.finditer(pattern, para_text):
                    occurrence += 1
                    if occurrence != match_index:
                        continue
                    start, end = m.start(), m.end()
                    overlapping = [r for (r, rs, rend) in run_spans if rs < end and rend > start]
                    if len(overlapping) == 1:
                        run, rs, rend = next(t for t in run_spans if t[0] is overlapping[0])
                        return para, overlapping, (run, start - rs, end - rs), partial
                    return para, overlapping, None, partial

        return None, None, None, False

    def _docx_split_run_at(self, run, start: int, end: int):
        """Split `run` into up to 3 sibling runs at [start:end), preserving formatting.

        Returns the run containing exactly text[start:end] -- the piece to anchor a
        comment to. If the span already covers the whole run, returns it unchanged.

        Runs holding non-text content (inline images, line breaks, field/footnote refs)
        can't be safely split via text reassignment -- `Run.text`'s setter calls
        `clear_content()`, which would silently delete that non-text content along with
        the text. Such a run is returned unchanged instead (widening the anchored span
        to the whole run rather than destroying its content).
        """
        from copy import deepcopy
        from docx.text.run import Run
        from docx.oxml.ns import qn

        text = run.text
        if start <= 0 and end >= len(text):
            return run

        preserved_tags = {qn("w:t"), qn("w:rPr")}
        if any(child.tag not in preserved_tags for child in run._r):
            return run

        before, middle, after = text[:start], text[start:end], text[end:]
        src_element = run._r
        parent = run._parent

        middle_element = deepcopy(src_element)
        src_element.addnext(middle_element)

        if after:
            after_element = deepcopy(src_element)
            middle_element.addnext(after_element)
            Run(after_element, parent).text = after

        Run(middle_element, parent).text = middle

        if before:
            Run(src_element, parent).text = before
        else:
            src_element.getparent().remove(src_element)

        return Run(middle_element, parent)

    def _docx_add_comment(self, doc, text: str, author: str, excerpt: str = "", match_index: int = 1, paragraph_index=None):
        """Attach a comment to an already-open python-docx Document (not yet saved).

        With `excerpt`: locates the exact quoted text (paragraphs + table cells) and
        anchors the comment to just that span when it fits in one run, or to the
        minimal set of overlapping runs when the quote crosses a formatting boundary.
        If the excerpt isn't found and `paragraph_index` was explicitly given, falls back
        to that whole-paragraph anchor with a warning. If the excerpt isn't found and NO
        `paragraph_index` was given, returns an error instead of guessing -- placing an
        unmatched comment on paragraph 0 by default would silently stack every unmatched
        entry from a batch onto the title, which is worse than reporting the miss.
        Without `excerpt`: behaves exactly as before (whole-paragraph anchor).

        Returns (error: Optional[str], warning: Optional[str]).
        """
        if excerpt:
            para, runs, single_offsets, partial = self._docx_find_excerpt_runs(doc, excerpt, match_index)
            if runs:
                if single_offsets:
                    run, local_start, local_end = single_offsets
                    target_run = self._docx_split_run_at(run, local_start, local_end)
                    doc.add_comment([target_run], text=text, author=author)
                    warning = "Matched only the text before the '...' in the excerpt -- verify placement." if partial else None
                else:
                    doc.add_comment(runs, text=text, author=author)
                    warning = "Excerpt spans a formatting boundary -- comment covers a slightly wider span than the exact quote."
                    if partial:
                        warning = "Matched only the text before the '...', and the match spans a formatting boundary -- verify placement."
                return None, warning
            if paragraph_index is None:
                return "Excerpt not found anywhere in the document, and no paragraph_index was given as a fallback.", None

        idx = paragraph_index if paragraph_index is not None else 0
        if not doc.paragraphs:
            doc.add_paragraph("")
        if idx >= len(doc.paragraphs) or idx < 0:
            idx = 0
        para = doc.paragraphs[idx]
        if not para.runs:
            para.add_run("")
        doc.add_comment(para.runs, text=text, author=author)
        fallback_warning = "Excerpt not found in the document -- fell back to paragraph_index." if excerpt else None
        return None, fallback_warning

    async def add_comments(self, file_id: str, comments: list, __user__=None, __request__=None) -> str:
        """Add MULTIPLE review comments to a PDF or DOCX in ONE call -- this is the right choice
        whenever there is more than one comment to add, producing a single output file instead
        of one file per comment.

        Use this instead of calling add_comment() repeatedly -- each add_comment call starts
        fresh from the original file, so 28 calls would produce 28 separate single-comment files
        instead of one file with 28 comments. This opens the document once, applies every
        comment, and saves once.

        Args:
            file_id: File ID of the PDF or DOCX to comment on
            comments: List of dicts. Each entry supports:
                text: Comment text (required)
                author: Reviewer name (optional, default "Reviewer")
                excerpt: Exact quoted text to locate and anchor the comment to (recommended).
                    May contain "..." to mark omitted text between two quoted spans.
                match_index: If excerpt appears more than once, which occurrence to use
                    (1-based, default 1). For PDF, counts across the whole document when
                    page_num is omitted, or on just that page when page_num is given.
                page_num: PDF only -- optional when excerpt is given (every page is searched
                    in order). Required when excerpt is empty (nowhere else to place the
                    sticky note), and used as the fallback position if a given excerpt isn't
                    found anywhere. Use find_text() first if you want to confirm a page.
                paragraph_index: DOCX only -- optional. DOCX excerpt search scans the whole
                    document, so this is only used as a fallback if excerpt is omitted or
                    not found anywhere. If excerpt IS given but not found, the entry is
                    reported as an error rather than guessed at (no silent mis-placement).
        """
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes:
            return json.dumps({"error": self._err("file_not_found")})
        if ftype not in ("pdf", "docx"):
            return json.dumps({"error": f"add_comments only supports PDF and DOCX. Got: {ftype}. Use add_comment for other formats."})
        if not comments:
            return json.dumps({"error": "comments list is empty"})

        if ftype == "pdf":
            import fitz
            pdf = fitz.open(stream=file_bytes, filetype="pdf")
            applied = 0
            errors = []
            warnings = []
            placements = []
            for i, c in enumerate(comments):
                page_num = c.get("page_num")
                text = c.get("text")
                author = c.get("author", "Reviewer")
                excerpt = c.get("excerpt", "")
                match_index = c.get("match_index", 1)
                if text is None:
                    errors.append(f"entry {i}: missing text")
                    continue
                if page_num is None and not excerpt:
                    errors.append(f"entry {i}: missing page_num (required when excerpt is not given)")
                    continue
                err, warning, page_used = self._pdf_add_comment_annot(pdf, page_num, text, author, excerpt=excerpt, match_index=match_index)
                if err:
                    errors.append(f"entry {i}: {err}")
                    continue
                applied += 1
                placements.append(f"entry {i}: page {page_used}")
                if warning:
                    warnings.append(f"entry {i}: {warning}")

            if applied == 0:
                pdf.close()
                return json.dumps({"error": "No comments applied", "details": errors})

            buf = io.BytesIO()
            pdf.save(buf)
            pdf.close()
            buf.seek(0)

        else:  # docx
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            applied = 0
            errors = []
            warnings = []
            placements = []
            for i, c in enumerate(comments):
                text = c.get("text")
                author = c.get("author", "Reviewer")
                excerpt = c.get("excerpt", "")
                match_index = c.get("match_index", 1)
                paragraph_index = c.get("paragraph_index")
                if text is None:
                    errors.append(f"entry {i}: missing text")
                    continue
                err, warning = self._docx_add_comment(doc, text, author, excerpt=excerpt, match_index=match_index, paragraph_index=paragraph_index)
                if err:
                    errors.append(f"entry {i}: {err}")
                    continue
                applied += 1
                if warning:
                    warnings.append(f"entry {i}: {warning}")

            if applied == 0:
                return json.dumps({"error": "No comments applied", "details": errors})

            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)

        url, fname = await self._save_and_link(buf.getvalue(), filename, __request__, __user__=__user__)
        if not url:
            return json.dumps({"error": self._err("could_not_save")})
        result = f"Added {applied} comment(s) to [{fname}]({url})"
        if placements:
            result += f"\nPlacements: {'; '.join(placements)}"
        if errors:
            result += f"\n{len(errors)} entries skipped: {'; '.join(errors)}"
        if warnings:
            result += f"\n{len(warnings)} warning(s): {'; '.join(warnings)}"
        return result

    async def version_diff(self, file_id: str, version_label: str = "") -> str:
        """Show differences between current file and a previous version saved via version_file().
        Diff is a heuristic added/removed line COUNT (lines present in one version but not the
        other by exact-match membership) -- not a true positional/word-level diff, and doesn't
        show which specific lines changed. Supports xlsx, xls, docx, pptx. Not supported: pdf,
        odt/ods/odp."""
        file_bytes, filename, ftype = self._resolve_file(file_id)
        if not file_bytes: return json.dumps({"error": self._err("file_not_found")})
        base = os.path.splitext(filename)[0]
        ext = os.path.splitext(filename)[1]
        import time as _time
        # version_file() saves through _save_and_link, which writes the file to disk as a
        # bare UUID with NO name/extension -- the human filename only lives in the DB
        # `filename` column (base64-encoded). Globbing the upload directory for
        # "{base}_v*{ext}" (the old code) can never match anything on disk; versions have to
        # be found by decoding each DB filename instead.
        conn = sqlite3.connect(f"file:{_DB_PATH}?mode=ro", uri=True)
        try:
            rows = conn.execute("SELECT id, filename, created_at FROM file").fetchall()
        finally:
            conn.close()
        candidates = []
        for vid, encoded_name, created_at in rows:
            decoded = _decode_filename(encoded_name) if encoded_name else ""
            if decoded.startswith(f"{base}_v") and decoded.endswith(ext):
                candidates.append((vid, decoded, created_at))
        if not candidates:
            return f"No previous versions found for {filename}. Use version_file() to create versions."
        if version_label:
            candidates = [c for c in candidates if version_label in c[1]]
            if not candidates: return f"No version matching '{version_label}' found."
        candidates.sort(key=lambda c: c[2], reverse=True)
        vid, vname, created_at = candidates[0]
        vbytes = _read_file_bytes(vid)
        if not vbytes:
            return json.dumps({"error": f"Version record found ({vname}) but its file content is missing on disk."})
        vtime = _time.strftime("%Y-%m-%d %H:%M", _time.localtime(created_at))
        if ftype in ("xlsx","xls"):
            curr = self._read_xlsx(file_bytes, filename) if ftype == "xlsx" else self._read_xls(file_bytes, filename)
            prev = self._read_xlsx(vbytes, vname) if ftype == "xlsx" else self._read_xls(vbytes, vname)
        elif ftype == "docx":
            curr = self._read_docx(file_bytes, filename)
            prev = self._read_docx(vbytes, vname)
        elif ftype == "pptx":
            curr = self._read_pptx(file_bytes, filename)
            prev = self._read_pptx(vbytes, vname)
        else:
            return json.dumps({"error": f"Diff not supported for {ftype}"})
        cl = curr.split('\n'); pl = prev.split('\n')
        added = sum(1 for l in cl if l not in pl)
        removed = sum(1 for l in pl if l not in cl)
        return f"**Version Diff: {filename} vs {vname}** ({vtime})\n- Lines added: {added}\n- Lines removed: {removed}\n- Current: {len(cl)} lines\n- Previous: {len(pl)} lines"

    async def webhook_trigger(self, event: str = "test", url: str = "", file_id: str = "") -> str:
        """Trigger a webhook on document events. event: test, created, edited, deleted."""
        if not url:
            return json.dumps({"error": "url parameter required"})
        url_err = _validate_outbound_url(url)
        if url_err:
            return json.dumps({"error": url_err})
        import urllib.request as _urllib
        payload = json.dumps({"event": event, "file_id": file_id, "timestamp": __import__('time').time()}).encode()
        try:
            req = _urllib.Request(url, data=payload, headers={"Content-Type": "application/json"})
            resp = _urllib.urlopen(req, timeout=10)
            return f"Webhook sent to {url}: HTTP {resp.getcode()}"
        except Exception as e:
            return json.dumps({"error": f"Webhook failed: {str(e)}"})

    async def import_from_api(self, url: str, data_path: str = "", output_filename: str = "api_data", __user__=None, __request__=None) -> str:
        """Import data from a REST API (GET request, expects a JSON response) and export to Excel.

        Args:
            url: API endpoint to GET.
            data_path: Dot-notation path into the JSON response to reach the list/object to
                export, e.g. "results.items" or "data.0.rows" (numeric segments index into
                lists). Leave empty to use the whole response body.
            output_filename: Output filename (".xlsx" appended if missing).
        """
        url_err = _validate_outbound_url(url)
        if url_err:
            return json.dumps({"error": url_err})
        import urllib.request as _urllib, io
        try:
            req = _urllib.Request(url, headers={"Accept": "application/json", "User-Agent": "OpenWebUI/1.0"})
            resp = _urllib.urlopen(req, timeout=15)
            raw = json.loads(resp.read())
        except Exception as e:
            return json.dumps({"error": f"API request failed: {str(e)}"})
        data = raw
        if data_path:
            for key in data_path.split("."):
                if isinstance(data, dict):
                    data = data.get(key, [])
                elif isinstance(data, list) and key.isdigit():
                    data = data[int(key)]
        if not isinstance(data, list):
            data = [data] if isinstance(data, dict) else [{"value": str(data)}]
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        wb = Workbook(); ws = wb.active
        if data and isinstance(data[0], dict):
            # Union of keys across all records -- a single record's keys would silently
            # drop columns only present on later, heterogeneous records.
            headers = []
            for item in data:
                if isinstance(item, dict):
                    for k in item.keys():
                        if k not in headers:
                            headers.append(k)
        elif data:
            # List of scalars (e.g. data_path="ids" -> [1,2,3]) -- one "value" column,
            # instead of writing nothing and reporting a false "N records imported".
            headers = ["value"]
            data = [{"value": item} for item in data]
        else:
            headers = []
        for j, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=j, value=h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        for i, item in enumerate(data, 2):
            for j, key in enumerate(headers, 1):
                ws.cell(row=i, column=j, value=str(item.get(key, "")) if isinstance(item, dict) else str(item))
        buf = io.BytesIO(); wb.save(buf); buf.seek(0)
        url_out, fname = await self._save_and_link(buf.getvalue(), _ensure_ext(output_filename, "xlsx"), __request__, __user__=__user__)
        if not url_out:
            return json.dumps({"error": self._err("could_not_save")})
        _PREVIEW_CAP = 50
        preview_lines = ["| " + " | ".join(headers) + " |", "|" + "---|" * len(headers)] if headers else []
        for item in data[:_PREVIEW_CAP]:
            preview_lines.append("| " + " | ".join(str(item.get(h, "")) for h in headers) + " |")
        preview = "\n".join(preview_lines)
        if len(data) > _PREVIEW_CAP:
            preview += f"\n...({len(data) - _PREVIEW_CAP} more records in the saved file)"
        return f"Imported {len(data)} records from API: [{fname}]({url_out})\n\n{preview}"

